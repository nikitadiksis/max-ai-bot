from __future__ import annotations

import asyncio
import base64
import csv
from contextlib import asynccontextmanager, suppress
from collections import deque
from dataclasses import dataclass
from datetime import date, datetime, timedelta
import hmac
import hashlib
import html
from io import StringIO, BytesIO
import json
import logging
from logging.handlers import RotatingFileHandler
import os
from pathlib import Path
import re
import secrets
import sqlite3
from typing import Any
from urllib.parse import parse_qsl, quote, urlencode, urlsplit, urlunsplit

import aiohttp
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Request
from fastapi.staticfiles import StaticFiles
from fastapi.responses import PlainTextResponse
from fastapi.responses import FileResponse
from fastapi.responses import HTMLResponse
from fastapi.responses import RedirectResponse
from fastapi.responses import Response
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
import uvicorn
from runtime_support import (
    RuntimeSupportDeps,
    backups_dir as backups_dir_impl,
    create_db_backup as create_db_backup_impl,
    format_timedelta_short as format_timedelta_short_impl,
    latest_backup_file as latest_backup_file_impl,
    service_status_report as service_status_report_impl,
    smoke_check_report as smoke_check_report_impl,
    system_resource_snapshot as system_resource_snapshot_impl,
)
from storage_backend import StorageBackend, create_storage_backend

load_dotenv()


BASE_DIR = Path(__file__).resolve().parent
CONFIG_PATH = BASE_DIR / "models.json"
LOGS_DIR = BASE_DIR / "logs"
DATA_DIR = BASE_DIR / "data"
SITE_DIR = BASE_DIR / "site"
LOGS_DIR.mkdir(exist_ok=True)
DATA_DIR.mkdir(exist_ok=True)

MAX_API = "https://platform-api.max.ru"
OPENROUTER_CHAT_API = "https://openrouter.ai/api/v1/chat/completions"
KIE_GEMINI_CHAT_API_DEFAULT = "https://api.kie.ai/gemini-2.5-flash/v1/chat/completions"
KIE_GPT54_API_DEFAULT = "https://api.kie.ai/codex/v1/responses"
DEFAULT_MAX_MESSAGE_LEN = 3900

MAX_TOKEN = os.getenv("MAX_TOKEN", "").strip()
OPENROUTER_KEY = os.getenv("OPENROUTER_KEY", "").strip()
KIE_API_KEY = os.getenv("KIE_API_KEY", "").strip()
RUN_MODE = os.getenv("RUN_MODE", "polling").strip().lower()
HOST = os.getenv("HOST", "0.0.0.0").strip()
PORT = int(os.getenv("PORT", "8000"))
LOG_LEVEL = os.getenv("LOG_LEVEL", "INFO").upper()
LOG_FILE = os.getenv("LOG_FILE", str(LOGS_DIR / "bot.log"))
WEBHOOK_SECRET = os.getenv("WEBHOOK_SECRET", "").strip()
DEFAULT_TEXT_MODEL_ALIAS = os.getenv("DEFAULT_TEXT_MODEL", "gpt").strip().lower() or "gpt"
DEFAULT_IMAGE_MODEL_ALIAS = os.getenv("DEFAULT_IMAGE_MODEL", "image").strip().lower() or "image"
HISTORY_LIMIT = int(os.getenv("HISTORY_LIMIT", "20"))
MAX_MESSAGE_LEN = int(os.getenv("MAX_MESSAGE_LEN", str(DEFAULT_MAX_MESSAGE_LEN)))
DEDUP_CACHE_SIZE = int(os.getenv("DEDUP_CACHE_SIZE", "300"))
PROCESSED_UPDATE_TTL_HOURS = int(os.getenv("PROCESSED_UPDATE_TTL_HOURS", "72"))
DB_PATH = Path(os.getenv("DB_PATH", str(DATA_DIR / "bot.sqlite3")))
DATABASE_URL = os.getenv("DATABASE_URL", "").strip()
DB_BACKEND = create_storage_backend(DB_PATH, DATABASE_URL)
MAX_TEXT_INPUT_CHARS = int(os.getenv("MAX_TEXT_INPUT_CHARS", "2500"))
MAX_IMAGE_PROMPT_CHARS = int(os.getenv("MAX_IMAGE_PROMPT_CHARS", "800"))
MAX_ASSISTANT_OUTPUT_CHARS = int(os.getenv("MAX_ASSISTANT_OUTPUT_CHARS", "1400"))
MAX_CONTEXT_CHARS = int(os.getenv("MAX_CONTEXT_CHARS", "7000"))
MESSAGE_COOLDOWN_SECONDS = int(os.getenv("MESSAGE_COOLDOWN_SECONDS", "1"))
IMAGE_COOLDOWN_SECONDS = int(os.getenv("IMAGE_COOLDOWN_SECONDS", "20"))
MAX_API_CONCURRENCY = int(os.getenv("MAX_API_CONCURRENCY", "12"))
OPENROUTER_TEXT_CONCURRENCY = int(os.getenv("OPENROUTER_TEXT_CONCURRENCY", "8"))
OPENROUTER_IMAGE_CONCURRENCY = int(os.getenv("OPENROUTER_IMAGE_CONCURRENCY", "2"))
KIE_TEXT_CONCURRENCY = int(os.getenv("KIE_TEXT_CONCURRENCY", "6"))
TBANK_API_CONCURRENCY = int(os.getenv("TBANK_API_CONCURRENCY", "4"))
HTTP_RETRY_ATTEMPTS = int(os.getenv("HTTP_RETRY_ATTEMPTS", "3"))
HTTP_RETRY_BASE_MS = int(os.getenv("HTTP_RETRY_BASE_MS", "400"))
KIE_TEXT_ALIASES = {
    alias.strip().lower()
    for alias in os.getenv("KIE_TEXT_ALIASES", "gemini,gpt54").split(",")
    if alias.strip()
}
KIE_GEMINI_CHAT_API = os.getenv("KIE_GEMINI_CHAT_API", KIE_GEMINI_CHAT_API_DEFAULT).strip() or KIE_GEMINI_CHAT_API_DEFAULT
KIE_GPT54_API = os.getenv("KIE_GPT54_API", KIE_GPT54_API_DEFAULT).strip() or KIE_GPT54_API_DEFAULT
KIE_GPT54_MODEL = os.getenv("KIE_GPT54_MODEL", "gpt-5.4").strip() or "gpt-5.4"
LITE_PLAN_PRICE_RUB = int(os.getenv("LITE_PLAN_PRICE_RUB", "390"))
START_PLAN_PRICE_RUB = int(os.getenv("START_PLAN_PRICE_RUB", "990"))
PRO_PLAN_PRICE_RUB = int(os.getenv("PRO_PLAN_PRICE_RUB", "2490"))
LITE_PLAN_DAYS = int(os.getenv("LITE_PLAN_DAYS", "30"))
START_PLAN_DAYS = int(os.getenv("START_PLAN_DAYS", "30"))
PRO_PLAN_DAYS = int(os.getenv("PRO_PLAN_DAYS", "30"))
START_DAILY_GPT54_LIMIT = int(os.getenv("START_DAILY_GPT54_LIMIT", "3"))
PRO_DAILY_GPT54_LIMIT = int(os.getenv("PRO_DAILY_GPT54_LIMIT", "0"))
FREE_DAILY_CREDITS = int(os.getenv("FREE_DAILY_CREDITS", "40"))
FREE_FILE_COOLDOWN_DAYS = int(os.getenv("FREE_FILE_COOLDOWN_DAYS", "14"))
MAX_COMPLETION_TOKENS_FREE = int(os.getenv("MAX_COMPLETION_TOKENS_FREE", "500"))
MAX_COMPLETION_TOKENS_LITE = int(os.getenv("MAX_COMPLETION_TOKENS_LITE", "550"))
MAX_COMPLETION_TOKENS_START = int(os.getenv("MAX_COMPLETION_TOKENS_START", "650"))
MAX_COMPLETION_TOKENS_PRO = int(os.getenv("MAX_COMPLETION_TOKENS_PRO", "800"))
LITE_PLAN_CREDITS = int(os.getenv("LITE_PLAN_CREDITS", "5500"))
START_PLAN_CREDITS = int(os.getenv("START_PLAN_CREDITS", "15000"))
PRO_PLAN_CREDITS = int(os.getenv("PRO_PLAN_CREDITS", "40000"))
CREDIT_COST_DEEPSEEK = int(os.getenv("CREDIT_COST_DEEPSEEK", "1"))
CREDIT_COST_GPT = int(os.getenv("CREDIT_COST_GPT", "3"))
CREDIT_COST_GPTO = int(os.getenv("CREDIT_COST_GPTO", "4"))
CREDIT_COST_GEMINI = int(os.getenv("CREDIT_COST_GEMINI", "5"))
CREDIT_COST_GPT54 = int(os.getenv("CREDIT_COST_GPT54", "20"))
CREDIT_COST_IMAGE = int(os.getenv("CREDIT_COST_IMAGE", "35"))
CREDIT_COST_IMAGE_EDIT = int(os.getenv("CREDIT_COST_IMAGE_EDIT", "55"))
FILE_DOC_REQUEST_COST = max(1, int(os.getenv("FILE_DOC_REQUEST_COST", "10")))
FILE_PPT_REQUEST_COST = max(1, int(os.getenv("FILE_PPT_REQUEST_COST", "20")))
FILE_SHEET_REQUEST_COST = max(1, int(os.getenv("FILE_SHEET_REQUEST_COST", "12")))
PUBLIC_REQUEST_UNIT_CREDITS = max(1, int(os.getenv("PUBLIC_REQUEST_UNIT_CREDITS", "5")))
VAR_CREDITS_PER_1K_DEEPSEEK = int(os.getenv("VAR_CREDITS_PER_1K_DEEPSEEK", "0"))
VAR_CREDITS_PER_1K_GPT = int(os.getenv("VAR_CREDITS_PER_1K_GPT", "1"))
VAR_CREDITS_PER_1K_GPTO = int(os.getenv("VAR_CREDITS_PER_1K_GPTO", "1"))
VAR_CREDITS_PER_1K_GEMINI = int(os.getenv("VAR_CREDITS_PER_1K_GEMINI", "2"))
VAR_CREDITS_PER_1K_GPT54 = int(os.getenv("VAR_CREDITS_PER_1K_GPT54", "4"))
MAX_VARIABLE_CREDITS_PER_TEXT = int(os.getenv("MAX_VARIABLE_CREDITS_PER_TEXT", "4"))
TOPUP_SMALL_PRICE_RUB = int(os.getenv("TOPUP_SMALL_PRICE_RUB", "199"))
TOPUP_SMALL_CREDITS = int(os.getenv("TOPUP_SMALL_CREDITS", "1200"))
TOPUP_MEDIUM_PRICE_RUB = int(os.getenv("TOPUP_MEDIUM_PRICE_RUB", "499"))
TOPUP_MEDIUM_CREDITS = int(os.getenv("TOPUP_MEDIUM_CREDITS", "3200"))
TOPUP_LARGE_PRICE_RUB = int(os.getenv("TOPUP_LARGE_PRICE_RUB", "990"))
TOPUP_LARGE_CREDITS = int(os.getenv("TOPUP_LARGE_CREDITS", "7000"))
TOPUP_QUICK_CODE = os.getenv("TOPUP_QUICK_CODE", "medium").strip().lower() or "medium"
LOW_CREDITS_NUDGE_THRESHOLD_FREE = int(os.getenv("LOW_CREDITS_NUDGE_THRESHOLD_FREE", "10"))
LOW_CREDITS_NUDGE_THRESHOLD_PAID = int(os.getenv("LOW_CREDITS_NUDGE_THRESHOLD_PAID", "120"))
LOW_CREDITS_NUDGE_COOLDOWN_HOURS = int(os.getenv("LOW_CREDITS_NUDGE_COOLDOWN_HOURS", "18"))
ANALYTICS_USD_TO_RUB = float(os.getenv("ANALYTICS_USD_TO_RUB", "95"))
ANALYTICS_PAYMENT_FEE_PCT = float(os.getenv("ANALYTICS_PAYMENT_FEE_PCT", "2.5"))
ANALYTICS_RECEIPT_FEE_PCT = float(os.getenv("ANALYTICS_RECEIPT_FEE_PCT", "1.5"))
ANALYTICS_TAX_PCT = float(os.getenv("ANALYTICS_TAX_PCT", "6.0"))
ANALYTICS_EXPECTED_COST_PER_CREDIT_RUB = float(os.getenv("ANALYTICS_EXPECTED_COST_PER_CREDIT_RUB", "0.03"))
TRANSIENT_HTTP_STATUSES = {408, 409, 425, 429, 500, 502, 503, 504}
PAYMENT_DETAILS_TEXT = os.getenv(
    "PAYMENT_DETAILS_TEXT",
    "Реквизиты не настроены. Напиши администратору для оплаты.",
).replace("\\n", "\n").strip()
PUBLIC_BASE_URL = os.getenv("PUBLIC_BASE_URL", "").strip().rstrip("/")
TBANK_TERMINAL_KEY = os.getenv("TBANK_TERMINAL_KEY", "").strip()
TBANK_PASSWORD = os.getenv("TBANK_PASSWORD", "").strip()
TBANK_INIT_URL = os.getenv("TBANK_INIT_URL", "https://securepay.tinkoff.ru/v2/Init").strip()
TBANK_GET_STATE_URL = os.getenv("TBANK_GET_STATE_URL", "https://securepay.tinkoff.ru/v2/GetState").strip()
TBANK_NOTIFICATION_URL = os.getenv("TBANK_NOTIFICATION_URL", "").strip()
TBANK_SUCCESS_URL = os.getenv("TBANK_SUCCESS_URL", "").strip()
TBANK_FAIL_URL = os.getenv("TBANK_FAIL_URL", "").strip()
TBANK_RECEIPT_TAXATION = os.getenv("TBANK_RECEIPT_TAXATION", "usn_income").strip()
TBANK_RECEIPT_TAX = os.getenv("TBANK_RECEIPT_TAX", "none").strip()
TBANK_RECEIPT_PAYMENT_METHOD = os.getenv("TBANK_RECEIPT_PAYMENT_METHOD", "full_prepayment").strip()
TBANK_RECEIPT_PAYMENT_OBJECT = os.getenv("TBANK_RECEIPT_PAYMENT_OBJECT", "service").strip()
TBANK_RECEIPT_FFD_VERSION = os.getenv("TBANK_RECEIPT_FFD_VERSION", "1.05").strip()
TBANK_CANCEL_STATUSES = {"REJECTED", "CANCELED", "DEADLINE_EXPIRED"}
TBANK_REFUND_STATUSES = {"REFUNDED", "REVERSED", "PARTIAL_REVERSED", "PARTIAL_REFUNDED", "CHARGEDBACK"}
TBANK_SUCCESS_STATUSES = {"CONFIRMED"}
PAYMENT_FINAL_STATUSES = {"paid", "canceled", "refunded"}
PAYMENT_STATUS_LABELS = {
    "pending": "Ожидает оплату",
    "claimed": "Проверка оплаты",
    "paid": "Оплачено",
    "canceled": "Отменено",
    "refunded": "Возврат",
    "unknown": "Неизвестно",
}
SUPPORT_URL = os.getenv("SUPPORT_URL", "").strip()
SUPPORT_TEXT = os.getenv("SUPPORT_TEXT", "Поддержка: напиши нам, поможем быстро.").strip()
CONTACT_EMAIL = os.getenv("CONTACT_EMAIL", "support@aimaxbots.ru").strip()
CONTACT_PHONE = os.getenv("CONTACT_PHONE", "").strip()
BOT_PUBLIC_URL = os.getenv("BOT_PUBLIC_URL", "").strip()
CHANNEL_URL = os.getenv("CHANNEL_URL", "https://max.ru/id231128398751_biz").strip()
CHANNEL_GATE_ENABLED = os.getenv("CHANNEL_GATE_ENABLED", "1").strip().lower() not in {"0", "false", "no", "off"}
CHANNEL_CHAT_ID = os.getenv("CHANNEL_CHAT_ID", "").strip()
CHANNEL_MEMBERSHIP_CACHE_HOURS = int(os.getenv("CHANNEL_MEMBERSHIP_CACHE_HOURS", "12"))
REFERRAL_BONUS_CREDITS = int(os.getenv("REFERRAL_BONUS_CREDITS", "70"))
PROMO_WELCOME_CREDITS = int(os.getenv("PROMO_WELCOME_CREDITS", "0"))
PROMO_CODES_RAW = os.getenv("PROMO_CODES", "").strip()
ADMIN_PANEL_TOKEN = os.getenv("ADMIN_PANEL_TOKEN", "").strip()
ADMIN_SESSION_COOKIE = "aimax_admin_session"
ADMIN_SESSION_MAX_AGE = 60 * 60 * 24 * 14
ADMIN_LOGIN_WINDOW_SECONDS = int(os.getenv("ADMIN_LOGIN_WINDOW_SECONDS", "300"))
ADMIN_LOGIN_MAX_ATTEMPTS = int(os.getenv("ADMIN_LOGIN_MAX_ATTEMPTS", "8"))
ADMIN_LOGIN_BLOCK_SECONDS = int(os.getenv("ADMIN_LOGIN_BLOCK_SECONDS", "900"))
PAYMENT_STATUS_TOKEN_TTL_SECONDS = int(os.getenv("PAYMENT_STATUS_TOKEN_TTL_SECONDS", str(60 * 60 * 24 * 7)))
BACKUP_KEEP_FILES = int(os.getenv("BACKUP_KEEP_FILES", "12"))
AUTO_BACKUP_ENABLED = os.getenv("AUTO_BACKUP_ENABLED", "1").strip().lower() not in {"0", "false", "no"}
AUTO_BACKUP_INTERVAL_HOURS = int(os.getenv("AUTO_BACKUP_INTERVAL_HOURS", "24"))
ERROR_ALERT_COOLDOWN_SEC = int(os.getenv("ERROR_ALERT_COOLDOWN_SEC", "120"))
ERROR_ALERTS_ENABLED = os.getenv("ERROR_ALERTS_ENABLED", "1").strip().lower() not in {"0", "false", "no"}
SERVICE_MONITOR_ENABLED = os.getenv("SERVICE_MONITOR_ENABLED", "1").strip().lower() not in {"0", "false", "no"}
SERVICE_MONITOR_INTERVAL_MINUTES = int(os.getenv("SERVICE_MONITOR_INTERVAL_MINUTES", "15"))
ALERT_HIGH_ERRORS_WINDOW_MINUTES = int(os.getenv("ALERT_HIGH_ERRORS_WINDOW_MINUTES", "30"))
ALERT_HIGH_ERRORS_THRESHOLD = int(os.getenv("ALERT_HIGH_ERRORS_THRESHOLD", "8"))
ALERT_LOW_PAYMENTS_LOOKBACK_HOURS = int(os.getenv("ALERT_LOW_PAYMENTS_LOOKBACK_HOURS", "24"))
ALERT_LOW_PAYMENTS_MIN_ACTIVE_USERS = int(os.getenv("ALERT_LOW_PAYMENTS_MIN_ACTIVE_USERS", "15"))
ALERT_LOW_PAYMENTS_MAX_PAYMENTS = int(os.getenv("ALERT_LOW_PAYMENTS_MAX_PAYMENTS", "0"))
ALERT_SPEND_SPIKE_LOOKBACK_HOURS = int(os.getenv("ALERT_SPEND_SPIKE_LOOKBACK_HOURS", "1"))
ALERT_SPEND_SPIKE_MIN_RUB = float(os.getenv("ALERT_SPEND_SPIKE_MIN_RUB", "150"))
ALERT_SPEND_SPIKE_MULTIPLIER = float(os.getenv("ALERT_SPEND_SPIKE_MULTIPLIER", "3.0"))
ALERT_CPU_LOAD_PER_CORE_THRESHOLD = float(os.getenv("ALERT_CPU_LOAD_PER_CORE_THRESHOLD", "1.2"))
ALERT_MEMORY_USED_PCT_THRESHOLD = float(os.getenv("ALERT_MEMORY_USED_PCT_THRESHOLD", "92"))
ALERT_DISK_USED_PCT_THRESHOLD = float(os.getenv("ALERT_DISK_USED_PCT_THRESHOLD", "90"))
ALERT_BACKUP_STALE_HOURS = int(os.getenv("ALERT_BACKUP_STALE_HOURS", "36"))
REENGAGE_DORMANT_DAYS = int(os.getenv("REENGAGE_DORMANT_DAYS", "5"))
REENGAGE_BATCH_LIMIT = int(os.getenv("REENGAGE_BATCH_LIMIT", "30"))
SENTRY_DSN = os.getenv("SENTRY_DSN", "").strip()
SENTRY_ENVIRONMENT = os.getenv("SENTRY_ENVIRONMENT", "production").strip() or "production"
REFERENCE_IMAGE_TTL_MINUTES = int(os.getenv("REFERENCE_IMAGE_TTL_MINUTES", "180"))
FILE_KIND_LABELS = {
    "doc": "документ",
    "ppt": "презентация",
    "sheet": "таблица",
}
FILE_KIND_EXTENSIONS = {
    "doc": "docx",
    "ppt": "pptx",
    "sheet": "xlsx",
}
FILE_KIND_MIME_TYPES = {
    "doc": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "ppt": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "sheet": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}
FILE_PROFILE_LABELS = {
    "short": "Короткая",
    "medium": "Средняя",
    "full": "Полная",
}
ADMIN_IDS = {
    int(value.strip())
    for value in os.getenv("ADMIN_IDS", "").split(",")
    if value.strip().isdigit()
}
ADMIN_MAX_USER_IDS = {
    int(value.strip())
    for value in os.getenv("ADMIN_MAX_USER_IDS", "").split(",")
    if value.strip().isdigit()
}

SYSTEM_PROMPT_BASE = (
    "Ты полезный AI-ассистент в мессенджере MAX. "
    "Отвечай по-русски, если пользователь не попросил иначе. "
    "Не упоминай внутренние технические детали без необходимости. "
    "По умолчанию отвечай кратко и легко для чтения: 2-4 коротких абзаца или 3-5 коротких пунктов. "
    "Не используй длинные нумерованные списки и тяжелое оформление без необходимости. "
    "Старайся писать простым живым языком, без воды. "
    "Если тема большая, дай сжатый ответ и не обрывай фразу на полуслове."
)

STYLE_PROMPTS = {
    "gpt": "Стиль: четко структурируй ответ, хорошо объясняй шаги и варианты.",
    "gpt4o": "Стиль: отвечай живо, понятно и по делу, с хорошими примерами.",
    "gemini": "Стиль: отвечай быстро, дружелюбно и с упором на практический результат.",
    "deepseek": "Стиль: делай упор на рассуждение, анализ и техническую точность.",
    "gpt54": "Стиль: отвечай как эксперт-консультант, глубоко и обоснованно.",
}

PLAN_ORDER = {"free": 0, "lite": 1, "start": 2, "pro": 3}
PAID_PLANS = {"lite", "start", "pro"}
BUYABLE_PLANS = {"lite", "start", "pro"}
IMAGE_STYLE_OPTIONS: dict[str, tuple[str, str]] = {
    "auto": ("✨ Авто", ""),
    "photo": ("📷 Фото", "photorealistic style, natural lighting"),
    "anime": ("🌸 Аниме", "anime style, clean line art"),
    "art": ("🖼 Арт", "digital art illustration, cinematic composition"),
}
IMAGE_ASPECT_OPTIONS: dict[str, tuple[str, str]] = {
    "square": ("1:1", "square composition"),
    "portrait": ("9:16", "vertical composition"),
    "landscape": ("16:9", "horizontal composition"),
}
IMAGE_PRESET_OPTIONS: dict[str, dict[str, str]] = {
    "avatar": {
        "label": "👤 Аватар",
        "style": "photo",
        "aspect": "square",
        "hint": "аватар, портрет, лицо крупно",
        "prompt": "avatar portrait, subject centered, clean background, expressive face",
    },
    "art_portrait": {
        "label": "🎨 Арт-портрет",
        "style": "art",
        "aspect": "square",
        "hint": "арт-портрет, яркий стиль, атмосферно",
        "prompt": "stylized art portrait, cinematic mood, detailed illustration",
    },
    "product": {
        "label": "📦 Товар",
        "style": "photo",
        "aspect": "square",
        "hint": "товар, чистый фон, акцент на продукте",
        "prompt": "commercial product shot, clean background, focus on product details",
    },
    "poster": {
        "label": "🎬 Постер",
        "style": "art",
        "aspect": "portrait",
        "hint": "постер, заглавный кадр, выразительная композиция",
        "prompt": "poster composition, dramatic framing, bold focal point",
    },
}
IMAGE_EDIT_PRESET_OPTIONS: dict[str, dict[str, str]] = {
    "anime_ref": {
        "label": "🌸 Аниме по фото",
        "style": "anime",
        "aspect": "square",
        "hint": "перерисовать в стиле аниме",
        "prompt": "convert the reference photo into clean anime artwork while preserving identity",
    },
    "background_ref": {
        "label": "🫧 Сменить фон",
        "style": "photo",
        "aspect": "portrait",
        "hint": "заменить фон или перенести в новую сцену",
        "prompt": "replace the background cleanly while keeping the main subject natural",
    },
    "enhance_ref": {
        "label": "✨ Улучшить фото",
        "style": "photo",
        "aspect": "square",
        "hint": "улучшить качество, свет, детали",
        "prompt": "improve photo quality, lighting, clarity and skin tones while keeping it natural",
    },
    "art_ref": {
        "label": "🖼 Арт по фото",
        "style": "art",
        "aspect": "square",
        "hint": "сделать художественную версию фото",
        "prompt": "turn the reference photo into polished digital art with strong composition",
    },
}
ADS_MEDIA_CHANNELS: list[tuple[str, str, int, str]] = [
    ("AD01", "КиберПоток | ИИ & Нейросети", 601, "https://telega.in/channels/Kiber_potok/card"),
    ("AD02", "Технологичка | Dev/IT", 825, "https://telega.in/channels/technologichka/card"),
    ("AD03", "Code Learning", 839, "https://telega.in/channels/codelearning_tg/card"),
    ("AD04", "Техносплит – Нейросети, Технологии, Новости IT", 1259, "https://telega.in/channels/technosplit/card"),
    ("AD05", "AI Simplify", 1399, "https://telega.in/channels/simplify_ai/card"),
    ("AD06", "Machine Learning | Нейронные сети, ИИ, Big Data", 3357, "https://telega.in/channels/ML_secrets/card"),
    ("AD07", "ChatGpt | Нейросеть", 4196, "https://telega.in/channels/gpt_chat1/card"),
    ("AD08", "Техконтент | Ai | ChatGPT", 4895, "https://telega.in/channels/tech_contents/card"),
    ("AD09", "Нейро Лептик", 4895, "https://telega.in/channels/neiro_leptik/card"),
    ("AD10", "Bard AI | Нейросети & IT", 5594, "https://telega.in/channels/NeuralToday/card"),
    ("AD11", "Дневник ChatGPT", 8392, "https://telega.in/channels/chatgptmachine/card"),
    ("AD12", "StudGPT (ChatGPT) Ai | ИИ", 11189, "https://telega.in/channels/studgpt/card"),
    ("AD13", "Библиотека нейротекста | ChatGPT, Gemini, Bing", 12028, "https://telega.in/channels/neuro_text/card"),
    ("AD14", "Искусственный интеллект. Высокие технологии", 20280, "https://telega.in/channels/vistehno/card"),
    ("AD15", "ИИволюция", 27972, "https://telega.in/channels/ai_volution/card"),
    ("AD16", "Tips AI | IT & AI", 39161, "https://telega.in/channels/tips_ai/card"),
]
DEFAULT_IMAGE_STYLE = "auto"
DEFAULT_IMAGE_ASPECT = "square"

WELCOME_TEXT = (
    "Привет. Это твой AI-бот в MAX.\n\n"
    "Что умею:\n"
    "• ответы через GPT, Gemini и DeepSeek\n"
    "• генерация картинок\n"
    "• сохранение контекста диалога\n\n"
    "Выбери действие кнопками или просто напиши вопрос."
)

MENU_TEXT = (
    "Кнопки ниже помогают быстро открыть нужный режим, тарифы и картинки.\n"
    "Новости и обновления — в кнопке «📣 Канал».\n"
    "Можно просто написать вопрос в чат — бот ответит."
)

HELP_TEXT = (
    "Справка\n\n"
    "Основной способ пользоваться ботом — кнопки ниже.\n"
    "Команды пригодятся, если нужно быстрое действие:\n\n"
    "/start или /menu — меню\n"
    "/id — твой chat_id и user_id\n"
    "/models — версии и описание моделей\n"
    "/plan — твой тариф и остатки\n"
    "/preset <fast|balanced|quality|expert> — выбрать режим\n"
    "/model <alias> — выбрать модель вручную\n"
    "/gpt /gpt4o /gemini /deepseek /gpt54 — быстрый выбор модели\n"
    "/image <описание> — сгенерировать картинку\n"
    "/image_ref <описание> — сгенерировать по последнему фото\n"
    "/tariffs — тарифы\n"
    "/topup — пакеты запросов\n"
    "/buy <lite|start|pro> — заявка на подписку\n"
    "/payments — мои заявки\n"
    "/ref [код] — реферальный код и активация\n"
    "/promo <код> — активировать промокод\n"
    "/channel — наш канал\n"
    "/credits или /requests — остаток запросов\n"
    "/support — помощь по оплате и работе бота\n"
    "/clear — очистить контекст"
)

HELP_TEXT += "\n/files — меню файлов: документ, презентация, таблица"

ADMIN_HELP_TEXT = (
    "\n\nАдмин:\n"
    "/admin help\n"
    "/admin user <chat_id>\n"
    "/admin block <chat_id> <on|off>\n"
    "/admin templates\n"
    "/admin backup\n"
    "/admin nudge [days] [limit]\n"
    "Изменение тарифов/платежей — через /admin/panel\n"
    "/costs — модели и цены\n"
    "/id — твой chat_id и user_id"
)

@dataclass(slots=True)
class PlanInfo:
    name: str
    daily_gpt54_limit: int


PLAN_CONFIGS = {
    "free": PlanInfo(
        name="free",
        daily_gpt54_limit=0,
    ),
    "lite": PlanInfo(
        name="lite",
        daily_gpt54_limit=0,
    ),
    "start": PlanInfo(
        name="start",
        daily_gpt54_limit=START_DAILY_GPT54_LIMIT,
    ),
    "pro": PlanInfo(
        name="pro",
        daily_gpt54_limit=PRO_DAILY_GPT54_LIMIT,
    ),
}

PLAN_CREDITS = {
    "free": 0,
    "lite": LITE_PLAN_CREDITS,
    "start": START_PLAN_CREDITS,
    "pro": PRO_PLAN_CREDITS,
}

MODEL_CREDIT_COSTS = {
    "deepseek": CREDIT_COST_DEEPSEEK,
    "gpt": CREDIT_COST_GPT,
    "gpt4o": CREDIT_COST_GPTO,
    "gemini": CREDIT_COST_GEMINI,
    "gpt54": CREDIT_COST_GPT54,
}

MODEL_VAR_CREDITS_PER_1K = {
    "deepseek": VAR_CREDITS_PER_1K_DEEPSEEK,
    "gpt": VAR_CREDITS_PER_1K_GPT,
    "gpt4o": VAR_CREDITS_PER_1K_GPTO,
    "gemini": VAR_CREDITS_PER_1K_GEMINI,
    "gpt54": VAR_CREDITS_PER_1K_GPT54,
}

MODEL_PRESETS: dict[str, dict[str, Any]] = {
    "fast": {
        "label": "⚡ Быстро",
        "description": "короткие и быстрые ответы",
        "aliases": ["deepseek", "gpt"],
    },
    "balanced": {
        "label": "⚖ Баланс",
        "description": "ежедневные задачи и диалог",
        "aliases": ["gpt4o", "gpt", "deepseek"],
    },
    "quality": {
        "label": "🧠 Качество",
        "description": "подробно и аккуратно",
        "aliases": ["gemini", "gpt4o", "gpt"],
    },
    "expert": {
        "label": "🚀 Эксперт",
        "description": "максимум качества для сложных задач",
        "aliases": ["gpt54", "gemini", "gpt4o", "gpt"],
    },
}

TOPUP_PACKS = {
    "small": {
        "label": "Small",
        "price_rub": TOPUP_SMALL_PRICE_RUB,
        "credits": TOPUP_SMALL_CREDITS,
    },
    "medium": {
        "label": "Medium",
        "price_rub": TOPUP_MEDIUM_PRICE_RUB,
        "credits": TOPUP_MEDIUM_CREDITS,
    },
    "large": {
        "label": "Large",
        "price_rub": TOPUP_LARGE_PRICE_RUB,
        "credits": TOPUP_LARGE_CREDITS,
    },
}


@dataclass(slots=True)
class ImageResult:
    image_bytes: bytes
    mime_type: str


@dataclass(slots=True)
class TextAnswerResult:
    text: str
    prompt_tokens: int
    completion_tokens: int
    total_tokens: int


@dataclass(slots=True)
class ImageJob:
    kind: str
    chat_id: int
    user_prompt: str
    model_prompt: str
    credits_spent: int
    details: str
    reference_image_data_url: str = ""


@dataclass(slots=True)
class ModelInfo:
    alias: str
    provider: str
    model: str
    label: str
    kind: str
    version: str
    input_price_usd_per_m: str
    output_price_usd_per_m: str
    min_plan: str
    description: str
    prompt_style: str


def safe_float(value: Any, default: float = 0.0) -> float:
    try:
        return float(str(value).replace(",", "."))
    except (TypeError, ValueError):
        return default


def parse_usage_details_blob(details: str) -> tuple[int, int]:
    value = str(details or "")
    prompt_match = re.search(r"(?:^|;)prompt=(\d+)", value)
    completion_match = re.search(r"(?:^|;)completion=(\d+)", value)
    prompt_tokens = int(prompt_match.group(1)) if prompt_match else 0
    completion_tokens = int(completion_match.group(1)) if completion_match else 0
    return max(0, prompt_tokens), max(0, completion_tokens)


def parse_detail_field(details: str, key: str) -> str:
    if not details or not key:
        return ""
    match = re.search(rf"(?:^|;){re.escape(key)}=([^;]+)", str(details))
    return match.group(1).strip() if match else ""


def model_info_for_alias(alias: str) -> tuple[str, str, ModelInfo | None]:
    raw_alias = str(alias or "").strip()
    if not raw_alias:
        return "-", "—", None
    base_alias, _, suffix = raw_alias.partition(":")
    text_model = TEXT_MODELS.get(base_alias)
    if text_model:
        return text_model.label, "Текст", text_model
    image_model = IMAGE_MODELS.get(base_alias)
    if image_model:
        label = image_model.label
        if suffix == "edit":
            label = f"{label} (ред. фото)"
        return label, "Картинка", image_model
    return raw_alias, "—", None


def estimate_text_cost_usd(model: ModelInfo | None, prompt_tokens: int, completion_tokens: int, total_tokens: int) -> float:
    if not model or model.kind != "text":
        return 0.0
    input_price = safe_float(model.input_price_usd_per_m)
    output_price = safe_float(model.output_price_usd_per_m)
    prompt = max(0, int(prompt_tokens or 0))
    completion = max(0, int(completion_tokens or 0))
    total = max(0, int(total_tokens or 0))
    if prompt <= 0 and completion <= 0 and total > 0:
        prompt = int(total * 0.65)
        completion = max(0, total - prompt)
    return (prompt / 1_000_000.0) * input_price + (completion / 1_000_000.0) * output_price


def expected_unit_economics(price_rub: int, credits: int) -> dict[str, float]:
    revenue = float(max(0, int(price_rub or 0)))
    expected_cost = float(max(0, int(credits or 0))) * ANALYTICS_EXPECTED_COST_PER_CREDIT_RUB
    payment_fee = revenue * (ANALYTICS_PAYMENT_FEE_PCT / 100.0)
    receipt_fee = revenue * (ANALYTICS_RECEIPT_FEE_PCT / 100.0)
    tax = revenue * (ANALYTICS_TAX_PCT / 100.0)
    margin = revenue - expected_cost - payment_fee - receipt_fee - tax
    margin_pct = (margin * 100.0 / revenue) if revenue > 0 else 0.0
    return {
        "revenue_rub": revenue,
        "expected_cost_rub": expected_cost,
        "payment_fee_rub": payment_fee,
        "receipt_fee_rub": receipt_fee,
        "tax_rub": tax,
        "margin_rub": margin,
        "margin_pct": margin_pct,
    }


def configure_logging() -> logging.Logger:
    logger = logging.getLogger("max_ai_agent")
    logger.setLevel(getattr(logging, LOG_LEVEL, logging.INFO))
    logger.handlers.clear()

    formatter = logging.Formatter("%(asctime)s [%(levelname)s] %(message)s")
    console_handler = logging.StreamHandler()
    console_handler.setFormatter(formatter)
    logger.addHandler(console_handler)

    file_handler = RotatingFileHandler(LOG_FILE, maxBytes=2_000_000, backupCount=5, encoding="utf-8")
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)

    logger.propagate = False
    return logger


log = configure_logging()


def to_base36(value: int) -> str:
    alphabet = "0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    if value <= 0:
        return "0"
    out: list[str] = []
    n = value
    while n:
        n, rem = divmod(n, 36)
        out.append(alphabet[rem])
    return "".join(reversed(out))


def referral_code_for_chat(chat_id: int) -> str:
    return f"RF{to_base36(chat_id).zfill(6)}"


def normalize_referral_code(value: str) -> str:
    return re.sub(r"[^A-Z0-9]", "", value.strip().upper())


def is_referral_code(value: str) -> bool:
    code = normalize_referral_code(value)
    return code.startswith("RF") and len(code) >= 8


def parse_start_payload(update: dict[str, Any]) -> str:
    candidates: list[Any] = [
        update.get("payload"),
        update.get("start_payload"),
        update.get("start_parameter"),
        ((update.get("message") or {}).get("body") or {}).get("payload") if isinstance(update.get("message"), dict) else None,
        ((update.get("message") or {}).get("body") or {}).get("start_payload") if isinstance(update.get("message"), dict) else None,
        ((update.get("message") or {}).get("body") or {}).get("start_parameter") if isinstance(update.get("message"), dict) else None,
        ((update.get("chat") or {}).get("payload")) if isinstance(update.get("chat"), dict) else None,
        ((update.get("chat") or {}).get("start_payload")) if isinstance(update.get("chat"), dict) else None,
        ((update.get("chat") or {}).get("start_parameter")) if isinstance(update.get("chat"), dict) else None,
    ]
    for candidate in candidates:
        if isinstance(candidate, str) and candidate.strip():
            return candidate.strip()
    return ""


def referral_code_from_start_payload(update: dict[str, Any]) -> str:
    payload = parse_start_payload(update)
    if not payload:
        return ""
    raw = payload.strip()
    if "=" in raw or "&" in raw:
        params = dict(parse_qsl(raw, keep_blank_values=True))
        for key in ("ref", "referral", "code", "r"):
            value = params.get(key)
            if isinstance(value, str) and is_referral_code(value):
                return normalize_referral_code(value)
    for prefix in ("ref:", "ref=", "referral:", "referral=", "r:", "r="):
        if raw.lower().startswith(prefix):
            raw = raw[len(prefix):]
            break
    return normalize_referral_code(raw) if is_referral_code(raw) else ""


def acquisition_meta_from_start_payload(update: dict[str, Any]) -> tuple[str, str]:
    payload = parse_start_payload(update)
    source = ""
    campaign = ""
    if payload:
        raw = payload.strip()
        if "=" in raw or "&" in raw:
            params = dict(parse_qsl(raw, keep_blank_values=True))
            source = str(params.get("src") or params.get("source") or params.get("utm_source") or "").strip().lower()
            campaign = str(params.get("campaign") or params.get("utm_campaign") or "").strip().lower()
        else:
            for chunk in raw.split(";"):
                part = chunk.strip()
                if not part or ":" not in part:
                    continue
                key, value = part.split(":", 1)
                key = key.strip().lower()
                value = value.strip().lower()
                if key in {"src", "source", "utm_source"} and not source:
                    source = value
                elif key in {"campaign", "utm_campaign"} and not campaign:
                    campaign = value
    code = referral_code_from_start_payload(update)
    if code and not source:
        source = "referral"
    source = re.sub(r"[^a-z0-9_-]", "", source)[:32]
    campaign = re.sub(r"[^a-z0-9_-]", "", campaign)[:48]
    return source, campaign


def max_share_url(text: str) -> str:
    return f"https://max.ru/:share?text={quote(str(text or ''), safe='')}"


def bot_public_url_value() -> str:
    return (BOT_PUBLIC_URL or "").strip()


def campaign_start_payload(campaign: str, source: str = "ads") -> str:
    campaign_value = re.sub(r"[^a-z0-9_-]", "", str(campaign or "").strip().lower())[:48]
    source_value = re.sub(r"[^a-z0-9_-]", "", str(source or "").strip().lower())[:32]
    params: list[tuple[str, str]] = []
    if source_value:
        params.append(("source", source_value))
    if campaign_value:
        params.append(("campaign", campaign_value))
    return urlencode(params)


def campaign_deep_link(campaign: str, source: str = "ads") -> str:
    base_url = bot_public_url_value().rstrip("/")
    if not base_url:
        return ""
    payload = campaign_start_payload(campaign, source=source)
    if not payload:
        return base_url
    separator = "&" if "?" in base_url else "?"
    return f"{base_url}{separator}start={payload}"


def referral_share_message_v2(code: str) -> str:
    normalized = normalize_referral_code(code)
    return (
        "Попробуй моего AI-бота в MAX.\n"
        f"Канал с обновлениями: {channel_url_value()}\n\n"
        "Как активировать бонус:\n"
        "1. Перейди в бота\n"
        "2. Открой Меню → Бонусы\n"
        "3. Нажми «Ввести реф-код»\n"
        "4. Вставь код друга:\n"
        f"`{normalized}`\n\n"
        f"После активации тебе и мне начислят по +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов."
    )


def channel_promo_meta(today: date | None = None) -> dict[str, Any]:
    # Legacy channel promo is intentionally disabled.
    return {
        "enabled": False,
        "active": False,
        "code": "",
        "credits": 0,
        "start": (today or datetime.utcnow().date()),
        "end_exclusive": (today or datetime.utcnow().date()),
        "days_left": 0,
        "bonus_ttl_days": 0,
    }


def promo_offer_for_code(code: str) -> tuple[int, int, str]:
    promo_code = normalize_referral_code(code)
    if not promo_code:
        return 0, 0, "Пустой промокод."

    if promo_code == channel["code"] and channel["enabled"]:
        if channel["active"]:
            return int(channel["credits"]), int(channel["bonus_ttl_days"]), ""
        if datetime.utcnow().date() < channel["start"]:
            return 0, 0, f"Акция еще не началась (старт: {channel['start'].isoformat()})."
        return 0, 0, "Акция по этому промокоду завершена."

    credits = int(promo_catalog().get(promo_code, 0) or 0)
    if credits > 0:
        return credits, 0, ""
    return 0, 0, "Такого промокода нет или он выключен."


def promo_catalog() -> dict[str, int]:
    catalog: dict[str, int] = {}
    if PROMO_WELCOME_CREDITS > 0:
        catalog["WELCOME"] = int(PROMO_WELCOME_CREDITS)
    raw = PROMO_CODES_RAW
    if not raw:
        return {k: v for k, v in catalog.items() if v > 0}
    for chunk in raw.split(","):
        part = chunk.strip()
        if not part or ":" not in part:
            continue
        key_raw, val_raw = part.split(":", 1)
        key = normalize_referral_code(key_raw)
        if not key:
            continue
        try:
            value = int(val_raw.strip())
        except Exception:
            continue
        if value > 0:
            catalog[key] = value
    return {k: v for k, v in catalog.items() if v > 0}


class UserStore:
    def __init__(self, backend: StorageBackend) -> None:
        self.backend = backend
        self._init()

    def _connect(self) -> Any:
        return self.backend.connect()

    def _init(self) -> None:
        with self._connect() as conn:
            users_pk = "chat_id BIGINT PRIMARY KEY" if self.backend.kind == "postgres" else "chat_id INTEGER PRIMARY KEY"
            payment_pk = "id BIGSERIAL PRIMARY KEY" if self.backend.kind == "postgres" else "id INTEGER PRIMARY KEY AUTOINCREMENT"
            usage_pk = "id BIGSERIAL PRIMARY KEY" if self.backend.kind == "postgres" else "id INTEGER PRIMARY KEY AUTOINCREMENT"
            promo_pk = "id BIGSERIAL PRIMARY KEY" if self.backend.kind == "postgres" else "id INTEGER PRIMARY KEY AUTOINCREMENT"
            grant_pk = "id BIGSERIAL PRIMARY KEY" if self.backend.kind == "postgres" else "id INTEGER PRIMARY KEY AUTOINCREMENT"
            conn.execute(
                f"""
                CREATE TABLE IF NOT EXISTS users (
                    {users_pk},
                    max_user_id INTEGER NOT NULL DEFAULT 0,
                    acquisition_source TEXT NOT NULL DEFAULT '',
                    acquisition_campaign TEXT NOT NULL DEFAULT '',
                    acquired_at TEXT NOT NULL DEFAULT '',
                    channel_subscribed_at TEXT NOT NULL DEFAULT '',
                    channel_subscription_checked_at TEXT NOT NULL DEFAULT '',
                    plan TEXT NOT NULL DEFAULT 'free',
                    is_blocked INTEGER NOT NULL DEFAULT 0,
                    onboarding_done INTEGER NOT NULL DEFAULT 0,
                    referral_code TEXT NOT NULL DEFAULT '',
                    referred_by_chat_id INTEGER NOT NULL DEFAULT 0,
                    referrals_invited INTEGER NOT NULL DEFAULT 0,
                    receipt_email TEXT NOT NULL DEFAULT '',
                    receipt_phone TEXT NOT NULL DEFAULT '',
                    selected_model_alias TEXT NOT NULL DEFAULT '',
                    selected_preset TEXT NOT NULL DEFAULT '',
                    subscription_expires_at TEXT NOT NULL DEFAULT '',
                    recurring_enabled INTEGER NOT NULL DEFAULT 0,
                    recurring_cancel_from TEXT NOT NULL DEFAULT '',
                    recurring_canceled_at TEXT NOT NULL DEFAULT '',
                    usage_date TEXT NOT NULL DEFAULT '',
                    daily_messages_used INTEGER NOT NULL DEFAULT 0,
                    daily_images_used INTEGER NOT NULL DEFAULT 0,
                    daily_gpt54_used INTEGER NOT NULL DEFAULT 0,
                    free_image_week_key TEXT NOT NULL DEFAULT '',
                    free_image_week_used INTEGER NOT NULL DEFAULT 0,
                    free_image_last_used_at TEXT NOT NULL DEFAULT '',
                    free_file_last_used_at TEXT NOT NULL DEFAULT '',
                    credits_balance INTEGER NOT NULL DEFAULT 0,
                    credits_spent_total INTEGER NOT NULL DEFAULT 0,
                    last_active_at TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL DEFAULT '',
                    updated_at TEXT NOT NULL DEFAULT ''
                )
                """
            )
            conn.execute(
                f"""
                CREATE TABLE IF NOT EXISTS payment_requests (
                    {payment_pk},
                    chat_id INTEGER NOT NULL,
                    plan TEXT NOT NULL,
                    days INTEGER NOT NULL,
                    amount_rub INTEGER NOT NULL,
                    recurring_consent INTEGER NOT NULL DEFAULT 0,
                    recurring_consent_at TEXT NOT NULL DEFAULT '',
                    recurring_consent_text TEXT NOT NULL DEFAULT '',
                    receipt_email TEXT NOT NULL DEFAULT '',
                    receipt_phone TEXT NOT NULL DEFAULT '',
                    status TEXT NOT NULL DEFAULT 'pending',
                    provider TEXT NOT NULL DEFAULT 'manual',
                    provider_ref TEXT NOT NULL DEFAULT '',
                    payment_url TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL DEFAULT '',
                    paid_at TEXT NOT NULL DEFAULT '',
                    activated_at TEXT NOT NULL DEFAULT ''
                )
                """
            )
            conn.execute(
                f"""
                CREATE TABLE IF NOT EXISTS usage_events (
                    {usage_pk},
                    chat_id INTEGER NOT NULL,
                    event_type TEXT NOT NULL,
                    plan TEXT NOT NULL DEFAULT '',
                    model_alias TEXT NOT NULL DEFAULT '',
                    credits_spent INTEGER NOT NULL DEFAULT 0,
                    rub_amount INTEGER NOT NULL DEFAULT 0,
                    tokens_total INTEGER NOT NULL DEFAULT 0,
                    details TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL DEFAULT ''
                )
                """
            )
            conn.execute(
                f"""
                CREATE TABLE IF NOT EXISTS promo_activations (
                    {promo_pk},
                    chat_id INTEGER NOT NULL,
                    promo_code TEXT NOT NULL,
                    credits INTEGER NOT NULL DEFAULT 0,
                    created_at TEXT NOT NULL DEFAULT '',
                    UNIQUE(chat_id, promo_code)
                )
                """
            )
            conn.execute(
                f"""
                CREATE TABLE IF NOT EXISTS promo_bonus_grants (
                    {grant_pk},
                    chat_id INTEGER NOT NULL,
                    promo_code TEXT NOT NULL,
                    amount_total INTEGER NOT NULL DEFAULT 0,
                    amount_remaining INTEGER NOT NULL DEFAULT 0,
                    expires_at TEXT NOT NULL DEFAULT '',
                    created_at TEXT NOT NULL DEFAULT '',
                    expired_at TEXT NOT NULL DEFAULT ''
                )
                """
            )
            conn.execute(
                """
                CREATE TABLE IF NOT EXISTS processed_updates (
                    fingerprint TEXT PRIMARY KEY,
                    created_at TEXT NOT NULL DEFAULT ''
                )
                """
            )
            self._ensure_column(conn, "users", "subscription_expires_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "max_user_id", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "acquisition_source", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "acquisition_campaign", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "acquired_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "channel_subscribed_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "channel_subscription_checked_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "recurring_enabled", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "recurring_cancel_from", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "recurring_canceled_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "receipt_email", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "receipt_phone", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "selected_preset", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "onboarding_done", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "referral_code", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "referred_by_chat_id", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "referrals_invited", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "last_active_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "daily_gpt54_used", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "free_image_week_key", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "free_image_week_used", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "free_image_last_used_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "free_file_last_used_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "users", "credits_balance", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "users", "credits_spent_total", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "payment_requests", "recurring_consent", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(conn, "payment_requests", "recurring_consent_at", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "payment_requests", "recurring_consent_text", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "payment_requests", "receipt_email", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "payment_requests", "receipt_phone", "TEXT NOT NULL DEFAULT ''")
            self._ensure_column(conn, "payment_requests", "payment_url", "TEXT NOT NULL DEFAULT ''")
            conn.execute(
                "CREATE UNIQUE INDEX IF NOT EXISTS idx_promo_activations_unique ON promo_activations(chat_id, promo_code)"
            )
            conn.execute(
                "CREATE UNIQUE INDEX IF NOT EXISTS idx_users_max_user_id_unique ON users(max_user_id) WHERE max_user_id > 0"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_users_plan_updated_at ON users(plan, updated_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_users_acquisition_source_campaign ON users(acquisition_source, acquisition_campaign)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_promo_bonus_grants_expiry ON promo_bonus_grants(chat_id, expires_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_processed_updates_created_at ON processed_updates(created_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_payment_requests_chat_created_at ON payment_requests(chat_id, created_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_payment_requests_status_created_at ON payment_requests(status, created_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_payment_requests_provider_ref ON payment_requests(provider_ref)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_usage_events_chat_created_at ON usage_events(chat_id, created_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_usage_events_type_created_at ON usage_events(event_type, created_at)"
            )
            conn.execute(
                "CREATE INDEX IF NOT EXISTS idx_usage_events_model_created_at ON usage_events(model_alias, created_at)"
            )
            conn.commit()

    def _ensure_column(
        self,
        conn: Any,
        table: str,
        column: str,
        sqlite_spec: str,
        postgres_spec: str | None = None,
    ) -> None:
        if self.backend.kind == "postgres":
            exists_row = conn.execute(
                """
                SELECT 1
                FROM information_schema.columns
                WHERE table_schema = current_schema()
                  AND table_name = ?
                  AND column_name = ?
                LIMIT 1
                """,
                (table, column),
            ).fetchone()
            if exists_row:
                return
            spec = postgres_spec or sqlite_spec
            conn.execute(f'ALTER TABLE "{table}" ADD COLUMN "{column}" {spec}')
            return

        info = conn.execute(f"PRAGMA table_info({table})").fetchall()
        names = {row["name"] for row in info}
        if column not in names:
            conn.execute(f"ALTER TABLE {table} ADD COLUMN {column} {sqlite_spec}")

    def _today(self) -> str:
        return date.today().isoformat()

    def _week_key(self) -> str:
        today = date.today()
        year, week, _ = today.isocalendar()
        return f"{year}-W{week:02d}"

    def _expire_bonus_grants_if_needed(self, conn: sqlite3.Connection, chat_id: int, row: sqlite3.Row, now: datetime) -> sqlite3.Row:
        now_iso = now.isoformat()
        grants = conn.execute(
            """
            SELECT id, amount_remaining
            FROM promo_bonus_grants
            WHERE chat_id = ? AND amount_remaining > 0 AND expires_at <> '' AND expires_at <= ?
            ORDER BY expires_at ASC, id ASC
            """,
            (chat_id, now_iso),
        ).fetchall()
        if not grants:
            return row

        balance = int(row["credits_balance"] or 0)
        for grant in grants:
            remaining = int(grant["amount_remaining"] or 0)
            if remaining <= 0:
                continue
            deduction = min(balance, remaining)
            balance -= deduction
            conn.execute(
                "UPDATE promo_bonus_grants SET amount_remaining = 0, expired_at = ? WHERE id = ?",
                (now_iso, int(grant["id"])),
            )

        conn.execute(
            "UPDATE users SET credits_balance = ?, updated_at = ? WHERE chat_id = ?",
            (max(0, balance), now_iso, chat_id),
        )
        conn.commit()
        refreshed = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
        return refreshed if refreshed is not None else row

    def _insert_new_user(
        self,
        conn: sqlite3.Connection,
        chat_id: int,
        default_model_alias: str,
        now: str,
        today: str,
        max_user_id: int = 0,
    ) -> sqlite3.Row:
        referral_code = referral_code_for_chat(chat_id)
        conn.execute(
            """
            INSERT INTO users (
                chat_id, max_user_id, plan, is_blocked, onboarding_done, referral_code, referred_by_chat_id, referrals_invited,
                acquisition_source, acquisition_campaign, acquired_at,
                selected_model_alias, selected_preset, usage_date,
                daily_messages_used, daily_images_used, daily_gpt54_used,
                free_image_week_key, free_image_week_used, free_image_last_used_at,
                free_file_last_used_at,
                credits_balance, credits_spent_total,
                last_active_at,
                created_at, updated_at
            ) VALUES (?, ?, 'free', 0, 0, ?, 0, 0, '', '', '', ?, '', ?, 0, 0, 0, '', 0, '', '', ?, 0, ?, ?, ?)
            """,
            (
                chat_id,
                max(0, int(max_user_id or 0)),
                referral_code,
                default_model_alias,
                today,
                FREE_DAILY_CREDITS,
                now,
                now,
                now,
            ),
        )
        return conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()

    def _update_chat_references(self, conn: sqlite3.Connection, old_chat_id: int, new_chat_id: int, now: str) -> None:
        conn.execute(
            "UPDATE users SET chat_id = ?, updated_at = ? WHERE chat_id = ?",
            (new_chat_id, now, old_chat_id),
        )
        conn.execute(
            "UPDATE users SET referred_by_chat_id = ? WHERE referred_by_chat_id = ?",
            (new_chat_id, old_chat_id),
        )
        for table in ("payment_requests", "usage_events", "promo_activations", "promo_bonus_grants"):
            conn.execute(f"UPDATE {table} SET chat_id = ? WHERE chat_id = ?", (new_chat_id, old_chat_id))

    def _merge_rebound_user_rows(
        self,
        existing: sqlite3.Row,
        current: sqlite3.Row | None,
        max_user_id: int,
        now: str,
    ) -> dict[str, Any]:
        merged = dict(existing)
        merged["max_user_id"] = max(0, int(max_user_id or 0))
        if current is None:
            merged["updated_at"] = now
            return merged

        current_dict = dict(current)
        existing_plan = str(merged.get("plan", "free") or "free")
        current_plan = str(current_dict.get("plan", "free") or "free")
        existing_rank = PLAN_ORDER.get(existing_plan, 0)
        current_rank = PLAN_ORDER.get(current_plan, 0)

        def later_iso(left: Any, right: Any) -> str:
            left_s = str(left or "").strip()
            right_s = str(right or "").strip()
            if not left_s:
                return right_s
            if not right_s:
                return left_s
            return max(left_s, right_s)

        merged["is_blocked"] = 1 if int(merged.get("is_blocked", 0) or 0) or int(current_dict.get("is_blocked", 0) or 0) else 0
        merged["onboarding_done"] = 1 if int(merged.get("onboarding_done", 0) or 0) or int(current_dict.get("onboarding_done", 0) or 0) else 0
        merged["referral_code"] = str(merged.get("referral_code", "") or "").strip() or str(current_dict.get("referral_code", "") or "").strip() or referral_code_for_chat(int(existing["chat_id"]))
        merged["referred_by_chat_id"] = int(merged.get("referred_by_chat_id", 0) or 0) or int(current_dict.get("referred_by_chat_id", 0) or 0)
        merged["referrals_invited"] = max(int(merged.get("referrals_invited", 0) or 0), int(current_dict.get("referrals_invited", 0) or 0))
        merged["acquisition_source"] = str(merged.get("acquisition_source", "") or "").strip() or str(current_dict.get("acquisition_source", "") or "").strip()
        merged["acquisition_campaign"] = str(merged.get("acquisition_campaign", "") or "").strip() or str(current_dict.get("acquisition_campaign", "") or "").strip()
        merged["acquired_at"] = later_iso(merged.get("acquired_at", ""), current_dict.get("acquired_at", ""))
        merged["channel_subscribed_at"] = later_iso(merged.get("channel_subscribed_at", ""), current_dict.get("channel_subscribed_at", ""))
        merged["channel_subscription_checked_at"] = later_iso(merged.get("channel_subscription_checked_at", ""), current_dict.get("channel_subscription_checked_at", ""))
        merged["receipt_email"] = str(merged.get("receipt_email", "") or "").strip() or str(current_dict.get("receipt_email", "") or "").strip()
        merged["receipt_phone"] = str(merged.get("receipt_phone", "") or "").strip() or str(current_dict.get("receipt_phone", "") or "").strip()
        merged["last_active_at"] = later_iso(merged.get("last_active_at", ""), current_dict.get("last_active_at", ""))
        merged["created_at"] = min(str(merged.get("created_at", "") or "").strip() or now, str(current_dict.get("created_at", "") or "").strip() or now)

        if current_rank > existing_rank:
            for key in (
                "plan",
                "selected_model_alias",
                "selected_preset",
                "subscription_expires_at",
                "recurring_enabled",
                "recurring_cancel_from",
                "recurring_canceled_at",
            ):
                merged[key] = current_dict.get(key, merged.get(key))
            merged["credits_balance"] = int(current_dict.get("credits_balance", 0) or 0)
        elif current_rank == existing_rank and existing_plan != "free":
            merged["credits_balance"] = max(int(merged.get("credits_balance", 0) or 0), int(current_dict.get("credits_balance", 0) or 0))
            merged["subscription_expires_at"] = later_iso(merged.get("subscription_expires_at", ""), current_dict.get("subscription_expires_at", ""))
            merged["recurring_enabled"] = 1 if int(merged.get("recurring_enabled", 0) or 0) or int(current_dict.get("recurring_enabled", 0) or 0) else 0
            merged["recurring_cancel_from"] = later_iso(merged.get("recurring_cancel_from", ""), current_dict.get("recurring_cancel_from", ""))
            merged["recurring_canceled_at"] = later_iso(merged.get("recurring_canceled_at", ""), current_dict.get("recurring_canceled_at", ""))
        else:
            merged["credits_balance"] = min(int(merged.get("credits_balance", 0) or 0), int(current_dict.get("credits_balance", 0) or 0))

        if not str(merged.get("selected_model_alias", "") or "").strip():
            merged["selected_model_alias"] = str(current_dict.get("selected_model_alias", "") or "").strip()
        if not str(merged.get("selected_preset", "") or "").strip():
            merged["selected_preset"] = str(current_dict.get("selected_preset", "") or "").strip()

        merged["credits_spent_total"] = max(int(merged.get("credits_spent_total", 0) or 0), int(current_dict.get("credits_spent_total", 0) or 0))
        merged["free_image_last_used_at"] = later_iso(merged.get("free_image_last_used_at", ""), current_dict.get("free_image_last_used_at", ""))
        merged["free_image_week_key"] = later_iso(merged.get("free_image_week_key", ""), current_dict.get("free_image_week_key", ""))
        merged["free_image_week_used"] = max(int(merged.get("free_image_week_used", 0) or 0), int(current_dict.get("free_image_week_used", 0) or 0))
        merged["free_file_last_used_at"] = later_iso(merged.get("free_file_last_used_at", ""), current_dict.get("free_file_last_used_at", ""))

        existing_usage_date = str(merged.get("usage_date", "") or "")
        current_usage_date = str(current_dict.get("usage_date", "") or "")
        if current_usage_date > existing_usage_date:
            merged["usage_date"] = current_usage_date
            for key in ("daily_messages_used", "daily_images_used", "daily_gpt54_used"):
                merged[key] = int(current_dict.get(key, 0) or 0)
        elif current_usage_date == existing_usage_date and current_usage_date:
            for key in ("daily_messages_used", "daily_images_used", "daily_gpt54_used"):
                merged[key] = max(int(merged.get(key, 0) or 0), int(current_dict.get(key, 0) or 0))

        merged["updated_at"] = now
        return merged

    def ensure_user_binding(self, chat_id: int, max_user_id: int | None, default_model_alias: str) -> int | None:
        identity = max(0, int(max_user_id or 0))
        now = datetime.utcnow().isoformat()
        today = self._today()
        with self._connect() as conn:
            current = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
            if identity <= 0:
                if current is None:
                    self._insert_new_user(conn, chat_id, default_model_alias, now, today, 0)
                    conn.commit()
                return None

            existing = conn.execute("SELECT * FROM users WHERE max_user_id = ? LIMIT 1", (identity,)).fetchone()
            if current is None and existing is None:
                self._insert_new_user(conn, chat_id, default_model_alias, now, today, identity)
                conn.commit()
                return None

            if existing is None:
                conn.execute(
                    "UPDATE users SET max_user_id = ?, updated_at = ? WHERE chat_id = ?",
                    (identity, now, chat_id),
                )
                conn.commit()
                return None

            existing_chat_id = int(existing["chat_id"])
            if existing_chat_id == chat_id:
                if int(existing["max_user_id"] or 0) != identity:
                    conn.execute(
                        "UPDATE users SET max_user_id = ?, updated_at = ? WHERE chat_id = ?",
                        (identity, now, chat_id),
                    )
                    conn.commit()
                return None

            merged = self._merge_rebound_user_rows(existing, current, identity, now)
            if current is not None:
                conn.execute("DELETE FROM users WHERE chat_id = ?", (chat_id,))
            self._update_chat_references(conn, existing_chat_id, chat_id, now)
            conn.execute(
                """
                UPDATE users
                SET max_user_id = ?, acquisition_source = ?, acquisition_campaign = ?, acquired_at = ?,
                    channel_subscribed_at = ?, channel_subscription_checked_at = ?,
                    plan = ?, is_blocked = ?, onboarding_done = ?, referral_code = ?,
                    referred_by_chat_id = ?, referrals_invited = ?, receipt_email = ?, receipt_phone = ?,
                    selected_model_alias = ?, selected_preset = ?, subscription_expires_at = ?,
                    recurring_enabled = ?, recurring_cancel_from = ?, recurring_canceled_at = ?,
                    usage_date = ?, daily_messages_used = ?, daily_images_used = ?, daily_gpt54_used = ?,
                    free_image_week_key = ?, free_image_week_used = ?, free_image_last_used_at = ?, free_file_last_used_at = ?,
                    credits_balance = ?, credits_spent_total = ?, last_active_at = ?, created_at = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (
                    int(merged.get("max_user_id", 0) or 0),
                    str(merged.get("acquisition_source", "") or ""),
                    str(merged.get("acquisition_campaign", "") or ""),
                    str(merged.get("acquired_at", "") or ""),
                    str(merged.get("channel_subscribed_at", "") or ""),
                    str(merged.get("channel_subscription_checked_at", "") or ""),
                    str(merged.get("plan", "free") or "free"),
                    int(merged.get("is_blocked", 0) or 0),
                    int(merged.get("onboarding_done", 0) or 0),
                    str(merged.get("referral_code", "") or ""),
                    int(merged.get("referred_by_chat_id", 0) or 0),
                    int(merged.get("referrals_invited", 0) or 0),
                    str(merged.get("receipt_email", "") or ""),
                    str(merged.get("receipt_phone", "") or ""),
                    str(merged.get("selected_model_alias", "") or ""),
                    str(merged.get("selected_preset", "") or ""),
                    str(merged.get("subscription_expires_at", "") or ""),
                    int(merged.get("recurring_enabled", 0) or 0),
                    str(merged.get("recurring_cancel_from", "") or ""),
                    str(merged.get("recurring_canceled_at", "") or ""),
                    str(merged.get("usage_date", "") or ""),
                    int(merged.get("daily_messages_used", 0) or 0),
                    int(merged.get("daily_images_used", 0) or 0),
                    int(merged.get("daily_gpt54_used", 0) or 0),
                    str(merged.get("free_image_week_key", "") or ""),
                    int(merged.get("free_image_week_used", 0) or 0),
                    str(merged.get("free_image_last_used_at", "") or ""),
                    str(merged.get("free_file_last_used_at", "") or ""),
                    int(merged.get("credits_balance", 0) or 0),
                    int(merged.get("credits_spent_total", 0) or 0),
                    str(merged.get("last_active_at", "") or ""),
                    str(merged.get("created_at", "") or now),
                    now,
                    chat_id,
                ),
            )
            conn.commit()
            return existing_chat_id

    def get_or_create_user(self, chat_id: int, default_model_alias: str) -> dict[str, Any]:
        now = datetime.utcnow().isoformat()
        today = self._today()
        with self._connect() as conn:
            row = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
            if row is None:
                row = self._insert_new_user(conn, chat_id, default_model_alias, now, today, 0)
                conn.commit()

            row = self._expire_bonus_grants_if_needed(conn, chat_id, row, datetime.utcnow())

            day_changed = row["usage_date"] != today
            if day_changed:
                conn.execute(
                    """
                    UPDATE users
                    SET usage_date = ?, daily_messages_used = 0, daily_images_used = 0, daily_gpt54_used = 0, updated_at = ?
                    WHERE chat_id = ?
                    """,
                    (today, now, chat_id),
                )
                conn.commit()
                row = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()

            # Daily free bonus should be topped up only once per day.
            if day_changed and row["plan"] == "free" and int(row["credits_balance"] or 0) < FREE_DAILY_CREDITS:
                conn.execute(
                    """
                    UPDATE users
                    SET credits_balance = ?, updated_at = ?
                    WHERE chat_id = ?
                    """,
                    (FREE_DAILY_CREDITS, now, chat_id),
                )
                conn.commit()
                row = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()

            if not str(row["referral_code"] or "").strip():
                conn.execute(
                    "UPDATE users SET referral_code = ?, updated_at = ? WHERE chat_id = ?",
                    (referral_code_for_chat(chat_id), now, chat_id),
                )
                conn.commit()
                row = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()

            return dict(row)

    def set_selected_model(self, chat_id: int, alias: str) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET selected_model_alias = ?, updated_at = ? WHERE chat_id = ?",
                (alias, datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def set_selected_preset(self, chat_id: int, preset: str) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET selected_preset = ?, updated_at = ? WHERE chat_id = ?",
                (preset.strip(), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def set_receipt_contact(self, chat_id: int, email: str, phone: str) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET receipt_email = ?, receipt_phone = ?, updated_at = ? WHERE chat_id = ?",
                (email.strip(), phone.strip(), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def touch_last_active(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET last_active_at = ?, updated_at = ? WHERE chat_id = ?",
                (datetime.utcnow().isoformat(), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def set_onboarding_done(self, chat_id: int, done: bool = True) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET onboarding_done = ?, updated_at = ? WHERE chat_id = ?",
                (1 if done else 0, datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def mark_channel_subscription(self, chat_id: int, subscribed: bool) -> None:
        now = datetime.utcnow().replace(microsecond=0).isoformat()
        subscribed_at = now if subscribed else ""
        with self._connect() as conn:
            conn.execute(
                """
                UPDATE users
                SET channel_subscribed_at = ?, channel_subscription_checked_at = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (subscribed_at, now, now, chat_id),
            )
            conn.commit()

    def apply_referral_code(self, chat_id: int, referral_code: str, bonus_credits: int) -> tuple[bool, str]:
        code = normalize_referral_code(referral_code)
        if not code:
            return False, "Пустой реферальный код."
        with self._connect() as conn:
            conn.execute("BEGIN IMMEDIATE")
            user = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
            if not user:
                conn.rollback()
                return False, "Пользователь не найден."
            if int(user["referred_by_chat_id"] or 0) > 0:
                conn.rollback()
                return False, "Реферальный код уже был использован."
            owner = conn.execute("SELECT * FROM users WHERE referral_code = ? LIMIT 1", (code,)).fetchone()
            if not owner:
                conn.rollback()
                return False, "Такого реферального кода нет."
            owner_chat_id = int(owner["chat_id"])
            if owner_chat_id == int(chat_id):
                conn.rollback()
                return False, "Нельзя активировать свой код."

            now = datetime.utcnow().isoformat()
            conn.execute(
                """
                UPDATE users
                SET referred_by_chat_id = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (owner_chat_id, now, chat_id),
            )
            conn.execute(
                """
                UPDATE users
                SET referrals_invited = referrals_invited + 1, updated_at = ?
                WHERE chat_id = ?
                """,
                (now, owner_chat_id),
            )
            if bonus_credits > 0:
                conn.execute(
                    "UPDATE users SET credits_balance = credits_balance + ?, updated_at = ? WHERE chat_id = ?",
                    (bonus_credits, now, chat_id),
                )
                conn.execute(
                    "UPDATE users SET credits_balance = credits_balance + ?, updated_at = ? WHERE chat_id = ?",
                    (bonus_credits, now, owner_chat_id),
                )
            conn.commit()
            return True, str(owner_chat_id)

    def redeem_promo_code(self, chat_id: int, promo_code: str, credits: int, bonus_ttl_days: int = 0) -> tuple[bool, str]:
        code = normalize_referral_code(promo_code)
        if not code:
            return False, "Пустой промокод."
        if credits <= 0:
            return False, "Промокод сейчас недоступен."
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            try:
                conn.execute(
                    """
                    INSERT INTO promo_activations (chat_id, promo_code, credits, created_at)
                    VALUES (?, ?, ?, ?)
                    """,
                    (chat_id, code, credits, now),
                )
            except sqlite3.IntegrityError:
                return False, "Этот промокод уже активирован."
            conn.execute(
                "UPDATE users SET credits_balance = credits_balance + ?, updated_at = ? WHERE chat_id = ?",
                (credits, now, chat_id),
            )
            ttl_days = max(0, int(bonus_ttl_days))
            if ttl_days > 0:
                expires_at = (datetime.utcnow() + timedelta(days=ttl_days)).replace(microsecond=0).isoformat()
                conn.execute(
                    """
                    INSERT INTO promo_bonus_grants (chat_id, promo_code, amount_total, amount_remaining, expires_at, created_at, expired_at)
                    VALUES (?, ?, ?, ?, ?, ?, '')
                    """,
                    (chat_id, code, credits, credits, expires_at, now),
                )
            conn.commit()
        return True, str(credits)

    def active_bonus_credits_summary(self, chat_id: int) -> tuple[int, str]:
        now_iso = datetime.utcnow().isoformat()
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT amount_remaining, expires_at
                FROM promo_bonus_grants
                WHERE chat_id = ? AND amount_remaining > 0 AND expires_at > ?
                ORDER BY expires_at ASC
                """,
                (chat_id, now_iso),
            ).fetchall()
        if not rows:
            return 0, ""
        total = sum(int(row["amount_remaining"] or 0) for row in rows)
        nearest = str(rows[0]["expires_at"] or "")
        return total, nearest

    def list_recent_users(self, limit: int = 50) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT chat_id, max_user_id, plan, is_blocked, onboarding_done, credits_balance,
                       recurring_enabled, subscription_expires_at,
                       daily_messages_used, daily_images_used, last_active_at, updated_at
                FROM users
                ORDER BY updated_at DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
            return [dict(row) for row in rows]

    def search_users(self, query: str, limit: int = 50) -> list[dict[str, Any]]:
        needle = str(query or "").strip()
        if not needle:
            return self.list_recent_users(limit=limit)
        like_value = f"%{needle.lower()}%"
        with self._connect() as conn:
            if needle.isdigit():
                rows = conn.execute(
                    """
                    SELECT chat_id, max_user_id, plan, credits_balance,
                           recurring_enabled, subscription_expires_at,
                           daily_messages_used, daily_images_used, last_active_at, updated_at
                    FROM users
                    WHERE chat_id = ? OR max_user_id = ?
                    ORDER BY updated_at DESC
                    LIMIT ?
                    """,
                    (int(needle), int(needle), limit),
                ).fetchall()
                if rows:
                    return [dict(row) for row in rows]
            rows = conn.execute(
                """
                SELECT chat_id, max_user_id, plan, credits_balance,
                       recurring_enabled, subscription_expires_at,
                       daily_messages_used, daily_images_used, last_active_at, updated_at
                FROM users
                WHERE lower(plan) LIKE ?
                   OR lower(referral_code) LIKE ?
                ORDER BY updated_at DESC
                LIMIT ?
                """,
                (like_value, like_value, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def suspicious_users_report(self, limit: int = 25) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT
                    u.chat_id,
                    u.max_user_id,
                    u.plan,
                    u.credits_balance,
                    u.subscription_expires_at,
                    u.recurring_enabled,
                    u.last_active_at,
                    CASE
                        WHEN u.plan != 'free'
                             AND NOT EXISTS (
                                SELECT 1
                                FROM payment_requests p
                                WHERE p.chat_id = u.chat_id
                                  AND p.status = 'paid'
                                  AND p.plan = u.plan
                             )
                        THEN 'Платный план без оплаченной заявки'
                        WHEN u.plan != 'free'
                             AND EXISTS (
                                SELECT 1
                                FROM usage_events e
                                WHERE e.chat_id = u.chat_id
                                  AND e.event_type = 'payment'
                                  AND e.details LIKE '%source=admin panel%'
                             )
                        THEN 'Оплата подтверждена вручную через админку'
                        ELSE ''
                    END AS risk_reason
                FROM users u
                WHERE u.plan != 'free'
                  AND (
                    NOT EXISTS (
                        SELECT 1
                        FROM payment_requests p
                        WHERE p.chat_id = u.chat_id
                          AND p.status = 'paid'
                          AND p.plan = u.plan
                    )
                    OR EXISTS (
                        SELECT 1
                        FROM usage_events e
                        WHERE e.chat_id = u.chat_id
                          AND e.event_type = 'payment'
                          AND e.details LIKE '%source=admin panel%'
                    )
                  )
                ORDER BY u.updated_at DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
            return [dict(row) for row in rows]

    def top_referrers_report(self, limit: int = 10) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT
                    u.chat_id,
                    u.max_user_id,
                    u.referral_code,
                    u.referrals_invited,
                    SUM(CASE WHEN r.plan != 'free' THEN 1 ELSE 0 END) AS paid_referrals,
                    MAX(COALESCE(r.created_at, '')) AS last_referral_at
                FROM users u
                LEFT JOIN users r ON r.referred_by_chat_id = u.chat_id
                WHERE u.referrals_invited > 0
                GROUP BY u.chat_id, u.max_user_id, u.referral_code, u.referrals_invited
                ORDER BY u.referrals_invited DESC, paid_referrals DESC, last_referral_at DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
            return [dict(row) for row in rows]

    def suspicious_referral_report(self, limit: int = 10) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT
                    u.chat_id,
                    u.max_user_id,
                    u.referral_code,
                    u.referrals_invited,
                    SUM(CASE WHEN r.plan != 'free' THEN 1 ELSE 0 END) AS paid_referrals,
                    SUM(CASE WHEN r.last_active_at >= ? THEN 1 ELSE 0 END) AS active_recent_referrals
                FROM users u
                LEFT JOIN users r ON r.referred_by_chat_id = u.chat_id
                WHERE u.referrals_invited >= 3
                GROUP BY u.chat_id, u.max_user_id, u.referral_code, u.referrals_invited
                HAVING paid_referrals = 0
                ORDER BY u.referrals_invited DESC, active_recent_referrals ASC, u.chat_id DESC
                LIMIT ?
                """,
                ((datetime.utcnow() - timedelta(days=14)).replace(microsecond=0).isoformat(), limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def get_user(self, chat_id: int) -> dict[str, Any] | None:
        with self._connect() as conn:
            row = conn.execute("SELECT * FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
            return dict(row) if row else None

    def get_user_by_max_user_id(self, max_user_id: int) -> dict[str, Any] | None:
        identity = max(0, int(max_user_id or 0))
        if identity <= 0:
            return None
        with self._connect() as conn:
            row = conn.execute("SELECT * FROM users WHERE max_user_id = ? LIMIT 1", (identity,)).fetchone()
            return dict(row) if row else None

    def list_recent_payments(self, limit: int = 50) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT id, chat_id, plan, amount_rub, status, provider_ref, created_at, paid_at, activated_at
                FROM payment_requests
                ORDER BY id DESC
                LIMIT ?
                """,
                (limit,),
            ).fetchall()
            return [dict(row) for row in rows]

    def search_payments(self, query: str = "", status: str = "", limit: int = 50) -> list[dict[str, Any]]:
        needle = str(query or "").strip()
        status_filter = str(status or "").strip().lower()
        conditions: list[str] = []
        params: list[Any] = []
        if needle:
            if needle.isdigit():
                conditions.append("(id = ? OR chat_id = ?)")
                params.extend([int(needle), int(needle)])
            else:
                like_value = f"%{needle.lower()}%"
                conditions.append("lower(plan) LIKE ?")
                params.append(like_value)
        if status_filter:
            conditions.append("lower(status) = ?")
            params.append(status_filter)
        where_sql = f"WHERE {' AND '.join(conditions)}" if conditions else ""
        with self._connect() as conn:
            rows = conn.execute(
                f"""
                SELECT id, chat_id, plan, amount_rub, status, provider_ref, created_at, paid_at, activated_at
                FROM payment_requests
                {where_sql}
                ORDER BY id DESC
                LIMIT ?
                """,
                (*params, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def adjust_credits(self, chat_id: int, delta: int) -> int:
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET credits_balance = credits_balance + ?, updated_at = ? WHERE chat_id = ?",
                (delta, now, chat_id),
            )
            conn.commit()
            row = conn.execute("SELECT credits_balance FROM users WHERE chat_id = ?", (chat_id,)).fetchone()
            return int((row["credits_balance"] if row else 0) or 0)

    def reset_daily_counters(self, chat_id: int) -> None:
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            conn.execute(
                """
                UPDATE users
                SET daily_messages_used = 0, daily_images_used = 0, daily_gpt54_used = 0, updated_at = ?
                WHERE chat_id = ?
                """,
                (now, chat_id),
            )
            conn.commit()

    def list_reengage_candidates(self, dormant_days: int, limit: int) -> list[dict[str, Any]]:
        if dormant_days <= 0:
            dormant_days = 1
        if limit <= 0:
            limit = 1
        cutoff = (datetime.utcnow() - timedelta(days=dormant_days)).replace(microsecond=0).isoformat()
        free_image_ready_cutoff = (datetime.utcnow() - timedelta(days=7)).replace(microsecond=0).isoformat()
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT chat_id, plan, free_image_last_used_at, last_active_at, credits_balance
                FROM users
                WHERE plan = 'free'
                  AND (last_active_at = '' OR last_active_at <= ?)
                  AND (free_image_last_used_at = '' OR free_image_last_used_at <= ?)
                ORDER BY last_active_at ASC
                LIMIT ?
                """,
                (cutoff, free_image_ready_cutoff, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def set_plan(self, chat_id: int, plan: str) -> None:
        expires = "" if plan == "free" else None
        with self._connect() as conn:
            if expires is None:
                conn.execute(
                    "UPDATE users SET plan = ?, updated_at = ? WHERE chat_id = ?",
                    (plan, datetime.utcnow().isoformat(), chat_id),
                )
            else:
                conn.execute(
                    """
                    UPDATE users
                    SET plan = ?, subscription_expires_at = ?, recurring_enabled = 0,
                        recurring_cancel_from = '', recurring_canceled_at = '',
                        credits_balance = CASE WHEN credits_balance < ? THEN ? ELSE credits_balance END,
                        updated_at = ?
                    WHERE chat_id = ?
                    """,
                    (plan, expires, FREE_DAILY_CREDITS, FREE_DAILY_CREDITS, datetime.utcnow().isoformat(), chat_id),
                )
            conn.commit()

    def set_subscription(
        self,
        chat_id: int,
        plan: str,
        days: int,
        selected_model_alias: str,
        recurring_enabled: bool | None = None,
    ) -> str:
        expires_at = (datetime.utcnow() + timedelta(days=days)).replace(microsecond=0).isoformat()
        with self._connect() as conn:
            if recurring_enabled is None:
                conn.execute(
                    """
                    UPDATE users
                    SET plan = ?, subscription_expires_at = ?, selected_model_alias = ?,
                        credits_balance = ?, updated_at = ?
                    WHERE chat_id = ?
                    """,
                    (
                        plan,
                        expires_at,
                        selected_model_alias,
                        PLAN_CREDITS.get(plan, 0),
                        datetime.utcnow().isoformat(),
                        chat_id,
                    ),
                )
            else:
                conn.execute(
                    """
                    UPDATE users
                    SET plan = ?, subscription_expires_at = ?, selected_model_alias = ?,
                        recurring_enabled = ?, recurring_cancel_from = '', recurring_canceled_at = '',
                        credits_balance = ?, updated_at = ?
                    WHERE chat_id = ?
                    """,
                    (
                        plan,
                        expires_at,
                        selected_model_alias,
                        1 if recurring_enabled else 0,
                        PLAN_CREDITS.get(plan, 0),
                        datetime.utcnow().isoformat(),
                        chat_id,
                    ),
                )
            conn.commit()
        return expires_at

    def cancel_recurring(self, chat_id: int, cancel_from: str) -> None:
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            conn.execute(
                """
                UPDATE users
                SET recurring_enabled = 0, recurring_cancel_from = ?, recurring_canceled_at = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (cancel_from, now, now, chat_id),
            )
            conn.commit()

    def create_payment_request(
        self,
        chat_id: int,
        plan: str,
        days: int,
        amount_rub: int,
        recurring_consent: bool = False,
        recurring_consent_text: str = "",
        receipt_email: str = "",
        receipt_phone: str = "",
    ) -> int:
        now = datetime.utcnow().isoformat()
        consent_at = now if recurring_consent else ""
        with self._connect() as conn:
            cur = conn.execute(
                """
                INSERT INTO payment_requests (
                    chat_id, plan, days, amount_rub, recurring_consent,
                    recurring_consent_at, recurring_consent_text, receipt_email, receipt_phone, status, created_at
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, 'pending', ?)
                """,
                (
                    chat_id,
                    plan,
                    days,
                    amount_rub,
                    1 if recurring_consent else 0,
                    consent_at,
                    recurring_consent_text,
                    receipt_email.strip(),
                    receipt_phone.strip(),
                    now,
                ),
            )
            conn.commit()
            return int(cur.lastrowid)

    def set_payment_provider_ref(self, request_id: int, provider_ref: str) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE payment_requests SET provider_ref = ? WHERE id = ?",
                (provider_ref, request_id),
            )
            conn.commit()

    def set_payment_url(self, request_id: int, payment_url: str) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE payment_requests SET payment_url = ? WHERE id = ?",
                (payment_url.strip(), request_id),
            )
            conn.commit()

    def list_user_payments(self, chat_id: int, limit: int = 5) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT id, plan, days, amount_rub, status, created_at
                FROM payment_requests
                WHERE chat_id = ?
                ORDER BY id DESC
                LIMIT ?
                """,
                (chat_id, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def list_user_usage_events(self, chat_id: int, limit: int = 12) -> list[dict[str, Any]]:
        with self._connect() as conn:
            rows = conn.execute(
                """
                SELECT event_type, plan, model_alias, credits_spent, rub_amount, tokens_total, details, created_at
                FROM usage_events
                WHERE chat_id = ?
                ORDER BY id DESC
                LIMIT ?
                """,
                (chat_id, limit),
            ).fetchall()
            return [dict(row) for row in rows]

    def get_payment(self, request_id: int) -> dict[str, Any] | None:
        with self._connect() as conn:
            row = conn.execute("SELECT * FROM payment_requests WHERE id = ?", (request_id,)).fetchone()
            return dict(row) if row else None

    def set_payment_status(self, request_id: int, status: str) -> None:
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            if status == "paid":
                conn.execute(
                    "UPDATE payment_requests SET status = ?, paid_at = ? WHERE id = ?",
                    (status, now, request_id),
                )
            else:
                conn.execute(
                    "UPDATE payment_requests SET status = ? WHERE id = ?",
                    (status, request_id),
                )
            conn.commit()

    def mark_payment_activated(self, request_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE payment_requests SET activated_at = ? WHERE id = ?",
                (datetime.utcnow().isoformat(), request_id),
            )
            conn.commit()

    def set_blocked(self, chat_id: int, blocked: bool) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET is_blocked = ?, updated_at = ? WHERE chat_id = ?",
                (1 if blocked else 0, datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def increment_message_usage(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET daily_messages_used = daily_messages_used + 1, updated_at = ? WHERE chat_id = ?",
                (datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def increment_image_usage(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET daily_images_used = daily_images_used + 1, updated_at = ? WHERE chat_id = ?",
                (datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def increment_free_week_image_usage(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                """
                UPDATE users
                SET free_image_week_used = 1, free_image_last_used_at = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (datetime.utcnow().isoformat(), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def increment_free_file_usage(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET free_file_last_used_at = ?, updated_at = ? WHERE chat_id = ?",
                (datetime.utcnow().isoformat(), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def increment_gpt54_usage(self, chat_id: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET daily_gpt54_used = daily_gpt54_used + 1, updated_at = ? WHERE chat_id = ?",
                (datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def consume_credits(self, chat_id: int, amount: int) -> bool:
        if amount <= 0:
            return True
        with self._connect() as conn:
            cur = conn.execute(
                """
                UPDATE users
                SET credits_balance = credits_balance - ?, credits_spent_total = credits_spent_total + ?, updated_at = ?
                WHERE chat_id = ? AND credits_balance >= ?
                """,
                (amount, amount, datetime.utcnow().isoformat(), chat_id, amount),
            )
            conn.commit()
            return cur.rowcount > 0

    def refund_credits(self, chat_id: int, amount: int) -> None:
        if amount <= 0:
            return
        with self._connect() as conn:
            conn.execute(
                """
                UPDATE users
                SET credits_balance = credits_balance + ?,
                    credits_spent_total = CASE WHEN credits_spent_total >= ? THEN credits_spent_total - ? ELSE 0 END,
                    updated_at = ?
                WHERE chat_id = ?
                """,
                (amount, amount, amount, datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def set_credits(self, chat_id: int, amount: int) -> None:
        with self._connect() as conn:
            conn.execute(
                "UPDATE users SET credits_balance = ?, updated_at = ? WHERE chat_id = ?",
                (max(0, int(amount)), datetime.utcnow().isoformat(), chat_id),
            )
            conn.commit()

    def set_acquisition_meta(self, chat_id: int, source: str = "", campaign: str = "") -> None:
        source_value = re.sub(r"[^a-z0-9_-]", "", str(source or "").strip().lower())[:32]
        campaign_value = re.sub(r"[^a-z0-9_-]", "", str(campaign or "").strip().lower())[:48]
        if not source_value and not campaign_value:
            return
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            row = conn.execute(
                "SELECT acquisition_source, acquisition_campaign, acquired_at FROM users WHERE chat_id = ?",
                (chat_id,),
            ).fetchone()
            if not row:
                return
            current_source = str(row["acquisition_source"] or "").strip().lower()
            current_campaign = str(row["acquisition_campaign"] or "").strip().lower()
            acquired_at = str(row["acquired_at"] or "").strip()
            next_source = current_source
            next_campaign = current_campaign
            next_acquired_at = acquired_at
            if source_value and (not current_source or current_source == "direct"):
                next_source = source_value
                next_acquired_at = acquired_at or now
            if campaign_value and not current_campaign:
                next_campaign = campaign_value
                next_acquired_at = acquired_at or now
            if not next_source and source_value:
                next_source = source_value
                next_acquired_at = acquired_at or now
            if next_source == current_source and next_campaign == current_campaign and next_acquired_at == acquired_at:
                return
            conn.execute(
                """
                UPDATE users
                SET acquisition_source = ?, acquisition_campaign = ?, acquired_at = ?, updated_at = ?
                WHERE chat_id = ?
                """,
                (next_source, next_campaign, next_acquired_at, now, chat_id),
            )
            conn.commit()

    def record_usage_event(
        self,
        chat_id: int,
        event_type: str,
        plan: str = "",
        model_alias: str = "",
        credits_spent: int = 0,
        rub_amount: int = 0,
        tokens_total: int = 0,
        details: str = "",
    ) -> None:
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            conn.execute(
                """
                INSERT INTO usage_events (
                    chat_id, event_type, plan, model_alias, credits_spent, rub_amount, tokens_total, details, created_at
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    int(chat_id),
                    str(event_type),
                    str(plan),
                    str(model_alias),
                    int(credits_spent),
                    int(rub_amount),
                    int(tokens_total),
                    str(details)[:500],
                    now,
                ),
            )
            conn.commit()

    def remember_processed_update(self, fingerprint: str) -> bool:
        if not fingerprint:
            return False
        now = datetime.utcnow().isoformat()
        with self._connect() as conn:
            row = conn.execute(
                "SELECT 1 FROM processed_updates WHERE fingerprint = ?",
                (fingerprint,),
            ).fetchone()
            if row:
                return False
            conn.execute(
                "INSERT INTO processed_updates (fingerprint, created_at) VALUES (?, ?)",
                (fingerprint, now),
            )
            cutoff = (datetime.utcnow() - timedelta(hours=max(1, PROCESSED_UPDATE_TTL_HOURS))).isoformat()
            conn.execute(
                "DELETE FROM processed_updates WHERE created_at < ?",
                (cutoff,),
            )
            conn.commit()
        return True

    def kpi_report(self, days: int = 30) -> dict[str, Any]:
        period_days = max(1, min(int(days), 365))
        now_dt = datetime.utcnow()
        since_dt = now_dt - timedelta(days=period_days)
        since = since_dt.isoformat()
        with self._connect() as conn:
            summary = conn.execute(
                """
                SELECT
                    COUNT(*) AS events_total,
                    COUNT(DISTINCT chat_id) AS active_users,
                    SUM(CASE WHEN event_type = 'payment' THEN rub_amount ELSE 0 END) AS revenue_rub,
                    SUM(CASE WHEN event_type = 'refund' THEN rub_amount ELSE 0 END) AS refunds_rub,
                    SUM(CASE WHEN event_type = 'payment' AND rub_amount > 0 THEN 1 ELSE 0 END) AS payments_count,
                    SUM(CASE WHEN event_type = 'refund' THEN 1 ELSE 0 END) AS refunds_count,
                    SUM(CASE WHEN event_type = 'text_request' THEN 1 ELSE 0 END) AS text_requests,
                    SUM(CASE WHEN event_type = 'image_request' THEN 1 ELSE 0 END) AS image_requests,
                    SUM(CASE WHEN event_type = 'referral_activation' THEN 1 ELSE 0 END) AS referral_activations,
                    SUM(CASE WHEN event_type IN ('referral_activation', 'referral_reward') THEN credits_spent ELSE 0 END) AS referral_bonus_credits,
                    SUM(CASE WHEN event_type = 'text_request' THEN credits_spent ELSE 0 END) AS text_credits,
                    SUM(CASE WHEN event_type = 'image_request' THEN credits_spent ELSE 0 END) AS image_credits,
                    SUM(CASE WHEN event_type IN ('text_request','image_request') THEN credits_spent ELSE 0 END) AS total_credits_spent,
                    SUM(CASE WHEN event_type = 'text_request' THEN tokens_total ELSE 0 END) AS text_tokens
                FROM usage_events
                WHERE created_at >= ?
                """,
                (since,),
            ).fetchone()

            payers_row = conn.execute(
                """
                SELECT COUNT(DISTINCT chat_id) AS payers
                FROM usage_events
                WHERE created_at >= ? AND event_type = 'payment' AND rub_amount > 0
                """,
                (since,),
            ).fetchone()

            referred_payers_row = conn.execute(
                """
                SELECT COUNT(DISTINCT e.chat_id) AS referred_payers
                FROM usage_events e
                JOIN users u ON u.chat_id = e.chat_id
                WHERE e.created_at >= ?
                  AND e.event_type = 'payment'
                  AND e.rub_amount > 0
                  AND u.referred_by_chat_id > 0
                """,
                (since,),
            ).fetchone()

            request_rows = conn.execute(
                """
                SELECT event_type, plan, model_alias, credits_spent, tokens_total, details
                FROM usage_events
                WHERE created_at >= ? AND event_type IN ('text_request', 'image_request')
                """,
                (since,),
            ).fetchall()

            plan_rows = conn.execute(
                """
                SELECT plan, COUNT(*) AS payments, SUM(rub_amount) AS revenue
                FROM usage_events
                WHERE created_at >= ? AND event_type = 'payment' AND rub_amount > 0
                GROUP BY plan
                ORDER BY revenue DESC, payments DESC
                """
                ,
                (since,),
            ).fetchall()

            daily_rows = conn.execute(
                """
                SELECT
                    substr(created_at, 1, 10) AS day,
                    SUM(CASE WHEN event_type = 'payment' THEN rub_amount ELSE 0 END) AS revenue,
                    COUNT(DISTINCT chat_id) AS active_users
                FROM usage_events
                WHERE created_at >= ?
                GROUP BY substr(created_at, 1, 10)
                ORDER BY day DESC
                LIMIT 14
                """,
                (since,),
            ).fetchall()
            top_referrers = self.top_referrers_report(limit=10)
            suspicious_referrals = self.suspicious_referral_report(limit=10)
            source_rows = conn.execute(
                """
                WITH payer_revenue AS (
                    SELECT chat_id, SUM(rub_amount) AS revenue_rub
                    FROM usage_events
                    WHERE created_at >= ? AND event_type = 'payment' AND rub_amount > 0
                    GROUP BY chat_id
                )
                SELECT
                    COALESCE(NULLIF(lower(u.acquisition_source), ''), 'direct') AS source,
                    COALESCE(NULLIF(lower(u.acquisition_campaign), ''), '-') AS campaign,
                    COUNT(*) AS users_count,
                    SUM(CASE WHEN u.plan != 'free' THEN 1 ELSE 0 END) AS paid_users,
                    SUM(COALESCE(pr.revenue_rub, 0)) AS revenue_rub
                FROM users u
                LEFT JOIN payer_revenue pr ON pr.chat_id = u.chat_id
                WHERE COALESCE(NULLIF(u.acquired_at, ''), u.created_at) >= ?
                GROUP BY source, campaign
                ORDER BY revenue_rub DESC, paid_users DESC, users_count DESC
                LIMIT 12
                """,
                (since, since),
            ).fetchall()
            promo_code_rows = conn.execute(
                """
                WITH promo_base AS (
                    SELECT chat_id, promo_code, credits, created_at
                    FROM promo_activations
                    WHERE created_at >= ?
                ),
                promo_payments AS (
                    SELECT
                        pb.promo_code,
                        pb.chat_id,
                        SUM(CASE WHEN e.event_type = 'payment' AND e.rub_amount > 0 THEN e.rub_amount ELSE 0 END) AS revenue_rub
                    FROM promo_base pb
                    LEFT JOIN usage_events e
                      ON e.chat_id = pb.chat_id
                     AND e.created_at >= pb.created_at
                    GROUP BY pb.promo_code, pb.chat_id
                )
                SELECT
                    pb.promo_code,
                    COUNT(DISTINCT pb.chat_id) AS activations,
                    SUM(pb.credits) AS credits,
                    SUM(CASE WHEN COALESCE(pp.revenue_rub, 0) > 0 THEN 1 ELSE 0 END) AS paid_users,
                    SUM(COALESCE(pp.revenue_rub, 0)) AS revenue_rub
                FROM promo_base pb
                LEFT JOIN promo_payments pp
                  ON pp.promo_code = pb.promo_code
                 AND pp.chat_id = pb.chat_id
                GROUP BY pb.promo_code
                ORDER BY revenue_rub DESC, paid_users DESC, activations DESC
                LIMIT 20
                """,
                (since,),
            ).fetchall()
            cohort_rows = conn.execute(
                """
                SELECT chat_id, COALESCE(NULLIF(acquired_at, ''), created_at) AS cohort_at
                FROM users
                WHERE COALESCE(NULLIF(acquired_at, ''), created_at) >= ?
                """,
                (since,),
            ).fetchall()
            behavior_rows = conn.execute(
                """
                SELECT chat_id, event_type, model_alias, details, created_at
                FROM usage_events
                WHERE created_at >= ?
                  AND event_type IN ('payment', 'text_request', 'image_request', 'screen_view', 'preset_select', 'model_select')
                ORDER BY id DESC
                """,
                (since,),
            ).fetchall()

            model_stats: dict[str, dict[str, Any]] = {}
            plan_cost_stats: dict[str, dict[str, Any]] = {}
            for raw_row in request_rows:
                alias = str(raw_row["model_alias"] or "").strip() or "-"
                plan_name = str(raw_row["plan"] or "").strip() or "free"
                label, kind_label, model_info = model_info_for_alias(alias)
                stat = model_stats.setdefault(
                    alias,
                    {
                        "model_alias": alias,
                        "label": label,
                        "kind": kind_label,
                        "requests": 0,
                        "credits": 0,
                        "tokens": 0,
                        "prompt_tokens": 0,
                        "completion_tokens": 0,
                        "estimated_cost_usd": 0.0,
                    },
                )
                stat["requests"] += 1
                credits = int(raw_row["credits_spent"] or 0)
                total_tokens = int(raw_row["tokens_total"] or 0)
                prompt_tokens, completion_tokens = parse_usage_details_blob(str(raw_row["details"] or ""))
                stat["credits"] += credits
                stat["tokens"] += total_tokens
                stat["prompt_tokens"] += prompt_tokens
                stat["completion_tokens"] += completion_tokens
                estimated_cost_usd = estimate_text_cost_usd(
                    model_info,
                    prompt_tokens,
                    completion_tokens,
                    total_tokens,
                )
                stat["estimated_cost_usd"] += estimated_cost_usd
                plan_stat = plan_cost_stats.setdefault(
                    plan_name,
                    {
                        "plan": plan_name,
                        "text_requests": 0,
                        "image_requests": 0,
                        "estimated_text_cost_usd": 0.0,
                    },
                )
                if kind_label == "Текст":
                    plan_stat["text_requests"] += 1
                    plan_stat["estimated_text_cost_usd"] += estimated_cost_usd
                elif kind_label == "Картинка":
                    plan_stat["image_requests"] += 1

            estimated_text_cost_usd = sum(float(row.get("estimated_cost_usd", 0.0) or 0.0) for row in model_stats.values())
            model_rows = sorted(
                model_stats.values(),
                key=lambda row: (
                    -float(row.get("estimated_cost_usd", 0.0) or 0.0),
                    -int(row.get("tokens", 0) or 0),
                    -int(row.get("requests", 0) or 0),
                ),
            )[:12]
            payment_revenue_by_plan = {
                str(row["plan"] or "").strip() or "-": int(row["revenue"] or 0)
                for row in plan_rows
            }
            margin_rows: list[dict[str, Any]] = []
            for plan_name in ("lite", "start", "pro", "small", "medium", "large"):
                revenue = int(payment_revenue_by_plan.get(plan_name, 0) or 0)
                plan_stat = plan_cost_stats.get(plan_name, {})
                text_cost_rub = float(plan_stat.get("estimated_text_cost_usd", 0.0) or 0.0) * ANALYTICS_USD_TO_RUB
                contribution_rub = revenue - text_cost_rub
                margin_pct = (contribution_rub * 100.0 / revenue) if revenue > 0 else 0.0
                margin_rows.append(
                    {
                        "plan": plan_name,
                        "revenue_rub": revenue,
                        "estimated_text_cost_rub": text_cost_rub,
                        "contribution_rub": contribution_rub,
                        "margin_pct": margin_pct,
                        "text_requests": int(plan_stat.get("text_requests", 0) or 0),
                        "image_requests": int(plan_stat.get("image_requests", 0) or 0),
                    }
                )

            new_users = 0
            new_paid_users = 0
            d1_eligible = 0
            d1_retained = 0
            d7_eligible = 0
            d7_retained = 0
            activity_dates_by_chat: dict[int, set[str]] = {}
            payment_datetimes_by_chat: dict[int, list[datetime]] = {}
            screen_counts: dict[str, int] = {}
            preset_counts: dict[str, int] = {}

            for row in behavior_rows:
                cid = int(row["chat_id"] or 0)
                event_type = str(row["event_type"] or "").strip().lower()
                created_at = parse_iso_datetime(str(row["created_at"] or ""))
                details = str(row["details"] or "")
                if created_at and event_type in {"text_request", "image_request", "screen_view"}:
                    activity_dates_by_chat.setdefault(cid, set()).add(created_at.date().isoformat())
                if created_at and event_type == "payment":
                    payment_datetimes_by_chat.setdefault(cid, []).append(created_at)
                if event_type == "screen_view":
                    screen = parse_detail_field(details, "screen") or "unknown"
                    screen_counts[screen] = screen_counts.get(screen, 0) + 1
                if event_type == "preset_select":
                    preset = parse_detail_field(details, "preset") or "unknown"
                    preset_counts[preset] = preset_counts.get(preset, 0) + 1

            for row in cohort_rows:
                cid = int(row["chat_id"] or 0)
                cohort_at = parse_iso_datetime(str(row["cohort_at"] or ""))
                if not cohort_at:
                    continue
                new_users += 1
                payment_dates = payment_datetimes_by_chat.get(cid, [])
                if any(payment_dt >= cohort_at for payment_dt in payment_dates):
                    new_paid_users += 1
                cohort_date = cohort_at.date()
                activity_dates = activity_dates_by_chat.get(cid, set())
                if cohort_at <= now_dt - timedelta(days=1):
                    d1_eligible += 1
                    if (cohort_date + timedelta(days=1)).isoformat() in activity_dates:
                        d1_retained += 1
                if cohort_at <= now_dt - timedelta(days=7):
                    d7_eligible += 1
                    if (cohort_date + timedelta(days=7)).isoformat() in activity_dates:
                        d7_retained += 1

            top_screens = [
                {"screen": screen, "views": count}
                for screen, count in sorted(screen_counts.items(), key=lambda item: (-item[1], item[0]))[:8]
            ]
            top_presets = [
                {"preset": preset, "uses": count}
                for preset, count in sorted(preset_counts.items(), key=lambda item: (-item[1], item[0]))[:8]
            ]

            return {
                "days": period_days,
                "since": since,
                "events_total": int((summary["events_total"] or 0) if summary else 0),
                "active_users": int((summary["active_users"] or 0) if summary else 0),
                "revenue_rub": int((summary["revenue_rub"] or 0) if summary else 0),
                "refunds_rub": int((summary["refunds_rub"] or 0) if summary else 0),
                "payments_count": int((summary["payments_count"] or 0) if summary else 0),
                "refunds_count": int((summary["refunds_count"] or 0) if summary else 0),
                "text_requests": int((summary["text_requests"] or 0) if summary else 0),
                "image_requests": int((summary["image_requests"] or 0) if summary else 0),
                "referral_activations": int((summary["referral_activations"] or 0) if summary else 0),
                "referral_bonus_credits": int((summary["referral_bonus_credits"] or 0) if summary else 0),
                "text_credits": int((summary["text_credits"] or 0) if summary else 0),
                "image_credits": int((summary["image_credits"] or 0) if summary else 0),
                "total_credits_spent": int((summary["total_credits_spent"] or 0) if summary else 0),
                "text_tokens": int((summary["text_tokens"] or 0) if summary else 0),
                "new_users": new_users,
                "new_paid_users": new_paid_users,
                "d1_eligible": d1_eligible,
                "d1_retained": d1_retained,
                "d7_eligible": d7_eligible,
                "d7_retained": d7_retained,
                "payers": int((payers_row["payers"] or 0) if payers_row else 0),
                "referred_payers": int((referred_payers_row["referred_payers"] or 0) if referred_payers_row else 0),
                "estimated_text_cost_usd": estimated_text_cost_usd,
                "estimated_text_cost_rub": estimated_text_cost_usd * ANALYTICS_USD_TO_RUB,
                "models": model_rows,
                "plans": [dict(row) for row in plan_rows],
                "margins": margin_rows,
                "daily": [dict(row) for row in daily_rows],
                "top_referrers": top_referrers,
                "suspicious_referrals": suspicious_referrals,
                "sources": [dict(row) for row in source_rows],
                "promo_codes": [dict(row) for row in promo_code_rows],
                "top_screens": top_screens,
                "top_presets": top_presets,
            }

    def export_logical_backup(self) -> dict[str, Any]:
        tables = (
            "users",
            "payment_requests",
            "usage_events",
            "promo_activations",
            "promo_bonus_grants",
            "processed_updates",
        )
        snapshot: dict[str, Any] = {
            "backend": self.backend.kind,
            "generated_at": datetime.utcnow().isoformat(),
            "tables": {},
        }
        with self._connect() as conn:
            for table in tables:
                rows = conn.execute(f"SELECT * FROM {table}").fetchall()
                snapshot["tables"][table] = [dict(row) for row in rows]
        return snapshot

    def service_monitor_report(
        self,
        payments_hours: int = 24,
        spend_hours: int = 1,
        baseline_hours: int = 24,
    ) -> dict[str, Any]:
        payments_window = max(1, int(payments_hours or 24))
        spend_window = max(1, int(spend_hours or 1))
        baseline_window = max(spend_window + 1, int(baseline_hours or 24))
        now = datetime.utcnow()
        payments_since = (now - timedelta(hours=payments_window)).isoformat()
        spend_since = (now - timedelta(hours=spend_window)).isoformat()
        baseline_since = (now - timedelta(hours=baseline_window)).isoformat()
        baseline_until = spend_since

        with self._connect() as conn:
            payments_row = conn.execute(
                """
                SELECT
                    COUNT(DISTINCT CASE WHEN event_type IN ('text_request', 'image_request') THEN chat_id END) AS active_users,
                    SUM(CASE WHEN event_type = 'payment' AND rub_amount > 0 THEN 1 ELSE 0 END) AS payments_count,
                    SUM(CASE WHEN event_type = 'payment' THEN rub_amount ELSE 0 END) AS revenue_rub,
                    SUM(CASE WHEN event_type = 'refund' THEN rub_amount ELSE 0 END) AS refunds_rub
                FROM usage_events
                WHERE created_at >= ?
                """,
                (payments_since,),
            ).fetchone()

            recent_rows = conn.execute(
                """
                SELECT model_alias, tokens_total, details
                FROM usage_events
                WHERE created_at >= ? AND event_type = 'text_request'
                """,
                (spend_since,),
            ).fetchall()
            baseline_rows = conn.execute(
                """
                SELECT model_alias, tokens_total, details
                FROM usage_events
                WHERE created_at >= ? AND created_at < ? AND event_type = 'text_request'
                """,
                (baseline_since, baseline_until),
            ).fetchall()

        def rows_cost_rub(rows: list[sqlite3.Row]) -> float:
            total_usd = 0.0
            for row in rows:
                alias = str(row["model_alias"] or "").strip()
                _, _, model_info = model_info_for_alias(alias)
                prompt_tokens, completion_tokens = parse_usage_details_blob(str(row["details"] or ""))
                total_tokens = int(row["tokens_total"] or 0)
                total_usd += estimate_text_cost_usd(model_info, prompt_tokens, completion_tokens, total_tokens)
            return total_usd * ANALYTICS_USD_TO_RUB

        recent_cost_rub = rows_cost_rub(recent_rows)
        baseline_total_cost_rub = rows_cost_rub(baseline_rows)
        baseline_hourly_cost_rub = baseline_total_cost_rub / float(max(1, baseline_window - spend_window))
        baseline_cost_for_window_rub = baseline_hourly_cost_rub * spend_window
        return {
            "payments_window_hours": payments_window,
            "spend_window_hours": spend_window,
            "active_users": int((payments_row["active_users"] or 0) if payments_row else 0),
            "payments_count": int((payments_row["payments_count"] or 0) if payments_row else 0),
            "revenue_rub": int((payments_row["revenue_rub"] or 0) if payments_row else 0),
            "refunds_rub": int((payments_row["refunds_rub"] or 0) if payments_row else 0),
            "recent_text_cost_rub": recent_cost_rub,
            "baseline_text_cost_rub": baseline_cost_for_window_rub,
            "baseline_hourly_text_cost_rub": baseline_hourly_cost_rub,
        }


class BotState:
    def __init__(self) -> None:
        self.started_at = datetime.utcnow()
        self.user_histories: dict[int, deque[dict[str, str]]] = {}
        self.pending_receipt_plan: dict[int, str] = {}
        self.pending_image_prompt: set[int] = set()
        self.pending_image_ref_prompt: set[int] = set()
        self.pending_promo_code_input: set[int] = set()
        self.pending_referral_code_input: set[int] = set()
        self.pending_file_kind: dict[int, str] = {}
        self.pending_file_profile: dict[int, str] = {}
        self.image_request_prefs: dict[int, dict[str, str]] = {}
        self.pending_image_jobs: set[int] = set()
        self.image_job_queue: asyncio.Queue[ImageJob] = asyncio.Queue()
        self.last_reference_image_data_url: dict[int, str] = {}
        self.last_reference_image_at: dict[int, datetime] = {}
        self.processed_updates: deque[str] = deque()
        self.processed_lookup: set[str] = set()
        self.last_message_at: dict[int, datetime] = {}
        self.last_image_at: dict[int, datetime] = {}
        self.last_low_credits_nudge_at: dict[int, datetime] = {}
        self.error_alert_last_at: dict[str, datetime] = {}
        self.runtime_error_events: deque[datetime] = deque()
        self.channel_chat_id_cache: str = ""
        self.ui_message_mid: dict[int, str] = {}
        self.onboarding_message_mid: dict[int, str] = {}
        self.ui_current_page: dict[int, str] = {}
        self.ui_back_stack: dict[int, list[str]] = {}
        self.ui_forward_stack: dict[int, list[str]] = {}
        self.admin_sessions: dict[str, datetime] = {}
        self.admin_login_attempts: dict[str, deque[datetime]] = {}
        self.admin_login_blocked_until: dict[str, datetime] = {}
        self.max_api_semaphore = asyncio.Semaphore(max(1, MAX_API_CONCURRENCY))
        self.openrouter_text_semaphore = asyncio.Semaphore(max(1, OPENROUTER_TEXT_CONCURRENCY))
        self.openrouter_image_semaphore = asyncio.Semaphore(max(1, OPENROUTER_IMAGE_CONCURRENCY))
        self.kie_text_semaphore = asyncio.Semaphore(max(1, KIE_TEXT_CONCURRENCY))
        self.tbank_api_semaphore = asyncio.Semaphore(max(1, TBANK_API_CONCURRENCY))
        self.session: aiohttp.ClientSession | None = None
        self.polling_task: asyncio.Task[None] | None = None
        self.backup_task: asyncio.Task[None] | None = None
        self.monitor_task: asyncio.Task[None] | None = None
        self.image_worker_tasks: list[asyncio.Task[None]] = []
        self.user_store = UserStore(DB_BACKEND)

    def history(self, chat_id: int) -> deque[dict[str, str]]:
        if chat_id not in self.user_histories:
            self.user_histories[chat_id] = deque(maxlen=HISTORY_LIMIT)
        return self.user_histories[chat_id]


state = BotState()


def clear_growth_pending_inputs(chat_id: int) -> None:
    state.pending_referral_code_input.discard(chat_id)
    state.pending_promo_code_input.discard(chat_id)


def clear_file_pending_input(chat_id: int) -> None:
    state.pending_file_kind.pop(chat_id, None)
    state.pending_file_profile.pop(chat_id, None)


def migrate_runtime_chat_state(old_chat_id: int, new_chat_id: int) -> None:
    if old_chat_id == new_chat_id:
        return

    history = state.user_histories.pop(old_chat_id, None)
    if history is not None:
        target = state.user_histories.get(new_chat_id)
        if target is None:
            state.user_histories[new_chat_id] = history
        else:
            for item in history:
                target.append(item)

    dict_attrs = (
        "pending_receipt_plan",
        "pending_file_kind",
        "pending_file_profile",
        "image_request_prefs",
        "last_reference_image_data_url",
        "last_reference_image_at",
        "last_message_at",
        "last_image_at",
        "last_low_credits_nudge_at",
        "ui_message_mid",
        "onboarding_message_mid",
        "ui_current_page",
        "ui_back_stack",
        "ui_forward_stack",
    )
    for attr in dict_attrs:
        mapping = getattr(state, attr)
        if old_chat_id not in mapping:
            continue
        value = mapping.pop(old_chat_id)
        if new_chat_id not in mapping:
            mapping[new_chat_id] = value

    set_attrs = (
        "pending_image_prompt",
        "pending_image_ref_prompt",
        "pending_promo_code_input",
        "pending_referral_code_input",
    )
    for attr in set_attrs:
        items = getattr(state, attr)
        if old_chat_id in items:
            items.discard(old_chat_id)
            items.add(new_chat_id)


def ensure_update_user_binding(chat_id: int, update: dict[str, Any]) -> None:
    if chat_id <= 0:
        return
    actor_user_id = parse_actor_user_id(update)
    rebound_from = state.user_store.ensure_user_binding(chat_id, actor_user_id, best_default_alias_for_plan("free"))
    if rebound_from and rebound_from != chat_id:
        log.info("Rebound MAX user_id=%s from chat_id=%s to chat_id=%s", actor_user_id, rebound_from, chat_id)
        migrate_runtime_chat_state(rebound_from, chat_id)


def handoff_onboarding_to_ui(chat_id: int, source_mid: str | None) -> None:
    onboarding_mid = state.onboarding_message_mid.get(chat_id)
    target_mid = onboarding_mid or source_mid
    if target_mid:
        state.ui_message_mid[chat_id] = target_mid
    state.onboarding_message_mid.pop(chat_id, None)


def looks_like_bonus_code(text: str) -> bool:
    value = text.strip()
    if not value:
        return False
    return bool(re.fullmatch(r"[A-Za-zА-Яа-я0-9_-]{3,40}", value))


def init_sentry_if_enabled() -> None:
    if not SENTRY_DSN:
        return
    try:
        import sentry_sdk  # type: ignore

        sentry_sdk.init(
            dsn=SENTRY_DSN,
            environment=SENTRY_ENVIRONMENT,
            traces_sample_rate=0.0,
        )
        log.info("Sentry initialized")
    except Exception:
        log.exception("Failed to initialize Sentry")


def capture_exception_safe(exc: Exception) -> None:
    if not SENTRY_DSN:
        return
    try:
        import sentry_sdk  # type: ignore

        sentry_sdk.capture_exception(exc)
    except Exception:
        pass


async def notify_admin_alert(key: str, text: str) -> None:
    if not ERROR_ALERTS_ENABLED:
        return
    now = datetime.utcnow()
    last = state.error_alert_last_at.get(key)
    if last and (now - last).total_seconds() < max(10, ERROR_ALERT_COOLDOWN_SEC):
        return
    state.error_alert_last_at[key] = now
    for admin_id in admin_target_chat_ids():
        with suppress(Exception):
            await max_send_message(admin_id, f"⚠️ ALERT [{key}]\n{text}", notify=False)


def record_runtime_error() -> None:
    now = datetime.utcnow()
    state.runtime_error_events.append(now)
    cutoff = now - timedelta(hours=24)
    while state.runtime_error_events and state.runtime_error_events[0] < cutoff:
        state.runtime_error_events.popleft()


def require_env() -> None:
    missing = []
    if not MAX_TOKEN:
        missing.append("MAX_TOKEN")
    if not OPENROUTER_KEY:
        missing.append("OPENROUTER_KEY")
    if missing:
        raise RuntimeError(f"Missing required env vars: {', '.join(missing)}")


def validate_pricing_sanity() -> None:
    if PRO_PLAN_PRICE_RUB <= 0 or credits_for_plan("pro") <= 0:
        log.warning("Pricing sanity skipped: invalid Pro plan price/credits.")
        return

    pro_ratio = credits_for_plan("pro") / float(PRO_PLAN_PRICE_RUB)
    for code, pack in TOPUP_PACKS.items():
        price = int(pack.get("price_rub", 0) or 0)
        credits = int(pack.get("credits", 0) or 0)
        if price <= 0 or credits <= 0:
            log.warning("Top-up pack %s has invalid config: price=%s credits=%s", code, price, credits)
            continue
        ratio = credits / float(price)
        if ratio >= pro_ratio:
            log.warning(
                "Top-up pack %s is too generous: %.4f credits/RUB >= Pro %.4f credits/RUB",
                code,
                ratio,
                pro_ratio,
            )


def load_model_registry() -> tuple[dict[str, ModelInfo], dict[str, ModelInfo]]:
    with CONFIG_PATH.open("r", encoding="utf-8") as fh:
        raw = json.load(fh)

    text_models: dict[str, ModelInfo] = {}
    image_models: dict[str, ModelInfo] = {}

    for alias, payload in raw.get("text_models", {}).items():
        text_models[alias] = ModelInfo(
            alias=alias,
            provider=payload["provider"],
            model=payload["model"],
            label=payload["label"],
            kind="text",
            version=payload.get("version", ""),
            input_price_usd_per_m=payload.get("input_price_usd_per_m", ""),
            output_price_usd_per_m=payload.get("output_price_usd_per_m", ""),
            min_plan=payload.get("min_plan", "free"),
            description=payload.get("description", ""),
            prompt_style=payload.get("prompt_style", ""),
        )

    for alias, payload in raw.get("image_models", {}).items():
        image_models[alias] = ModelInfo(
            alias=alias,
            provider=payload["provider"],
            model=payload["model"],
            label=payload["label"],
            kind="image",
            version=payload.get("version", ""),
            input_price_usd_per_m=payload.get("input_price_usd_per_m", ""),
            output_price_usd_per_m=payload.get("output_price_usd_per_m", ""),
            min_plan=payload.get("min_plan", "free"),
            description=payload.get("description", ""),
            prompt_style=payload.get("prompt_style", ""),
        )

    return text_models, image_models


TEXT_MODELS, IMAGE_MODELS = load_model_registry()
DEFAULT_TEXT_MODEL = TEXT_MODELS.get(DEFAULT_TEXT_MODEL_ALIAS, TEXT_MODELS["gpt"])
DEFAULT_IMAGE_MODEL = IMAGE_MODELS.get(DEFAULT_IMAGE_MODEL_ALIAS, IMAGE_MODELS["image"])


def max_headers() -> dict[str, str]:
    return {"Authorization": MAX_TOKEN, "Content-Type": "application/json"}


def openrouter_headers() -> dict[str, str]:
    return {
        "Authorization": f"Bearer {OPENROUTER_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": "https://max.ru",
        "X-Title": "MAX Multi AI Bot",
    }


def kie_headers() -> dict[str, str]:
    return {
        "Authorization": f"Bearer {KIE_API_KEY}",
        "Content-Type": "application/json",
    }


def kie_enabled_for_alias(alias: str) -> bool:
    return bool(KIE_API_KEY and alias.strip().lower() in KIE_TEXT_ALIASES)


def kie_error_message(data: Any) -> str:
    if not isinstance(data, dict):
        return ""
    code = data.get("code")
    if isinstance(code, int) and code >= 400:
        message = data.get("msg") or data.get("message") or data.get("error")
        return str(message or f"Kie API error code {code}")
    if isinstance(code, str) and code.isdigit() and int(code) >= 400:
        message = data.get("msg") or data.get("message") or data.get("error")
        return str(message or f"Kie API error code {code}")
    return ""


def tbank_enabled() -> bool:
    return bool(TBANK_TERMINAL_KEY and TBANK_PASSWORD)


def tbank_order_id(request_id: int) -> str:
    suffix = datetime.utcnow().strftime("%y%m%d%H%M%S")
    return f"MAXBOT-{request_id}-{suffix}"


def parse_request_id_from_order_id(order_id: str) -> int | None:
    value = order_id.strip().upper()
    if not value.startswith("MAXBOT-"):
        return None
    payload = value.split("-", 1)[1]
    request_part = payload.split("-", 1)[0]
    return int(request_part) if request_part.isdigit() else None


def add_request_id_to_url(url: str, request_id: int) -> str:
    if not url:
        return ""
    parsed = urlsplit(url)
    query = dict(parse_qsl(parsed.query, keep_blank_values=True))
    query["request_id"] = str(request_id)
    expires_at = int((datetime.utcnow() + timedelta(seconds=max(60, PAYMENT_STATUS_TOKEN_TTL_SECONDS))).timestamp())
    status_sig = payment_status_signature(request_id, expires_at)
    if status_sig:
        query["status_ts"] = str(expires_at)
        query["status_sig"] = status_sig
    return urlunsplit(
        (
            parsed.scheme,
            parsed.netloc,
            parsed.path,
            urlencode(query),
            parsed.fragment,
        )
    )


def tbank_payment_id_from_provider_ref(provider_ref: str) -> str:
    value = provider_ref.strip()
    if not value:
        return ""
    if value.startswith("tbank:"):
        return value.split(":", 1)[1].strip()
    return value


def build_tbank_receipt(
    amount_rub: int,
    description: str,
    receipt_email: str = "",
    receipt_phone: str = "",
) -> dict[str, Any]:
    amount_kop = int(amount_rub) * 100
    item_name = (description or "Подписка").strip()[:128] or "Подписка"
    receipt: dict[str, Any] = {
        "Taxation": TBANK_RECEIPT_TAXATION or "usn_income",
        "Items": [
            {
                "Name": item_name,
                "Price": amount_kop,
                "Quantity": 1,
                "Amount": amount_kop,
                "Tax": TBANK_RECEIPT_TAX or "none",
                "PaymentMethod": TBANK_RECEIPT_PAYMENT_METHOD or "full_prepayment",
                "PaymentObject": TBANK_RECEIPT_PAYMENT_OBJECT or "service",
            }
        ],
    }
    if TBANK_RECEIPT_FFD_VERSION:
        receipt["FfdVersion"] = TBANK_RECEIPT_FFD_VERSION
    if receipt_email:
        receipt["Email"] = receipt_email
    if receipt_phone:
        receipt["Phone"] = receipt_phone
    return receipt


def scalar_string(value: Any) -> str:
    if isinstance(value, bool):
        return "true" if value else "false"
    if value is None:
        return ""
    return str(value)


def tbank_token_from_payload(payload: dict[str, Any], password: str) -> str:
    # T-Bank token uses only root scalar fields (+ Password), excluding Token and nested objects.
    values: dict[str, str] = {}
    for key, value in payload.items():
        if key == "Token":
            continue
        if isinstance(value, (dict, list)):
            continue
        values[key] = scalar_string(value)
    values["Password"] = password
    raw = "".join(values[key] for key in sorted(values.keys()))
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def tbank_notification_is_valid(payload: dict[str, Any]) -> bool:
    if TBANK_TERMINAL_KEY:
        terminal_key = scalar_string(payload.get("TerminalKey"))
        if terminal_key != TBANK_TERMINAL_KEY:
            return False
    if not TBANK_PASSWORD:
        return True
    received = scalar_string(payload.get("Token")).lower()
    if not received:
        return False
    expected = tbank_token_from_payload(payload, TBANK_PASSWORD).lower()
    return received == expected


def tbank_is_cancel_status(status: str) -> bool:
    return status.strip().upper() in TBANK_CANCEL_STATUSES


def tbank_is_refund_status(status: str) -> bool:
    return status.strip().upper() in TBANK_REFUND_STATUSES


def site_file(name: str) -> Path:
    path = SITE_DIR / name
    if not path.exists():
        raise HTTPException(status_code=404, detail="page not found")
    return path


def support_url_value() -> str:
    if SUPPORT_URL:
        return SUPPORT_URL
    if PUBLIC_BASE_URL:
        return f"{PUBLIC_BASE_URL}/support"
    return "/support"


def support_help_text() -> str:
    return (
        "Помощь\n\n"
        "Если что-то не работает:\n"
        "1. Открой «Тарифы» и создай новую заявку\n"
        "2. После оплаты подожди 1-2 минуты\n"
        "3. Открой «Мой план» или «Мои оплаты» и проверь статус\n\n"
        f"{SUPPORT_TEXT}\n"
        f"Email: {CONTACT_EMAIL}\n"
        f"FAQ: {support_url_value()}\n\n"
        "Что прислать в поддержку:\n"
        "• номер заявки\n"
        "• что произошло\n"
        "• скрин ошибки или оплаты, если он есть"
    )


def channel_url_value() -> str:
    return CHANNEL_URL or "https://max.ru/id231128398751_biz"


def channel_chat_id_value() -> str:
    raw = CHANNEL_CHAT_ID or channel_url_value()
    value = str(raw or "").strip()
    if value.startswith(("http://", "https://")):
        value = urlsplit(value).path.strip("/").rsplit("/", 1)[-1].strip()
    if re.fullmatch(r"-?\d+", value):
        return value
    match = re.fullmatch(r"id(-?\d+)(?:[_-].*)?", value)
    if match:
        return match.group(1)
    return value


def normalize_channel_link(value: str) -> str:
    raw = str(value or "").strip()
    if not raw:
        return ""
    parsed = urlsplit(raw if raw.startswith(("http://", "https://")) else f"https://max.ru/{raw.strip('/')}")
    host = parsed.netloc.lower().removeprefix("www.")
    path = parsed.path.strip("/").lower()
    return f"{host}/{path}".rstrip("/")


def extract_chats_from_response(data: Any) -> tuple[list[dict[str, Any]], str]:
    if isinstance(data, list):
        return [item for item in data if isinstance(item, dict)], ""
    if not isinstance(data, dict):
        return [], ""
    for key in ("chats", "items", "data"):
        value = data.get(key)
        if isinstance(value, list):
            marker = str(data.get("marker") or data.get("next_marker") or data.get("nextMarker") or "")
            return [item for item in value if isinstance(item, dict)], marker
    return [], ""


def chat_id_from_chat_item(item: dict[str, Any]) -> str:
    for key in ("chat_id", "chatId", "id"):
        value = item.get(key)
        if isinstance(value, (int, str)) and str(value).strip():
            return str(value).strip()
    return ""


def chat_item_link_values(item: dict[str, Any]) -> list[str]:
    values: list[str] = []
    for key in ("link", "url", "invite_link", "inviteLink", "public_link", "publicLink"):
        value = item.get(key)
        if isinstance(value, str) and value.strip():
            values.append(value)
    for key in ("description", "payload", "settings"):
        value = item.get(key)
        if isinstance(value, dict):
            values.extend(chat_item_link_values(value))
    return values


async def resolve_channel_chat_id() -> str:
    if CHANNEL_CHAT_ID:
        return channel_chat_id_value()
    if state.channel_chat_id_cache:
        return state.channel_chat_id_cache
    if not MAX_TOKEN:
        return ""

    marker = ""
    target_link = normalize_channel_link(channel_url_value())
    channel_candidates: list[dict[str, Any]] = []
    for _ in range(5):
        params = {"count": "100"}
        if marker:
            params["marker"] = marker
        status, data, _ = await http_json_request_with_retries(
            "GET",
            f"{MAX_API}/chats",
            headers=max_headers(),
            params=params,
            semaphore=state.max_api_semaphore,
            request_name="max_chats_resolve",
        )
        if status >= 400:
            await notify_admin_alert("channel_gate_resolve", f"MAX chats resolve failed: status={status}, body={str(data)[:300]}")
            return ""
        chats, marker = extract_chats_from_response(data)
        for item in chats:
            chat_type = str(item.get("type") or item.get("chat_type") or item.get("chatType") or "").lower()
            status = str(item.get("status") or "").lower()
            if chat_type and chat_type not in {"channel", "chat"}:
                continue
            if status and status not in {"active", "enabled"}:
                continue
            item_id = chat_id_from_chat_item(item)
            if not item_id:
                continue
            normalized_links = [normalize_channel_link(value) for value in chat_item_link_values(item)]
            if target_link and target_link in normalized_links:
                state.channel_chat_id_cache = item_id
                return item_id
            if chat_type == "channel":
                channel_candidates.append(item)
        if not marker:
            break

    if len(channel_candidates) == 1:
        item_id = chat_id_from_chat_item(channel_candidates[0])
        if item_id:
            state.channel_chat_id_cache = item_id
            await notify_admin_alert(
                "channel_gate_resolve",
                f"CHANNEL_CHAT_ID не задан. Использую единственный найденный канал chat_id={item_id}.",
            )
            return item_id

    await notify_admin_alert(
        "channel_gate_resolve",
        "Не удалось найти канал бота через /chats. Добавь бота в канал как админа или задай CHANNEL_CHAT_ID вручную.",
    )
    return ""


def channel_subscription_cache_valid(row: dict[str, Any]) -> bool:
    subscribed_at = parse_iso_datetime(str(row.get("channel_subscribed_at", "") or ""))
    checked_at = parse_iso_datetime(str(row.get("channel_subscription_checked_at", "") or ""))
    if subscribed_at is None or checked_at is None:
        return False
    cache_hours = max(1, CHANNEL_MEMBERSHIP_CACHE_HOURS)
    return datetime.utcnow() - checked_at <= timedelta(hours=cache_hours)


def response_contains_user_id(node: Any, user_id: int) -> bool:
    if isinstance(node, dict):
        for key in ("user_id", "userId", "id"):
            value = node.get(key)
            if isinstance(value, int) and value == user_id:
                return True
            if isinstance(value, str) and value.isdigit() and int(value) == user_id:
                return True
        return any(response_contains_user_id(value, user_id) for value in node.values())
    if isinstance(node, list):
        return any(response_contains_user_id(item, user_id) for item in node)
    return False


def response_has_member_items(node: Any) -> bool:
    if isinstance(node, list):
        return len(node) > 0
    if not isinstance(node, dict):
        return False
    for key in ("members", "items", "participants", "users"):
        value = node.get(key)
        if isinstance(value, list):
            return len(value) > 0
        if isinstance(value, dict):
            if response_has_member_items(value):
                return True
    return False


def channel_membership_response_is_positive(data: Any, max_user_id: int) -> bool:
    return response_contains_user_id(data, max_user_id) or response_has_member_items(data)


async def check_channel_subscription(chat_id: int, force: bool = False) -> tuple[bool, str]:
    if not CHANNEL_GATE_ENABLED:
        return True, "disabled"
    if is_admin(chat_id):
        return True, "admin"

    row = user_profile(chat_id)
    if not force and channel_subscription_cache_valid(row):
        return True, "cached"

    max_user_id = int(row.get("max_user_id", 0) or 0)
    if max_user_id <= 0:
        await notify_admin_alert("channel_gate_user_id", f"Не удалось определить max_user_id для chat_id={chat_id}")
        return False, "no_user_id"

    channel_id = await resolve_channel_chat_id()
    if not MAX_TOKEN or not channel_id:
        await notify_admin_alert(
            "channel_gate_config",
            f"Проверка подписки включена, но не найден chat_id канала. chat_id={chat_id}",
        )
        return False, "config_missing"

    url = f"{MAX_API}/chats/{quote(channel_id, safe='')}/members"
    try:
        status, data, _ = await http_json_request_with_retries(
            "GET",
            url,
            headers=max_headers(),
            params={"user_ids": str(max_user_id)},
            semaphore=state.max_api_semaphore,
            request_name="channel_members_check",
        )
        if status >= 400:
            await notify_admin_alert(
                "channel_gate_api",
                f"MAX members check failed: status={status}, channel_id={channel_id}, user_id={max_user_id}, body={str(data)[:300]}",
            )
            if status == 400 and "dialogs" in str(data).lower():
                state.channel_chat_id_cache = ""
            return False, f"api_error_{status}"
    except Exception as exc:
        await notify_admin_alert(
            "channel_gate_api",
            f"MAX members check exception: channel_id={channel_id}, user_id={max_user_id}, error={exc}",
        )
        return False, "api_exception"

    subscribed = channel_membership_response_is_positive(data, max_user_id)
    state.user_store.mark_channel_subscription(chat_id, subscribed)
    state.user_store.record_usage_event(
        chat_id=chat_id,
        event_type="channel_subscription_check",
        plan=str(row.get("plan", "")),
        details=f"subscribed={1 if subscribed else 0};force={1 if force else 0};channel_id={channel_id}",
    )
    return subscribed, "subscribed" if subscribed else "not_subscribed"


def support_admin_templates_text() -> str:
    return (
        "Шаблоны для спорных кейсов\n\n"
        "1) Оплата есть, доступа нет:\n"
        "«Проверили оплату по заявке #{request_id}. Статус в банке: {bank_status}. "
        "Доступ обновлен, проверь раздел “Мой план”. Если не обновилось — напишите нам, проверим вручную.»\n\n"
        "2) Оплата не прошла:\n"
        "«По заявке #{request_id} банк вернул статус {bank_status}. Оплата не завершена. "
        "Создайте новую заявку в разделе “Тарифы” и повторите оплату.»\n\n"
        "3) Возврат подтвержден:\n"
        "«По заявке #{request_id} оформлен возврат. Подписка переведена на free, автопродление отключено. "
        "Срок зачисления средств зависит от банка-эмитента.»\n\n"
        "4) Чарджбэк/спор:\n"
        "«Получили запрос на оспаривание платежа по заявке #{request_id}. "
        "Для проверки пришлите дату/время оплаты и скрин подтверждения операции.»"
    )


def admin_panel_authorized(token: str) -> bool:
    return bool(ADMIN_PANEL_TOKEN) and token.strip() == ADMIN_PANEL_TOKEN


def request_client_ip(request: Request) -> str:
    forwarded = request.headers.get("X-Forwarded-For", "")
    if forwarded:
        first = forwarded.split(",", 1)[0].strip()
        if first:
            return first
    client = request.client
    return client.host if client and client.host else "unknown"


def prune_admin_sessions() -> None:
    now = datetime.utcnow()
    expired = [sid for sid, exp in state.admin_sessions.items() if exp <= now]
    for sid in expired:
        state.admin_sessions.pop(sid, None)


def issue_admin_session() -> str:
    sid = secrets.token_urlsafe(32)
    state.admin_sessions[sid] = datetime.utcnow() + timedelta(seconds=ADMIN_SESSION_MAX_AGE)
    return sid


def admin_session_valid(session_id: str) -> bool:
    if not session_id:
        return False
    prune_admin_sessions()
    expires_at = state.admin_sessions.get(session_id)
    if not expires_at:
        return False
    if expires_at <= datetime.utcnow():
        state.admin_sessions.pop(session_id, None)
        return False
    return True


def resolve_admin_session(request: Request, token: str = "") -> str:
    cookie_sid = request.cookies.get(ADMIN_SESSION_COOKIE, "").strip()
    if admin_session_valid(cookie_sid):
        return cookie_sid
    provided = token.strip()
    if admin_panel_authorized(provided):
        return issue_admin_session()
    return ""


def set_admin_cookie(response: HTMLResponse | RedirectResponse, session_id: str) -> None:
    secure = PUBLIC_BASE_URL.startswith("https://")
    response.set_cookie(
        key=ADMIN_SESSION_COOKIE,
        value=session_id,
        max_age=ADMIN_SESSION_MAX_AGE,
        httponly=True,
        secure=secure,
        samesite="lax",
        path="/",
    )


def clear_admin_cookie(response: HTMLResponse | RedirectResponse) -> None:
    response.delete_cookie(key=ADMIN_SESSION_COOKIE, path="/")


def admin_csrf_token(session_id: str) -> str:
    if not session_id:
        return ""
    payload = f"{session_id}:csrf:v1"
    return hmac.new(ADMIN_PANEL_TOKEN.encode("utf-8"), payload.encode("utf-8"), hashlib.sha256).hexdigest()


def admin_csrf_valid(session_id: str, csrf_token: str) -> bool:
    if not session_id or not csrf_token:
        return False
    expected = admin_csrf_token(session_id)
    return bool(expected) and hmac.compare_digest(expected, csrf_token.strip())


def admin_login_allowed(request: Request) -> tuple[bool, str]:
    ip = request_client_ip(request)
    now = datetime.utcnow()
    blocked_until = state.admin_login_blocked_until.get(ip)
    if blocked_until and blocked_until > now:
        left = int((blocked_until - now).total_seconds())
        return False, f"Слишком много попыток. Повтори через {max(1, left)} сек."
    attempts = state.admin_login_attempts.setdefault(ip, deque())
    cutoff = now - timedelta(seconds=max(30, ADMIN_LOGIN_WINDOW_SECONDS))
    while attempts and attempts[0] < cutoff:
        attempts.popleft()
    if len(attempts) >= max(1, ADMIN_LOGIN_MAX_ATTEMPTS):
        state.admin_login_blocked_until[ip] = now + timedelta(seconds=max(60, ADMIN_LOGIN_BLOCK_SECONDS))
        return False, "Слишком много попыток входа. Попробуй позже."
    return True, ""


def admin_login_register_failure(request: Request) -> None:
    ip = request_client_ip(request)
    now = datetime.utcnow()
    attempts = state.admin_login_attempts.setdefault(ip, deque())
    cutoff = now - timedelta(seconds=max(30, ADMIN_LOGIN_WINDOW_SECONDS))
    while attempts and attempts[0] < cutoff:
        attempts.popleft()
    attempts.append(now)


def admin_login_register_success(request: Request) -> None:
    ip = request_client_ip(request)
    state.admin_login_attempts.pop(ip, None)
    state.admin_login_blocked_until.pop(ip, None)


def admin_url(path: str, token: str = "", **params: Any) -> str:
    query: dict[str, str] = {}
    if token:
        query["token"] = token
    for key, value in params.items():
        if value is None or value == "":
            continue
        query[key] = str(value)
    if not query:
        return path
    return f"{path}?{urlencode(query)}"


def repair_mojibake(text: str) -> str:
    if not text or ("Р" not in text and "вЂ" not in text):
        return text
    try:
        fixed = text.encode("latin1", errors="ignore").decode("utf-8", errors="ignore")
    except Exception:
        return text
    return fixed or text


def runtime_support_deps() -> RuntimeSupportDeps:
    return RuntimeSupportDeps(
        data_dir=DATA_DIR,
        db_path=DB_PATH,
        backup_keep_files=BACKUP_KEEP_FILES,
        max_token=MAX_TOKEN,
        openrouter_key=OPENROUTER_KEY,
        run_mode=RUN_MODE,
        auto_backup_enabled=AUTO_BACKUP_ENABLED,
        service_monitor_enabled=SERVICE_MONITOR_ENABLED,
        config_path=CONFIG_PATH,
        site_index_path=site_file("index.html"),
        tbank_terminal_key=TBANK_TERMINAL_KEY,
        tbank_password=TBANK_PASSWORD,
        channel_gate_enabled=CHANNEL_GATE_ENABLED,
        channel_chat_id=CHANNEL_CHAT_ID,
        channel_membership_cache_hours=CHANNEL_MEMBERSHIP_CACHE_HOURS,
        alert_high_errors_window_minutes=ALERT_HIGH_ERRORS_WINDOW_MINUTES,
        alert_low_payments_lookback_hours=ALERT_LOW_PAYMENTS_LOOKBACK_HOURS,
        alert_spend_spike_lookback_hours=ALERT_SPEND_SPIKE_LOOKBACK_HOURS,
        channel_url_resolver=channel_url_value,
        system_snapshot_provider=system_resource_snapshot,
    )


def backups_dir() -> Path:
    return backups_dir_impl(DATA_DIR)


def create_db_backup() -> Path:
    return create_db_backup_impl(
        data_dir=DATA_DIR,
        db_path=DB_PATH,
        backup_keep_files=BACKUP_KEEP_FILES,
        user_store=state.user_store,
    )


def latest_backup_file() -> Path | None:
    return latest_backup_file_impl(DATA_DIR)


def format_timedelta_short(delta: timedelta) -> str:
    return format_timedelta_short_impl(delta)


def smoke_check_report() -> list[dict[str, str]]:
    return smoke_check_report_impl(
        deps=runtime_support_deps(),
        user_store=state.user_store,
        state=state,
    )


def service_status_report() -> dict[str, Any]:
    return service_status_report_impl(
        deps=runtime_support_deps(),
        user_store=state.user_store,
        state=state,
    )


def system_resource_snapshot() -> dict[str, Any]:
    return system_resource_snapshot_impl(DATA_DIR)


def payment_status_label(status: str) -> str:
    return PAYMENT_STATUS_LABELS.get(status.strip().lower(), "Статус обновляется")


def payment_status_title_message(status: str) -> tuple[str, str]:
    key = status.strip().lower()
    titles = {
        "pending": "Платеж создан",
        "claimed": "Проверяем оплату",
        "paid": "Оплата подтверждена",
        "canceled": "Оплата не завершена",
        "refunded": "Оформлен возврат",
    }
    messages = {
        "pending": "Платеж в обработке. Обычно подтверждение приходит в течение 1-2 минут.",
        "claimed": "Банк прислал сигнал, подтверждаем оплату. Обычно это занимает до 1-2 минут.",
        "paid": "Подписка уже активирована. Можно возвращаться в бот.",
        "canceled": "Оплата не была завершена. Вернись в бот и попробуй еще раз.",
        "refunded": "Возврат подтвержден. Подписка отключена, действует тариф free.",
    }
    return (
        titles.get(key, "Статус обновляется"),
        messages.get(key, "Подожди немного и обнови страницу."),
    )


def payment_user_status_text(payment: dict[str, Any], bank_status: str = "") -> str:
    request_id = int(payment.get("id", 0) or 0)
    status = str(payment.get("status", "pending")).lower()
    status_human = payment_status_label(status)
    plan = str(payment.get("plan", ""))
    amount = int(payment.get("amount_rub", 0) or 0)
    payment_url = str(payment.get("payment_url", "") or "").strip()
    bank_line = f"\nСтатус банка: {bank_status}" if bank_status else ""
    pay_line = f"\nСсылка на оплату: {payment_url}" if payment_url and status in {"pending", "claimed"} else ""

    if status == "paid":
        return (
            f"✅ Заявка #{request_id}: {status_human}\n"
            f"Тариф/пакет: {plan}, сумма: {amount} ₽.\n"
            "Доступ уже активирован. Можно возвращаться в меню."
        )
    if status == "refunded":
        return (
            f"↩️ Заявка #{request_id}: {status_human}\n"
            f"Тариф/пакет: {plan}, сумма: {amount} ₽.{bank_line}\n"
            "Возврат подтвержден. Если деньги долго не приходят, напиши в поддержку."
        )
    if status == "canceled":
        return (
            f"❌ Заявка #{request_id}: {status_human}\n"
            f"Тариф/пакет: {plan}, сумма: {amount} ₽.{bank_line}\n"
            "Оплата не завершена. Можно создать новую заявку в «Тарифы»."
        )
    if status == "claimed":
        return (
            f"🕒 Заявка #{request_id}: {status_human}\n"
            f"Тариф/пакет: {plan}, сумма: {amount} ₽.{bank_line}{pay_line}\n"
            "Платеж уже отправлен на проверку. Обычно это занимает 1-2 минуты."
        )
    return (
        f"⏳ Заявка #{request_id}: {status_human}\n"
        f"Тариф/пакет: {plan}, сумма: {amount} ₽.{bank_line}{pay_line}\n"
        "Открой ссылку на оплату. Если уже оплатил — нажми «Проверить статус»."
    )


def plan_allowed(plan: str, min_plan: str) -> bool:
    if min_plan not in PLAN_ORDER:
        log.warning("Unknown min_plan=%r in model config; denying access", min_plan)
        return False
    if plan not in PLAN_ORDER:
        return False
    return PLAN_ORDER[plan] >= PLAN_ORDER[min_plan]


def text_model_allowed_for_plan(plan: str, model_alias: str) -> bool:
    info = TEXT_MODELS.get(model_alias)
    if not info:
        return False
    if plan_allowed(plan, info.min_plan):
        return True
    if model_alias == "gpt54" and plan == "start" and PLAN_CONFIGS["start"].daily_gpt54_limit > 0:
        return True
    return False


def is_admin(chat_id: int) -> bool:
    row = state.user_store.get_user(chat_id)
    max_user_id = int((row or {}).get("max_user_id", 0) or 0)
    if ADMIN_MAX_USER_IDS:
        return max_user_id > 0 and max_user_id in ADMIN_MAX_USER_IDS
    return chat_id in ADMIN_IDS


def admin_target_chat_ids() -> set[int]:
    targets = {int(value) for value in ADMIN_IDS if int(value) > 0}
    for max_user_id in ADMIN_MAX_USER_IDS:
        row = state.user_store.get_user_by_max_user_id(max_user_id)
        if not row:
            continue
        chat_id = int(row.get("chat_id", 0) or 0)
        if chat_id > 0:
            targets.add(chat_id)
    return targets


def best_default_alias_for_plan(plan: str) -> str:
    preferred = ["gpt4o", DEFAULT_TEXT_MODEL.alias, "gpt", "deepseek"]
    for alias in preferred:
        if text_model_allowed_for_plan(plan, alias):
            return alias
    for alias in TEXT_MODELS:
        if text_model_allowed_for_plan(plan, alias):
            return alias
    return DEFAULT_TEXT_MODEL.alias


def resolve_preset_alias_for_plan(plan: str, preset: str) -> str:
    preset_cfg = MODEL_PRESETS.get(preset)
    if not preset_cfg:
        return best_default_alias_for_plan(plan)
    for alias in preset_cfg.get("aliases", []):
        if text_model_allowed_for_plan(plan, alias):
            return alias
    return best_default_alias_for_plan(plan)


def resolve_preset_alias_for_chat(chat_id: int, preset: str) -> str:
    plan = str(user_profile(chat_id).get("plan", "free"))
    return resolve_preset_alias_for_plan(plan, preset)


def allowed_text_aliases_for_plan(plan: str) -> list[str]:
    aliases: list[str] = []
    seen: set[str] = set()
    preferred = ["deepseek", "gpt", "gpt4o", "gemini", "gpt54"]
    for alias in preferred:
        if alias in seen or not text_model_allowed_for_plan(plan, alias):
            continue
        seen.add(alias)
        aliases.append(alias)
    for alias in TEXT_MODELS:
        if alias in seen or not text_model_allowed_for_plan(plan, alias):
            continue
        seen.add(alias)
        aliases.append(alias)
    return aliases


def preset_available_aliases_for_plan(plan: str, preset: str) -> list[str]:
    preset_cfg = MODEL_PRESETS.get(preset) or {}
    aliases: list[str] = []
    seen: set[str] = set()
    for raw_alias in preset_cfg.get("aliases", []):
        alias = str(raw_alias).strip()
        if not alias or alias in seen:
            continue
        info = TEXT_MODELS.get(alias)
        if not info or not text_model_allowed_for_plan(plan, alias):
            continue
        seen.add(alias)
        aliases.append(alias)
    if not aliases:
        aliases.append(best_default_alias_for_plan(plan))
    return aliases


def preset_choice_enabled_for_plan(plan: str, preset: str) -> bool:
    if plan == "free":
        return False
    return len(preset_available_aliases_for_plan(plan, preset)) > 1


def current_preset_for_chat(chat_id: int) -> str:
    row = user_profile(chat_id)
    preset = str(row.get("selected_preset", "") or "").strip().lower()
    return preset if preset in MODEL_PRESETS else ""


def can_pick_models_for_current_preset(chat_id: int) -> bool:
    plan = str(user_profile(chat_id).get("plan", "free"))
    return len(allowed_text_aliases_for_plan(plan)) > 1


def preset_model_hint(alias: str) -> str:
    return {
        "deepseek": "самый быстрый и выгодный",
        "gpt": "аккуратный и лёгкий для everyday-задач",
        "gpt4o": "живой универсальный диалог",
        "gemini": "лучше для длинных и подробных задач",
        "gpt54": "максимум качества для сложных запросов",
    }.get(alias, "подходит для этого режима")


def build_preset_block(plan: str) -> str:
    lines = ["🎛 Режимы ответов для вашего тарифа:"]
    for key in ("fast", "balanced", "quality", "expert"):
        cfg = MODEL_PRESETS[key]
        alias = resolve_preset_alias_for_plan(plan, key)
        label = TEXT_MODELS.get(alias, DEFAULT_TEXT_MODEL).label
        lines.append(f"• {cfg['label']} — {cfg['description']} ({label})")
    lines.append("• 🎨 Картинка — отдельный режим для генерации и редактирования")
    return "\n".join(lines)


def record_ui_page_view(chat_id: int, page: str | None) -> None:
    if page not in UI_PAGE_KEYS:
        return
    row = user_profile(chat_id)
    state.user_store.record_usage_event(
        chat_id=chat_id,
        event_type="screen_view",
        plan=str(row.get("plan", "")),
        details=f"screen={page}",
    )


def normalize_text_content(content: Any) -> str:
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, dict) and item.get("type") == "text":
                text = item.get("text")
                if isinstance(text, str):
                    parts.append(text)
            elif isinstance(item, dict) and item.get("type") == "output_text":
                text = item.get("text")
                if isinstance(text, str):
                    parts.append(text)
        return "\n".join(part for part in parts if part).strip()
    return ""


def normalize_response_output_text(data: dict[str, Any]) -> str:
    choices = data.get("choices")
    if isinstance(choices, list) and choices:
        message = choices[0].get("message") if isinstance(choices[0], dict) else None
        if isinstance(message, dict):
            text = normalize_text_content(message.get("content"))
            if text:
                return text

    output = data.get("output")
    if isinstance(output, list):
        parts: list[str] = []
        for item in output:
            if not isinstance(item, dict):
                continue
            content = item.get("content")
            text = normalize_text_content(content)
            if text:
                parts.append(text)
        if parts:
            return "\n".join(parts).strip()

    output_text = data.get("output_text")
    if isinstance(output_text, str) and output_text.strip():
        return output_text.strip()
    return ""


def split_message(text: str, limit: int = MAX_MESSAGE_LEN) -> list[str]:
    content = text.strip()
    if not content:
        return []
    if len(content) <= limit:
        return [content]

    chunks: list[str] = []
    remaining = content
    while remaining:
        if len(remaining) <= limit:
            chunks.append(remaining)
            break
        split_at = remaining.rfind("\n", 0, limit)
        if split_at < limit // 2:
            split_at = remaining.rfind(" ", 0, limit)
        if split_at < limit // 2:
            split_at = limit
        chunk = remaining[:split_at].strip()
        if chunk:
            chunks.append(chunk)
        remaining = remaining[split_at:].strip()
    return chunks


def truncate_text(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    candidate = text[:limit].rstrip()
    sentence_marks = [candidate.rfind(mark) for mark in (". ", "! ", "? ", ".\n", "!\n", "?\n")]
    split_at = max(sentence_marks)
    if split_at >= int(limit * 0.65):
        return candidate[: split_at + 1].rstrip()
    split_at = candidate.rfind("\n")
    if split_at >= int(limit * 0.65):
        return candidate[:split_at].rstrip() + "…"
    split_at = candidate.rfind(" ")
    if split_at >= int(limit * 0.65):
        return candidate[:split_at].rstrip() + "…"
    return candidate.rstrip() + "…"


def extract_first_http_url(text: str) -> str:
    match = re.search(r"https?://[^\s]+", text or "")
    if not match:
        return ""
    return match.group(0).rstrip(").,]>")


def trim_history_by_chars(history: list[dict[str, str]], budget_chars: int) -> list[dict[str, str]]:
    if budget_chars <= 0:
        return []
    kept: list[dict[str, str]] = []
    used = 0
    for item in reversed(history):
        content = str(item.get("content", ""))
        role = str(item.get("role", ""))
        size = len(content) + len(role) + 8
        if used + size > budget_chars:
            break
        kept.append({"role": role, "content": content})
        used += size
    kept.reverse()
    return kept


def completion_tokens_for_plan(plan: str) -> int:
    if plan == "pro":
        return MAX_COMPLETION_TOKENS_PRO
    if plan == "start":
        return MAX_COMPLETION_TOKENS_START
    if plan == "lite":
        return MAX_COMPLETION_TOKENS_LITE
    return MAX_COMPLETION_TOKENS_FREE


def credits_for_plan(plan: str) -> int:
    return int(PLAN_CREDITS.get(plan, 0))


def public_requests_from_credits(credits: int) -> int:
    return max(0, int(credits or 0) // PUBLIC_REQUEST_UNIT_CREDITS)


def public_request_cost_from_credits(credits: int) -> int:
    credits = max(0, int(credits or 0))
    if credits <= 0:
        return 0
    return (credits + PUBLIC_REQUEST_UNIT_CREDITS - 1) // PUBLIC_REQUEST_UNIT_CREDITS


def credits_for_public_requests(requests: int) -> int:
    return max(0, int(requests or 0)) * PUBLIC_REQUEST_UNIT_CREDITS


def normalize_public_request_credit_cost(credits: int) -> int:
    return credits_for_public_requests(public_request_cost_from_credits(credits))


def request_balance_text(credits: int) -> str:
    return str(public_requests_from_credits(credits))


def request_cost_text(credits: int) -> str:
    return str(public_request_cost_from_credits(credits))


def topup_plan_code(code: str) -> str:
    return f"topup_{code}"


def topup_code_from_plan(plan: str) -> str:
    if not plan.startswith("topup_"):
        return ""
    return plan.split("_", 1)[1]


def is_topup_plan(plan: str) -> bool:
    return topup_code_from_plan(plan) in TOPUP_PACKS


def topup_spec(code: str) -> dict[str, Any] | None:
    return TOPUP_PACKS.get(code)


def text_credit_cost(alias: str) -> int:
    return int(MODEL_CREDIT_COSTS.get(alias, CREDIT_COST_GPT))


def image_credit_cost() -> int:
    return normalize_public_request_credit_cost(CREDIT_COST_IMAGE)


def image_edit_credit_cost() -> int:
    return normalize_public_request_credit_cost(CREDIT_COST_IMAGE_EDIT)


def file_request_cost(kind: str) -> int:
    requests_map = {
        "doc": FILE_DOC_REQUEST_COST,
        "ppt": FILE_PPT_REQUEST_COST,
        "sheet": FILE_SHEET_REQUEST_COST,
    }
    return credits_for_public_requests(requests_map.get(kind, FILE_DOC_REQUEST_COST))


def text_var_credits_per_1k(alias: str) -> int:
    return int(MODEL_VAR_CREDITS_PER_1K.get(alias, 0))


def estimate_tokens_from_messages(messages: list[dict[str, Any]]) -> int:
    total_chars = 0
    for msg in messages:
        content = msg.get("content")
        if isinstance(content, str):
            total_chars += len(content)
            continue
        if isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    text = item.get("text")
                    if isinstance(text, str):
                        total_chars += len(text)
    # Safe-side approximation: 1 token ~= 3 chars for mixed RU/EN chats.
    return max(1, (total_chars + 2) // 3)


def variable_text_credits(alias: str, total_tokens: int) -> int:
    rate = text_var_credits_per_1k(alias)
    if rate <= 0 or total_tokens <= 0:
        return 0
    value = (total_tokens * rate + 999) // 1000
    if MAX_VARIABLE_CREDITS_PER_TEXT > 0:
        value = min(value, MAX_VARIABLE_CREDITS_PER_TEXT)
    return max(0, value)


def parse_usage_tokens(data: dict[str, Any]) -> tuple[int, int, int]:
    usage = data.get("usage")
    if not isinstance(usage, dict):
        return 0, 0, 0
    prompt_tokens = int(usage.get("prompt_tokens", usage.get("input_tokens", 0)) or 0)
    completion_tokens = int(usage.get("completion_tokens", usage.get("output_tokens", 0)) or 0)
    total_tokens = int(usage.get("total_tokens", 0) or 0)
    if total_tokens <= 0:
        total_tokens = max(0, prompt_tokens) + max(0, completion_tokens)
    return max(0, prompt_tokens), max(0, completion_tokens), max(0, total_tokens)


def build_keyboard(chat_id: int | None = None) -> list[dict[str, Any]]:
    plan_buttons = [
        {"type": "callback", "text": "Тарифы", "payload": "action:tariffs"},
        {"type": "callback", "text": "Мой план", "payload": "action:plan"},
    ]
    if chat_id is not None and can_pick_models_for_current_preset(chat_id):
        plan_buttons.append({"type": "callback", "text": "⚙ Модели", "payload": "action:preset_models"})
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "⚡ Быстро", "payload": "set_preset:fast"},
                        {"type": "callback", "text": "⚖ Баланс", "payload": "set_preset:balanced"},
                    ],
                    [
                        {"type": "callback", "text": "🧠 Качество", "payload": "set_preset:quality"},
                        {"type": "callback", "text": "🚀 Эксперт", "payload": "set_preset:expert"},
                    ],
                    [
                        *plan_buttons,
                    ],
                    [
                        {"type": "callback", "text": "🎨 Картинка", "payload": "action:image_menu"},
                        {"type": "callback", "text": "🎁 Бонусы", "payload": "action:growth"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Сброс", "payload": "action:clear"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                    [
                        {"type": "callback", "text": "📣 Канал", "payload": "action:channel"},
                    ],
                ]
            },
        }
    ]


def build_channel_gate_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [{"type": "link", "text": "📣 Подписаться на канал", "url": channel_url_value()}],
                    [{"type": "callback", "text": "✅ Проверить подписку", "payload": "channel_gate:check"}],
                    [{"type": "callback", "text": "Помощь", "payload": "action:support"}],
                ]
            },
        }
    ]


def channel_gate_text() -> str:
    return (
        "👋 Добро пожаловать!\n\n"
        "Это AI-бот в MAX: помогает с текстами, вопросами, идеями, кодом и картинками.\n\n"
        "Чтобы пользоваться ботом, подпишись на канал проекта.\n\n"
        "1. Нажми «Подписаться на канал».\n"
        "2. Вернись сюда и нажми «Проверить подписку».\n\n"
        "Если ты уже подписан, просто нажми проверку."
    )


def channel_gate_setup_text(reason: str) -> str:
    return (
        "Канал почти подключен, но бот пока не может проверить подписку.\n\n"
        "Мы уже видим проблему и настраиваем доступ. Попробуй нажать «Проверить подписку» чуть позже.\n\n"
        f"Техническая причина: {reason}"
    )


def channel_gate_allows_payload(payload: str) -> bool:
    return payload in {"channel_gate:check", "action:channel", "action:support"}


def channel_gate_allows_text(text: str) -> bool:
    value = text.strip()
    if not value.startswith("/"):
        return value.lower() in {"старт", "start", "начать"}
    command = value.split(maxsplit=1)[0].lower()
    return command in {"/start", "/id", "/support", "/channel"}


def build_reply_shortcuts_keyboard(chat_id: int, include_share: bool = False) -> list[dict[str, Any]]:
    row = user_profile(chat_id)
    buttons: list[list[dict[str, Any]]] = [
        [
            {"type": "callback", "text": "Меню", "payload": "reply_action:menu"},
            {"type": "callback", "text": "🎨 Картинка", "payload": "reply_action:image_menu"},
        ]
    ]
    if str(row.get("plan", "free")) == "free":
        buttons[0].append({"type": "callback", "text": "Тарифы", "payload": "reply_action:tariffs"})
    buttons.append([{"type": "callback", "text": "Сброс", "payload": "reply_action:clear"}])
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def build_growth_keyboard(chat_id: int | None = None) -> list[dict[str, Any]]:
    share_button: dict[str, Any]
    if chat_id:
        row = user_profile(chat_id)
        code = str(row.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
        share_button = {"type": "link", "text": "🔗 Поделиться", "url": max_share_url(referral_share_message_v2(code))}
    else:
        share_button = {"type": "callback", "text": "🔗 Поделиться", "payload": "growth:ref_share"}
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "👥 Мой реф-код", "payload": "growth:ref_show"},
                        share_button,
                    ],
                    [
                        {"type": "callback", "text": "🎟 Ввести реф-код", "payload": "growth:ref_enter"},
                    ],
                    [
                        {"type": "callback", "text": "📣 Канал", "payload": "action:channel"},
                        {"type": "callback", "text": "🎁 Промокод", "payload": "growth:promo_enter"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                    ],
                ]
            },
        }
    ]


def build_growth_input_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [{"type": "callback", "text": "Отмена", "payload": "growth:input_cancel"}],
                    [{"type": "callback", "text": "Меню", "payload": "action:menu"}],
                ]
            },
        }
    ]


def build_receipt_contact_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [{"type": "callback", "text": "Отмена", "payload": "payment:input_cancel"}],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_onboarding_keyboard(step: int) -> list[dict[str, Any]]:
    buttons = [[{"type": "callback", "text": "Начать!", "payload": "onboard:done"}]]
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def get_image_prefs(chat_id: int) -> dict[str, str]:
    prefs = state.image_request_prefs.get(chat_id)
    if not isinstance(prefs, dict):
        prefs = {}
    mode = str(prefs.get("mode", "")).strip().lower()
    panel = str(prefs.get("panel", "root")).strip().lower()
    style = str(prefs.get("style", DEFAULT_IMAGE_STYLE)).strip().lower()
    aspect = str(prefs.get("aspect", DEFAULT_IMAGE_ASPECT)).strip().lower()
    preset = str(prefs.get("preset", "")).strip().lower()
    edit_preset = str(prefs.get("edit_preset", "")).strip().lower()
    if mode not in {"", "generate", "edit"}:
        mode = ""
    if panel not in {"root", "scenario", "style"}:
        panel = "root"
    if style not in IMAGE_STYLE_OPTIONS:
        style = DEFAULT_IMAGE_STYLE
    if aspect not in IMAGE_ASPECT_OPTIONS:
        aspect = DEFAULT_IMAGE_ASPECT
    if preset not in IMAGE_PRESET_OPTIONS:
        preset = ""
    if edit_preset not in IMAGE_EDIT_PRESET_OPTIONS:
        edit_preset = ""
    normalized = {"mode": mode, "panel": panel, "style": style, "aspect": aspect, "preset": preset, "edit_preset": edit_preset}
    state.image_request_prefs[chat_id] = normalized
    return normalized


def set_image_mode(chat_id: int, mode: str) -> dict[str, str]:
    prefs = get_image_prefs(chat_id)
    prefs["mode"] = mode if mode in {"generate", "edit"} else ""
    prefs["panel"] = "scenario" if prefs["mode"] else "root"
    if prefs["mode"] != "generate":
        prefs["preset"] = ""
    if prefs["mode"] != "edit":
        prefs["edit_preset"] = ""
    state.image_request_prefs[chat_id] = prefs
    return prefs


def set_image_panel(chat_id: int, panel: str) -> dict[str, str]:
    prefs = get_image_prefs(chat_id)
    prefs["panel"] = panel if panel in {"root", "scenario", "style"} else "root"
    state.image_request_prefs[chat_id] = prefs
    return prefs


def clear_image_preset(chat_id: int) -> dict[str, str]:
    prefs = get_image_prefs(chat_id)
    prefs["preset"] = ""
    prefs["edit_preset"] = ""
    state.image_request_prefs[chat_id] = prefs
    return prefs


def apply_image_preset(chat_id: int, preset_key: str) -> dict[str, str]:
    prefs = get_image_prefs(chat_id)
    preset = IMAGE_PRESET_OPTIONS[preset_key]
    prefs["mode"] = "generate"
    prefs["style"] = preset["style"]
    prefs["aspect"] = preset["aspect"]
    prefs["preset"] = preset_key
    prefs["edit_preset"] = ""
    state.image_request_prefs[chat_id] = prefs
    return prefs


def apply_image_edit_preset(chat_id: int, preset_key: str) -> dict[str, str]:
    prefs = get_image_prefs(chat_id)
    preset = IMAGE_EDIT_PRESET_OPTIONS[preset_key]
    prefs["mode"] = "edit"
    prefs["style"] = preset["style"]
    prefs["aspect"] = preset["aspect"]
    prefs["preset"] = ""
    prefs["edit_preset"] = preset_key
    state.image_request_prefs[chat_id] = prefs
    return prefs


def image_preset_summary_lines(chat_id: int) -> list[str]:
    prefs = get_image_prefs(chat_id)
    lines: list[str] = []
    preset_key = prefs.get("preset", "")
    edit_preset_key = prefs.get("edit_preset", "")
    if preset_key in IMAGE_PRESET_OPTIONS:
        lines.append(f"Сценарий: {IMAGE_PRESET_OPTIONS[preset_key]['label']}")
    if edit_preset_key in IMAGE_EDIT_PRESET_OPTIONS:
        lines.append(f"По фото: {IMAGE_EDIT_PRESET_OPTIONS[edit_preset_key]['label']}")
    return lines


def build_image_menu_keyboard(chat_id: int) -> list[dict[str, Any]]:
    prefs = get_image_prefs(chat_id)
    current_mode = prefs.get("mode", "")
    current_panel = prefs.get("panel", "root")
    current_style = prefs["style"]
    current_aspect = prefs["aspect"]
    current_preset = prefs.get("preset", "")
    current_edit_preset = prefs.get("edit_preset", "")

    style_buttons: list[dict[str, Any]] = []
    for key in ("auto", "photo", "anime", "art"):
        label = IMAGE_STYLE_OPTIONS[key][0]
        prefix = "● " if key == current_style else ""
        style_buttons.append({"type": "callback", "text": f"{prefix}{label}", "payload": f"image_style:{key}"})

    aspect_buttons: list[dict[str, Any]] = []
    for key in ("square", "portrait", "landscape"):
        label = IMAGE_ASPECT_OPTIONS[key][0]
        prefix = "● " if key == current_aspect else ""
        aspect_buttons.append({"type": "callback", "text": f"{prefix}{label}", "payload": f"image_aspect:{key}"})

    preset_buttons: list[dict[str, Any]] = []
    for key in ("avatar", "art_portrait", "product", "poster"):
        label = IMAGE_PRESET_OPTIONS[key]["label"]
        prefix = "● " if key == current_preset else ""
        preset_buttons.append({"type": "callback", "text": f"{prefix}{label}", "payload": f"image_preset:{key}"})

    edit_preset_buttons: list[dict[str, Any]] = []
    for key in ("anime_ref", "background_ref", "enhance_ref", "art_ref"):
        label = IMAGE_EDIT_PRESET_OPTIONS[key]["label"]
        prefix = "● " if key == current_edit_preset else ""
        edit_preset_buttons.append({"type": "callback", "text": f"{prefix}{label}", "payload": f"image_edit_preset:{key}"})

    if not current_mode:
        buttons = [
            [
                {"type": "callback", "text": "✅ Сгенерировать", "payload": "image_mode:generate"},
                {"type": "callback", "text": "🖼 Редактировать фото", "payload": "image_mode:edit"},
            ],
            [
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
            ],
        ]
        return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]

    if current_mode == "generate":
        if current_panel == "style":
            buttons = [
                style_buttons[:2],
                style_buttons[2:],
                aspect_buttons,
                [
                    {"type": "callback", "text": "✅ Сгенерировать", "payload": "image_prompt:start"},
                ],
                [
                    {"type": "callback", "text": "◀ Назад к сценариям", "payload": "image_panel:scenario"},
                    {"type": "callback", "text": "Меню", "payload": "action:menu"},
                ],
                [
                    {"type": "callback", "text": "Помощь", "payload": "action:support"},
                ],
            ]
            return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]

        buttons = [
            preset_buttons[:2],
            preset_buttons[2:],
            [
                {"type": "callback", "text": "Без сценария", "payload": "image_preset:none"},
            ],
            [
                {"type": "callback", "text": "◀ Выбор режима", "payload": "image_mode:back"},
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
            ],
            [
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
            ],
        ]
        return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]

    if current_panel == "style":
        buttons = [
            style_buttons[:2],
            style_buttons[2:],
            aspect_buttons,
            [
                {"type": "callback", "text": "🖼 Редактировать фото", "payload": "image_ref:start"},
            ],
            [
                {"type": "callback", "text": "◀ Назад к сценариям", "payload": "image_panel:scenario"},
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
            ],
            [
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
            ],
        ]
        return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]

    buttons = [
        edit_preset_buttons[:2],
        edit_preset_buttons[2:],
        [
            {"type": "callback", "text": "Без сценария", "payload": "image_edit_preset:none"},
        ],
        [
            {"type": "callback", "text": "◀ Выбор режима", "payload": "image_mode:back"},
            {"type": "callback", "text": "Меню", "payload": "action:menu"},
        ],
        [
            {"type": "callback", "text": "Помощь", "payload": "action:support"},
        ],
    ]
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def build_image_prompt_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "Отмена", "payload": "image_prompt:cancel"},
                    ],
                ]
            },
        }
    ]


def image_availability_text(chat_id: int) -> str:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free")).strip().lower()
    if plan_name != "free":
        return ""
    if plan_name == "free":
        if not free_image_is_available(row):
            next_at = free_image_next_available_at(row)
            return (
                f"На тарифе Free: не более 1 действия с картинкой каждые 7 дней. "
                f"Осталось {format_remaining_time(next_at)}. "
                f"Новая будет доступна с {format_msk_datetime(next_at)}."
        )
        return "На тарифе Free: доступно 1 действие с картинкой каждые 7 дней."


def build_image_menu_text(chat_id: int) -> str:
    prefs = get_image_prefs(chat_id)
    row = user_profile(chat_id)
    is_free_plan = str(row.get("plan", "free")).strip().lower() == "free"
    mode = prefs.get("mode", "")
    panel = prefs.get("panel", "root")
    availability = image_availability_text(chat_id)
    availability_block = f"{availability}\n" if availability else ""
    edit_cost_line = "" if is_free_plan else f"Редактирование фото: {request_cost_text(image_edit_credit_cost())} запросов.\n\n"
    generation_line = (
        "Лимит Free: 1 генерация или редактирование фото в 7 дней с момента последнего использования.\n"
        if is_free_plan
        else f"Генерация: {request_cost_text(image_credit_cost())} запросов.\n"
    )
    if not mode:
        return (
            "🎨 Картинки\n\n"
            f"{availability_block}"
            f"{generation_line}"
            f"{edit_cost_line}"
            "Сначала выбери, что хочешь сделать:\n"
            "• сгенерировать новую картинку\n"
            "• изменить фото"
        )
    if mode == "generate":
        if panel == "style":
            return (
                "🎛 Стиль и формат\n\n"
                "Здесь можно настроить внешний вид картинки вручную.\n\n"
                f"{image_params_summary(chat_id)}"
            )
        return (
            "✨ Сценарии генерации\n\n"
            "Выбери готовый сценарий или не выбирай его.\n"
            "После этого откроется шаг со стилем и форматом.\n\n"
            f"{availability_block}"
            + (
                "Лимит Free: 1 генерация или редактирование фото в 7 дней с момента последнего использования."
                if is_free_plan
                else f"Стоимость: {request_cost_text(image_credit_cost())} запросов."
            )
        )
    if panel == "style":
        return (
            "🎛 Стиль и формат фото\n\n"
            "Здесь можно вручную выбрать стиль результата и формат кадра.\n\n"
            f"{image_params_summary(chat_id)}"
        )
    edit_intro_line = (
        "Лимит Free: 1 генерация или редактирование фото в 7 дней с момента последнего использования.\n"
        if is_free_plan
        else f"Стоимость: {request_cost_text(image_edit_credit_cost())} запросов.\n"
    )
    availability_line = "" if is_free_plan else f"Доступно с тарифа {DEFAULT_IMAGE_MODEL.min_plan}."
    return (
        "🖼 Сценарии для фото\n\n"
        "Выбери готовый сценарий для редактирования фото или не выбирай его.\n"
        "После этого откроется шаг со стилем и форматом.\n\n"
        f"{edit_intro_line}"
        f"{availability_line}"
    )


def image_params_summary(chat_id: int) -> str:
    prefs = get_image_prefs(chat_id)
    style_label = IMAGE_STYLE_OPTIONS[prefs["style"]][0]
    aspect_label = IMAGE_ASPECT_OPTIONS[prefs["aspect"]][0]
    lines = image_preset_summary_lines(chat_id)
    lines.append(f"Стиль: {style_label}")
    lines.append(f"Формат: {aspect_label}")
    return "\n".join(lines)


def build_image_prompt(user_text: str, chat_id: int) -> str:
    prefs = get_image_prefs(chat_id)
    style_instruction = IMAGE_STYLE_OPTIONS[prefs["style"]][1]
    aspect_instruction = IMAGE_ASPECT_OPTIONS[prefs["aspect"]][1]
    preset_instruction = ""
    preset_key = prefs.get("preset", "")
    edit_preset_key = prefs.get("edit_preset", "")
    if preset_key in IMAGE_PRESET_OPTIONS:
        preset_instruction = IMAGE_PRESET_OPTIONS[preset_key]["prompt"]
    elif edit_preset_key in IMAGE_EDIT_PRESET_OPTIONS:
        preset_instruction = IMAGE_EDIT_PRESET_OPTIONS[edit_preset_key]["prompt"]
    instructions = [part for part in (style_instruction, aspect_instruction, preset_instruction) if part]
    if not instructions:
        return user_text
    return f"{user_text}\n\nStyle constraints: {', '.join(instructions)}."


def should_intercept_image_flow_text(chat_id: int, text: str) -> bool:
    if not text or text.strip().startswith("/"):
        return False
    if chat_id in state.pending_image_prompt or chat_id in state.pending_image_ref_prompt:
        return False
    if state.ui_current_page.get(chat_id) != UI_PAGE_IMAGE_MENU:
        return False
    prefs = get_image_prefs(chat_id)
    return prefs.get("panel") == "style" and prefs.get("mode") in {"generate", "edit"}


def build_quick_topup_keyboard(code: str) -> list[dict[str, Any]]:
    pack = topup_spec(code) or topup_spec("medium")
    payload_code = code if topup_spec(code) else "medium"
    label = str(pack.get("label", "Medium")) if pack else "Medium"
    price_rub = int(pack.get("price_rub", TOPUP_MEDIUM_PRICE_RUB)) if pack else TOPUP_MEDIUM_PRICE_RUB
    credits = int(pack.get("credits", TOPUP_MEDIUM_CREDITS)) if pack else TOPUP_MEDIUM_CREDITS
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {
                            "type": "callback",
                            "text": f"⚡ Быстро докупить {label} ({request_balance_text(credits)} запросов / {price_rub}₽)",
                            "payload": f"topup_quick:{payload_code}",
                        }
                    ],
                    [
                        {"type": "callback", "text": "⭐ Все пакеты", "payload": "action:topups"},
                        {"type": "callback", "text": "Тарифы", "payload": "action:tariffs"},
                    ],
                ]
            },
        }
    ]


def build_consent_keyboard(plan: str) -> list[dict[str, Any]]:
    amount, days = plan_price_and_days(plan)
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {
                            "type": "callback",
                            "text": f"Согласен: {amount}₽ каждые {days} дн",
                            "payload": f"buy_consent:{plan}",
                        },
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_tariffs_keyboard_pricing() -> list[dict[str, Any]]:
    buy_row_1: list[dict[str, Any]] = [
        {"type": "callback", "text": f"🍬 Lite {LITE_PLAN_PRICE_RUB}₽", "payload": "buy:lite"},
        {"type": "callback", "text": f"👌 Start {START_PLAN_PRICE_RUB}₽", "payload": "buy:start"},
    ]
    buy_row_2: list[dict[str, Any]] = [
        {"type": "callback", "text": f"🚀 Pro {PRO_PLAN_PRICE_RUB}₽", "payload": "buy:pro"},
    ]
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    buy_row_1,
                    buy_row_2,
                    [
                        {"type": "callback", "text": "⭐ Пакеты запросов", "payload": "action:topups"},
                        {"type": "callback", "text": "Мои оплаты", "payload": "action:payments"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_topups_keyboard() -> list[dict[str, Any]]:
    small = TOPUP_PACKS["small"]
    medium = TOPUP_PACKS["medium"]
    large = TOPUP_PACKS["large"]
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {
                            "type": "callback",
                            "text": f"🪙 Small {request_balance_text(int(small['credits']))} запросов • {small['price_rub']}₽",
                            "payload": "topup:small",
                        },
                    ],
                    [
                        {
                            "type": "callback",
                            "text": f"💎 Medium {request_balance_text(int(medium['credits']))} запросов • {medium['price_rub']}₽",
                            "payload": "topup:medium",
                        },
                    ],
                    [
                        {
                            "type": "callback",
                            "text": f"🚀 Large {request_balance_text(int(large['credits']))} запросов • {large['price_rub']}₽",
                            "payload": "topup:large",
                        },
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_topup_consent_keyboard(code: str) -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "Подтверждаю покупку", "payload": f"topup_consent:{code}"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_plan_keyboard(row: dict[str, Any]) -> list[dict[str, Any]]:
    keyboard = build_keyboard()
    if row.get("plan") != "free" and recurring_enabled_for_row(row):
        keyboard[0]["payload"]["buttons"].append(
            [
                {"type": "callback", "text": "Отменить подписку", "payload": "sub_cancel:start"},
            ]
        )
    return keyboard


def build_cancel_subscription_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "Подтвердить отмену", "payload": "sub_cancel:confirm"},
                    ],
                    [
                        {"type": "callback", "text": "Не отменять", "payload": "sub_cancel:back"},
                    ],
                ]
            },
        }
    ]


def plan_access_human(min_plan: str) -> str:
    if min_plan == "free":
        return "free/lite/start/pro"
    if min_plan == "lite":
        return "lite/start/pro"
    if min_plan == "start":
        return "start/pro"
    if min_plan == "pro":
        return "pro"
    return f"{min_plan}+"


def model_line(model: ModelInfo, include_prices: bool) -> str:
    lines = [
        f"{model.alias} — {model.label} ({model.provider})",
        f"версия: {model.version}, доступ: {plan_access_human(model.min_plan)}",
        f"для чего: {model.description}",
    ]
    if model.kind == "text" and model.alias in MODEL_CREDIT_COSTS:
        fixed_requests = public_request_cost_from_credits(text_credit_cost(model.alias))
        if text_var_credits_per_1k(model.alias) > 0:
            lines.append(f"списание: обычно {fixed_requests} запрос, длинный ответ до {fixed_requests + 1}")
        else:
            lines.append(f"списание: {fixed_requests} запрос")
    if model.kind == "image":
        lines.append(
            f"списание: {request_cost_text(image_credit_cost())} запросов/картинка, "
            f"{request_cost_text(image_edit_credit_cost())} запросов/редактирование фото"
        )
    if include_prices:
        lines.append(f"цена: in ${model.input_price_usd_per_m}/M, out ${model.output_price_usd_per_m}/M")
    return "\n".join(lines)


def build_models_text(user_plan: str, include_prices: bool = False) -> str:
    lines = [f"Текстовые модели (твой план: {user_plan}):"]
    start_gpt54_limit = PLAN_CONFIGS["start"].daily_gpt54_limit
    for alias, model in TEXT_MODELS.items():
        if alias == "gpt54" and user_plan == "start" and start_gpt54_limit > 0:
            prefix = f"✅ до {start_gpt54_limit}/день"
        else:
            prefix = "✅" if text_model_allowed_for_plan(user_plan, alias) else f"нужно {plan_access_human(model.min_plan)}"
        lines.append(f"\n[{prefix}]\n{model_line(model, include_prices)}")
    image_model = DEFAULT_IMAGE_MODEL
    if user_plan == "free":
        image_prefix = "✅ 1 раз в 7 дней"
    else:
        image_prefix = "✅" if plan_allowed(user_plan, image_model.min_plan) else f"нужно {plan_access_human(image_model.min_plan)}"
    lines.append("\nКартинки:")
    lines.append(f"\n[{image_prefix}]\n{model_line(image_model, include_prices)}")
    return "\n".join(lines)


def build_update_fingerprint(update: dict[str, Any]) -> str:
    try:
        payload = json.dumps(update, sort_keys=True, ensure_ascii=False).encode("utf-8")
    except Exception:
        payload = repr(update).encode("utf-8", errors="ignore")
    return hashlib.sha1(payload).hexdigest()


def remember_update(update: dict[str, Any]) -> bool:
    fingerprint = build_update_fingerprint(update)
    if fingerprint in state.processed_lookup:
        return False
    if not state.user_store.remember_processed_update(fingerprint):
        return False
    state.processed_updates.append(fingerprint)
    state.processed_lookup.add(fingerprint)
    while len(state.processed_updates) > DEDUP_CACHE_SIZE:
        old = state.processed_updates.popleft()
        state.processed_lookup.discard(old)
    return True


def image_extension(mime_type: str) -> str:
    return {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/webp": "webp",
    }.get(mime_type, "png")


def decode_data_url(data_url: str) -> ImageResult:
    header, encoded = data_url.split(",", 1)
    mime_type = "image/png"
    if ";" in header and ":" in header:
        mime_type = header.split(":", 1)[1].split(";", 1)[0]
    return ImageResult(image_bytes=base64.b64decode(encoded), mime_type=mime_type)


def encode_data_url(image: ImageResult) -> str:
    encoded = base64.b64encode(image.image_bytes).decode("ascii")
    return f"data:{image.mime_type};base64,{encoded}"


def parse_incoming_text(update: dict[str, Any]) -> tuple[int | None, str]:
    message = update.get("message") or {}
    recipient = message.get("recipient") or {}
    body = message.get("body") or {}
    chat_id = update.get("chat_id")
    if not isinstance(chat_id, int):
        chat_id = recipient.get("chat_id")
    text = body.get("text")
    if not isinstance(chat_id, int):
        return None, ""
    if not isinstance(text, str):
        return chat_id, ""
    return chat_id, text.strip()


def _walk_for_string_value(node: Any, keys: set[str]) -> str:
    if isinstance(node, dict):
        for key, value in node.items():
            if key in keys and isinstance(value, str) and value.strip():
                return value.strip()
        for value in node.values():
            found = _walk_for_string_value(value, keys)
            if found:
                return found
    elif isinstance(node, list):
        for item in node:
            found = _walk_for_string_value(item, keys)
            if found:
                return found
    return ""


def is_channel_update(update: dict[str, Any]) -> bool:
    chat_type = _walk_for_string_value(update, {"type", "chat_type", "chatType"}).lower()
    if chat_type == "channel":
        return True

    callback = update.get("callback")
    if isinstance(callback, dict):
        cb_message = callback.get("message")
        if isinstance(cb_message, dict):
            recipient = cb_message.get("recipient")
            if isinstance(recipient, dict):
                if str(recipient.get("type") or "").lower() == "channel":
                    return True
                if int(recipient.get("chat_id") or 0) < 0 and not recipient.get("dialog_with_user"):
                    return True
            cb_chat = cb_message.get("chat")
            if isinstance(cb_chat, dict):
                if str(cb_chat.get("type") or cb_chat.get("chat_type") or cb_chat.get("chatType") or "").lower() == "channel":
                    return True
                if int(cb_chat.get("chat_id") or 0) < 0 and not cb_chat.get("dialog_with_user"):
                    return True

    message = update.get("message")
    if isinstance(message, dict):
        recipient = message.get("recipient")
        if isinstance(recipient, dict):
            if str(recipient.get("type") or "").lower() == "channel":
                return True
            if int(recipient.get("chat_id") or 0) < 0 and not recipient.get("dialog_with_user"):
                return True
    chat = update.get("chat")
    if isinstance(chat, dict):
        if str(chat.get("type") or chat.get("chat_type") or chat.get("chatType") or "").lower() == "channel":
            return True
        if int(chat.get("chat_id") or 0) < 0 and not chat.get("dialog_with_user"):
            return True
    return False


def _extract_nested_int(node: Any, *path: str) -> int | None:
    current = node
    for key in path:
        if not isinstance(current, dict):
            return None
        current = current.get(key)
    return current if isinstance(current, int) else None


def parse_actor_user_id(update: dict[str, Any]) -> int | None:
    candidates = [
        _extract_nested_int(update, "user", "user_id"),
        _extract_nested_int(update, "sender", "user_id"),
        _extract_nested_int(update, "chat", "dialog_with_user", "user_id"),
        _extract_nested_int(update, "message", "sender", "user_id"),
        _extract_nested_int(update, "callback", "user", "user_id"),
        _extract_nested_int(update, "callback", "sender", "user_id"),
        _extract_nested_int(update, "callback", "initiator", "user_id"),
        _extract_nested_int(update, "callback", "message", "recipient", "dialog_with_user", "user_id"),
    ]
    for candidate in candidates:
        if isinstance(candidate, int) and candidate > 0:
            return candidate
    return None


def _walk_for_image_urls(node: Any, found: list[str]) -> None:
    if isinstance(node, dict):
        node_type = str(node.get("type", "")).lower()
        payload = node.get("payload")
        if node_type in {"image", "photo"} and isinstance(payload, dict):
            for key in ("url", "image_url", "imageUrl", "file_url", "src"):
                value = payload.get(key)
                if isinstance(value, str) and value.startswith(("http://", "https://")):
                    found.append(value)
            return

        # Explicit image fields in alternative schemas.
        for key in ("image_url", "imageUrl", "photo_url"):
            value = node.get(key)
            if isinstance(value, str) and value.startswith(("http://", "https://")):
                found.append(value)
                return

        attachments = node.get("attachments")
        if isinstance(attachments, list):
            for item in attachments:
                _walk_for_image_urls(item, found)
            return

        # Fallback for nested image payloads only.
        if "photo" in node:
            _walk_for_image_urls(node.get("photo"), found)
            return

        if "image" in node:
            _walk_for_image_urls(node.get("image"), found)
            return
    elif isinstance(node, list):
        for item in node:
            _walk_for_image_urls(item, found)


def parse_incoming_image_url(update: dict[str, Any]) -> str:
    found: list[str] = []
    message = update.get("message") if isinstance(update, dict) else None
    if isinstance(message, dict):
        _walk_for_image_urls(message.get("body"), found)
        _walk_for_image_urls(message.get("attachments"), found)
    else:
        _walk_for_image_urls(update, found)
    for url in found:
        low = url.lower()
        if any(ext in low for ext in (".jpg", ".jpeg", ".png", ".webp", "/image", "content-type=image")):
            return url
    return found[0] if found else ""


def parse_callback_source_mid(update: dict[str, Any]) -> str | None:
    callback = update.get("callback") if isinstance(update, dict) else None
    if not isinstance(callback, dict):
        return None
    message = callback.get("message")
    if not isinstance(message, dict):
        return None

    candidates: list[Any] = [
        message.get("mid"),
        (message.get("body") or {}).get("mid") if isinstance(message.get("body"), dict) else None,
        message.get("message_id"),
        (message.get("body") or {}).get("message_id") if isinstance(message.get("body"), dict) else None,
    ]
    for candidate in candidates:
        if isinstance(candidate, (int, str)) and str(candidate).strip():
            return str(candidate).strip()
    return None


def parse_callback_payload(update: dict[str, Any]) -> tuple[int | None, str | None, str | None, str | None]:
    chat_id, _ = parse_incoming_text(update)
    callback = update.get("callback") or {}
    callback_id = callback.get("callback_id")
    payload = callback.get("payload")
    source_mid = parse_callback_source_mid(update)
    if isinstance(payload, dict):
        payload = payload.get("value") or payload.get("payload") or payload.get("data")
    if not isinstance(callback_id, str):
        callback_id = None
    if not isinstance(payload, str):
        payload = None
    return chat_id, callback_id, payload, source_mid


def is_supported_update(update: dict[str, Any]) -> bool:
    return update.get("update_type") in {
        "message_created",
        "message_callback",
        "bot_started",
        "user_added",
        "bot_added",
    }


def parse_admin_target(value: str) -> int | None:
    value = value.strip()
    if value.startswith("@"):
        value = value[1:]
    if value.isdigit():
        return int(value)
    return None


def normalize_receipt_phone(raw: str) -> str:
    value = raw.strip()
    has_plus = value.startswith("+")
    digits = "".join(ch for ch in value if ch.isdigit())
    if not (10 <= len(digits) <= 15):
        return ""
    return f"+{digits}" if has_plus else digits


def normalize_receipt_email(raw: str) -> str:
    value = raw.strip().lower()
    if not value:
        return ""
    if len(value) > 254:
        return ""
    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", value):
        return ""
    return value


def parse_receipt_contact(raw: str) -> tuple[str, str]:
    text = raw.strip()
    if not text:
        return "", ""
    if "@" in text:
        email = normalize_receipt_email(text)
        return email, ""
    phone = normalize_receipt_phone(text)
    return "", phone


def parse_iso_datetime(value: str) -> datetime | None:
    if not value:
        return None
    try:
        return datetime.fromisoformat(value)
    except Exception:
        return None


def format_msk_datetime(value: datetime | None) -> str:
    if value is None:
        return "-"
    return (value + timedelta(hours=3)).strftime("%Y-%m-%d %H:%M МСК")


def format_remaining_time(target: datetime | None) -> str:
    if target is None:
        return "0 ч."
    delta = target - datetime.utcnow()
    total_seconds = max(0, int(delta.total_seconds()))
    if total_seconds <= 0:
        return "0 ч."
    days, rem = divmod(total_seconds, 86400)
    hours, rem = divmod(rem, 3600)
    minutes = rem // 60
    parts: list[str] = []
    if days > 0:
        parts.append(f"{days} д.")
    if hours > 0 or days > 0:
        parts.append(f"{hours} ч.")
    elif minutes > 0:
        parts.append(f"{minutes} мин.")
    else:
        parts.append("1 мин.")
    return " ".join(parts[:2])


def free_image_next_available_at(row: dict[str, Any]) -> datetime | None:
    last_used = parse_iso_datetime(str(row.get("free_image_last_used_at", "") or ""))
    if last_used is None:
        return None
    return last_used + timedelta(days=7)


def free_image_is_available(row: dict[str, Any]) -> bool:
    next_at = free_image_next_available_at(row)
    return next_at is None or datetime.utcnow() >= next_at


def free_file_next_available_at(row: dict[str, Any]) -> datetime | None:
    last_used = parse_iso_datetime(str(row.get("free_file_last_used_at", "") or ""))
    if not last_used:
        return None
    return last_used + timedelta(days=FREE_FILE_COOLDOWN_DAYS)


def free_file_is_available(row: dict[str, Any]) -> bool:
    next_at = free_file_next_available_at(row)
    return next_at is None or datetime.utcnow() >= next_at


def plan_price_and_days(plan: str) -> tuple[int, int]:
    if plan == "lite":
        return LITE_PLAN_PRICE_RUB, LITE_PLAN_DAYS
    if plan == "start":
        return START_PLAN_PRICE_RUB, START_PLAN_DAYS
    if plan == "pro":
        return PRO_PLAN_PRICE_RUB, PRO_PLAN_DAYS
    return 0, 0


def build_tariffs_text() -> str:
    start_cfg = PLAN_CONFIGS["start"]
    pro_cfg = PLAN_CONFIGS["pro"]
    start_gpt54_line = f"+ GPT-5.4 в «Эксперт» до {start_cfg.daily_gpt54_limit}/день" if start_cfg.daily_gpt54_limit > 0 else "как Lite"
    pro_gpt54_line = f"+ GPT-5.4 без дневного лимита" if pro_cfg.daily_gpt54_limit <= 0 else f"+ GPT-5.4 до {pro_cfg.daily_gpt54_limit}/день"
    free_requests = public_requests_from_credits(FREE_DAILY_CREDITS)
    lite_requests = public_requests_from_credits(credits_for_plan("lite"))
    start_requests = public_requests_from_credits(credits_for_plan("start"))
    pro_requests = public_requests_from_credits(credits_for_plan("pro"))
    gpt54_fixed_requests = request_cost_text(text_credit_cost("gpt54"))
    gpt54_long_requests = request_cost_text(text_credit_cost("gpt54") + MAX_VARIABLE_CREDITS_PER_TEXT)
    return (
        "💠 Тарифы:\n"
        f"• 🆓 **Free (бесплатный)**: {free_requests} запросов/день + 1 генерация или редактирование фото / 7 дней\n"
        f"• 🍬 **Lite**: {LITE_PLAN_PRICE_RUB} ₽ / {LITE_PLAN_DAYS} дней, {lite_requests} запросов\n"
        f"• 👌 **Start**: {START_PLAN_PRICE_RUB} ₽ / {START_PLAN_DAYS} дней, {start_requests} запросов\n"
        f"• 🚀 **Pro**: {PRO_PLAN_PRICE_RUB} ₽ / {PRO_PLAN_DAYS} дней, {pro_requests} запросов\n\n"
        "Модели по тарифам:\n"
        "• **Free**: DeepSeek V4 Flash, GPT-4.1 Nano, Gemini 2.5 Flash Image (1 раз в 7 дней)\n"
        "• **Lite**: + GPT-4o Mini, Gemini 2.5 Flash, Gemini 2.5 Flash Image\n"
        f"• **Start**: {start_gpt54_line}\n"
        f"• **Pro**: {pro_gpt54_line}\n\n"
        "🪙 Обычно списывается:\n"
        "• Текст: 1 запрос\n"
        "• Длинный ответ: до 2 запросов\n"
        f"• GPT-5.4: {gpt54_fixed_requests}-{gpt54_long_requests} запросов\n"
        f"• Картинка: {request_cost_text(image_credit_cost())} запросов\n"
        f"• Редактирование фото: {request_cost_text(image_edit_credit_cost())} запросов\n\n"
        "Для платных тарифов действует автопродление.\n"
        "Перед оплатой мы отдельно попросим согласие с суммой и периодичностью.\n"
        "Отменить автопродление можно в разделе «Мой план»."
    )


def recurring_terms_for_plan(plan: str) -> str:
    amount, days = plan_price_and_days(plan)
    return (
        f"Подписка {plan}: {amount} ₽ каждые {days} дней. "
        "Пользователь подтверждает регулярные списания до отмены. "
        "Отменить автопродление можно в разделе «Мой план»."
    )


async def get_session() -> aiohttp.ClientSession:
    if state.session is None or state.session.closed:
        timeout = aiohttp.ClientTimeout(total=180)
        connector_limit = max(
            16,
            MAX_API_CONCURRENCY + OPENROUTER_TEXT_CONCURRENCY + OPENROUTER_IMAGE_CONCURRENCY + KIE_TEXT_CONCURRENCY + TBANK_API_CONCURRENCY + 4,
        )
        connector = aiohttp.TCPConnector(limit=connector_limit, limit_per_host=max(4, connector_limit // 2), ttl_dns_cache=300)
        state.session = aiohttp.ClientSession(timeout=timeout, connector=connector)
    return state.session


def retry_backoff_seconds(attempt: int) -> float:
    base_ms = max(100, HTTP_RETRY_BASE_MS)
    return (base_ms / 1000.0) * max(1, 2 ** max(0, attempt - 1))


def should_retry_http_status(status: int) -> bool:
    return int(status) in TRANSIENT_HTTP_STATUSES


async def http_json_request_with_retries(
    method: str,
    url: str,
    *,
    headers: dict[str, str] | None = None,
    params: Any = None,
    json_payload: dict[str, Any] | None = None,
    data: Any = None,
    semaphore: asyncio.Semaphore | None = None,
    request_name: str = "http",
) -> tuple[int, Any, str]:
    session = await get_session()
    attempts = max(1, HTTP_RETRY_ATTEMPTS)
    for attempt in range(1, attempts + 1):
        try:
            async def _perform() -> tuple[int, Any, str]:
                async with session.request(
                    method.upper(),
                    url,
                    headers=headers,
                    params=params,
                    json=json_payload,
                    data=data,
                ) as resp:
                    raw_text = await resp.text()
                    parsed: Any = raw_text
                    if raw_text:
                        with suppress(Exception):
                            parsed = json.loads(raw_text)
                    if should_retry_http_status(resp.status) and attempt < attempts:
                        raise RuntimeError(f"retryable_http_status:{resp.status}:{raw_text[:300]}")
                    return resp.status, parsed, raw_text

            if semaphore is not None:
                async with semaphore:
                    return await _perform()
            return await _perform()
        except RuntimeError as exc:
            if not str(exc).startswith("retryable_http_status:") or attempt >= attempts:
                raise
            log.warning("%s retry %s/%s after transient status: %s", request_name, attempt, attempts, exc)
        except (aiohttp.ClientError, asyncio.TimeoutError) as exc:
            if attempt >= attempts:
                raise RuntimeError(f"{request_name} transport error: {exc}") from exc
            log.warning("%s retry %s/%s after transport error: %s", request_name, attempt, attempts, exc)
        await asyncio.sleep(retry_backoff_seconds(attempt))
    raise RuntimeError(f"{request_name} failed after retries")


async def http_bytes_request_with_retries(
    method: str,
    url: str,
    *,
    headers: dict[str, str] | None = None,
    params: Any = None,
    semaphore: asyncio.Semaphore | None = None,
    request_name: str = "http-bytes",
) -> tuple[int, bytes, str]:
    session = await get_session()
    attempts = max(1, HTTP_RETRY_ATTEMPTS)
    for attempt in range(1, attempts + 1):
        try:
            async def _perform() -> tuple[int, bytes, str]:
                async with session.request(
                    method.upper(),
                    url,
                    headers=headers,
                    params=params,
                ) as resp:
                    body = await resp.read()
                    mime_type = resp.headers.get("Content-Type", "application/octet-stream").split(";", 1)[0]
                    if should_retry_http_status(resp.status) and attempt < attempts:
                        preview = body[:300].decode("utf-8", errors="ignore")
                        raise RuntimeError(f"retryable_http_status:{resp.status}:{preview}")
                    return resp.status, body, mime_type

            if semaphore is not None:
                async with semaphore:
                    return await _perform()
            return await _perform()
        except RuntimeError as exc:
            if not str(exc).startswith("retryable_http_status:") or attempt >= attempts:
                raise
            log.warning("%s retry %s/%s after transient status: %s", request_name, attempt, attempts, exc)
        except (aiohttp.ClientError, asyncio.TimeoutError) as exc:
            if attempt >= attempts:
                raise RuntimeError(f"{request_name} transport error: {exc}") from exc
            log.warning("%s retry %s/%s after transport error: %s", request_name, attempt, attempts, exc)
        await asyncio.sleep(retry_backoff_seconds(attempt))
    raise RuntimeError(f"{request_name} failed after retries")


def user_profile(chat_id: int) -> dict[str, Any]:
    row = state.user_store.get_or_create_user(chat_id, best_default_alias_for_plan("free"))
    plan = row["plan"]
    expires_at = parse_iso_datetime(row.get("subscription_expires_at", ""))
    if plan != "free" and (expires_at is None or expires_at <= datetime.utcnow()):
        state.user_store.set_plan(chat_id, "free")
        state.user_store.set_selected_model(chat_id, best_default_alias_for_plan("free"))
        row = state.user_store.get_or_create_user(chat_id, best_default_alias_for_plan("free"))
        plan = row["plan"]

    selected = row["selected_model_alias"] or best_default_alias_for_plan(plan)
    info = TEXT_MODELS.get(selected)
    if info is None or not plan_allowed(plan, info.min_plan):
        selected = best_default_alias_for_plan(plan)
        state.user_store.set_selected_model(chat_id, selected)
        state.user_store.set_selected_preset(chat_id, "")
        row = state.user_store.get_or_create_user(chat_id, selected)
    return row


def usage_text(row: dict[str, Any]) -> str:
    plan_name = row["plan"]
    cfg = PLAN_CONFIGS[plan_name]
    gpt54_used = int(row.get("daily_gpt54_used", 0) or 0)
    gpt54_left = max(0, cfg.daily_gpt54_limit - gpt54_used)
    expires_raw = row.get("subscription_expires_at", "")
    expires_at = parse_iso_datetime(expires_raw)
    expires_text = "-"
    if plan_name != "free":
        expires_text = expires_at.strftime("%Y-%m-%d %H:%M UTC") if expires_at else "не задан"
    balance = int(row.get("credits_balance", 0) or 0)
    bonus_total, bonus_expires = state.user_store.active_bonus_credits_summary(int(row.get("chat_id", 0) or 0))
    text = (
        f"План: {plan_name}\n"
        f"Подписка до: {expires_text}\n"
        f"Запросы: {request_balance_text(balance)}"
    )
    if bonus_total > 0 and bonus_expires:
        bonus_dt = parse_iso_datetime(bonus_expires)
        bonus_until = bonus_dt.strftime("%Y-%m-%d %H:%M UTC") if bonus_dt else bonus_expires
        text += f"\n🎁 Временный бонус: {request_balance_text(bonus_total)} запросов (сгорит {bonus_until})"
    if plan_name == "free":
        text += f"\nДневной бонус free: {request_balance_text(FREE_DAILY_CREDITS)} запросов"
        next_at = free_image_next_available_at(row)
        if free_image_is_available(row):
            text += "\nКартинки на free: 1/1 доступна сейчас"
        else:
            text += (
                f"\nКартинки на free: лимит исчерпан. "
                f"Осталось {format_remaining_time(next_at)}. "
                f"Новая будет доступна с {format_msk_datetime(next_at)}."
            )
    if cfg.daily_gpt54_limit > 0:
        text += f"\nGPT-5.4 сегодня: {gpt54_used}/{cfg.daily_gpt54_limit} (осталось {gpt54_left})"
    return text


def recurring_enabled_for_row(row: dict[str, Any]) -> bool:
    return int(row.get("recurring_enabled", 0) or 0) == 1


def recurring_status_text(row: dict[str, Any]) -> str:
    plan_name = str(row.get("plan", "free"))
    if plan_name == "free":
        return ""

    expires_at = parse_iso_datetime(str(row.get("subscription_expires_at", "")))
    expires_text = expires_at.strftime("%Y-%m-%d %H:%M UTC") if expires_at else "не задано"

    if recurring_enabled_for_row(row):
        return (
            "\n\nАвтопродление: включено.\n"
            f"Текущий период действует до {expires_text}.\n"
            "Если хочешь отключить списания — нажми кнопку «Отменить подписку»."
        )

    canceled_at_raw = str(row.get("recurring_canceled_at", ""))
    if not canceled_at_raw:
        return "\n\nАвтопродление: не подключено."

    cancel_from = parse_iso_datetime(str(row.get("recurring_cancel_from", "")))
    cancel_text = cancel_from.strftime("%Y-%m-%d %H:%M UTC") if cancel_from else expires_text
    return (
        "\n\nАвтопродление: отключено.\n"
        f"Подписка действует до конца оплаченного периода: до {expires_text}.\n"
        f"Отмена автопродления с {cancel_text}."
    )


def low_credits_threshold_for_plan(plan_name: str) -> int:
    if plan_name == "free":
        return max(1, LOW_CREDITS_NUDGE_THRESHOLD_FREE)
    return max(1, LOW_CREDITS_NUDGE_THRESHOLD_PAID)


def should_nudge_low_credits(row: dict[str, Any]) -> bool:
    plan_name = str(row.get("plan", "free"))
    balance = int(row.get("credits_balance", 0) or 0)
    threshold = low_credits_threshold_for_plan(plan_name)
    if balance > threshold:
        return False
    chat_id = int(row.get("chat_id", 0) or 0)
    if chat_id <= 0:
        return False
    cooldown = timedelta(hours=max(1, LOW_CREDITS_NUDGE_COOLDOWN_HOURS))
    last = state.last_low_credits_nudge_at.get(chat_id)
    if last and (datetime.utcnow() - last) < cooldown:
        return False
    return True


async def maybe_send_low_credits_nudge(chat_id: int) -> None:
    row = user_profile(chat_id)
    if not should_nudge_low_credits(row):
        return
    plan_name = str(row.get("plan", "free"))
    balance = int(row.get("credits_balance", 0) or 0)
    threshold = low_credits_threshold_for_plan(plan_name)
    state.last_low_credits_nudge_at[chat_id] = datetime.utcnow()
    balance_requests = request_balance_text(balance)
    threshold_requests = request_balance_text(threshold)
    await max_send_message(
        chat_id,
        (
            f"⚠️ Осталось мало запросов: {balance_requests} (порог уведомления: {threshold_requests}).\n"
            "Чтобы не прерывать диалог, можно быстро докупить пакет в 1 тап."
        ),
        attachments=build_quick_topup_keyboard(TOPUP_QUICK_CODE),
        notify=False,
    )


def purchase_help_keyboard_for_row(row: dict[str, Any]) -> list[dict[str, Any]]:
    plan_name = str(row.get("plan", "free"))
    if plan_name == "free":
        return build_tariffs_keyboard_pricing()
    return build_quick_topup_keyboard(TOPUP_QUICK_CODE)


def can_use_model(plan: str, model_alias: str) -> tuple[bool, str]:
    info = TEXT_MODELS.get(model_alias)
    if not info:
        return False, "Неизвестная модель. Используй /models"
    if not text_model_allowed_for_plan(plan, model_alias):
        return False, f"Модель {info.label} доступна с тарифа {info.min_plan}."
    return True, ""


def check_and_consume_credits(chat_id: int, amount: int, operation_name: str) -> tuple[bool, str]:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free")).strip().lower()
    if plan_name == "free" and operation_name.strip().lower() in {"картинка", "картинка по фото"}:
        return True, ""
    if amount <= 0:
        return True, ""
    balance = int(row.get("credits_balance", 0) or 0)
    if balance < amount:
        needed_requests = request_cost_text(amount)
        available_requests = request_balance_text(balance)
        if plan_name == "free":
            return (
                False,
                f"На сегодня free-запросы закончились для операции «{operation_name}». Нужно {needed_requests}, доступно {available_requests}. Завтра бонус обновится, либо открой «Тарифы».",
            )
        return (
            False,
            f"Недостаточно запросов для операции «{operation_name}». Нужно {needed_requests}, доступно {available_requests}. Открой «Тарифы».",
        )
    ok = state.user_store.consume_credits(chat_id, amount)
    if not ok:
        row = user_profile(chat_id)
        balance = int(row.get("credits_balance", 0) or 0)
        return (
            False,
            f"Недостаточно запросов для операции «{operation_name}». Нужно {request_cost_text(amount)}, доступно {request_balance_text(balance)}. Открой «Тарифы».",
        )
    return True, ""


def check_limit_only(chat_id: int, limit_type: str) -> tuple[bool, str]:
    row = user_profile(chat_id)
    if row["is_blocked"]:
        return False, "Доступ к боту временно ограничен. Напиши в поддержку."
    plan_name = row["plan"]
    cfg = PLAN_CONFIGS[plan_name]

    if limit_type == "messages":
        selected_alias = str(row.get("selected_model_alias") or "")
        if selected_alias == "gpt54" and cfg.daily_gpt54_limit > 0:
            gpt54_used = int(row.get("daily_gpt54_used", 0) or 0)
            if gpt54_used >= cfg.daily_gpt54_limit:
                return False, "Лимит GPT-5.4 на сегодня исчерпан. Выбери другую модель."
        return True, ""

    if limit_type == "images":
        if plan_name == "free":
            if not free_image_is_available(row):
                next_at = free_image_next_available_at(row)
                return (
                    False,
                    f"На free доступно 1 действие с картинкой каждые 7 дней. Следующая генерация или обработка фото будет доступна с {format_msk_datetime(next_at)}. "
                    "Хочешь больше — выбери тариф ниже или пакет запросов.",
                )
            return True, ""

        return True, ""

    if limit_type == "files":
        if plan_name == "free":
            if not free_file_is_available(row):
                next_at = free_file_next_available_at(row)
                return (
                    False,
                    f"На Free доступен 1 короткий файл каждые {FREE_FILE_COOLDOWN_DAYS} дней. "
                    f"Следующий документ, презентация или таблица будут доступны с {format_msk_datetime(next_at)}. "
                    "Хочешь больше — открой «Тарифы».",
                )
            return True, ""
        return True, ""

    return True, ""


def consume_limit(chat_id: int, limit_type: str) -> None:
    row = user_profile(chat_id)
    if limit_type == "messages":
        cfg = PLAN_CONFIGS[row["plan"]]
        if str(row.get("selected_model_alias") or "") == "gpt54" and cfg.daily_gpt54_limit > 0:
            state.user_store.increment_gpt54_usage(chat_id)
        return

    if limit_type == "images":
        if str(row.get("plan", "free")) == "free":
            state.user_store.increment_free_week_image_usage(chat_id)
            return
        state.user_store.increment_image_usage(chat_id)
        return

    if limit_type == "files":
        if str(row.get("plan", "free")) == "free":
            state.user_store.increment_free_file_usage(chat_id)
        return


def check_and_consume_limit(chat_id: int, limit_type: str) -> tuple[bool, str]:
    ok, reason = check_limit_only(chat_id, limit_type)
    if not ok:
        return False, reason
    consume_limit(chat_id, limit_type)
    return True, ""


def check_cooldown(chat_id: int, kind: str) -> tuple[bool, str]:
    now = datetime.utcnow()
    if kind == "message":
        last = state.last_message_at.get(chat_id)
        if last and (now - last).total_seconds() < MESSAGE_COOLDOWN_SECONDS:
            return False, f"Слишком часто. Подожди {MESSAGE_COOLDOWN_SECONDS} сек."
        state.last_message_at[chat_id] = now
        return True, ""

    last = state.last_image_at.get(chat_id)
    if last and (now - last).total_seconds() < IMAGE_COOLDOWN_SECONDS:
        return False, f"Картинки можно запрашивать раз в {IMAGE_COOLDOWN_SECONDS} сек."
    state.last_image_at[chat_id] = now
    return True, ""


async def max_send_message(
    chat_id: int,
    text: str,
    attachments: list[dict[str, Any]] | None = None,
    notify: bool = True,
    text_format: str | None = None,
) -> str | None:
    attachments = normalize_channel_link_buttons(attachments)
    chunks = split_message(text)
    if not chunks and attachments:
        chunks = [""]

    first_mid: str | None = None
    for index, chunk in enumerate(chunks):
        payload: dict[str, Any] = {
            "type": "text",
            "text": chunk,
            "notify": notify,
        }
        if text_format:
            payload["format"] = text_format
        if attachments and index == 0:
            payload["attachments"] = attachments

        status, body_json, body_text = await http_json_request_with_retries(
            "POST",
            f"{MAX_API}/messages",
            headers=max_headers(),
            params={"chat_id": str(chat_id)},
            json_payload=payload,
            semaphore=state.max_api_semaphore,
            request_name="max_send_message",
        )
        if status >= 400:
            raise RuntimeError(f"MAX send error {status}: {body_text[:500]}")
        if index == 0 and isinstance(body_json, dict):
            first_mid = extract_message_mid(body_json)
    if first_mid and attachments and any(
        isinstance(item, dict) and item.get("type") == "inline_keyboard" for item in attachments
    ):
        state.ui_message_mid[chat_id] = first_mid
    return first_mid


def extract_message_mid(node: Any) -> str | None:
    if isinstance(node, dict):
        for key in ("mid", "message_id", "messageId"):
            value = node.get(key)
            if isinstance(value, (str, int)) and str(value).strip():
                return str(value).strip()
        for key in ("message", "body", "payload", "result", "data"):
            if key in node:
                found = extract_message_mid(node.get(key))
                if found:
                    return found
    elif isinstance(node, list):
        for item in node:
            found = extract_message_mid(item)
            if found:
                return found
    return None


def is_expected_max_delivery_error(exc: Exception | str) -> bool:
    text = str(exc).lower()
    return (
        "max send error 403" in text
        and ("chat.denied" in text or "error.dialog.suspended" in text)
    )


def is_max_attachment_not_ready_error(exc: Exception | str) -> bool:
    text = str(exc).lower()
    return (
        "attachment.not.ready" in text
        or "file.not.processed" in text
        or "errors.process.attachment.file.not.processed" in text
    )


async def max_edit_message(
    chat_id: int,
    message_mid: str,
    text: str,
    attachments: list[dict[str, Any]] | None = None,
    text_format: str | None = None,
) -> bool:
    attachments = normalize_channel_link_buttons(attachments)
    payload: dict[str, Any] = {
        "type": "text",
        "text": text,
    }
    if text_format:
        payload["format"] = text_format
    if attachments is not None:
        payload["attachments"] = attachments

    status, _, body_text = await http_json_request_with_retries(
        "PUT",
        f"{MAX_API}/messages",
        headers=max_headers(),
        params={"chat_id": str(chat_id), "message_id": str(message_mid)},
        json_payload=payload,
        semaphore=state.max_api_semaphore,
        request_name="max_edit_message",
    )
    if status >= 400:
        log.warning("MAX edit error %s (chat_id=%s, mid=%s): %s", status, chat_id, message_mid, body_text[:300])
        return False
    return True


def normalize_channel_link_buttons(attachments: list[dict[str, Any]] | None) -> list[dict[str, Any]] | None:
    if not attachments:
        return attachments
    channel_url = channel_url_value()
    for attachment in attachments:
        if not isinstance(attachment, dict):
            continue
        if attachment.get("type") != "inline_keyboard":
            continue
        payload = attachment.get("payload")
        if not isinstance(payload, dict):
            continue
        rows = payload.get("buttons")
        if not isinstance(rows, list):
            continue
        for row in rows:
            if not isinstance(row, list):
                continue
            for button in row:
                if not isinstance(button, dict):
                    continue
                if button.get("type") == "callback" and str(button.get("payload", "")).strip() == "action:channel":
                    button["type"] = "link"
                    button["url"] = channel_url
                    button.pop("payload", None)
    return attachments


async def answer_callback(
    callback_id: str,
    notification: str,
    text: str | None = None,
    attachments: list[dict[str, Any]] | None = None,
) -> bool:
    payload: dict[str, Any] = {"notification": notification}
    if text is not None:
        message_payload: dict[str, Any] = {"text": text}
        if attachments:
            message_payload["attachments"] = attachments
        payload["message"] = message_payload
    status, response_json, body_text = await http_json_request_with_retries(
        "POST",
        f"{MAX_API}/answers",
        headers=max_headers(),
        params={"callback_id": callback_id},
        json_payload=payload,
        semaphore=state.max_api_semaphore,
        request_name="max_answer_callback",
    )
    if status >= 400:
        log.warning("Callback answer failed %s: %s", status, body_text[:300])
        return False
    if isinstance(response_json, dict) and response_json.get("success") is False:
        log.warning("Callback answer returned success=false: %s", response_json)
        return False
    return True


async def get_upload_url(upload_type: str = "image") -> str:
    status, body, _ = await http_json_request_with_retries(
        "POST",
        f"{MAX_API}/uploads",
        headers=max_headers(),
        params={"type": upload_type},
        semaphore=state.max_api_semaphore,
        request_name="max_get_upload_url",
    )
    if status >= 400:
        raise RuntimeError(f"MAX uploads error {status}: {body}")
    upload_url = body.get("url") if isinstance(body, dict) else None
    if not upload_url:
        raise RuntimeError(f"MAX uploads response has no url: {body}")
    return upload_url


async def upload_image_to_max(image_bytes: bytes, mime_type: str) -> dict[str, Any]:
    return await upload_binary_to_max(
        data_bytes=image_bytes,
        mime_type=mime_type,
        filename=f"generated.{image_extension(mime_type)}",
        upload_type="image",
    )


async def upload_binary_to_max(
    *,
    data_bytes: bytes,
    mime_type: str,
    filename: str,
    upload_type: str = "file",
) -> dict[str, Any]:
    upload_url = await get_upload_url(upload_type)
    form = aiohttp.FormData()
    form.add_field("data", BytesIO(data_bytes), filename=filename, content_type=mime_type)

    status, body, _ = await http_json_request_with_retries(
        "POST",
        upload_url,
        data=form,
        semaphore=state.max_api_semaphore,
        request_name="max_upload_image",
    )
    if status >= 400:
        raise RuntimeError(f"MAX file upload error {status}: {body}")
    if not isinstance(body, dict):
        raise RuntimeError(f"MAX file upload invalid response: {body}")
    return body


async def send_generated_image(chat_id: int, prompt: str, image: ImageResult, display_prompt: str | None = None) -> None:
    attachment_payload = await upload_image_to_max(image.image_bytes, image.mime_type)
    attachment = {"type": "image", "payload": attachment_payload}
    shown_prompt = display_prompt or prompt
    await max_send_message(
        chat_id,
        f"Готово. Вот картинка по запросу:\n{shown_prompt}",
        attachments=[attachment, *build_reply_shortcuts_keyboard(chat_id, include_share=True)],
    )


async def send_generated_file(
    chat_id: int,
    *,
    kind: str,
    title: str,
    file_bytes: bytes,
    filename: str,
) -> None:
    attachment_payload = await upload_binary_to_max(
        data_bytes=file_bytes,
        mime_type=FILE_KIND_MIME_TYPES.get(kind, "application/octet-stream"),
        filename=filename,
        upload_type="file",
    )
    attachment = {"type": "file", "payload": attachment_payload}
    kind_label = FILE_KIND_LABELS.get(kind, "файл")
    await max_send_message(
        chat_id,
        f"Готово. Собрал {kind_label}: {title}",
        attachments=[attachment, *build_reply_shortcuts_keyboard(chat_id)],
    )


async def send_generated_file(
    chat_id: int,
    *,
    kind: str,
    title: str,
    file_bytes: bytes,
    filename: str,
) -> None:
    attachment_payload = await upload_binary_to_max(
        data_bytes=file_bytes,
        mime_type=FILE_KIND_MIME_TYPES.get(kind, "application/octet-stream"),
        filename=filename,
        upload_type="file",
    )
    attachment = {"type": "file", "payload": attachment_payload}
    text = f"Готово. Собрал {FILE_KIND_LABELS.get(kind, 'файл')}: {title}"
    attachments = [attachment, *build_reply_shortcuts_keyboard(chat_id)]
    last_exc: Exception | None = None
    for attempt in range(5):
        try:
            await max_send_message(chat_id, text, attachments=attachments)
            return
        except Exception as exc:
            if not is_max_attachment_not_ready_error(exc) or attempt == 4:
                raise
            last_exc = exc
            await asyncio.sleep(1.2 + attempt * 0.8)
    if last_exc:
        raise last_exc


def friendly_image_generation_error(exc: Exception, *, is_edit: bool = False) -> str:
    text = str(exc).strip()
    if "Image was not returned by the selected model." in text or "Edited image was not returned by the selected model." in text:
        if is_edit:
            return "Не удалось получить готовую картинку по фото. Попробуй ещё раз чуть позже или измени описание."
        return "Не удалось получить картинку от модели. Попробуй ещё раз чуть позже или немного измени запрос."
    if is_edit:
        return f"Не удалось обработать фото: {text}" if text else "Не удалось обработать фото. Попробуй ещё раз позже."
    return f"Не удалось сгенерировать картинку: {text}" if text else "Не удалось сгенерировать картинку. Попробуй ещё раз позже."


async def fetch_image_bytes(url: str, use_max_auth: bool = False) -> ImageResult:
    headers = max_headers() if use_max_auth else None
    status, data, mime_type = await http_bytes_request_with_retries(
        "GET",
        url,
        headers=headers,
        semaphore=state.openrouter_image_semaphore,
        request_name="fetch_image_bytes",
    )
    if status >= 400:
        raise RuntimeError(f"Image fetch error {status}")
    return ImageResult(image_bytes=data, mime_type=mime_type)


def remember_reference_image(chat_id: int, data_url: str) -> None:
    state.last_reference_image_data_url[chat_id] = data_url
    state.last_reference_image_at[chat_id] = datetime.utcnow()


def get_recent_reference_image(chat_id: int) -> str:
    data_url = state.last_reference_image_data_url.get(chat_id, "")
    ts = state.last_reference_image_at.get(chat_id)
    if not data_url or not ts:
        return ""
    age = datetime.utcnow() - ts
    if age > timedelta(minutes=max(1, REFERENCE_IMAGE_TTL_MINUTES)):
        state.last_reference_image_data_url.pop(chat_id, None)
        state.last_reference_image_at.pop(chat_id, None)
        return ""
    return data_url


def image_job_already_pending(chat_id: int) -> bool:
    return chat_id in state.pending_image_jobs


async def enqueue_image_job(job: ImageJob) -> int:
    queue_size_before = state.image_job_queue.qsize()
    state.pending_image_jobs.add(job.chat_id)
    await state.image_job_queue.put(job)
    return queue_size_before


async def process_image_job(job: ImageJob) -> None:
    chat_id = int(job.chat_id)
    try:
        if job.kind == "edit":
            image = await generate_image_from_reference(job.model_prompt, job.reference_image_data_url)
            await send_generated_image(chat_id, job.model_prompt, image, display_prompt=job.user_prompt)
            final_row = user_profile(chat_id)
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="image_request",
                plan=str(final_row.get("plan", "")),
                model_alias=f"{DEFAULT_IMAGE_MODEL.alias}:edit",
                credits_spent=int(job.credits_spent),
                tokens_total=0,
                details=job.details,
            )
        else:
            image = await generate_image(job.model_prompt)
            await send_generated_image(chat_id, job.model_prompt, image, display_prompt=job.user_prompt)
            final_row = user_profile(chat_id)
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="image_request",
                plan=str(final_row.get("plan", "")),
                model_alias=DEFAULT_IMAGE_MODEL.alias,
                credits_spent=int(job.credits_spent),
                tokens_total=0,
                details=job.details,
            )
        with suppress(Exception):
            await maybe_send_low_credits_nudge(chat_id)
    except RuntimeError as exc:
        if int(job.credits_spent) > 0:
            state.user_store.refund_credits(chat_id, int(job.credits_spent))
        log.warning("Image job failed for chat_id=%s kind=%s: %s", chat_id, job.kind, exc)
        await max_send_message(chat_id, friendly_image_generation_error(exc, is_edit=job.kind == "edit"), attachments=build_image_menu_keyboard(chat_id))
    except Exception:
        if int(job.credits_spent) > 0:
            state.user_store.refund_credits(chat_id, int(job.credits_spent))
        raise
    finally:
        state.pending_image_jobs.discard(chat_id)


async def image_worker_loop(worker_index: int) -> None:
    log.info("Image worker started: #%s", worker_index)
    while True:
        try:
            job = await state.image_job_queue.get()
            try:
                await process_image_job(job)
            finally:
                state.image_job_queue.task_done()
        except asyncio.CancelledError:
            raise
        except Exception as exc:
            log.exception("Image worker loop error")
            capture_exception_safe(exc)
            record_runtime_error()
            await notify_admin_alert("image_worker", f"Image worker #{worker_index} error: {exc}")
            await asyncio.sleep(1)


def looks_like_image_ref_request(text: str) -> bool:
    value = text.strip().lower()
    if not value:
        return False
    verb = re.search(r"\b(нарисуй|перерисуй|сделай|измени|отрисуй|сгенерируй)\b", value)
    ref = re.search(r"\b(фото|фотк|картинк|из нее|из неё|ее|её|эту|этого человека)\b", value)
    return bool(verb and ref)


def build_text_request(chat_id: int, user_text: str, selected_alias: str | None = None) -> tuple[str, str, ModelInfo, list[dict[str, Any]]]:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free"))
    alias = selected_alias or row["selected_model_alias"] or best_default_alias_for_plan(row["plan"])
    model_info = TEXT_MODELS.get(alias, DEFAULT_TEXT_MODEL)
    history = trim_history_by_chars(list(state.history(chat_id)), MAX_CONTEXT_CHARS)
    messages: list[dict[str, Any]] = [{"role": "system", "content": f"{SYSTEM_PROMPT_BASE} {STYLE_PROMPTS.get(alias, '')}".strip()}]
    messages.extend(history)
    messages.append({"role": "user", "content": user_text})
    return plan_name, alias, model_info, messages


def detect_file_kind_from_text(text: str) -> str:
    value = (text or "").strip().lower()
    if not value:
        return ""
    if any(token in value for token in ("презентац", "слайды", "pptx", "ppt ")):
        return "ppt"
    if any(token in value for token in ("таблиц", "xlsx", "excel", "эксель", "смет", "бюджет", "медиаплан", "контент-план")):
        return "sheet"
    if any(token in value for token in ("документ", "docx", "word", "резюме", "коммерческ", "кп ", "инструкц", "бриф", "тз ")):
        return "doc"
    return ""


def looks_like_file_request(text: str) -> str:
    value = (text or "").strip().lower()
    if not value or value.startswith("/"):
        return ""
    kind = detect_file_kind_from_text(value)
    if not kind:
        return ""
    if re.search(r"\b(сделай|создай|подготовь|собери|оформи|сгенерируй|сформируй)\b", value):
        return kind
    return ""


def file_prompt_has_enough_detail(text: str) -> bool:
    words = [part for part in re.split(r"\s+", (text or "").strip()) if part]
    return len(words) >= 4 or len((text or "").strip()) >= 28


def default_file_profile_for_plan(plan_name: str) -> str:
    return "short" if str(plan_name or "").strip().lower() == "free" else "medium"


def detect_file_profile_from_text(text: str, default_profile: str = "medium") -> str:
    value = (text or "").strip().lower()
    if not value:
        return default_profile
    if any(token in value for token in ("кратк", "коротк", "быстро", "на 3 слайда", "на 4 слайда", "на 5 слайдов", "до 5 слайдов")):
        return "short"
    if any(token in value for token in ("средн", "обычн", "стандарт")):
        return "medium"
    if any(token in value for token in ("подроб", "полный", "развернут", "на 10 слайдов", "на 12 слайдов", "до 12 слайдов", "расширен")):
        return "full"
    return default_profile


def sanitize_filename_part(value: str, fallback: str) -> str:
    normalized = re.sub(r"[^a-zA-Z0-9а-яА-ЯёЁ_-]+", "-", (value or "").strip(), flags=re.UNICODE)
    normalized = normalized.strip("-_")
    return normalized[:60] or fallback


def extract_json_object(text: str) -> dict[str, Any] | None:
    raw = (text or "").strip()
    if not raw:
        return None
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
    try:
        loaded = json.loads(raw)
        if isinstance(loaded, dict):
            return loaded
    except Exception:
        pass
    start = raw.find("{")
    end = raw.rfind("}")
    if start >= 0 and end > start:
        try:
            loaded = json.loads(raw[start : end + 1])
            if isinstance(loaded, dict):
                return loaded
        except Exception:
            return None
    return None


def normalized_file_topic(user_prompt: str) -> str:
    topic = (user_prompt or "").strip()
    if not topic:
        return "Новый файл"
    topic = re.sub(
        r"^(подготовь|сделай|собери|оформи|создай|напиши)\s+",
        "",
        topic,
        flags=re.IGNORECASE,
    )
    topic = re.sub(
        r"^(документ|презентацию|презентация|таблицу|таблица|файл)\s+",
        "",
        topic,
        flags=re.IGNORECASE,
    )
    return topic.strip(" .:-") or (user_prompt or "").strip() or "Новый файл"


def default_file_title(kind: str, topic: str) -> str:
    topic_lc = topic.lower()
    if kind == "doc":
        if "анализ" in topic_lc and "рын" in topic_lc:
            return "Анализ рынка"
        if "коммерческ" in topic_lc or topic_lc.startswith("кп"):
            return "Коммерческое предложение"
        return "Документ"
    if kind == "ppt":
        return "Презентация"
    if kind == "sheet":
        return "Таблица"
    return "Файл"


def doc_request_mode(user_prompt: str) -> str:
    topic_lc = normalized_file_topic(user_prompt).lower()
    analysis_markers = (
        "анализ",
        "обзор",
        "исслед",
        "рын",
        "конкурент",
        "спрос",
        "ниша",
        "перспектив",
        "оценк",
    )
    if any(marker in topic_lc for marker in analysis_markers):
        return "analysis"
    return "general"


def analysis_doc_headings(profile: str) -> list[str]:
    if profile == "short":
        return [
            "Резюме",
            "Спрос и клиенты",
            "Конкуренты и альтернативы",
            "Вывод",
        ]
    if profile == "full":
        return [
            "Резюме",
            "Спрос и сегменты клиентов",
            "Конкуренты и альтернативы",
            "Локация и формат предложения",
            "Цены и экономика",
            "Риски и ограничения",
            "Итоговый вывод и рекомендации",
        ]
    return [
        "Резюме",
        "Спрос и сегменты клиентов",
        "Конкуренты и альтернативы",
        "Цены и экономика",
        "Риски и ограничения",
        "Итоговый вывод",
    ]


def build_doc_generation_rule(profile: str, user_prompt: str) -> str:
    base_rule = (
        'Верни только JSON объект вида {"title":"...", "subtitle":"...", '
        '"sections":[{"heading":"...", "paragraphs":["..."], "bullets":["..."]}], '
        '"table":{"title":"...", "columns":["..."], "rows":[["..."]]}}. '
        "sections обязателен. table добавляй только если она реально нужна. "
        "Для документа пиши содержательные абзацы и конкретные пункты, а не шаблонные заглушки вроде "
        '"Ключевая цель", "Основной результат" или "Следующий шаг". '
        "Не дублируй запрос пользователя как единственное содержание документа. "
    )
    if doc_request_mode(user_prompt) == "analysis":
        headings = ", ".join(analysis_doc_headings(profile))
        return (
            base_rule
            + "Пользователь просит именно результат анализа, а не текст о том, как его проводить. "
            "Собери документ как готовый аналитический вывод по теме. "
            "Не пиши про методику, этапы анализа, необходимость что-то еще оценить или общие заглушки. "
            "Если точных цифр нет, давай содержательные вероятностные выводы и помечай их как оценочные допущения, "
            "но все равно формулируй выводы по рынку, спросу, конкурентам, ценам, рискам и перспективе ниши. "
            f"Желательная структура разделов: {headings}. "
            "В каждом разделе должны быть именно наблюдения, интерпретация и практический вывод."
        )
    return base_rule


def placeholder_doc_sections(user_prompt: str, profile: str = "medium") -> list[dict[str, Any]]:
    topic = normalized_file_topic(user_prompt)
    topic_lc = topic.lower()
    is_market = any(token in topic_lc for token in ("анализ", "рын", "конкурент", "спрос", "ниша", "обзор"))
    is_warehouse = any(token in topic_lc for token in ("склад", "контейнер", "storage", "логист"))

    if is_market:
        sections = [
            {
                "heading": "Резюме",
                "paragraphs": [
                    f"Документ содержит предварительный анализ по теме: {topic}. "
                    "По характеру спроса ниша выглядит рабочей там, где клиенту важны быстрый запуск, "
                    "доступность и более низкий порог входа по сравнению с классическим складом."
                ],
                "bullets": [
                    "Наиболее вероятный спрос формируют микро-бизнес, стройка, e-commerce и сезонное хранение",
                    "Формат выигрывает в скорости запуска и гибкости, но уступает капитальным складам по комфорту и инфраструктуре",
                    "Наиболее реалистичный сценарий входа — пилот на понятной локации с быстрой проверкой заполняемости",
                ],
            },
            {
                "heading": "Спрос и сегменты клиентов",
                "paragraphs": [
                    "Основной спрос на такой продукт обычно возникает не у крупных складских операторов, "
                    "а у клиентов, которым нужен недорогой, быстрый и локальный формат хранения без сложного договора и длинного входа."
                ],
                "bullets": [
                    "Малый бизнес и торговля с нехваткой площади под запас",
                    "Строительные и подрядные компании для инструмента и материалов",
                    "Локальные продавцы и e-commerce для временного и сезонного хранения",
                ],
            },
            {
                "heading": "Конкуренты и альтернативы",
                "paragraphs": [
                    "Прямыми конкурентами будут другие контейнерные и модульные решения, а косвенными — гаражи, боксы, ячейки и классические небольшие склады. "
                    "Конкуренция обычно идет не за “лучший склад вообще”, а за сочетание цены, доступности и удобства доступа."
                ],
                "bullets": [
                    "Прямые конкуренты: контейнеры и модульные склады",
                    "Косвенные конкуренты: гаражи, складские боксы, мини-склады, ячейки",
                    "Ключевое преимущество контейнера — быстрый запуск и относительная дешевизна",
                ],
            },
            {
                "heading": "Цены и экономика",
                "paragraphs": [
                    "Экономика такого проекта обычно держится на недорогом входе и быстрой оборачиваемости. "
                    "Если локация подобрана удачно, контейнерный формат может зайти быстрее, чем капитальный склад, "
                    "но сильно зависит от загрузки и понятного клиентского сценария."
                ],
                "bullets": [
                    "Ключевые драйверы экономики: ставка аренды, загрузка, локация, стоимость доработки контейнера",
                    "Слабое место модели — простой объекта и низкая повторяемость спроса",
                    "На старте безопаснее проверять спрос через пилот, а не сразу через масштабирование",
                ],
            },
            {
                "heading": "Риски и ограничения",
                "paragraphs": [
                    "Главные риски лежат не в самом контейнере, а в локации, юридических ограничениях, уровне охраны и том, насколько понятен клиенту формат хранения."
                ],
                "bullets": [
                    "Недостаточный трафик и слабый спрос в конкретной точке",
                    "Непрозрачные ограничения по размещению и эксплуатации площадки",
                    "Сильное ценовое давление со стороны гаражей, боксов и небольших складов",
                ],
            },
            {
                "heading": "Выводы и рекомендации",
                "paragraphs": [
                    "Ниша выглядит жизнеспособной как локальный формат недорогого хранения, особенно там, где клиенту важны скорость входа и простой доступ. "
                    "Наиболее разумный сценарий — запуск через короткий пилот с конкретной локацией и проверкой заполняемости."
                ],
                "bullets": [
                    "Тестировать не “город в целом”, а 2-3 конкретные точки с понятным спросом",
                    "Сразу сравнивать ставку и условия с альтернативами: гаражи, боксы, мини-склады",
                    "Считать экономику через пилотный сценарий, а не через оптимистичную загрузку",
                ],
            },
        ]
        if profile == "short":
            return sections[:3] + [sections[-1]]
        if profile == "full":
            if is_warehouse:
                sections.insert(
                    4,
                    {
                        "heading": "Локация и формат предложения",
                        "paragraphs": [
                            "Для контейнеров под склад критично, чтобы формат был не просто дешевым, а реально удобным для клиента. "
                            "Именно локация, подъезд и режим доступа часто определяют, будет ли объект сдаваться стабильно."
                        ],
                        "bullets": [
                            "Подъезд легкового и грузового транспорта",
                            "Освещение, охрана и предсказуемый доступ",
                            "Понятный формат хранения для бизнеса и частного клиента",
                        ],
                    },
                )
            return sections
        return sections

    sections = [
        {
            "heading": "Контекст",
            "paragraphs": [
                f"Документ подготовлен по теме: {topic}. Ниже собрана рабочая структура, "
                "которую можно использовать как черновик и доработать под конкретную цель."
            ],
            "bullets": [],
        },
        {
            "heading": "Основные тезисы",
            "paragraphs": ["Ниже перечислены ключевые вопросы и направления, которые стоит раскрыть по теме."],
            "bullets": [
                "Что именно нужно получить на выходе",
                "Какие ограничения и вводные влияют на результат",
                "Какие данные или решения важны в первую очередь",
            ],
        },
        {
            "heading": "Практические выводы",
            "paragraphs": [
                "После уточнения деталей этот документ можно быстро превратить в полноценную рабочую версию."
            ],
            "bullets": ["Уточнить цель", "Добавить фактуру", "Согласовать следующий шаг"],
        },
    ]
    if profile == "full":
        sections.insert(
            2,
            {
                "heading": "Риски и вопросы для проверки",
                "paragraphs": ["Перед финализацией важно проверить допущения и закрыть недостающие данные."],
                "bullets": ["Проверить факты", "Согласовать критерии успеха", "Уточнить формат использования документа"],
            },
        )
    return sections


def doc_spec_from_plain_text(user_prompt: str, text: str, profile: str = "medium") -> dict[str, Any] | None:
    raw = (text or "").strip()
    if not raw:
        return None
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
    raw = raw.strip()
    if not raw:
        return None

    topic = normalized_file_topic(user_prompt)
    blocks = [block.strip() for block in re.split(r"\n\s*\n", raw) if block.strip()]
    sections: list[dict[str, Any]] = []
    heading_hints = (
        "введение",
        "обзор",
        "резюме",
        "вывод",
        "рекоменда",
        "риски",
        "спрос",
        "предлож",
        "цены",
        "эконом",
        "конкур",
        "рын",
        "сегмент",
    )

    for index, block in enumerate(blocks[:8], start=1):
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        heading = ""
        content_lines = lines
        first_line = re.sub(r"^[#>\-\*\d\.\)\s]+", "", lines[0]).strip()
        if first_line and (
            lines[0].startswith("#")
            or lines[0].endswith(":")
            or (len(first_line) <= 70 and any(first_line.lower().startswith(hint) for hint in heading_hints))
        ):
            heading = first_line.rstrip(":")
            content_lines = lines[1:]
        paragraphs: list[str] = []
        bullets: list[str] = []
        for line in content_lines:
            cleaned = line.strip()
            if not cleaned:
                continue
            if re.match(r"^[-*•]\s+", cleaned):
                bullets.append(re.sub(r"^[-*•]\s+", "", cleaned).strip())
            elif re.match(r"^\d+[\.\)]\s+", cleaned):
                bullets.append(re.sub(r"^\d+[\.\)]\s+", "", cleaned).strip())
            else:
                paragraphs.append(cleaned)
        if not heading:
            heading = "Краткий вывод" if index == 1 else f"Раздел {index}"
        if paragraphs or bullets:
            sections.append(
                {
                    "heading": heading,
                    "paragraphs": paragraphs[:4],
                    "bullets": bullets[:6],
                }
            )

    if not sections:
        return None

    return {
        "title": default_file_title("doc", topic),
        "subtitle": topic,
        "sections": sections,
    }


def looks_like_broken_json_payload(text: str) -> bool:
    raw = (text or "").strip()
    if not raw:
        return False
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\s*", "", raw)
        raw = re.sub(r"\s*```$", "", raw)
    raw = raw.strip()
    if not raw:
        return False
    if raw.startswith("{") or raw.startswith("["):
        return True
    json_markers = ('"title"', '"subtitle"', '"sections"', '"slides"', '"columns"', '"rows"', '"table"')
    return sum(1 for marker in json_markers if marker in raw) >= 2


def doc_spec_needs_enrichment(spec: dict[str, Any], user_prompt: str) -> bool:
    sections = spec.get("sections")
    if not isinstance(sections, list) or not sections:
        return True
    if len(sections) != 1:
        return False
    section = sections[0] if isinstance(sections[0], dict) else {}
    heading = str(section.get("heading") or "").strip().lower()
    paragraphs = [str(item).strip().lower() for item in (section.get("paragraphs", []) or []) if str(item).strip()]
    bullets = [str(item).strip().lower() for item in (section.get("bullets", []) or []) if str(item).strip()]
    placeholder_bullets = {"ключевая цель", "основной результат", "следующий шаг"}
    prompt_lc = (user_prompt or "").strip().lower()
    return (
        heading in {"задача", "документ", "описание"}
        and set(bullets).issubset(placeholder_bullets)
        and (not paragraphs or paragraphs == [prompt_lc])
    )


async def ask_openrouter_text_model(
    *,
    alias: str,
    model_info: ModelInfo,
    plan_name: str,
    messages: list[dict[str, Any]],
) -> tuple[str, int, int, int]:
    payload = {
        "model": model_info.model,
        "messages": messages,
        "max_tokens": completion_tokens_for_plan(plan_name),
    }
    status, data, _ = await http_json_request_with_retries(
        "POST",
        OPENROUTER_CHAT_API,
        headers=openrouter_headers(),
        json_payload=payload,
        semaphore=state.openrouter_text_semaphore,
        request_name=f"openrouter_text:{alias}",
    )
    if status >= 400:
        message = data.get("error", {}).get("message", "Unknown OpenRouter error") if isinstance(data, dict) else str(data)
        raise RuntimeError(message)

    answer = normalize_response_output_text(data) or "РќРµ СѓРґР°Р»РѕСЃСЊ РїРѕР»СѓС‡РёС‚СЊ С‚РµРєСЃС‚РѕРІС‹Р№ РѕС‚РІРµС‚."
    prompt_tokens, completion_tokens, total_tokens = parse_usage_tokens(data)
    return truncate_text(answer, MAX_ASSISTANT_OUTPUT_CHARS), prompt_tokens, completion_tokens, total_tokens


async def ask_kie_gemini_text_model(
    *,
    alias: str,
    plan_name: str,
    messages: list[dict[str, Any]],
) -> tuple[str, int, int, int]:
    payload = {
        "messages": messages,
        "max_tokens": completion_tokens_for_plan(plan_name),
    }
    status, data, body_text = await http_json_request_with_retries(
        "POST",
        KIE_GEMINI_CHAT_API,
        headers=kie_headers(),
        json_payload=payload,
        semaphore=state.kie_text_semaphore,
        request_name=f"kie_text:{alias}",
    )
    message = kie_error_message(data)
    if status >= 400 or message:
        message = message or (data.get("msg") if isinstance(data, dict) else body_text)
        raise RuntimeError(message or "Unknown Kie Gemini error")

    answer = normalize_response_output_text(data) or "РќРµ СѓРґР°Р»РѕСЃСЊ РїРѕР»СѓС‡РёС‚СЊ С‚РµРєСЃС‚РѕРІС‹Р№ РѕС‚РІРµС‚."
    prompt_tokens, completion_tokens, total_tokens = parse_usage_tokens(data)
    return truncate_text(answer, MAX_ASSISTANT_OUTPUT_CHARS), prompt_tokens, completion_tokens, total_tokens


async def ask_kie_gpt54_text_model(
    *,
    alias: str,
    plan_name: str,
    messages: list[dict[str, Any]],
) -> tuple[str, int, int, int]:
    user_parts: list[str] = []
    system_parts: list[str] = []
    for item in messages:
        if not isinstance(item, dict):
            continue
        role = str(item.get("role") or "").strip().lower()
        content = normalize_text_content(item.get("content"))
        if not content:
            continue
        if role == "system":
            system_parts.append(content)
        elif role == "user":
            user_parts.append(content)
        else:
            user_parts.append(f"{role}: {content}")

    input_parts: list[str] = []
    if system_parts:
        input_parts.append("System instructions:\n" + "\n\n".join(system_parts))
    if user_parts:
        input_parts.append("\n\n".join(user_parts))

    payload = {
        "model": KIE_GPT54_MODEL,
        "input": "\n\n".join(part for part in input_parts if part).strip(),
        "max_output_tokens": completion_tokens_for_plan(plan_name),
    }
    status, data, body_text = await http_json_request_with_retries(
        "POST",
        KIE_GPT54_API,
        headers=kie_headers(),
        json_payload=payload,
        semaphore=state.kie_text_semaphore,
        request_name=f"kie_text:{alias}",
    )
    message = kie_error_message(data)
    if status >= 400 or message:
        message = message or (data.get("msg") if isinstance(data, dict) else body_text)
        raise RuntimeError(message or "Unknown Kie GPT-5.4 error")

    answer = normalize_response_output_text(data) or "РќРµ СѓРґР°Р»РѕСЃСЊ РїРѕР»СѓС‡РёС‚СЊ С‚РµРєСЃС‚РѕРІС‹Р№ РѕС‚РІРµС‚."
    prompt_tokens, completion_tokens, total_tokens = parse_usage_tokens(data)
    return truncate_text(answer, MAX_ASSISTANT_OUTPUT_CHARS), prompt_tokens, completion_tokens, total_tokens


async def _ask_text_model_openrouter_legacy(chat_id: int, user_text: str, selected_alias: str | None = None) -> TextAnswerResult:
    plan_name, alias, model_info, messages = build_text_request(chat_id, user_text, selected_alias=selected_alias)

    payload = {
        "model": model_info.model,
        "messages": messages,
        "max_tokens": completion_tokens_for_plan(plan_name),
    }
    status, data, _ = await http_json_request_with_retries(
        "POST",
        OPENROUTER_CHAT_API,
        headers=openrouter_headers(),
        json_payload=payload,
        semaphore=state.openrouter_text_semaphore,
        request_name=f"openrouter_text:{alias}",
    )
    if status >= 400:
        message = data.get("error", {}).get("message", "Unknown OpenRouter error") if isinstance(data, dict) else str(data)
        raise RuntimeError(message)

    choice = data["choices"][0]["message"]
    answer = normalize_text_content(choice.get("content")) or "Не удалось получить текстовый ответ."
    answer = truncate_text(answer, MAX_ASSISTANT_OUTPUT_CHARS)
    state.history(chat_id).append({"role": "user", "content": user_text})
    state.history(chat_id).append({"role": "assistant", "content": answer})

    prompt_tokens, completion_tokens, total_tokens = parse_usage_tokens(data)
    if total_tokens <= 0:
        estimated_prompt = estimate_tokens_from_messages(messages)
        total_tokens = estimated_prompt + max(0, completion_tokens)
        prompt_tokens = max(prompt_tokens, estimated_prompt)

    return TextAnswerResult(
        text=answer,
        prompt_tokens=prompt_tokens,
        completion_tokens=completion_tokens,
        total_tokens=total_tokens,
    )


async def complete_text_messages(
    *,
    alias: str,
    plan_name: str,
    messages: list[dict[str, Any]],
) -> TextAnswerResult:
    model_info = TEXT_MODELS.get(alias, DEFAULT_TEXT_MODEL)
    provider_name = "openrouter"
    try:
        if kie_enabled_for_alias(alias):
            if alias == "gemini":
                provider_name = "kie"
                answer, prompt_tokens, completion_tokens, total_tokens = await ask_kie_gemini_text_model(
                    alias=alias,
                    plan_name=plan_name,
                    messages=messages,
                )
            elif alias == "gpt54":
                provider_name = "kie"
                answer, prompt_tokens, completion_tokens, total_tokens = await ask_kie_gpt54_text_model(
                    alias=alias,
                    plan_name=plan_name,
                    messages=messages,
                )
            else:
                answer, prompt_tokens, completion_tokens, total_tokens = await ask_openrouter_text_model(
                    alias=alias,
                    model_info=model_info,
                    plan_name=plan_name,
                    messages=messages,
                )
        else:
            answer, prompt_tokens, completion_tokens, total_tokens = await ask_openrouter_text_model(
                alias=alias,
                model_info=model_info,
                plan_name=plan_name,
                messages=messages,
            )
    except Exception as exc:
        if provider_name == "kie":
            log.warning("Kie text fallback for %s: %s", alias, exc)
            answer, prompt_tokens, completion_tokens, total_tokens = await ask_openrouter_text_model(
                alias=alias,
                model_info=model_info,
                plan_name=plan_name,
                messages=messages,
            )
        else:
            raise

    if total_tokens <= 0:
        estimated_prompt = estimate_tokens_from_messages(messages)
        total_tokens = estimated_prompt + max(0, completion_tokens)
        prompt_tokens = max(prompt_tokens, estimated_prompt)

    return TextAnswerResult(
        text=answer,
        prompt_tokens=prompt_tokens,
        completion_tokens=completion_tokens,
        total_tokens=total_tokens,
    )


async def ask_text_model(chat_id: int, user_text: str, selected_alias: str | None = None) -> TextAnswerResult:
    plan_name, alias, _, messages = build_text_request(chat_id, user_text, selected_alias=selected_alias)
    result = await complete_text_messages(alias=alias, plan_name=plan_name, messages=messages)

    state.history(chat_id).append({"role": "user", "content": user_text})
    state.history(chat_id).append({"role": "assistant", "content": result.text})
    return result


def build_file_generation_messages(kind: str, user_prompt: str) -> list[dict[str, Any]]:
    format_rules = {
        "doc": (
            "Верни только JSON объект вида "
            '{"title":"...", "subtitle":"...", "sections":[{"heading":"...", "paragraphs":["..."], "bullets":["..."]}], '
            '"table":{"title":"...", "columns":["..."], "rows":[["..."]]}}. '
            "sections обязателен. table добавляй только если она реально нужна."
        ),
        "ppt": (
            "Верни только JSON объект вида "
            '{"title":"...", "subtitle":"...", "slides":[{"title":"...", "bullets":["..."], "note":"..."}]}. '
            "Сделай 6-8 содержательных слайдов, если пользователь не указал другое."
        ),
        "sheet": (
            "Верни только JSON объект вида "
            '{"title":"...", "sheet_name":"...", "columns":["..."], "rows":[["..."]], "summary":"..."}. '
            "columns и rows обязательны. Данные делай практичными и пригодными для работы."
        ),
    }
    system_text = (
        "Ты помогаешь собрать готовый офисный файл для пользователя. "
        "Отвечай только валидным JSON без markdown, пояснений и code fences. "
        "Пиши на русском, если пользователь пишет по-русски. "
        "Делай структуру компактной, практичной и пригодной для реального использования. "
        + format_rules.get(kind, "")
    )
    return [
        {"role": "system", "content": system_text},
        {"role": "user", "content": user_prompt.strip()},
    ]


def fallback_file_spec(kind: str, user_prompt: str) -> dict[str, Any]:
    topic = user_prompt.strip() or "Новый файл"
    if kind == "ppt":
        return {
            "title": "Презентация",
            "subtitle": topic,
            "slides": [
                {"title": "Тема и цель", "bullets": [topic, "Цель и ожидаемый результат"]},
                {"title": "Ключевые тезисы", "bullets": ["Главная идея", "Что важно учесть", "Следующий шаг"]},
                {"title": "Рекомендации", "bullets": ["Приоритет 1", "Приоритет 2", "Приоритет 3"]},
            ],
        }
    if kind == "sheet":
        return {
            "title": "Таблица",
            "sheet_name": "Данные",
            "columns": ["Параметр", "Значение", "Комментарий"],
            "rows": [
                ["Запрос", topic, "Исходная задача пользователя"],
                ["Статус", "Черновик", "Можно доработать"],
            ],
            "summary": "Черновая таблица по запросу пользователя.",
        }
    return {
        "title": "Документ",
        "subtitle": topic,
        "sections": [
            {
                "heading": "Задача",
                "paragraphs": [topic],
                "bullets": ["Ключевая цель", "Основной результат", "Следующий шаг"],
            }
        ],
    }


async def generate_file_spec(chat_id: int, kind: str, user_prompt: str) -> tuple[dict[str, Any], TextAnswerResult]:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free"))
    alias = str(row.get("selected_model_alias") or best_default_alias_for_plan(plan_name))
    messages = build_file_generation_messages(kind, user_prompt)
    result = await complete_text_messages(alias=alias, plan_name=plan_name, messages=messages)
    spec = extract_json_object(result.text) or fallback_file_spec(kind, user_prompt)
    return spec, result


def build_file_generation_messages(kind: str, profile: str, user_prompt: str) -> list[dict[str, Any]]:
    profile_rules = {
        "short": "Сделай короткую и плотную версию без лишней воды.",
        "medium": "Сделай стандартную рабочую версию с нормальной полнотой.",
        "full": "Сделай подробную и более полную версию с расширенным содержанием.",
    }
    ppt_rule = {
        "short": "Сделай до 5 содержательных слайдов.",
        "medium": "Сделай 6-8 содержательных слайдов, если пользователь не указал другое.",
        "full": "Сделай 9-12 содержательных слайдов, если пользователь не указал другое.",
    }.get(profile, "Сделай 6-8 содержательных слайдов, если пользователь не указал другое.")
    format_rules = {
        "doc": build_doc_generation_rule(profile, user_prompt),
        "ppt": (
            'Верни только JSON объект вида {"title":"...", "subtitle":"...", "slides":[{"title":"...", "bullets":["..."], "note":"..."}]}. '
            + ppt_rule
        ),
        "sheet": (
            'Верни только JSON объект вида {"title":"...", "sheet_name":"...", "columns":["..."], "rows":[["..."]], "summary":"..."}. '
            "columns и rows обязательны. Данные делай практичными и пригодными для работы."
        ),
    }
    system_text = (
        "Ты помогаешь собрать готовый офисный файл для пользователя. "
        "Отвечай только валидным JSON без markdown, пояснений и code fences. "
        "Пиши на русском, если пользователь пишет по-русски. "
        "Делай структуру практичной и пригодной для реального использования. "
        + profile_rules.get(profile, profile_rules["medium"]) + " "
        + format_rules.get(kind, "")
    )
    return [
        {"role": "system", "content": system_text},
        {"role": "user", "content": user_prompt.strip()},
    ]


def fallback_file_spec(kind: str, user_prompt: str, profile: str = "medium") -> dict[str, Any]:
    topic = normalized_file_topic(user_prompt)
    if kind == "ppt":
        slide_count = {"short": 3, "medium": 5, "full": 7}.get(profile, 5)
        base_slides = [
            {"title": "Тема и цель", "bullets": [topic, "Цель и ожидаемый результат"]},
            {"title": "Ключевые тезисы", "bullets": ["Главная идея", "Что важно учесть", "Следующий шаг"]},
            {"title": "Рекомендации", "bullets": ["Приоритет 1", "Приоритет 2", "Приоритет 3"]},
            {"title": "План действий", "bullets": ["Шаг 1", "Шаг 2", "Шаг 3"]},
            {"title": "Риски", "bullets": ["Что может помешать", "Как это снизить"]},
            {"title": "Ресурсы", "bullets": ["Люди", "Сроки", "Инструменты"]},
            {"title": "Итоги", "bullets": ["Главный вывод", "Что делать дальше"]},
        ]
        return {"title": "Презентация", "subtitle": topic, "slides": base_slides[:slide_count]}
    if kind == "sheet":
        extra_rows = {
            "short": [],
            "medium": [["Приоритет", "Средний", "Можно уточнить"]],
            "full": [["Приоритет", "Средний", "Можно уточнить"], ["Ответственный", "-", "Добавить при доработке"], ["Срок", "-", "Уточнить по задаче"]],
        }.get(profile, [])
        return {
            "title": "Таблица",
            "sheet_name": "Данные",
            "columns": ["Параметр", "Значение", "Комментарий"],
            "rows": [["Запрос", topic, "Исходная задача пользователя"], ["Статус", "Черновик", "Можно доработать"], *extra_rows],
            "summary": "Черновая таблица по запросу пользователя.",
        }
    return {
        "title": default_file_title("doc", topic),
        "subtitle": topic,
        "sections": placeholder_doc_sections(topic, profile),
    }


async def generate_file_spec(chat_id: int, kind: str, profile: str, user_prompt: str) -> tuple[dict[str, Any], TextAnswerResult]:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free"))
    alias = str(row.get("selected_model_alias") or best_default_alias_for_plan(plan_name))
    messages = build_file_generation_messages(kind, profile, user_prompt)
    result = await complete_text_messages(alias=alias, plan_name=plan_name, messages=messages)
    spec = extract_json_object(result.text)
    if kind == "doc":
        if not isinstance(spec, dict) and not looks_like_broken_json_payload(result.text):
            spec = doc_spec_from_plain_text(user_prompt, result.text, profile)
        elif isinstance(spec, dict) and doc_spec_needs_enrichment(spec, user_prompt):
            spec = fallback_file_spec(kind, user_prompt, profile)
    if not isinstance(spec, dict):
        spec = fallback_file_spec(kind, user_prompt, profile)
    return spec, result


def build_docx_bytes(spec: dict[str, Any]) -> bytes:
    document = Document()
    title = str(spec.get("title") or "Документ").strip()
    subtitle = str(spec.get("subtitle") or "").strip()
    title_par = document.add_paragraph()
    title_run = title_par.add_run(title)
    title_run.bold = True
    title_run.font.size = Pt(20)
    if subtitle:
        subtitle_par = document.add_paragraph()
        subtitle_run = subtitle_par.add_run(subtitle)
        subtitle_run.italic = True
        subtitle_run.font.size = Pt(11)
    sections = spec.get("sections")
    if not isinstance(sections, list) or not sections:
        sections = fallback_file_spec("doc", title).get("sections", [])
    for section in sections:
        if not isinstance(section, dict):
            continue
        heading = str(section.get("heading") or "").strip()
        if heading:
            document.add_heading(heading, level=1)
        for paragraph in section.get("paragraphs", []) or []:
            text = str(paragraph).strip()
            if text:
                document.add_paragraph(text)
        for bullet in section.get("bullets", []) or []:
            text = str(bullet).strip()
            if text:
                document.add_paragraph(text, style="List Bullet")
    table_spec = spec.get("table")
    if isinstance(table_spec, dict):
        columns = [str(col).strip() for col in table_spec.get("columns", []) if str(col).strip()]
        rows = table_spec.get("rows", []) or []
        if columns and rows:
            table_title = str(table_spec.get("title") or "").strip()
            if table_title:
                document.add_heading(table_title, level=2)
            table = document.add_table(rows=1, cols=len(columns))
            table.style = "Table Grid"
            for idx, col in enumerate(columns):
                table.rows[0].cells[idx].text = col
            for row in rows[:20]:
                values = row if isinstance(row, list) else [row]
                cells = table.add_row().cells
                for idx in range(len(columns)):
                    cells[idx].text = str(values[idx]).strip() if idx < len(values) else ""
    out = BytesIO()
    document.save(out)
    return out.getvalue()


def build_pptx_bytes(spec: dict[str, Any]) -> bytes:
    prs = Presentation()
    title = str(spec.get("title") or "Презентация").strip()
    subtitle = str(spec.get("subtitle") or "").strip()
    first_slide = prs.slides.add_slide(prs.slide_layouts[0])
    first_slide.shapes.title.text = title
    if len(first_slide.placeholders) > 1:
        first_slide.placeholders[1].text = subtitle
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        slides = fallback_file_spec("ppt", title).get("slides", [])
    for slide_spec in slides[:12]:
        if not isinstance(slide_spec, dict):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(slide_spec.get("title") or "Слайд").strip()
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = slide_spec.get("bullets", []) or []
        if bullets:
            for idx, bullet in enumerate(bullets[:6]):
                paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                paragraph.text = str(bullet).strip()
                paragraph.level = 0
                for run in paragraph.runs:
                    run.font.size = PptPt(20)
        note = str(slide_spec.get("note") or "").strip()
        if note:
            box = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.5), Inches(0.5))
            paragraph = box.text_frame.paragraphs[0]
            paragraph.text = note
            for run in paragraph.runs:
                run.font.size = PptPt(10)
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def build_xlsx_bytes(spec: dict[str, Any]) -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sanitize_filename_part(str(spec.get("sheet_name") or "Данные"), "Данные")[:31]
    summary = str(spec.get("summary") or "").strip()
    row_offset = 1
    if summary:
        sheet["A1"] = summary
        sheet["A1"].font = Font(bold=True, size=12)
        row_offset = 3
    columns = [str(col).strip() for col in spec.get("columns", []) if str(col).strip()]
    rows = spec.get("rows", []) or []
    if not columns:
        fallback = fallback_file_spec("sheet", str(spec.get("title") or "Таблица"))
        columns = fallback["columns"]
        rows = fallback["rows"]
    for col_idx, column in enumerate(columns, start=1):
        cell = sheet.cell(row=row_offset, column=col_idx, value=column)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E78")
        cell.alignment = Alignment(horizontal="center", vertical="center")
    for row_idx, row_values in enumerate(rows[:200], start=row_offset + 1):
        values = row_values if isinstance(row_values, list) else [row_values]
        for col_idx in range(1, len(columns) + 1):
            value = values[col_idx - 1] if col_idx - 1 < len(values) else ""
            cell = sheet.cell(row=row_idx, column=col_idx, value=str(value))
            cell.alignment = Alignment(vertical="top", wrap_text=True)
    sheet.freeze_panes = f"A{row_offset + 1}"
    for idx, column in enumerate(columns, start=1):
        max_len = len(column)
        for row_idx in range(row_offset + 1, min(sheet.max_row, row_offset + 25) + 1):
            max_len = max(max_len, len(str(sheet.cell(row=row_idx, column=idx).value or "")))
        sheet.column_dimensions[chr(64 + idx)].width = min(max(max_len + 3, 14), 36)
    out = BytesIO()
    workbook.save(out)
    return out.getvalue()


async def process_file_request(chat_id: int, kind: str, prompt: str, profile: str = "medium") -> bool:
    clean_prompt = prompt.strip()
    if not clean_prompt:
        return False
    if len(clean_prompt) > MAX_TEXT_INPUT_CHARS:
        await max_send_message(chat_id, f"Слишком длинное описание. Максимум {MAX_TEXT_INPUT_CHARS} символов.", attachments=build_keyboard())
        return True
    ok_cd, reason_cd = check_cooldown(chat_id, "message")
    if not ok_cd:
        await max_send_message(chat_id, reason_cd, attachments=build_keyboard())
        return True
    ok, reason = check_limit_only(chat_id, "messages")
    if not ok:
        await max_send_message(chat_id, reason, attachments=build_keyboard())
        return True
    cost = file_request_cost(kind)
    ok_credit, reason_credit = check_and_consume_credits(chat_id, cost, f"{FILE_KIND_LABELS.get(kind, 'файл')}")
    if not ok_credit:
        await max_send_message(chat_id, reason_credit, attachments=purchase_help_keyboard_for_row(user_profile(chat_id)))
        return True
    ok, reason = check_and_consume_limit(chat_id, "messages")
    if not ok:
        state.user_store.refund_credits(chat_id, cost)
        await max_send_message(chat_id, reason, attachments=build_keyboard())
        return True
    await max_send_message(chat_id, f"Собираю {FILE_KIND_LABELS.get(kind, 'файл')}, это может занять немного времени...", notify=False)
    try:
        spec, result = await generate_file_spec(chat_id, kind, profile, clean_prompt)
        title = str(spec.get("title") or clean_prompt[:60] or FILE_KIND_LABELS.get(kind, "Файл")).strip()
        if kind == "ppt":
            file_bytes = build_pptx_bytes(spec)
        elif kind == "sheet":
            file_bytes = build_xlsx_bytes(spec)
        else:
            file_bytes = build_docx_bytes(spec)
        filename = f"{sanitize_filename_part(title, FILE_KIND_LABELS.get(kind, 'file'))}.{FILE_KIND_EXTENSIONS.get(kind, 'bin')}"
        await send_generated_file(chat_id, kind=kind, title=title, file_bytes=file_bytes, filename=filename)
        final_row = user_profile(chat_id)
        state.user_store.record_usage_event(
            chat_id=chat_id,
            event_type="file_request",
            plan=str(final_row.get("plan", "")),
            model_alias=str(final_row.get("selected_model_alias") or ""),
            credits_spent=cost,
            tokens_total=int(result.total_tokens),
            details=f"kind={kind};profile={profile};title={title}",
        )
        return True
    except Exception:
        state.user_store.refund_credits(chat_id, cost)
        raise


async def process_file_request(chat_id: int, kind: str, prompt: str, profile: str = "medium") -> bool:
    clean_prompt = prompt.strip()
    if not clean_prompt:
        return False
    if len(clean_prompt) > MAX_TEXT_INPUT_CHARS:
        await max_send_message(chat_id, f"Слишком длинное описание. Максимум {MAX_TEXT_INPUT_CHARS} символов.", attachments=build_keyboard())
        return True

    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free")).strip().lower()
    is_free_plan = plan_name == "free"

    ok_cd, reason_cd = check_cooldown(chat_id, "message")
    if not ok_cd:
        await max_send_message(chat_id, reason_cd, attachments=build_keyboard())
        return True

    if is_free_plan:
        if profile != "short":
            await max_send_message(
                chat_id,
                "На Free доступен только короткий файл. Средняя и полная версии открываются на платных тарифах.",
                attachments=purchase_help_keyboard_for_row(row),
            )
            return True
        ok_limit, reason_limit = check_limit_only(chat_id, "files")
        if not ok_limit:
            await max_send_message(chat_id, reason_limit, attachments=purchase_help_keyboard_for_row(row))
            return True
        cost = 0
    else:
        ok_limit, reason_limit = check_limit_only(chat_id, "messages")
        if not ok_limit:
            await max_send_message(chat_id, reason_limit, attachments=build_keyboard())
            return True
        cost = file_request_cost(kind)
        ok_credit, reason_credit = check_and_consume_credits(chat_id, cost, f"{FILE_KIND_LABELS.get(kind, 'файл')}")
        if not ok_credit:
            await max_send_message(chat_id, reason_credit, attachments=purchase_help_keyboard_for_row(row))
            return True
        ok_limit, reason_limit = check_and_consume_limit(chat_id, "messages")
        if not ok_limit:
            state.user_store.refund_credits(chat_id, cost)
            await max_send_message(chat_id, reason_limit, attachments=build_keyboard())
            return True

    await max_send_message(chat_id, f"Собираю {FILE_KIND_LABELS.get(kind, 'файл')}, это может занять немного времени...", notify=False)
    try:
        spec, result = await generate_file_spec(chat_id, kind, profile, clean_prompt)
        title = str(spec.get("title") or clean_prompt[:60] or FILE_KIND_LABELS.get(kind, "Файл")).strip()
        if kind == "ppt":
            file_bytes = build_pptx_bytes(spec)
        elif kind == "sheet":
            file_bytes = build_xlsx_bytes(spec)
        else:
            file_bytes = build_docx_bytes(spec)
        filename = f"{sanitize_filename_part(title, FILE_KIND_LABELS.get(kind, 'file'))}.{FILE_KIND_EXTENSIONS.get(kind, 'bin')}"
        await send_generated_file(chat_id, kind=kind, title=title, file_bytes=file_bytes, filename=filename)
        if is_free_plan:
            consume_limit(chat_id, "files")
        final_row = user_profile(chat_id)
        state.user_store.record_usage_event(
            chat_id=chat_id,
            event_type="file_request",
            plan=str(final_row.get("plan", "")),
            model_alias=str(final_row.get("selected_model_alias") or ""),
            credits_spent=cost,
            tokens_total=int(result.total_tokens),
            details=f"kind={kind};profile={profile};title={title}",
        )
        return True
    except Exception:
        if cost > 0:
            state.user_store.refund_credits(chat_id, cost)
        raise


def build_files_menu_text(chat_id: int) -> str:
    row = user_profile(chat_id)
    if str(row.get("plan", "free")).strip().lower() == "free":
        return (
            "📄 Файлы\n\n"
            "Здесь можно собрать готовый файл и получить его прямо в чат.\n\n"
            "Что доступно:\n"
            "• Документ (.docx)\n"
            "• Презентация (.pptx)\n"
            "• Таблица (.xlsx)\n\n"
            "На Free:\n"
            "• 1 короткий файл каждые 14 дней\n"
            "• средняя и полная версии доступны на платных тарифах\n\n"
            "Можно нажать кнопку ниже или просто написать:\n"
            "«сделай презентацию ...», «подготовь документ ...», «собери таблицу ...»"
        )
    return (
        "📄 Файлы\n\n"
        "Здесь можно собрать готовый файл и получить его прямо в чат.\n\n"
        "Что доступно:\n"
        "• Документ (.docx)\n"
        "• Презентация (.pptx)\n"
        "• Таблица (.xlsx)\n\n"
        "Стоимость:\n"
        f"• Документ — {FILE_DOC_REQUEST_COST} запросов\n"
        f"• Презентация — {FILE_PPT_REQUEST_COST} запросов\n"
        f"• Таблица — {FILE_SHEET_REQUEST_COST} запросов\n\n"
        "Можно нажать кнопку ниже или просто написать:\n"
        "«сделай презентацию ...», «подготовь документ ...», «собери таблицу ...»"
    )


def build_file_size_text(kind: str) -> str:
    kind_label = FILE_KIND_LABELS.get(kind, "файл").capitalize()
    row = user_profile(0) if False else None
    current_row = None
    hint = ""
    if kind == "ppt":
        hint = "Короткая: до 5 слайдов\nСредняя: 6-8 слайдов\nПолная: 9-12 слайдов"
    elif kind == "sheet":
        hint = "Короткая: простая рабочая таблица\nСредняя: нормальная рабочая версия\nПолная: расширенная таблица"
    else:
        hint = "Короткая: краткий документ\nСредняя: стандартный рабочий документ\nПолная: подробная версия"
    return (
        f"{kind_label}\n\n"
        "Выбери объем перед генерацией.\n\n"
        f"{hint}\n\n"
        "После этого я попрошу описать задачу одним сообщением."
    )


def build_file_prompt_text(kind: str, profile: str = "medium", chat_id: int | None = None) -> str:
    profile_label = FILE_PROFILE_LABELS.get(profile, FILE_PROFILE_LABELS["medium"])
    is_free_plan = False
    if chat_id is not None:
        is_free_plan = str(user_profile(chat_id).get("plan", "free")).strip().lower() == "free"
    row = user_profile(0) if False else None
    if kind == "ppt":
        free_line = "Лимит Free: 1 короткий файл каждые 14 дней." if profile == "short" else f"Стоимость: {FILE_PPT_REQUEST_COST} запросов."
        return (
            f"Режим: Презентация — {profile_label}\n\n"
            "Напиши, какую презентацию собрать одним сообщением.\n\n"
            "Хороший пример: «сделай презентацию на 8 слайдов про запуск AI-бота для рекламы и продаж».\n\n"
            f"{free_line}"
        )
    if kind == "sheet":
        free_line = "Лимит Free: 1 короткий файл каждые 14 дней." if profile == "short" else f"Стоимость: {FILE_SHEET_REQUEST_COST} запросов."
        return (
            f"Режим: Таблица — {profile_label}\n\n"
            "Напиши, какую таблицу собрать одним сообщением.\n\n"
            "Хороший пример: «собери таблицу бюджета на рекламу по каналам на месяц».\n\n"
            f"{free_line}"
        )
    free_line = "Лимит Free: 1 короткий файл каждые 14 дней." if profile == "short" else f"Стоимость: {FILE_DOC_REQUEST_COST} запросов."
    return (
        f"Режим: Документ — {profile_label}\n\n"
        "Напиши, какой документ собрать одним сообщением.\n\n"
        "Хороший пример: «подготовь коммерческое предложение для клиента на разработку бота».\n\n"
        f"{free_line}"
    )


def build_file_prompt_text(kind: str, profile: str = "medium", chat_id: int | None = None) -> str:
    profile_label = FILE_PROFILE_LABELS.get(profile, FILE_PROFILE_LABELS["medium"])
    is_free_plan = False
    if chat_id is not None:
        is_free_plan = str(user_profile(chat_id).get("plan", "free")).strip().lower() == "free"

    if kind == "ppt":
        footer = (
            "Лимит Free: 1 короткий файл каждые 14 дней."
            if is_free_plan
            else f"Стоимость: {FILE_PPT_REQUEST_COST} запросов."
        )
        return (
            f"Режим: Презентация — {profile_label}\n\n"
            "Напиши, какую презентацию собрать одним сообщением.\n\n"
            "Хороший пример: «сделай презентацию на 8 слайдов про запуск AI-бота для рекламы и продаж».\n\n"
            f"{footer}"
        )
    if kind == "sheet":
        footer = (
            "Лимит Free: 1 короткий файл каждые 14 дней."
            if is_free_plan
            else f"Стоимость: {FILE_SHEET_REQUEST_COST} запросов."
        )
        return (
            f"Режим: Таблица — {profile_label}\n\n"
            "Напиши, какую таблицу собрать одним сообщением.\n\n"
            "Хороший пример: «собери таблицу бюджета на рекламу по каналам на месяц».\n\n"
            f"{footer}"
        )
    footer = (
        "Лимит Free: 1 короткий файл каждые 14 дней."
        if is_free_plan
        else f"Стоимость: {FILE_DOC_REQUEST_COST} запросов."
    )
    return (
        f"Режим: Документ — {profile_label}\n\n"
        "Напиши, какой документ собрать одним сообщением.\n\n"
        "Хороший пример: «подготовь коммерческое предложение для клиента на разработку бота».\n\n"
        f"{footer}"
    )


async def generate_image(prompt: str) -> ImageResult:
    payload = {
        "model": DEFAULT_IMAGE_MODEL.model,
        "messages": [{"role": "user", "content": prompt}],
        "modalities": ["image", "text"],
    }

    status, data, _ = await http_json_request_with_retries(
        "POST",
        OPENROUTER_CHAT_API,
        headers=openrouter_headers(),
        json_payload=payload,
        semaphore=state.openrouter_image_semaphore,
        request_name="openrouter_image_generate",
    )
    if status >= 400:
        message = data.get("error", {}).get("message", "Unknown OpenRouter error") if isinstance(data, dict) else str(data)
        raise RuntimeError(message)

    message = data["choices"][0]["message"]
    images = message.get("images") or []
    for image in images:
        if not isinstance(image, dict):
            continue
        image_url = image.get("image_url") or image.get("imageUrl") or {}
        if isinstance(image_url, dict):
            url = image_url.get("url")
            if isinstance(url, str) and url.startswith("data:"):
                return decode_data_url(url)
            if isinstance(url, str) and url.startswith("http"):
                return await fetch_image_bytes(url)
    raise RuntimeError("Image was not returned by the selected model.")


async def generate_image_from_reference(prompt: str, reference_image_data_url: str) -> ImageResult:
    payload = {
        "model": DEFAULT_IMAGE_MODEL.model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": reference_image_data_url}},
                ],
            }
        ],
        "modalities": ["image", "text"],
    }

    status, data, _ = await http_json_request_with_retries(
        "POST",
        OPENROUTER_CHAT_API,
        headers=openrouter_headers(),
        json_payload=payload,
        semaphore=state.openrouter_image_semaphore,
        request_name="openrouter_image_edit",
    )
    if status >= 400:
        message = data.get("error", {}).get("message", "Unknown OpenRouter error") if isinstance(data, dict) else str(data)
        raise RuntimeError(message)

    message = data["choices"][0]["message"]
    images = message.get("images") or []
    for image in images:
        if not isinstance(image, dict):
            continue
        image_url = image.get("image_url") or image.get("imageUrl") or {}
        if isinstance(image_url, dict):
            url = image_url.get("url")
            if isinstance(url, str) and url.startswith("data:"):
                return decode_data_url(url)
            if isinstance(url, str) and url.startswith("http"):
                return await fetch_image_bytes(url)
    raise RuntimeError("Edited image was not returned by the selected model.")


def current_model_label(chat_id: int) -> str:
    row = user_profile(chat_id)
    selected = row["selected_model_alias"] or best_default_alias_for_plan(row["plan"])
    model = TEXT_MODELS.get(selected, DEFAULT_TEXT_MODEL)
    return model.label


def current_model_display(chat_id: int) -> str:
    row = user_profile(chat_id)
    plan = str(row.get("plan", "free"))
    selected = str(row.get("selected_model_alias") or best_default_alias_for_plan(plan))
    model = TEXT_MODELS.get(selected, DEFAULT_TEXT_MODEL)
    preset = str(row.get("selected_preset", "") or "").strip().lower()
    preset_cfg = MODEL_PRESETS.get(preset)
    if preset_cfg and selected in preset_available_aliases_for_plan(plan, preset):
        preset_icon = str(preset_cfg["label"]).split()[0]
        return f"{preset_icon} {model.label}"
    return model.label


def build_files_menu_keyboard(chat_id: int) -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "📄 Документ", "payload": "files:doc"},
                        {"type": "callback", "text": "📊 Презентация", "payload": "files:ppt"},
                    ],
                    [
                        {"type": "callback", "text": "📈 Таблица", "payload": "files:sheet"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_keyboard(chat_id: int | None = None) -> list[dict[str, Any]]:
    plan_buttons = [
        {"type": "callback", "text": "Тарифы", "payload": "action:tariffs"},
        {"type": "callback", "text": "Мой план", "payload": "action:plan"},
    ]
    if chat_id is not None and can_pick_models_for_current_preset(chat_id):
        plan_buttons.append({"type": "callback", "text": "⚙ Модели", "payload": "action:preset_models"})
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "⚡ Быстро", "payload": "set_preset:fast"},
                        {"type": "callback", "text": "⚖ Баланс", "payload": "set_preset:balanced"},
                    ],
                    [
                        {"type": "callback", "text": "🧠 Качество", "payload": "set_preset:quality"},
                        {"type": "callback", "text": "🚀 Эксперт", "payload": "set_preset:expert"},
                    ],
                    [*plan_buttons],
                    [
                        {"type": "callback", "text": "🎨 Картинка", "payload": "action:image_menu"},
                        {"type": "callback", "text": "📄 Файлы", "payload": "action:files_menu"},
                        {"type": "callback", "text": "🎁 Бонусы", "payload": "action:growth"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Сброс", "payload": "action:clear"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                    [
                        {"type": "callback", "text": "📣 Канал", "payload": "action:channel"},
                    ],
                ]
            },
        }
    ]


def build_reply_shortcuts_keyboard(chat_id: int, include_share: bool = False) -> list[dict[str, Any]]:
    row = user_profile(chat_id)
    buttons: list[list[dict[str, Any]]] = [
        [
            {"type": "callback", "text": "Меню", "payload": "reply_action:menu"},
            {"type": "callback", "text": "🎨 Картинка", "payload": "reply_action:image_menu"},
            {"type": "callback", "text": "📄 Файлы", "payload": "reply_action:files_menu"},
        ]
    ]
    if str(row.get("plan", "free")) == "free":
        buttons[0].append({"type": "callback", "text": "Тарифы", "payload": "reply_action:tariffs"})
    buttons.append([{"type": "callback", "text": "Сброс", "payload": "reply_action:clear"}])
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def build_files_prompt_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {"buttons": [[{"type": "callback", "text": "Отмена", "payload": "files:cancel"}]]},
        }
    ]


def build_files_menu_text(chat_id: int) -> str:
    return (
        "📄 Файлы\n\n"
        "Здесь можно собрать готовый файл и получить его прямо в чат.\n\n"
        "Что доступно:\n"
        "• Документ (.docx)\n"
        "• Презентация (.pptx)\n"
        "• Таблица (.xlsx)\n\n"
        "Стоимость:\n"
        f"• Документ — {FILE_DOC_REQUEST_COST} запросов\n"
        f"• Презентация — {FILE_PPT_REQUEST_COST} запросов\n"
        f"• Таблица — {FILE_SHEET_REQUEST_COST} запросов\n\n"
        "Можно нажать кнопку ниже или просто написать:\n"
        "«сделай презентацию ...», «подготовь документ ...», «собери таблицу ...»"
    )


def build_file_prompt_text(kind: str) -> str:
    if kind == "ppt":
        return (
            "Напиши, какую презентацию собрать одним сообщением.\n\n"
            "Хороший пример: «сделай презентацию на 8 слайдов про запуск AI-бота для рекламы и продаж».\n\n"
            f"Стоимость: {FILE_PPT_REQUEST_COST} запросов."
        )
    if kind == "sheet":
        return (
            "Напиши, какую таблицу собрать одним сообщением.\n\n"
            "Хороший пример: «собери таблицу бюджета на рекламу по каналам на месяц».\n\n"
            f"Стоимость: {FILE_SHEET_REQUEST_COST} запросов."
        )
    return (
        "Напиши, какой документ собрать одним сообщением.\n\n"
        "Хороший пример: «подготовь коммерческое предложение для клиента на разработку бота».\n\n"
        f"Стоимость: {FILE_DOC_REQUEST_COST} запросов."
    )


async def send_image_menu(chat_id: int, notify: bool = False) -> None:
    set_image_mode(chat_id, "")
    await show_managed_content(
        chat_id,
        build_image_menu_text(chat_id),
        attachments=build_image_menu_keyboard(chat_id),
        page=UI_PAGE_IMAGE_MENU,
        push_history=False,
        force_new=False,
    )


def build_files_size_keyboard(kind: str) -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "Короткая", "payload": f"files:size:{kind}:short"},
                        {"type": "callback", "text": "Средняя", "payload": f"files:size:{kind}:medium"},
                        {"type": "callback", "text": "Полная", "payload": f"files:size:{kind}:full"},
                    ],
                    [
                        {"type": "callback", "text": "Назад", "payload": "action:files_menu"},
                        {"type": "callback", "text": "Отмена", "payload": "files:cancel"},
                    ],
                ]
            },
        }
    ]


def build_files_prompt_keyboard() -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {"buttons": [[{"type": "callback", "text": "Отмена", "payload": "files:cancel"}]]},
        }
    ]


def build_file_size_text(kind: str) -> str:
    kind_label = FILE_KIND_LABELS.get(kind, "файл").capitalize()
    if kind == "ppt":
        hint = "Короткая: до 5 слайдов\nСредняя: 6-8 слайдов\nПолная: 9-12 слайдов"
    elif kind == "sheet":
        hint = "Короткая: простая рабочая таблица\nСредняя: нормальная рабочая версия\nПолная: расширенная таблица"
    else:
        hint = "Короткая: краткий документ\nСредняя: стандартный рабочий документ\nПолная: подробная версия"
    return (
        f"{kind_label}\n\n"
        "Выбери объем перед генерацией.\n\n"
        f"{hint}\n\n"
        "После этого я попрошу описать задачу одним сообщением."
    )


def build_file_prompt_text(kind: str, profile: str = "medium") -> str:
    profile_label = FILE_PROFILE_LABELS.get(profile, FILE_PROFILE_LABELS["medium"])
    if kind == "ppt":
        return (
            f"Режим: Презентация — {profile_label}\n\n"
            "Напиши, какую презентацию собрать одним сообщением.\n\n"
            "Хороший пример: «сделай презентацию на 8 слайдов про запуск AI-бота для рекламы и продаж».\n\n"
            f"Стоимость: {FILE_PPT_REQUEST_COST} запросов."
        )
    if kind == "sheet":
        return (
            f"Режим: Таблица — {profile_label}\n\n"
            "Напиши, какую таблицу собрать одним сообщением.\n\n"
            "Хороший пример: «собери таблицу бюджета на рекламу по каналам на месяц».\n\n"
            f"Стоимость: {FILE_SHEET_REQUEST_COST} запросов."
        )
    return (
        f"Режим: Документ — {profile_label}\n\n"
        "Напиши, какой документ собрать одним сообщением.\n\n"
        "Хороший пример: «подготовь коммерческое предложение для клиента на разработку бота».\n\n"
        f"Стоимость: {FILE_DOC_REQUEST_COST} запросов."
    )


def build_files_menu_keyboard(chat_id: int) -> list[dict[str, Any]]:
    return [
        {
            "type": "inline_keyboard",
            "payload": {
                "buttons": [
                    [
                        {"type": "callback", "text": "📄 Документ", "payload": "files:sizepick:doc"},
                        {"type": "callback", "text": "📊 Презентация", "payload": "files:sizepick:ppt"},
                    ],
                    [
                        {"type": "callback", "text": "📈 Таблица", "payload": "files:sizepick:sheet"},
                    ],
                    [
                        {"type": "callback", "text": "Меню", "payload": "action:menu"},
                        {"type": "callback", "text": "Помощь", "payload": "action:support"},
                    ],
                ]
            },
        }
    ]


def build_file_prompt_text(kind: str, profile: str = "medium", chat_id: int | None = None) -> str:
    profile_label = FILE_PROFILE_LABELS.get(profile, FILE_PROFILE_LABELS["medium"])
    is_free_plan = False
    if chat_id is not None:
        is_free_plan = str(user_profile(chat_id).get("plan", "free")).strip().lower() == "free"

    if kind == "ppt":
        footer = "Лимит Free: 1 короткий файл каждые 14 дней." if is_free_plan else f"Стоимость: {FILE_PPT_REQUEST_COST} запросов."
        return (
            f"Режим: Презентация — {profile_label}\n\n"
            "Напиши, какую презентацию собрать одним сообщением.\n\n"
            "Хороший пример: «сделай презентацию на 8 слайдов про запуск AI-бота для рекламы и продаж».\n\n"
            f"{footer}"
        )
    if kind == "sheet":
        footer = "Лимит Free: 1 короткий файл каждые 14 дней." if is_free_plan else f"Стоимость: {FILE_SHEET_REQUEST_COST} запросов."
        return (
            f"Режим: Таблица — {profile_label}\n\n"
            "Напиши, какую таблицу собрать одним сообщением.\n\n"
            "Хороший пример: «собери таблицу бюджета на рекламу по каналам на месяц».\n\n"
            f"{footer}"
        )
    footer = "Лимит Free: 1 короткий файл каждые 14 дней." if is_free_plan else f"Стоимость: {FILE_DOC_REQUEST_COST} запросов."
    return (
        f"Режим: Документ — {profile_label}\n\n"
        "Напиши, какой документ собрать одним сообщением.\n\n"
        "Хороший пример: «подготовь коммерческое предложение для клиента на разработку бота».\n\n"
        f"{footer}"
    )


async def send_files_menu(chat_id: int, notify: bool = False) -> None:
    clear_file_pending_input(chat_id)
    await show_managed_content(
        chat_id,
        build_files_menu_text(chat_id),
        attachments=build_files_menu_keyboard(chat_id),
        page=UI_PAGE_FILES_MENU,
        push_history=False,
        force_new=False,
    )


async def process_image_generation(chat_id: int, user_prompt: str, model_prompt: str | None = None) -> bool:
    return await _process_image_generation_queued(chat_id, user_prompt, model_prompt=model_prompt)

async def process_image_edit_generation(chat_id: int, user_prompt: str, reference_image_data_url: str) -> bool:
    return await _process_image_edit_generation_queued(chat_id, user_prompt, reference_image_data_url)

async def _process_image_generation_queued(chat_id: int, user_prompt: str, model_prompt: str | None = None) -> bool:
    prompt = user_prompt.strip()
    if not prompt:
        await max_send_message(chat_id, "Опиши, что нужно сгенерировать.", attachments=build_image_prompt_keyboard())
        return True
    if len(prompt) > MAX_IMAGE_PROMPT_CHARS:
        await max_send_message(chat_id, f"Слишком длинный промпт. Максимум {MAX_IMAGE_PROMPT_CHARS} символов.", attachments=build_keyboard())
        return True
    if image_job_already_pending(chat_id):
        await max_send_message(chat_id, "У тебя уже есть активная задача по картинке. Дождись результата или ошибки, потом запускай следующую.", attachments=build_image_menu_keyboard(chat_id))
        return True

    ok_cd, reason_cd = check_cooldown(chat_id, "image")
    if not ok_cd:
        await max_send_message(chat_id, reason_cd, attachments=build_keyboard())
        return True

    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free")).strip().lower()
    if plan_name != "free" and not plan_allowed(plan_name, DEFAULT_IMAGE_MODEL.min_plan):
        await show_managed_content(
            chat_id,
            f"Картинки доступны с тарифа {DEFAULT_IMAGE_MODEL.min_plan}. Открой «Тарифы».",
            attachments=build_tariffs_keyboard_pricing(),
            page=UI_PAGE_TARIFFS,
            push_history=False,
            notification="Тарифы",
        )
        return True

    ok, reason = check_limit_only(chat_id, "images")
    if not ok:
        await show_managed_content(
            chat_id,
            reason,
            attachments=build_image_menu_keyboard(chat_id),
            page=UI_PAGE_IMAGE_MENU,
            push_history=False,
            notification="Лимит достигнут",
        )
        return True

    is_free_plan = plan_name == "free"
    img_cost = 0 if is_free_plan else image_credit_cost()
    if img_cost > 0:
        ok_credit, reason_credit = check_and_consume_credits(chat_id, img_cost, "картинка")
        if not ok_credit:
            await show_managed_content(
                chat_id,
                reason_credit,
                attachments=purchase_help_keyboard_for_row(row),
                page=UI_PAGE_TARIFFS,
                push_history=False,
                notification="Недостаточно запросов",
            )
            return True

    ok, reason = check_and_consume_limit(chat_id, "images")
    if not ok:
        if img_cost > 0:
            state.user_store.refund_credits(chat_id, img_cost)
        await show_managed_content(
            chat_id,
            reason,
            attachments=build_image_menu_keyboard(chat_id),
            page=UI_PAGE_IMAGE_MENU,
            push_history=False,
            notification="Лимит достигнут",
        )
        return True

    request_for_model = model_prompt or prompt
    queue_before = await enqueue_image_job(
        ImageJob(
            kind="generate",
            chat_id=chat_id,
            user_prompt=prompt,
            model_prompt=request_for_model,
            credits_spent=img_cost,
            details=f"style={get_image_prefs(chat_id).get('style','')};aspect={get_image_prefs(chat_id).get('aspect','')}",
        )
    )
    set_image_mode(chat_id, "")
    if queue_before > 0:
        await max_send_message(chat_id, f"Поставил картинку в очередь. Перед тобой задач: {queue_before}. Пришлю результат отдельным сообщением.", attachments=build_image_menu_keyboard(chat_id))
    else:
        await max_send_message(chat_id, "Генерирую картинку, это может занять немного времени...", attachments=build_image_menu_keyboard(chat_id))
    return True


async def _process_image_edit_generation_queued(chat_id: int, user_prompt: str, reference_image_data_url: str) -> bool:
    prompt = user_prompt.strip()
    if not prompt:
        await max_send_message(chat_id, "Опиши, что изменить на фото.", attachments=build_image_prompt_keyboard())
        return True
    if len(prompt) > MAX_IMAGE_PROMPT_CHARS:
        await max_send_message(chat_id, f"Слишком длинный промпт. Максимум {MAX_IMAGE_PROMPT_CHARS} символов.", attachments=build_keyboard())
        return True
    if image_job_already_pending(chat_id):
        await max_send_message(chat_id, "У тебя уже есть активная задача по картинке. Дождись результата или ошибки, потом запускай следующую.", attachments=build_image_menu_keyboard(chat_id))
        return True

    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free")).strip().lower()
    if plan_name != "free" and not plan_allowed(plan_name, DEFAULT_IMAGE_MODEL.min_plan):
        await show_managed_content(
            chat_id,
            f"Режим «по фото» доступен с тарифа {DEFAULT_IMAGE_MODEL.min_plan}. Открой «Тарифы».",
            attachments=build_tariffs_keyboard_pricing(),
            page=UI_PAGE_TARIFFS,
            push_history=False,
            notification="Тарифы",
        )
        return True

    ok_cd, reason_cd = check_cooldown(chat_id, "image")
    if not ok_cd:
        await max_send_message(chat_id, reason_cd, attachments=build_keyboard())
        return True

    ok, reason = check_limit_only(chat_id, "images")
    if not ok:
        await show_managed_content(
            chat_id,
            reason,
            attachments=build_image_menu_keyboard(chat_id),
            page=UI_PAGE_IMAGE_MENU,
            push_history=False,
            notification="Лимит достигнут",
        )
        return True

    is_free_plan = plan_name == "free"
    edit_cost = 0 if is_free_plan else image_edit_credit_cost()
    ok_credit, reason_credit = check_and_consume_credits(chat_id, edit_cost, "картинка по фото")
    if not ok_credit:
        await show_managed_content(
            chat_id,
            reason_credit,
            attachments=purchase_help_keyboard_for_row(row),
            page=UI_PAGE_TARIFFS,
            push_history=False,
            notification="Недостаточно запросов",
        )
        return True

    ok, reason = check_and_consume_limit(chat_id, "images")
    if not ok:
        state.user_store.refund_credits(chat_id, edit_cost)
        await show_managed_content(
            chat_id,
            reason,
            attachments=build_image_menu_keyboard(chat_id),
            page=UI_PAGE_IMAGE_MENU,
            push_history=False,
            notification="Лимит достигнут",
        )
        return True

    prepared_prompt = build_image_prompt(prompt, chat_id)
    queue_before = await enqueue_image_job(
        ImageJob(
            kind="edit",
            chat_id=chat_id,
            user_prompt=prompt,
            model_prompt=prepared_prompt,
            credits_spent=edit_cost,
            details=f"mode=image_edit;style={get_image_prefs(chat_id).get('style','')};aspect={get_image_prefs(chat_id).get('aspect','')}",
            reference_image_data_url=reference_image_data_url,
        )
    )
    set_image_mode(chat_id, "")
    if queue_before > 0:
        await max_send_message(chat_id, f"Поставил обработку фото в очередь. Перед тобой задач: {queue_before}. Пришлю результат отдельным сообщением.", attachments=build_image_menu_keyboard(chat_id))
    else:
        await max_send_message(chat_id, "Обрабатываю фото и готовлю вариант, это может занять немного времени...", attachments=build_image_menu_keyboard(chat_id))
    return True


async def handle_pending_image_prompt_input(chat_id: int, text: str) -> bool:
    if chat_id not in state.pending_image_prompt:
        return False

    lowered = text.strip().lower()
    if lowered in {"отмена", "cancel", "/cancel", "стоп", "/stop"}:
        state.pending_image_prompt.discard(chat_id)
        await max_send_message(chat_id, "Ок, генерацию отменил.", attachments=build_image_menu_keyboard(chat_id))
        return True

    if text.strip().startswith("/"):
        state.pending_image_prompt.discard(chat_id)
        return False

    state.pending_image_prompt.discard(chat_id)
    prepared_prompt = build_image_prompt(text.strip(), chat_id)
    await process_image_generation(chat_id, text.strip(), model_prompt=prepared_prompt)
    return True


async def handle_pending_image_ref_prompt_input(chat_id: int, text: str) -> bool:
    if chat_id not in state.pending_image_ref_prompt:
        return False

    lowered = text.strip().lower()
    if lowered in {"отмена", "cancel", "/cancel", "стоп", "/stop"}:
        state.pending_image_ref_prompt.discard(chat_id)
        await max_send_message(chat_id, "Ок, режим «по фото» отменил.", attachments=build_image_menu_keyboard(chat_id))
        return True

    if text.strip().startswith("/"):
        state.pending_image_ref_prompt.discard(chat_id)
        return False

    reference = get_recent_reference_image(chat_id)
    if not reference:
        await max_send_message(
            chat_id,
            "Сначала отправь фото, потом напиши, что с ним сделать.",
            attachments=build_image_menu_keyboard(chat_id),
        )
        return True

    state.pending_image_ref_prompt.discard(chat_id)
    await process_image_edit_generation(chat_id, text.strip(), reference)
    return True


async def handle_pending_file_input(chat_id: int, text: str) -> bool:
    kind = state.pending_file_kind.get(chat_id, "")
    profile = state.pending_file_profile.get(chat_id, "medium")
    if kind not in FILE_KIND_LABELS:
        return False

    lowered = text.strip().lower()
    if lowered in {"отмена", "cancel", "/cancel", "стоп", "/stop"}:
        clear_file_pending_input(chat_id)
        await show_ui_page(chat_id, UI_PAGE_FILES_MENU, push_history=False)
        return True

    if text.strip().startswith("/"):
        clear_file_pending_input(chat_id)
        return False

    clear_file_pending_input(chat_id)
    return await process_file_request(chat_id, kind, text.strip(), profile=profile)


async def send_help(chat_id: int) -> None:
    help_base = HELP_TEXT
    admin_part = ADMIN_HELP_TEXT if is_admin(chat_id) else ""
    text = (
        f"{help_base}"
        f"{admin_part}"
    )
    await send_managed_message(chat_id, text, attachments=build_keyboard(chat_id), page=UI_PAGE_MENU)


async def send_menu(chat_id: int) -> None:
    row = user_profile(chat_id)
    preset_block = build_preset_block(str(row.get("plan", "free")))
    text = (
        "Главное меню\n\n"
        "Выбери режим кнопками или просто напиши вопрос.\n\n"
        f"{preset_block}\n\n"
        f"{current_model_focus_block(chat_id)}\n"
        f"{usage_text(row)}\n\n"
        f"{MENU_TEXT}"
    )
    await send_managed_message(chat_id, text, attachments=build_keyboard(chat_id), page=UI_PAGE_MENU)


UI_PAGE_MENU = "menu"
UI_PAGE_MODELS = "models"
UI_PAGE_PLAN = "plan"
UI_PAGE_TARIFFS = "tariffs"
UI_PAGE_TOPUPS = "topups"
UI_PAGE_PAYMENTS = "payments"
UI_PAGE_GROWTH = "growth"
UI_PAGE_SUPPORT = "support"
UI_PAGE_IMAGE_MENU = "image_menu"
UI_PAGE_FILES_MENU = "files_menu"

UI_PAGE_KEYS = {
    UI_PAGE_MENU,
    UI_PAGE_MODELS,
    UI_PAGE_PLAN,
    UI_PAGE_TARIFFS,
    UI_PAGE_TOPUPS,
    UI_PAGE_PAYMENTS,
    UI_PAGE_GROWTH,
    UI_PAGE_SUPPORT,
    UI_PAGE_IMAGE_MENU,
    UI_PAGE_FILES_MENU,
}

UI_HISTORY_LIMIT = 20


def ui_back_stack(chat_id: int) -> list[str]:
    return state.ui_back_stack.setdefault(chat_id, [])


def ui_forward_stack(chat_id: int) -> list[str]:
    return state.ui_forward_stack.setdefault(chat_id, [])


def ui_can_go_back(chat_id: int) -> bool:
    return bool(ui_back_stack(chat_id))


def ui_can_go_forward(chat_id: int) -> bool:
    return bool(ui_forward_stack(chat_id))


def ui_set_page(chat_id: int, page: str, push_history: bool = True) -> None:
    current = state.ui_current_page.get(chat_id)
    if push_history and current is None and page != UI_PAGE_MENU:
        back = ui_back_stack(chat_id)
        back.append(UI_PAGE_MENU)
        ui_forward_stack(chat_id).clear()
    if push_history and current and current != page:
        back = ui_back_stack(chat_id)
        back.append(current)
        if len(back) > UI_HISTORY_LIMIT:
            del back[:-UI_HISTORY_LIMIT]
        ui_forward_stack(chat_id).clear()
    state.ui_current_page[chat_id] = page


def ui_nav_back(chat_id: int) -> str | None:
    back = ui_back_stack(chat_id)
    if not back:
        return None
    target = back.pop()
    current = state.ui_current_page.get(chat_id)
    if current:
        fwd = ui_forward_stack(chat_id)
        fwd.append(current)
        if len(fwd) > UI_HISTORY_LIMIT:
            del fwd[:-UI_HISTORY_LIMIT]
    state.ui_current_page[chat_id] = target
    return target


def ui_nav_forward(chat_id: int) -> str | None:
    fwd = ui_forward_stack(chat_id)
    if not fwd:
        return None
    target = fwd.pop()
    current = state.ui_current_page.get(chat_id)
    if current:
        back = ui_back_stack(chat_id)
        back.append(current)
        if len(back) > UI_HISTORY_LIMIT:
            del back[:-UI_HISTORY_LIMIT]
    state.ui_current_page[chat_id] = target
    return target


def add_ui_nav_buttons(chat_id: int, attachments: list[dict[str, Any]] | None) -> list[dict[str, Any]]:
    result = json.loads(json.dumps(attachments or []))
    nav_row: list[dict[str, Any]] = []
    if ui_can_go_back(chat_id):
        nav_row.append({"type": "callback", "text": "◀ Назад", "payload": "ui_nav:back"})
    if ui_can_go_forward(chat_id):
        nav_row.append({"type": "callback", "text": "↩ Откат", "payload": "ui_nav:forward"})
    if not nav_row:
        return result

    if result and isinstance(result[0], dict) and result[0].get("type") == "inline_keyboard":
        keyboard_payload = result[0].setdefault("payload", {})
        buttons = keyboard_payload.setdefault("buttons", [])
        if isinstance(buttons, list):
            buttons.append(nav_row)
            return result
    return [{"type": "inline_keyboard", "payload": {"buttons": [nav_row]}}]


def resolve_edit_target_mid(chat_id: int, source_mid: str | None, force_new: bool = False) -> str | None:
    if force_new:
        return None
    if source_mid:
        return source_mid
    managed_mid = state.ui_message_mid.get(chat_id)
    if managed_mid:
        return managed_mid
    return None


def current_model_focus_block(chat_id: int) -> str:
    return (
        "────────────\n"
        f"**Сейчас выбрана модель:** {current_model_display(chat_id)}\n"
        "────────────"
    )


def managed_page_text_format(page: str | None) -> str | None:
    if page in UI_PAGE_KEYS:
        return "markdown"
    return None


def build_model_picker_text(chat_id: int) -> str:
    row = user_profile(chat_id)
    plan = str(row.get("plan", "free"))
    current_alias = str(row.get("selected_model_alias") or best_default_alias_for_plan(plan))
    lines = ["⚙ Модели", "", "Здесь можно вручную выбрать текстовую модель."]
    if current_alias:
        current_label = TEXT_MODELS.get(current_alias, DEFAULT_TEXT_MODEL).label
        lines.extend(["", f"**Сейчас выбрана:** {current_label}"])
    for alias in allowed_text_aliases_for_plan(plan):
        info = TEXT_MODELS.get(alias, DEFAULT_TEXT_MODEL)
        prefix = "•"
        if alias == current_alias:
            prefix = "✅"
        lines.append(f"{prefix} {info.label} — {preset_model_hint(alias)}")
    lines.append("")
    lines.append("После выбора верну тебя в Главное меню.")
    return "\n".join(part for part in lines if part is not None)


def build_model_picker_keyboard(plan: str) -> list[dict[str, Any]]:
    rows: list[list[dict[str, Any]]] = []
    for alias in allowed_text_aliases_for_plan(plan):
        info = TEXT_MODELS.get(alias, DEFAULT_TEXT_MODEL)
        rows.append([{"type": "callback", "text": info.label, "payload": f"set_model:{alias}"}])
    rows.append(
        [
            {"type": "callback", "text": "Меню", "payload": "action:menu"},
            {"type": "callback", "text": "Помощь", "payload": "action:support"},
        ]
    )
    return [{"type": "inline_keyboard", "payload": {"buttons": rows}}]


def build_topups_text() -> str:
    small = TOPUP_PACKS["small"]
    medium = TOPUP_PACKS["medium"]
    large = TOPUP_PACKS["large"]

    def approx_images(credits: int) -> int:
        if image_credit_cost() <= 0:
            return 0
        return credits // image_credit_cost()

    return (
        "⭐ Пакеты запросов\n\n"
        f"• Small: {request_balance_text(int(small['credits']))} запросов за {small['price_rub']} ₽ (~{approx_images(int(small['credits']))} картинок)\n"
        f"• Medium: {request_balance_text(int(medium['credits']))} запросов за {medium['price_rub']} ₽ (~{approx_images(int(medium['credits']))} картинок)\n"
        f"• Large: {request_balance_text(int(large['credits']))} запросов за {large['price_rub']} ₽ (~{approx_images(int(large['credits']))} картинок)\n\n"
        "Запросы списываются за ответы моделей и генерацию картинок.\n"
        "Перед созданием оплаты бот попросит подтверждение покупки пакета."
    )


def active_promo_lines() -> list[str]:
    promo_items = sorted(promo_catalog().items())
    lines = [f"• {code}: +{request_balance_text(credits)} запросов" for code, credits in promo_items[:6]]
    channel = {"enabled": False, "active": False, "code": "", "credits": 0, "days_left": 0, "bonus_ttl_days": 0}
    if channel["enabled"] and channel["active"]:
        lines.append(
            f"• {channel['code']}: +{request_balance_text(int(channel['credits']))} запросов "
            f"(акция {channel['days_left']} дн, бонус на {channel['bonus_ttl_days']} дн)"
        )
    return lines


def build_growth_text(chat_id: int, row: dict[str, Any] | None = None) -> str:
    data = row or user_profile(chat_id)
    referral_code = str(data.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
    invited = int(data.get("referrals_invited", 0) or 0)
    promo_lines = active_promo_lines()
    promo_block = "\n".join(promo_lines) if promo_lines else "• Сейчас активных промокодов нет"
    welcome_line = (
        f"Базовый промокод: /promo WELCOME (+{request_balance_text(PROMO_WELCOME_CREDITS)} запросов, 1 раз)\n"
        if PROMO_WELCOME_CREDITS > 0
        else ""
    )
    return (
        "🎁 Бонусы и приглашения\n\n"
        f"Твой реф-код: {referral_code}\n"
        f"Приглашено друзей: {invited}\n"
        f"Бонус за друга: {request_balance_text(REFERRAL_BONUS_CREDITS)} запросов тебе и другу.\n"
        "Кнопка «Поделиться» откроет готовый текст приглашения.\n\n"
        "Промокоды из канала и акций:\n"
        f"{promo_block}\n\n"
        f"{welcome_line}"
        "Если в канале опубликован код, нажми «Промокод» и введи его здесь."
    )


def build_ui_page_payload(chat_id: int, page: str) -> tuple[str, list[dict[str, Any]]]:
    row = user_profile(chat_id)
    if page == UI_PAGE_MENU:
        preset_block = build_preset_block(str(row.get("plan", "free")))
        capabilities = (
            "Что умею:\n"
            "• ⚡ ответы через GPT, Gemini и DeepSeek\n"
            f"• 🎨 {image_capability_line().replace('• ', '')}\n"
            "• 🧠 сохранение контекста диалога"
        )
        text = (
            "Главное меню\n\n"
            "Выбери режим кнопками или просто напиши вопрос.\n\n"
            f"{capabilities}\n\n"
            f"{preset_block}\n\n"
            f"{current_model_focus_block(chat_id)}\n"
            f"{usage_text(row)}\n\n"
            f"{MENU_TEXT}"
        )
        return text, build_keyboard(chat_id)
    if page == UI_PAGE_MODELS:
        return build_models_text(row["plan"], include_prices=False), build_keyboard(chat_id)
    if page == UI_PAGE_PLAN:
        return f"{usage_text(row)}{recurring_status_text(row)}", build_plan_keyboard(row)
    if page == UI_PAGE_TARIFFS:
        return build_tariffs_text(), build_tariffs_keyboard_pricing()
    if page == UI_PAGE_TOPUPS:
        return build_topups_text(), build_topups_keyboard()
    if page == UI_PAGE_PAYMENTS:
        return build_payments_text(chat_id)
    if page == UI_PAGE_GROWTH:
        return build_growth_text(chat_id, row), build_growth_keyboard(chat_id)
    if page == UI_PAGE_SUPPORT:
        return support_help_text(), build_keyboard(chat_id)
    if page == UI_PAGE_IMAGE_MENU:
        return build_image_menu_text(chat_id), build_image_menu_keyboard(chat_id)
    return "Открой меню и выбери раздел.", build_keyboard(chat_id)


def build_ui_page_payload(chat_id: int, page: str) -> tuple[str, list[dict[str, Any]]]:
    row = user_profile(chat_id)
    if page == UI_PAGE_MENU:
        preset_block = build_preset_block(str(row.get("plan", "free")))
        capabilities = (
            "Что умею:\n"
            "• ⚡ ответы через GPT, Gemini и DeepSeek\n"
            f"• 🎨 {image_capability_line().replace('• ', '')}\n"
            "• 📄 документы, презентации и таблицы\n"
            "• 🧠 сохранение контекста диалога"
        )
        text = (
            "Главное меню\n\n"
            "Выбери режим кнопками или просто напиши вопрос.\n\n"
            f"{capabilities}\n\n"
            f"{preset_block}\n\n"
            f"{current_model_focus_block(chat_id)}\n"
            f"{usage_text(row)}\n\n"
            f"{MENU_TEXT}"
        )
        return text, build_keyboard(chat_id)
    if page == UI_PAGE_MODELS:
        return build_models_text(row["plan"], include_prices=False), build_keyboard(chat_id)
    if page == UI_PAGE_PLAN:
        return f"{usage_text(row)}{recurring_status_text(row)}", build_plan_keyboard(row)
    if page == UI_PAGE_TARIFFS:
        return build_tariffs_text(), build_tariffs_keyboard_pricing()
    if page == UI_PAGE_TOPUPS:
        return build_topups_text(), build_topups_keyboard()
    if page == UI_PAGE_PAYMENTS:
        return build_payments_text(chat_id)
    if page == UI_PAGE_GROWTH:
        return build_growth_text(chat_id, row), build_growth_keyboard(chat_id)
    if page == UI_PAGE_SUPPORT:
        return support_help_text(), build_keyboard(chat_id)
    if page == UI_PAGE_IMAGE_MENU:
        return build_image_menu_text(chat_id), build_image_menu_keyboard(chat_id)
    if page == UI_PAGE_FILES_MENU:
        return build_files_menu_text(chat_id), build_files_menu_keyboard(chat_id)
    return "Открой меню и выбери раздел.", build_keyboard(chat_id)


async def show_ui_page(
    chat_id: int,
    page: str,
    callback_id: str | None = None,
    source_mid: str | None = None,
    push_history: bool = True,
    notification: str = "Открываю",
    force_new: bool = False,
) -> None:
    text, attachments = build_ui_page_payload(chat_id, page)
    ui_set_page(chat_id, page, push_history=push_history)
    attachments = add_ui_nav_buttons(chat_id, attachments)

    target_mid = resolve_edit_target_mid(chat_id, source_mid, force_new=force_new)
    text_format = managed_page_text_format(page)
    if target_mid:
        ok = await max_edit_message(chat_id, target_mid, text, attachments=attachments, text_format=text_format)
        if ok:
            state.ui_message_mid[chat_id] = target_mid
            record_ui_page_view(chat_id, page)
            if callback_id:
                await answer_callback(callback_id, notification)
            return

    sent_mid = await max_send_message(chat_id, text, attachments=attachments, notify=False, text_format=text_format)
    if sent_mid:
        state.ui_message_mid[chat_id] = sent_mid
        record_ui_page_view(chat_id, page)
    if callback_id:
        await answer_callback(callback_id, notification)


async def show_managed_content(
    chat_id: int,
    text: str,
    attachments: list[dict[str, Any]] | None = None,
    callback_id: str | None = None,
    source_mid: str | None = None,
    page: str | None = None,
    push_history: bool = False,
    notification: str = "Открываю",
    force_new: bool = False,
) -> None:
    if page in UI_PAGE_KEYS:
        ui_set_page(chat_id, str(page), push_history=push_history)
    attachments = add_ui_nav_buttons(chat_id, attachments)

    target_mid = resolve_edit_target_mid(chat_id, source_mid, force_new=force_new)
    text_format = managed_page_text_format(page)
    if target_mid:
        ok = await max_edit_message(chat_id, target_mid, text, attachments=attachments, text_format=text_format)
        if ok:
            state.ui_message_mid[chat_id] = target_mid
            record_ui_page_view(chat_id, page)
            if callback_id:
                await answer_callback(callback_id, notification)
            return

    sent_mid = await max_send_message(chat_id, text, attachments=attachments, notify=False, text_format=text_format)
    if sent_mid:
        state.ui_message_mid[chat_id] = sent_mid
        record_ui_page_view(chat_id, page)
    if callback_id:
        await answer_callback(callback_id, notification)


async def show_channel_gate(
    chat_id: int,
    callback_id: str | None = None,
    source_mid: str | None = None,
    notification: str = "Нужна подписка",
    reason: str = "",
) -> None:
    text = channel_gate_setup_text(reason) if reason.startswith(("config_", "api_")) else channel_gate_text()
    await show_managed_content(
        chat_id,
        text,
        attachments=build_channel_gate_keyboard(),
        callback_id=callback_id,
        source_mid=source_mid,
        notification=notification,
        force_new=False,
    )


async def ensure_channel_access(
    chat_id: int,
    callback_id: str | None = None,
    source_mid: str | None = None,
) -> bool:
    ok, reason = await check_channel_subscription(chat_id, force=False)
    if ok:
        return True
    await show_channel_gate(chat_id, callback_id=callback_id, source_mid=source_mid, reason=reason)
    return False


async def send_managed_message(
    chat_id: int,
    text: str,
    attachments: list[dict[str, Any]] | None = None,
    notify: bool = False,
    page: str | None = None,
    push_history: bool = False,
) -> str | None:
    if page in UI_PAGE_KEYS:
        ui_set_page(chat_id, page, push_history=push_history)
    sent_mid = await max_send_message(chat_id, text, attachments=attachments, notify=notify, text_format=managed_page_text_format(page))
    if sent_mid:
        state.ui_message_mid[chat_id] = sent_mid
        record_ui_page_view(chat_id, page)
    return sent_mid


def build_onboarding_text(chat_id: int) -> str:
    row = user_profile(chat_id)
    return (
        "Как пользоваться ботом?\n"
        "────────────────────\n\n"
        "Можно ничего не выбирать и просто написать вопрос в чат. Бот ответит выбранной моделью и сохранит контекст диалога.\n\n"
        "Кнопки помогают быстрее выбрать режим:\n"
        "• ⚡ Быстро — короткие ответы и простые задачи\n"
        "• ⚖ Баланс — повседневные вопросы и диалог\n"
        "• 🧠 Качество — подробнее и аккуратнее\n"
        "• 🚀 Эксперт — сложные задачи, если доступно на тарифе\n\n"
        "Для картинок нажми «🎨 Картинка»: можно сгенерировать новую картинку или редактировать фото.\n\n"
        "Запросы списываются за ответы, картинки и модели. Баланс, тариф и дату следующей free-картинки можно смотреть в «Мой план».\n\n"
        f"Сейчас у тебя тариф: {str(row.get('plan', 'free')).title()}.\n"
        "Нажми «Начать!» — открою Главное меню."
    )


async def send_onboarding(chat_id: int, step: int = 1, notify: bool = False) -> None:
    text = build_onboarding_text(chat_id)
    sent_mid = await max_send_message(chat_id, text, attachments=build_onboarding_keyboard(step), notify=notify)
    if sent_mid:
        state.onboarding_message_mid[chat_id] = sent_mid


async def show_onboarding_step(
    chat_id: int,
    step: int,
    callback_id: str | None = None,
    source_mid: str | None = None,
    notification: str = "Открываю",
) -> None:
    text = build_onboarding_text(chat_id)
    target_mid = source_mid or state.onboarding_message_mid.get(chat_id)
    if target_mid:
        ok = await max_edit_message(chat_id, target_mid, text, attachments=build_onboarding_keyboard(step))
        if ok:
            state.onboarding_message_mid[chat_id] = target_mid
            if callback_id:
                await answer_callback(callback_id, notification)
            return
    sent_mid = await max_send_message(chat_id, text, attachments=build_onboarding_keyboard(step), notify=False)
    if sent_mid:
        state.onboarding_message_mid[chat_id] = sent_mid
    if callback_id:
        await answer_callback(callback_id, notification)


async def close_onboarding_message(chat_id: int, source_mid: str | None, text: str = " ") -> None:
    target_mid = source_mid or state.onboarding_message_mid.get(chat_id)
    if not target_mid:
        return
    with suppress(Exception):
        await max_edit_message(chat_id, target_mid, text, attachments=[])
    state.onboarding_message_mid.pop(chat_id, None)


async def finish_onboarding_to_page(
    chat_id: int,
    page: str,
    callback_id: str | None = None,
    source_mid: str | None = None,
    notification: str = "Готово",
) -> None:
    state.user_store.set_onboarding_done(chat_id, True)
    handoff_onboarding_to_ui(chat_id, source_mid)
    target_mid = state.ui_message_mid.get(chat_id) or source_mid
    await show_ui_page(
        chat_id,
        page,
        callback_id=callback_id,
        source_mid=target_mid,
        push_history=False,
        notification=notification,
    )

async def send_growth_menu(chat_id: int) -> None:
    await send_managed_message(
        chat_id,
        build_growth_text(chat_id),
        attachments=build_growth_keyboard(chat_id),
        page=UI_PAGE_GROWTH,
    )


async def send_channel(chat_id: int) -> None:
    if CHANNEL_GATE_ENABLED:
        subscribed, reason = await check_channel_subscription(chat_id, force=False)
        if not subscribed:
            await show_channel_gate(chat_id, notification="Канал", reason=reason)
            return
    await send_managed_message(
        chat_id,
        f"📣 Канал проекта:\n{channel_url_value()}",
        attachments=build_keyboard(),
    )


def log_referral_activation(friend_chat_id: int, owner_chat_id: int, code: str, source: str) -> None:
    friend_row = user_profile(friend_chat_id)
    owner_row = user_profile(owner_chat_id)
    normalized = normalize_referral_code(code)
    state.user_store.record_usage_event(
        chat_id=friend_chat_id,
        event_type="referral_activation",
        plan=str(friend_row.get("plan", "")),
        credits_spent=REFERRAL_BONUS_CREDITS,
        details=f"code={normalized};owner_chat_id={owner_chat_id};source={source}",
    )
    state.user_store.record_usage_event(
        chat_id=owner_chat_id,
        event_type="referral_reward",
        plan=str(owner_row.get("plan", "")),
        credits_spent=REFERRAL_BONUS_CREDITS,
        details=f"code={normalized};friend_chat_id={friend_chat_id};source={source}",
    )


async def notify_referral_success(friend_chat_id: int, owner_chat_id: int, source: str) -> None:
    source_tail = ""
    if source == "start":
        source_tail = " Код применился автоматически по стартовой ссылке."
    await show_managed_content(
        friend_chat_id,
        f"Готово! Реферальный код принят. Начислено +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов.{source_tail}",
        attachments=build_growth_keyboard(friend_chat_id),
        page=UI_PAGE_GROWTH,
    )
    with suppress(Exception):
        await max_send_message(
            owner_chat_id,
            f"🎉 По твоему коду зарегистрировался друг. Начислено +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов.",
            attachments=build_keyboard(),
            notify=False,
        )


async def maybe_apply_start_referral(chat_id: int, update: dict[str, Any]) -> bool:
    code = referral_code_from_start_payload(update)
    if not code:
        return False
    ok, info = state.user_store.apply_referral_code(chat_id, code, REFERRAL_BONUS_CREDITS)
    if not ok:
        log.info("Start referral ignored chat_id=%s code=%s reason=%s", chat_id, code, info)
        return False
    owner_chat_id = int(info)
    log_referral_activation(chat_id, owner_chat_id, code, "start")
    await notify_referral_success(chat_id, owner_chat_id, "start")
    return True


async def send_reengage_nudges(days: int, limit: int) -> tuple[int, int]:
    targets = state.user_store.list_reengage_candidates(days, limit)
    sent = 0
    total = len(targets)
    for item in targets:
        target = int(item.get("chat_id", 0) or 0)
        if target <= 0:
            continue
        with suppress(Exception):
            await max_send_message(
                target,
                (
                    "👋 Возвращайся, у тебя снова доступна 1 бесплатная картинка.\n"
                    "Нажми «🎨 Картинка» и отправь идею."
                ),
                attachments=build_keyboard(),
                notify=False,
            )
            sent += 1
    return sent, total


async def send_models(chat_id: int) -> None:
    row = user_profile(chat_id)
    await send_managed_message(chat_id, build_models_text(row["plan"], include_prices=False), attachments=build_keyboard(), page=UI_PAGE_MODELS)


async def send_costs(chat_id: int) -> None:
    row = user_profile(chat_id)
    await send_managed_message(chat_id, build_models_text(row["plan"], include_prices=True), attachments=build_keyboard())


async def send_plan(chat_id: int) -> None:
    row = user_profile(chat_id)
    text = f"{usage_text(row)}{recurring_status_text(row)}"
    await send_managed_message(chat_id, text, attachments=build_plan_keyboard(row), page=UI_PAGE_PLAN)


async def send_credits(chat_id: int) -> None:
    row = user_profile(chat_id)
    plan_name = str(row.get("plan", "free"))
    if plan_name not in PAID_PLANS:
        await send_managed_message(
            chat_id,
            f"🆓 На free каждый день доступно {request_balance_text(FREE_DAILY_CREDITS)} запросов.\nСейчас у тебя: {request_balance_text(int(row.get('credits_balance', 0) or 0))}.",
            attachments=build_tariffs_keyboard_pricing(),
            page=UI_PAGE_TARIFFS,
        )
        return
    text = (
        f"🪙 Твой баланс: {request_balance_text(int(row.get('credits_balance', 0) or 0))} запросов.\n\n"
        "Обычно списывается:\n"
        "• Текст: 1 запрос\n"
        "• Длинный ответ: до 2 запросов\n"
        f"• GPT-5.4: {request_cost_text(text_credit_cost('gpt54'))}-{request_cost_text(text_credit_cost('gpt54') + MAX_VARIABLE_CREDITS_PER_TEXT)} запросов\n"
        f"• Картинка: {request_cost_text(image_credit_cost())} запросов\n"
        f"• Редактирование фото: {request_cost_text(image_edit_credit_cost())} запросов"
    )
    await send_managed_message(chat_id, text, attachments=build_keyboard(), page=UI_PAGE_PLAN)


async def send_topups(chat_id: int) -> None:
    small = TOPUP_PACKS["small"]
    medium = TOPUP_PACKS["medium"]
    large = TOPUP_PACKS["large"]

    def approx_images(credits: int) -> int:
        if image_credit_cost() <= 0:
            return 0
        return credits // image_credit_cost()

    text = (
        "⭐ Пакеты запросов\n\n"
        f"• Small: {request_balance_text(int(small['credits']))} запросов за {small['price_rub']} ₽ (~{approx_images(int(small['credits']))} картинок)\n"
        f"• Medium: {request_balance_text(int(medium['credits']))} запросов за {medium['price_rub']} ₽ (~{approx_images(int(medium['credits']))} картинок)\n"
        f"• Large: {request_balance_text(int(large['credits']))} запросов за {large['price_rub']} ₽ (~{approx_images(int(large['credits']))} картинок)\n\n"
        "Запросы списываются за ответы моделей и генерацию картинок.\n"
        "Перед созданием оплаты бот попросит подтверждение покупки пакета."
    )
    await send_managed_message(chat_id, text, attachments=build_topups_keyboard(), page=UI_PAGE_TOPUPS)


async def send_topup_consent(
    chat_id: int,
    code: str,
    notify: bool = False,
    callback_id: str | None = None,
    source_mid: str | None = None,
) -> bool:
    pack = topup_spec(code)
    if not pack:
        await max_send_message(chat_id, "Неизвестный пакет запросов.", attachments=build_topups_keyboard(), notify=notify)
        return False
    text = (
        f"Пакет: {pack['label']}\n"
        f"Запросов: {request_balance_text(int(pack['credits']))}\n"
        f"Стоимость: {int(pack['price_rub'])} ₽\n\n"
        "Это разовая покупка без автосписаний.\n"
        "Подтверди покупку кнопкой ниже."
    )
    await show_managed_content(
        chat_id,
        text,
        attachments=build_topup_consent_keyboard(code),
        callback_id=callback_id,
        source_mid=source_mid,
        notification="Открываю покупку",
    )
    return True


async def send_payments(chat_id: int) -> None:
    rows = state.user_store.list_user_payments(chat_id, limit=8)
    if not rows:
        await send_managed_message(chat_id, "Заявок пока нет. Используй кнопку «Тарифы».", attachments=build_keyboard(), page=UI_PAGE_PAYMENTS)
        return
    lines = ["Твои последние заявки:"]
    for item in rows:
        status = str(item["status"]).lower()
        status_human = payment_status_label(status)
        lines.append(
            f"#{item['id']} | {item['plan']} | {item['days']} дн | {item['amount_rub']} RUB | {status_human} | {item['created_at'][:19]}"
        )
    lines.append("\nНажми «Проверить #...», чтобы обновить статус по банку.")
    await send_managed_message(chat_id, "\n".join(lines), attachments=build_payments_keyboard(rows), page=UI_PAGE_PAYMENTS)


def payment_item_human_name(plan: str) -> str:
    if is_topup_plan(plan):
        code = topup_code_from_plan(plan)
        pack = topup_spec(code)
        return f"Пакет {pack['label']}" if pack else "Пакет запросов"
    return f"Тариф {plan.capitalize()}"


def payment_return_target(plan: str) -> tuple[str, str]:
    if is_topup_plan(plan):
        return "action:topups", "К пакетам"
    return "action:tariffs", "К тарифам"


def build_payment_request_keyboard(request_id: int, payment_url: str = "", status: str = "pending") -> list[dict[str, Any]]:
    payment = state.user_store.get_payment(request_id) or {}
    plan = str(payment.get("plan", "") or "")
    return_payload, return_label = payment_return_target(plan)
    buttons: list[list[dict[str, Any]]] = []
    status_lc = str(status or "").lower()
    if payment_url and status_lc in {"pending", "claimed"}:
        buttons.append([{"type": "link", "text": "Оплатить", "url": payment_url}])
    if status_lc in {"pending", "claimed"}:
        buttons.append([{"type": "callback", "text": "Проверить статус", "payload": f"payment_status:{request_id}"}])
        buttons.append([{"type": "callback", "text": "Я оплатил", "payload": f"paid:{request_id}"}])
        buttons.append(
            [
                {"type": "callback", "text": "Мои оплаты", "payload": "action:payments"},
                {"type": "callback", "text": return_label, "payload": return_payload},
            ]
        )
        buttons.append(
            [
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
            ]
        )
    elif status_lc == "paid":
        primary_payload, primary_label = ("action:topups", "Пакеты запросов") if is_topup_plan(plan) else ("action:plan", "Мой план")
        buttons.append(
            [
                {"type": "callback", "text": primary_label, "payload": primary_payload},
                {"type": "callback", "text": "Мои оплаты", "payload": "action:payments"},
            ]
        )
        buttons.append(
            [
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
            ]
        )
    else:
        buttons.append(
            [
                {"type": "callback", "text": return_label, "payload": return_payload},
                {"type": "callback", "text": "Мои оплаты", "payload": "action:payments"},
            ]
        )
        buttons.append(
            [
                {"type": "callback", "text": "Меню", "payload": "action:menu"},
                {"type": "callback", "text": "Помощь", "payload": "action:support"},
            ]
        )
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def build_payments_keyboard(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    buttons: list[list[dict[str, Any]]] = []
    for item in rows[:5]:
        request_id = int(item.get("id", 0) or 0)
        if request_id <= 0:
            continue
        status = str(item.get("status", "")).lower()
        buttons.append(
            [
                {
                    "type": "callback",
                    "text": f"Открыть #{request_id}" if status in PAYMENT_FINAL_STATUSES else f"Проверить #{request_id}",
                    "payload": f"payment_status:{request_id}",
                }
            ]
        )
    buttons.append(
        [
            {"type": "callback", "text": "Мой план", "payload": "action:plan"},
            {"type": "callback", "text": "Тарифы", "payload": "action:tariffs"},
        ]
    )
    buttons.append(
        [
            {"type": "callback", "text": "Меню", "payload": "action:menu"},
            {"type": "callback", "text": "Помощь", "payload": "action:support"},
        ]
    )
    return [{"type": "inline_keyboard", "payload": {"buttons": buttons}}]


def build_payments_text(chat_id: int) -> tuple[str, list[dict[str, Any]]]:
    rows = state.user_store.list_user_payments(chat_id, limit=8)
    if not rows:
        return "Моих оплат пока нет.\nОткрой «Тарифы» или «Пакеты запросов», чтобы создать первую заявку.", build_keyboard()
    lines = ["💳 Мои оплаты"]
    for item in rows:
        status = str(item.get("status", "")).lower()
        status_human = payment_status_label(status)
        item_name = payment_item_human_name(str(item.get("plan", "") or ""))
        amount = int(item.get("amount_rub", 0) or 0)
        created_at = parse_iso_datetime(str(item.get("created_at", "") or ""))
        created_text = format_msk_datetime(created_at) if created_at else "-"
        lines.append(f"• #{item['id']} — {item_name} • {amount} ₽ • {status_human} • {created_text}")
    lines.append("\nНажми на заявку ниже, чтобы открыть её статус, оплату или вернуться к покупке.")
    return "\n".join(lines), build_payments_keyboard(rows)


def usage_text(row: dict[str, Any]) -> str:
    plan_name = str(row.get("plan", "free"))
    cfg = PLAN_CONFIGS.get(plan_name, PLAN_CONFIGS["free"])
    gpt54_used = int(row.get("daily_gpt54_used", 0) or 0)
    gpt54_left = max(0, cfg.daily_gpt54_limit - gpt54_used)
    expires_at = parse_iso_datetime(str(row.get("subscription_expires_at", "") or ""))
    balance = int(row.get("credits_balance", 0) or 0)
    bonus_total, bonus_expires = state.user_store.active_bonus_credits_summary(int(row.get("chat_id", 0) or 0))

    lines = [f"План: {plan_name}", f"Запросы: {request_balance_text(balance)}"]
    if plan_name != "free":
        lines.insert(1, f"Доступ до: {format_msk_datetime(expires_at)}")
    if bonus_total > 0 and bonus_expires:
        bonus_dt = parse_iso_datetime(bonus_expires)
        lines.append(f"🎁 Бонус: {request_balance_text(bonus_total)} запросов до {format_msk_datetime(bonus_dt)}")
    if plan_name == "free":
        lines.append(f"Free-бонус в день: {request_balance_text(FREE_DAILY_CREDITS)} запросов")
        next_at = free_image_next_available_at(row)
        if free_image_is_available(row):
            lines.append("Действие с картинкой на Free: доступно сейчас")
        else:
            lines.append(f"Действие с картинкой на Free: через {format_remaining_time(next_at)}")
            lines.append(f"Точное время: {format_msk_datetime(next_at)}")
    if cfg.daily_gpt54_limit > 0:
        lines.append(f"GPT-5.4 сегодня: {gpt54_used}/{cfg.daily_gpt54_limit} (осталось {gpt54_left})")
    return "\n".join(lines)


def recurring_status_text(row: dict[str, Any]) -> str:
    plan_name = str(row.get("plan", "free"))
    if plan_name == "free":
        return ""
    expires_at = parse_iso_datetime(str(row.get("subscription_expires_at", "") or ""))
    expires_text = format_msk_datetime(expires_at)
    if recurring_enabled_for_row(row):
        return (
            "\n\nАвтопродление: включено.\n"
            f"Следующее списание не раньше {expires_text}.\n"
            "Отключить можно кнопкой «Отменить подписку»."
        )
    canceled_at_raw = str(row.get("recurring_canceled_at", ""))
    if not canceled_at_raw:
        return "\n\nАвтопродление: не подключено."
    cancel_from = parse_iso_datetime(str(row.get("recurring_cancel_from", "") or ""))
    cancel_text = format_msk_datetime(cancel_from) if cancel_from else expires_text
    return (
        "\n\nАвтопродление: отключено.\n"
        f"Доступ сохранится до {expires_text}.\n"
        f"Повторных списаний не будет с {cancel_text}."
    )


def payment_user_status_text(payment: dict[str, Any], bank_status: str = "") -> str:
    request_id = int(payment.get("id", 0) or 0)
    status = str(payment.get("status", "pending")).lower()
    status_human = payment_status_label(status)
    plan = str(payment.get("plan", "") or "")
    amount = int(payment.get("amount_rub", 0) or 0)
    payment_url = str(payment.get("payment_url", "") or "").strip()
    item_name = payment_item_human_name(plan)
    bank_line = f"\nСтатус банка: {bank_status}" if bank_status else ""
    pay_line = f"\nСсылка на оплату: {payment_url}" if payment_url and status in {"pending", "claimed"} else ""
    if status == "paid":
        return (
            f"✅ Заявка #{request_id}: {status_human}\n"
            f"{item_name} • {amount} ₽\n"
            "Доступ уже активирован. Проверь «Мой план»."
        )
    if status == "refunded":
        return (
            f"↩️ Заявка #{request_id}: {status_human}\n"
            f"{item_name} • {amount} ₽{bank_line}\n"
            "Возврат подтвержден. Если деньги долго не приходят, напиши в поддержку."
        )
    if status == "canceled":
        return (
            f"❌ Заявка #{request_id}: {status_human}\n"
            f"{item_name} • {amount} ₽{bank_line}\n"
            "Оплата не завершена. Можно вернуться к покупке и создать новую заявку."
        )
    if status == "claimed":
        return (
            f"🕒 Заявка #{request_id}: {status_human}\n"
            f"{item_name} • {amount} ₽{bank_line}{pay_line}\n"
            "Платеж уже ушел на проверку. Обычно подтверждение занимает 1–2 минуты."
        )
    return (
        f"⏳ Заявка #{request_id}: {status_human}\n"
        f"{item_name} • {amount} ₽{bank_line}{pay_line}\n"
        "Открой оплату по кнопке ниже. После оплаты нажми «Проверить статус»."
    )


def effective_receipt_contact(row: dict[str, Any]) -> tuple[str, str]:
    email = str(row.get("receipt_email", "")).strip()
    phone = str(row.get("receipt_phone", "")).strip()
    return email, phone


def receipt_return_page(plan: str) -> str:
    return UI_PAGE_TOPUPS if plan.startswith("topup") else UI_PAGE_TARIFFS


async def request_receipt_contact(
    chat_id: int,
    plan: str,
    notify: bool = False,
    callback_id: str | None = None,
    source_mid: str | None = None,
) -> None:
    state.pending_receipt_plan[chat_id] = plan
    target_label = "подписки" if not plan.startswith("topup") else "пакета запросов"
    await show_managed_content(
        chat_id,
        (
            f"Перед оплатой {target_label} нужен контакт для отправки чека.\n"
            "Отправь одним сообщением email или телефон.\n\n"
            "Пример email: user@example.com\n"
            "Пример телефона: +79991234567\n\n"
            "Можно нажать «Отмена» ниже."
        ),
        attachments=build_receipt_contact_keyboard(),
        callback_id=callback_id,
        source_mid=source_mid,
        notification="Нужен контакт для чека",
        page=UI_PAGE_SUPPORT,
    )


async def start_buy_flow(
    chat_id: int,
    plan: str,
    notify: bool = False,
    callback_id: str | None = None,
    source_mid: str | None = None,
) -> bool:
    if plan not in BUYABLE_PLANS:
        await max_send_message(chat_id, "Доступно: lite, start или pro.", attachments=build_tariffs_keyboard_pricing(), notify=notify)
        return False
    row = user_profile(chat_id)
    email, phone = effective_receipt_contact(row)
    if not (email or phone):
        await request_receipt_contact(chat_id, plan, notify=notify, callback_id=callback_id, source_mid=source_mid)
        return False
    return await send_buy_consent(chat_id, plan, notify=notify, callback_id=callback_id, source_mid=source_mid)


async def send_buy_consent(
    chat_id: int,
    plan: str,
    notify: bool = False,
    callback_id: str | None = None,
    source_mid: str | None = None,
) -> bool:
    if plan not in BUYABLE_PLANS:
        await max_send_message(chat_id, "Доступно: lite, start или pro.", attachments=build_tariffs_keyboard_pricing(), notify=notify)
        return False

    amount, days = plan_price_and_days(plan)
    row = user_profile(chat_id)
    receipt_email, receipt_phone = effective_receipt_contact(row)
    receipt_label = receipt_email or receipt_phone or "не указан"
    text = (
        "Перед оплатой нужно согласие на автопродление.\n\n"
        f"Тариф: {plan}\n"
        f"Списание: {amount} ₽\n"
        f"Периодичность: каждые {days} дней\n\n"
        f"Чек придет на: {receipt_label}\n\n"
        "После согласия откроется оплата.\n"
        "Отменить автопродление можно в «Мой план»."
    )
    await show_managed_content(
        chat_id,
        text,
        attachments=build_consent_keyboard(plan),
        callback_id=callback_id,
        source_mid=source_mid,
        notification="Проверь условия",
    )
    return True


async def create_buy_request(chat_id: int, plan: str) -> str:
    if plan not in BUYABLE_PLANS:
        return "Доступно: lite, start или pro."
    amount, days = plan_price_and_days(plan)
    request_id = state.user_store.create_payment_request(chat_id, plan, days, amount)
    return (
        f"Заявка #{request_id} создана: {plan}, {days} дн, {amount} RUB.\n"
        "Сейчас подтверждение оплаты делает админ вручную."
    )


async def notify_admin_about_payment_claim(request_id: int, payment: dict[str, Any]) -> None:
    targets = admin_target_chat_ids()
    if not targets:
        return
    target = int(payment["chat_id"])
    plan = str(payment["plan"])
    amount = int(payment["amount_rub"])
    days = int(payment["days"])
    item = plan
    if is_topup_plan(plan):
        code = topup_code_from_plan(plan)
        pack = topup_spec(code)
        if pack:
            item = f"topup:{code} ({request_balance_text(int(pack['credits']))} requests)"
    text = (
        f"Пользователь подтвердил оплату по заявке #{request_id}.\n"
        f"user={target}, item={item}, amount={amount} RUB, days={days}\n"
        f"Проверка/отмена: /admin/panel?request_id={request_id}"
    )
    for admin_id in targets:
        with suppress(Exception):
            await max_send_message(admin_id, text)


def image_capability_line() -> str:
    label = DEFAULT_IMAGE_MODEL.label
    return f"• генерация картинок через {label}"


def resolve_tbank_notification_url() -> str:
    if TBANK_NOTIFICATION_URL:
        return TBANK_NOTIFICATION_URL
    if PUBLIC_BASE_URL:
        return f"{PUBLIC_BASE_URL}/webhook/tbank"
    return ""


def resolve_tbank_success_url() -> str:
    if TBANK_SUCCESS_URL:
        return TBANK_SUCCESS_URL
    if PUBLIC_BASE_URL:
        return f"{PUBLIC_BASE_URL}/payment/success"
    return ""


def resolve_tbank_fail_url() -> str:
    if TBANK_FAIL_URL:
        return TBANK_FAIL_URL
    if PUBLIC_BASE_URL:
        return f"{PUBLIC_BASE_URL}/payment/fail"
    return ""


async def tbank_init_payment(
    request_id: int,
    amount_rub: int,
    description: str,
    receipt_email: str = "",
    receipt_phone: str = "",
) -> tuple[str, str]:
    if not tbank_enabled():
        raise RuntimeError("T-Bank credentials are not configured.")

    order_id = tbank_order_id(request_id)
    payload: dict[str, Any] = {
        "TerminalKey": TBANK_TERMINAL_KEY,
        "Amount": int(amount_rub) * 100,
        "OrderId": order_id,
        "Description": description[:140],
        "PayType": "O",
    }
    final_email = receipt_email.strip()
    final_phone = receipt_phone.strip()
    if not (final_email or final_phone):
        raise RuntimeError("T-Bank Receipt required: buyer email or phone is missing")
    payload["Receipt"] = build_tbank_receipt(
        amount_rub=amount_rub,
        description=description,
        receipt_email=final_email,
        receipt_phone=final_phone,
    )
    notification_url = resolve_tbank_notification_url()
    if notification_url:
        payload["NotificationURL"] = notification_url
    success_url = add_request_id_to_url(resolve_tbank_success_url(), request_id)
    if success_url:
        payload["SuccessURL"] = success_url
    fail_url = add_request_id_to_url(resolve_tbank_fail_url(), request_id)
    if fail_url:
        payload["FailURL"] = fail_url

    payload["Token"] = tbank_token_from_payload(payload, TBANK_PASSWORD)
    status, data, _ = await http_json_request_with_retries(
        "POST",
        TBANK_INIT_URL,
        json_payload=payload,
        semaphore=state.tbank_api_semaphore,
        request_name="tbank_init",
    )
    if status >= 400:
        raise RuntimeError(f"T-Bank Init HTTP {status}: {data}")
    if not isinstance(data, dict):
        raise RuntimeError("Invalid T-Bank Init response.")
    if not data.get("Success"):
        error_code = scalar_string(data.get("ErrorCode"))
        message = scalar_string(data.get("Message"))
        details = scalar_string(data.get("Details"))
        parts = [part for part in [error_code and f"code={error_code}", message, details] if part]
        reason = " | ".join(parts) if parts else str(data)
        if error_code == "251":
            reason = f"{reason} | Сумма ниже минимальной для терминала."
        raise RuntimeError(f"T-Bank Init failed: {reason}")
    payment_url = data.get("PaymentURL")
    payment_id = data.get("PaymentId")
    if not isinstance(payment_url, str) or not payment_url:
        raise RuntimeError(f"T-Bank Init missing PaymentURL: {data}")
    return payment_url, scalar_string(payment_id)


async def tbank_get_state(payment_id: str) -> dict[str, Any]:
    if not tbank_enabled():
        raise RuntimeError("T-Bank credentials are not configured.")
    if not payment_id:
        raise RuntimeError("Empty PaymentId.")

    payload: dict[str, Any] = {
        "TerminalKey": TBANK_TERMINAL_KEY,
        "PaymentId": payment_id,
    }
    payload["Token"] = tbank_token_from_payload(payload, TBANK_PASSWORD)
    status, data, _ = await http_json_request_with_retries(
        "POST",
        TBANK_GET_STATE_URL,
        json_payload=payload,
        semaphore=state.tbank_api_semaphore,
        request_name="tbank_get_state",
    )
    if status >= 400:
        raise RuntimeError(f"T-Bank GetState HTTP {status}: {data}")
    if not isinstance(data, dict):
        raise RuntimeError("Invalid T-Bank GetState response.")
    return data


def cancel_payment_if_open(request_id: int) -> tuple[bool, str]:
    payment = state.user_store.get_payment(request_id)
    if not payment:
        return False, "not found"
    current_status = str(payment.get("status", "")).lower()
    if current_status in PAYMENT_FINAL_STATUSES:
        return False, f"already {current_status}"
    state.user_store.set_payment_status(request_id, "canceled")
    return True, "canceled"


async def refresh_payment_from_tbank(request_id: int, source: str) -> tuple[dict[str, Any] | None, str]:
    payment = state.user_store.get_payment(request_id)
    if not payment:
        return None, ""

    provider_ref = str(payment.get("provider_ref", ""))
    payment_id = tbank_payment_id_from_provider_ref(provider_ref)
    if not payment_id or not tbank_enabled():
        return payment, ""

    payload = await tbank_get_state(payment_id)
    bank_status = scalar_string(payload.get("Status")).upper()
    bank_success_raw = payload.get("Success")
    bank_success = bank_success_raw is True or scalar_string(bank_success_raw).lower() == "true"

    if bank_success and bank_status in TBANK_SUCCESS_STATUSES:
        await activate_payment_request(request_id, source=source)
    elif tbank_is_refund_status(bank_status):
        await process_refund_payment_request(request_id, source=source, bank_status=bank_status)
    elif tbank_is_cancel_status(bank_status):
        cancel_payment_if_open(request_id)

    return state.user_store.get_payment(request_id), bank_status


async def activate_payment_request(request_id: int, source: str) -> tuple[bool, str]:
    payment = state.user_store.get_payment(request_id)
    if not payment:
        return False, f"payment #{request_id} not found"

    status = str(payment.get("status", ""))
    if status in {"paid", "canceled", "refunded"}:
        return False, f"payment #{request_id} already {status}"

    target = int(payment["chat_id"])
    plan = str(payment["plan"])
    days = int(payment["days"])
    amount_rub = int(payment.get("amount_rub", 0) or 0)
    user_profile(target)
    state.user_store.set_payment_status(request_id, "paid")
    state.user_store.record_usage_event(
        chat_id=target,
        event_type="payment",
        plan=plan,
        rub_amount=amount_rub,
        details=f"request_id={request_id};source={source}",
    )

    if is_topup_plan(plan):
        code = topup_code_from_plan(plan)
        pack = topup_spec(code)
        if not pack:
            return False, f"unknown topup code in payment #{request_id}"
        credits = int(pack["credits"])
        row = user_profile(target)
        current = int(row.get("credits_balance", 0) or 0)
        state.user_store.set_credits(target, current + credits)
        state.user_store.mark_payment_activated(request_id)
        with suppress(Exception):
            await max_send_message(
                target,
                f"Оплата подтверждена ({source}). Зачислено {request_balance_text(credits)} запросов.",
                attachments=build_keyboard(),
            )
        return True, f"requests+{request_balance_text(credits)}"

    selected = best_default_alias_for_plan(plan)
    recurring_for_user = source.lower().startswith("t-bank")
    expires_at = state.user_store.set_subscription(
        target,
        plan,
        days,
        selected,
        recurring_enabled=recurring_for_user,
    )
    state.user_store.mark_payment_activated(request_id)

    with suppress(Exception):
        await max_send_message(
            target,
            f"Оплата подтверждена автоматически ({source}). Тариф {plan} активирован до {format_msk_datetime(parse_iso_datetime(expires_at))}.",
            attachments=build_keyboard(),
        )
    return True, expires_at


async def process_refund_payment_request(request_id: int, source: str, bank_status: str) -> tuple[bool, str]:
    payment = state.user_store.get_payment(request_id)
    if not payment:
        return False, f"payment #{request_id} not found"

    target = int(payment["chat_id"])
    plan = str(payment.get("plan", ""))
    amount_rub = int(payment.get("amount_rub", 0) or 0)
    current_status = str(payment.get("status", "")).lower()
    if current_status == "refunded":
        return False, "already refunded"
    if current_status != "refunded":
        state.user_store.set_payment_status(request_id, "refunded")
        state.user_store.record_usage_event(
            chat_id=target,
            event_type="refund",
            plan=plan,
            rub_amount=-abs(amount_rub),
            details=f"request_id={request_id};source={source};status={bank_status}",
        )

    if is_topup_plan(plan):
        code = topup_code_from_plan(plan)
        pack = topup_spec(code)
        if not pack:
            return False, "unknown topup for refund"
        credits = int(pack["credits"])
        row = user_profile(target)
        current = int(row.get("credits_balance", 0) or 0)
        state.user_store.set_credits(target, max(0, current - credits))
        with suppress(Exception):
            await max_send_message(
                target,
                (
                    f"Возврат подтвержден ({source}, статус {bank_status}).\n"
                    f"Пакет запросов отменен, списано {request_balance_text(credits)} запросов."
                ),
                attachments=build_keyboard(),
            )
        return True, f"topup requests-{request_balance_text(credits)}"

    row = user_profile(target)
    changed = False
    if row.get("plan") != "free" or recurring_enabled_for_row(row):
        state.user_store.set_plan(target, "free")
        state.user_store.set_selected_model(target, best_default_alias_for_plan("free"))
        changed = True

    if changed:
        with suppress(Exception):
            await max_send_message(
                target,
                (
                    f"Возврат подтвержден ({source}, статус {bank_status}).\n"
                    "Подписка переведена на free, автопродление отключено."
                ),
                attachments=build_keyboard(),
            )
        return True, "downgraded to free"
    return False, "already free or already canceled"


async def create_buy_request_v2(chat_id: int, plan: str, consent_text: str = "") -> tuple[int | None, str]:
    if plan not in BUYABLE_PLANS:
        return None, "Доступно: lite, start или pro."
    amount, days = plan_price_and_days(plan)
    row = user_profile(chat_id)
    receipt_email, receipt_phone = effective_receipt_contact(row)
    if not (receipt_email or receipt_phone):
        return None, "Нужен email или телефон для чека. Нажми «Тарифы» и начни оплату заново."
    request_id = state.user_store.create_payment_request(
        chat_id,
        plan,
        days,
        amount,
        recurring_consent=bool(consent_text),
        recurring_consent_text=consent_text,
        receipt_email=receipt_email,
        receipt_phone=receipt_phone,
    )
    payment_purpose = f"Оплата подписки, заказ #{request_id}"
    if tbank_enabled():
        try:
            payment_url, payment_id = await tbank_init_payment(
                request_id=request_id,
                amount_rub=amount,
                description=f"Подписка {plan}, заказ #{request_id}",
                receipt_email=receipt_email,
                receipt_phone=receipt_phone,
            )
            state.user_store.set_payment_provider_ref(request_id, f"tbank:{payment_id}")
            state.user_store.set_payment_url(request_id, payment_url)
            text = (
                f"Заявка #{request_id} создана\n"
                f"{payment_item_human_name(plan)}\n"
                f"Срок: {days} дней\n"
                f"Сумма: {amount} ₽\n"
                f"Чек: {receipt_email or receipt_phone}\n\n"
                "Открой оплату по кнопке ниже.\n"
                "После успешной оплаты тариф активируется автоматически."
            )
            return request_id, text
        except Exception as exc:
            log.exception("T-Bank Init failed for request %s", request_id)
            text = (
                f"Заявка #{request_id} создана\n"
                f"{payment_item_human_name(plan)}\n"
                f"Сумма: {amount} ₽\n"
                f"Автооплата сейчас недоступна ({exc}).\n\n"
                "Используй ручную оплату ниже."
            )
            return request_id, text
    text = (
        f"Заявка #{request_id} создана\n"
        f"{payment_item_human_name(plan)}\n"
        f"Срок: {days} дней\n"
        f"Сумма: {amount} ₽\n"
        f"Чек: {receipt_email or receipt_phone}\n\n"
        "Куда оплачивать:\n"
        f"{PAYMENT_DETAILS_TEXT}\n\nНазначение платежа: {payment_purpose}\nchat_id указывать не нужно.\n\n"
        "После оплаты нажми кнопку «Я оплатил»."
    )
    return request_id, text


async def create_topup_request_v2(chat_id: int, code: str) -> tuple[int | None, str]:
    pack = topup_spec(code)
    if not pack:
        return None, "Неизвестный пакет запросов."

    amount = int(pack["price_rub"])
    credits = int(pack["credits"])
    row = user_profile(chat_id)
    receipt_email, receipt_phone = effective_receipt_contact(row)
    if not (receipt_email or receipt_phone):
        return None, "Нужен email или телефон для чека. Нажми «Пакеты запросов» и начни покупку заново."

    request_id = state.user_store.create_payment_request(
        chat_id,
        topup_plan_code(code),
        0,
        amount,
        recurring_consent=False,
        recurring_consent_text="",
        receipt_email=receipt_email,
        receipt_phone=receipt_phone,
    )
    payment_purpose = f"Пакет запросов {pack['label']}, заказ #{request_id}"
    if tbank_enabled():
        try:
            payment_url, payment_id = await tbank_init_payment(
                request_id=request_id,
                amount_rub=amount,
                description=f"Пакет запросов {pack['label']}, заказ #{request_id}",
                receipt_email=receipt_email,
                receipt_phone=receipt_phone,
            )
            state.user_store.set_payment_provider_ref(request_id, f"tbank:{payment_id}")
            state.user_store.set_payment_url(request_id, payment_url)
            text = (
                f"Заявка #{request_id} создана\n"
                f"{payment_item_human_name(topup_plan_code(code))}\n"
                f"Запросов: {request_balance_text(credits)}\n"
                f"Сумма: {amount} ₽\n"
                f"Чек: {receipt_email or receipt_phone}\n\n"
                "Открой оплату по кнопке ниже.\n"
                "После успешной оплаты запросы зачислятся автоматически."
            )
            return request_id, text
        except Exception as exc:
            log.exception("T-Bank Init failed for topup request %s", request_id)
            text = (
                f"Заявка #{request_id} создана\n"
                f"{payment_item_human_name(topup_plan_code(code))}\n"
                f"Сумма: {amount} ₽\n"
                f"Автооплата сейчас недоступна ({exc}).\n\n"
                "Используй ручную оплату ниже."
            )
            return request_id, text

    text = (
        f"Заявка #{request_id} создана\n"
        f"{payment_item_human_name(topup_plan_code(code))}\n"
        f"Запросов: {request_balance_text(credits)}\n"
        f"Сумма: {amount} ₽\n"
        f"Чек: {receipt_email or receipt_phone}\n\n"
        "Куда оплачивать:\n"
        f"{PAYMENT_DETAILS_TEXT}\n\nНазначение платежа: {payment_purpose}\nchat_id указывать не нужно.\n\n"
        "После оплаты нажми кнопку «Я оплатил»."
    )
    return request_id, text


async def set_user_model(chat_id: int, alias: str) -> str:
    row = user_profile(chat_id)
    ok, reason = can_use_model(row["plan"], alias)
    if not ok:
        raise RuntimeError(reason)
    state.user_store.set_selected_model(chat_id, alias)
    return TEXT_MODELS[alias].label


async def handle_admin(chat_id: int, text: str) -> bool:
    if not is_admin(chat_id):
        return False

    parts = text.strip().split()
    if len(parts) < 2 or parts[1] == "help":
        await max_send_message(
            chat_id,
            "Админ-команды:\n"
            "/admin user <chat_id>\n"
            "/admin block <chat_id> <on|off>\n"
            "/admin templates\n"
            "/admin backup\n"
            "/admin nudge [days] [limit]\n\n"
            "Изменение тарифов и ручное подтверждение платежей — только через веб-админку.",
        )
        return True

    action = parts[1].lower()
    if action == "user" and len(parts) >= 3:
        target = parse_admin_target(parts[2])
        if target is None:
            await max_send_message(chat_id, "Некорректный chat_id")
            return True
        row = user_profile(target)
        await max_send_message(chat_id, f"user {target}\n{usage_text(row)}\nblocked={row['is_blocked']}")
        return True

    if action == "plan" and len(parts) >= 4:
        await max_send_message(chat_id, "Команда отключена. Изменение тарифов делай через веб-админку: /admin/panel")
        return True

    if action == "sub" and len(parts) >= 5:
        await max_send_message(chat_id, "Команда отключена. Изменение тарифов делай через веб-админку: /admin/panel")
        return True

    if action == "block" and len(parts) >= 4:
        target = parse_admin_target(parts[2])
        flag = parts[3].lower()
        if target is None or flag not in {"on", "off"}:
            await max_send_message(chat_id, "Используй: /admin block <chat_id> <on|off>")
            return True
        user_profile(target)
        state.user_store.set_blocked(target, flag == "on")
        await max_send_message(chat_id, f"Пользователь {target}: block={flag}")
        return True

    if action == "pay" and len(parts) >= 4:
        await max_send_message(chat_id, "Команда отключена. Проверку/подтверждение платежей делай через веб-админку: /admin/panel")
        return True

    if action == "templates":
        await max_send_message(chat_id, support_admin_templates_text())
        return True

    if action in {"panel", "kpi"}:
        await max_send_message(chat_id, "Веб-аналитика и админка теперь открываются только через сайт: /analytics")
        return True

    if action == "backup":
        try:
            backup_file = create_db_backup()
            await max_send_message(chat_id, f"Бэкап создан: {backup_file}")
        except Exception as exc:
            capture_exception_safe(exc)
            await max_send_message(chat_id, f"Не удалось создать бэкап: {exc}")
        return True

    if action == "nudge":
        days = REENGAGE_DORMANT_DAYS
        limit = REENGAGE_BATCH_LIMIT
        if len(parts) >= 3 and parts[2].isdigit():
            days = max(1, int(parts[2]))
        if len(parts) >= 4 and parts[3].isdigit():
            limit = max(1, int(parts[3]))
        sent, total = await send_reengage_nudges(days=days, limit=limit)
        await max_send_message(chat_id, f"Реактивация: отправлено {sent}/{total} (dormant_days={days}, limit={limit})")
        return True

    await max_send_message(chat_id, "Неизвестная админ-команда. Используй /admin help")
    return True


async def handle_callback(update: dict[str, Any]) -> bool:
    chat_id, callback_id, payload, source_mid = parse_callback_payload(update)
    if chat_id is None or not payload:
        return False
    ensure_update_user_binding(chat_id, update)

    if payload == "channel_gate:check":
        ok, reason = await check_channel_subscription(chat_id, force=True)
        if ok:
            row = user_profile(chat_id)
            if int(row.get("onboarding_done", 0) or 0) == 0:
                await show_onboarding_step(
                    chat_id,
                    step=1,
                    callback_id=callback_id,
                    source_mid=source_mid,
                    notification="Подписка подтверждена",
                )
            else:
                await show_ui_page(
                    chat_id,
                    UI_PAGE_MENU,
                    callback_id=callback_id,
                    source_mid=source_mid,
                    push_history=False,
                    notification="Подписка подтверждена" if reason in {"subscribed", "cached"} else "Открываю меню",
                )
        else:
            await show_channel_gate(
                chat_id,
                callback_id=callback_id,
                source_mid=source_mid,
                notification="Подписка не найдена",
                reason=reason,
            )
        return True

    if CHANNEL_GATE_ENABLED and not channel_gate_allows_payload(payload):
        if not await ensure_channel_access(chat_id, callback_id=callback_id, source_mid=source_mid):
            return True

    if payload.startswith("reply_action:"):
        reply_action = payload.split(":", 1)[1].strip()
        reply_page_map = {
            "menu": UI_PAGE_MENU,
            "image_menu": UI_PAGE_IMAGE_MENU,
            "files_menu": UI_PAGE_FILES_MENU,
            "tariffs": UI_PAGE_TARIFFS,
        }
        if reply_action == "clear":
            state.history(chat_id).clear()
            if callback_id:
                await answer_callback(callback_id, "Контекст очищен")
            await send_managed_message(chat_id, "Контекст диалога очищен.", attachments=build_keyboard(), page=UI_PAGE_MENU)
            return True
        target_page = reply_page_map.get(reply_action)
        if target_page:
            await show_ui_page(
                chat_id,
                target_page,
                callback_id=callback_id,
                source_mid=source_mid,
                push_history=True,
                force_new=True,
            )
            return True

    if payload.startswith("onboard:") and int(user_profile(chat_id).get("onboarding_done", 0) or 0) == 1:
        await show_ui_page(
            chat_id,
            UI_PAGE_MENU,
            callback_id=callback_id,
            source_mid=source_mid or state.ui_message_mid.get(chat_id),
            push_history=False,
            notification="Открываю меню",
        )
        return True

    if payload not in {"growth:ref_enter", "growth:promo_enter", "growth:input_cancel"}:
        clear_growth_pending_inputs(chat_id)
    if not (
        payload in {"files:doc", "files:ppt", "files:sheet", "files:cancel"}
        or payload.startswith("files:sizepick:")
        or payload.startswith("files:size:")
    ):
        clear_file_pending_input(chat_id)

    if payload == "ui_nav:back":
        target = ui_nav_back(chat_id)
        if target:
            await show_ui_page(chat_id, target, callback_id=callback_id, source_mid=source_mid, push_history=False)
        elif callback_id:
            await answer_callback(callback_id, "Назад недоступен")
        return True

    if payload == "ui_nav:forward":
        target = ui_nav_forward(chat_id)
        if target:
            await show_ui_page(chat_id, target, callback_id=callback_id, source_mid=source_mid, push_history=False)
        elif callback_id:
            await answer_callback(callback_id, "Откат недоступен")
        return True

    action_page_map = {
        "action:menu": UI_PAGE_MENU,
        "action:models": UI_PAGE_MODELS,
        "action:plan": UI_PAGE_PLAN,
        "action:tariffs": UI_PAGE_TARIFFS,
        "action:topups": UI_PAGE_TOPUPS,
        "action:payments": UI_PAGE_PAYMENTS,
        "action:growth": UI_PAGE_GROWTH,
        "action:support": UI_PAGE_SUPPORT,
        "action:image_menu": UI_PAGE_IMAGE_MENU,
        "action:files_menu": UI_PAGE_FILES_MENU,
    }
    if payload in action_page_map:
        await show_ui_page(chat_id, action_page_map[payload], callback_id=callback_id, source_mid=source_mid, push_history=True)
        return True

    if payload == "action:preset_models":
        plan = str(user_profile(chat_id).get("plan", "free"))
        if not can_pick_models_for_current_preset(chat_id):
            await show_ui_page(chat_id, UI_PAGE_MODELS, callback_id=callback_id, source_mid=source_mid, push_history=True)
            return True
        await show_managed_content(
            chat_id,
            build_model_picker_text(chat_id),
            attachments=build_model_picker_keyboard(plan),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_MENU,
            push_history=False,
            notification="Выбери модель",
        )
        return True

    if payload.startswith("set_preset:"):
        preset = payload.split(":", 1)[1].strip().lower()
        preset_cfg = MODEL_PRESETS.get(preset)
        if not preset_cfg:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный режим")
            return True
        try:
            plan = str(user_profile(chat_id).get("plan", "free"))
            aliases = preset_available_aliases_for_plan(plan, preset)
            alias = aliases[0]
            label = await set_user_model(chat_id, alias)
            state.user_store.set_selected_preset(chat_id, preset)
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="preset_select",
                plan=plan,
                model_alias=alias,
                details=f"preset={preset}",
            )
            await show_ui_page(
                chat_id,
                UI_PAGE_MENU,
                callback_id=callback_id,
                source_mid=source_mid,
                push_history=False,
                notification=f"{preset_cfg['label']} → {label}",
            )
        except Exception as exc:
            if callback_id:
                await answer_callback(callback_id, str(exc)[:120])
            await max_send_message(chat_id, f"Ошибка: {exc}", attachments=build_keyboard(), notify=False)
        return True

    if payload.startswith("preset_pick:"):
        _, preset, alias = (payload.split(":", 2) + ["", ""])[:3]
        preset = preset.strip().lower()
        alias = alias.strip().lower()
        preset_cfg = MODEL_PRESETS.get(preset)
        if not preset_cfg:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный режим")
            return True
        plan = str(user_profile(chat_id).get("plan", "free"))
        allowed_aliases = preset_available_aliases_for_plan(plan, preset)
        if alias not in allowed_aliases:
            if callback_id:
                await answer_callback(callback_id, "Модель недоступна")
            await show_managed_content(
                chat_id,
                build_model_picker_text(chat_id),
                attachments=build_model_picker_keyboard(plan),
                callback_id=None,
                source_mid=source_mid,
                page=UI_PAGE_MENU,
                push_history=False,
            )
            return True
        try:
            label = await set_user_model(chat_id, alias)
            state.user_store.set_selected_preset(chat_id, preset)
            await show_ui_page(
                chat_id,
                UI_PAGE_MENU,
                callback_id=callback_id,
                source_mid=source_mid,
                push_history=False,
                notification=f"{preset_cfg['label']} → {label}",
            )
        except Exception as exc:
            if callback_id:
                await answer_callback(callback_id, str(exc)[:120])
            await max_send_message(chat_id, f"Ошибка: {exc}", attachments=build_keyboard(), notify=False)
        return True

    if payload.startswith("set_model:"):
        alias = payload.split(":", 1)[1]
        try:
            label = await set_user_model(chat_id, alias)
            state.user_store.set_selected_preset(chat_id, "")
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="model_select",
                plan=str(user_profile(chat_id).get("plan", "free")),
                model_alias=alias,
                details="source=manual_picker",
            )
            target_page = state.ui_current_page.get(chat_id, UI_PAGE_MENU)
            if target_page not in UI_PAGE_KEYS:
                target_page = UI_PAGE_MENU
            await show_ui_page(
                chat_id,
                target_page,
                callback_id=callback_id,
                source_mid=source_mid,
                push_history=False,
                notification=f"Модель: {label}",
            )
        except Exception as exc:
            if callback_id:
                await answer_callback(callback_id, str(exc)[:120])
            await max_send_message(chat_id, f"Ошибка: {exc}", attachments=build_keyboard(), notify=False)
        return True

    if payload == "action:clear":
        state.history(chat_id).clear()
        if callback_id:
            await answer_callback(callback_id, "Контекст очищен")
        await send_managed_message(chat_id, "Контекст диалога очищен.", attachments=build_keyboard(), notify=False, page=UI_PAGE_MENU)
        return True

    if payload == "action:models":
        if callback_id:
            await answer_callback(callback_id, "Показываю модели")
        await send_models(chat_id)
        return True

    if payload == "action:plan":
        if callback_id:
            await answer_callback(callback_id, "Показываю план")
        await send_plan(chat_id)
        return True

    if payload == "action:tariffs":
        if callback_id:
            await answer_callback(callback_id, "Показываю тарифы")
        await max_send_message(chat_id, build_tariffs_text(), attachments=build_tariffs_keyboard_pricing(), notify=False)
        return True

    if payload == "action:topups":
        if callback_id:
            await answer_callback(callback_id, "Показываю пакеты")
        await send_topups(chat_id)
        return True

    if payload == "action:payments":
        if callback_id:
            await answer_callback(callback_id, "Показываю оплаты")
        await send_payments(chat_id)
        return True

    if payload == "action:menu":
        if callback_id:
            await answer_callback(callback_id, "Открываю меню")
        await send_menu(chat_id)
        return True

    if payload == "action:growth":
        if callback_id:
            await answer_callback(callback_id, "Бонусы")
        await send_growth_menu(chat_id)
        return True

    if payload == "action:channel":
        if CHANNEL_GATE_ENABLED:
            subscribed, reason = await check_channel_subscription(chat_id, force=False)
            if not subscribed:
                await show_channel_gate(
                    chat_id,
                    callback_id=callback_id,
                    source_mid=source_mid,
                    notification="Канал",
                    reason=reason,
                )
                return True
        await show_managed_content(
            chat_id,
            f"📣 Канал проекта:\n{channel_url_value()}",
            attachments=build_growth_keyboard(chat_id),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Канал",
        )
        return True

    if payload == "growth:ref_show":
        row = user_profile(chat_id)
        code = str(row.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
        invited = int(row.get("referrals_invited", 0) or 0)
        await show_managed_content(
            chat_id,
            (
                f"👥 Твой реф-код: {code}\n"
                f"Приглашено друзей: {invited}\n\n"
                f"Пригласи друга через кнопку «Поделиться» или попроси его ввести /ref {code}\n"
                f"После активации — бонус +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов вам обоим."
            ),
            attachments=build_growth_keyboard(chat_id),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Твой реф-код",
        )
        return True

    if payload == "growth:ref_share":
        row = user_profile(chat_id)
        code = str(row.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
        await show_managed_content(
            chat_id,
            (
                "🔗 Поделиться приглашением\n\n"
                "Для новых экранов кнопка «Поделиться» открывает шторку MAX сразу.\n"
                "Если ты нажал её на старом сообщении, просто нажми «Поделиться» ещё раз на обновлённой кнопке ниже."
            ),
            attachments=build_growth_keyboard(chat_id),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Поделиться",
        )
        return True

    if payload == "growth:ref_enter":
        state.pending_referral_code_input.add(chat_id)
        await show_managed_content(
            chat_id,
            "Введи реферальный код одним сообщением (пример: RFABC123). Для отмены отправь «отмена».",
            attachments=build_growth_input_keyboard(),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Жду код",
        )
        return True

    if payload == "growth:promo_enter":
        state.pending_promo_code_input.add(chat_id)
        await show_managed_content(
            chat_id,
            "Введи промокод одним сообщением. Для отмены отправь «отмена».",
            attachments=build_growth_input_keyboard(),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Жду промокод",
        )
        return True

    if payload == "growth:input_cancel":
        clear_growth_pending_inputs(chat_id)
        await show_ui_page(chat_id, UI_PAGE_GROWTH, callback_id=callback_id, source_mid=source_mid, push_history=False)
        return True

    if payload == "payment:input_cancel":
        pending_plan = state.pending_receipt_plan.pop(chat_id, "")
        await show_ui_page(
            chat_id,
            receipt_return_page(pending_plan or ""),
            callback_id=callback_id,
            source_mid=source_mid,
            push_history=False,
        )
        return True

    if payload == "growth:channel_bonus":
        await show_managed_content(
            chat_id,
            "Промокоды из канала активируются через кнопку «Промокод». Если в канале опубликован код, введи его здесь.",
            attachments=build_growth_keyboard(chat_id),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_GROWTH,
            notification="Промокоды",
        )
        return True

    if payload.startswith("onboard:"):
        target_page = UI_PAGE_MENU
        if payload == "onboard:scenario:image":
            target_page = UI_PAGE_IMAGE_MENU
        elif payload == "onboard:scenario:tariff":
            target_page = UI_PAGE_TARIFFS
        await finish_onboarding_to_page(chat_id, target_page, callback_id=callback_id, source_mid=source_mid)
        return True

    if payload == "action:image_menu":
        await send_image_menu(chat_id)
        return True

    if payload == "action:files_menu":
        await send_files_menu(chat_id)
        return True

    if payload.startswith("files:sizepick:"):
        kind = payload.split(":", 2)[2].strip()
        if kind not in FILE_KIND_LABELS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный тип файла")
            return True
        plan_name = str(user_profile(chat_id).get("plan", "free")).strip().lower()
        state.pending_file_kind[chat_id] = kind
        state.pending_file_profile[chat_id] = default_file_profile_for_plan(plan_name)
        await show_managed_content(
            chat_id,
            build_file_size_text(kind),
            attachments=build_files_size_keyboard(kind),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_FILES_MENU,
            notification="Выбери объем",
        )
        return True

    if payload.startswith("files:size:"):
        _, _, kind, profile = (payload.split(":", 3) + ["", "", "", ""])[:4]
        if kind not in FILE_KIND_LABELS or profile not in FILE_PROFILE_LABELS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный режим")
            return True
        if str(user_profile(chat_id).get("plan", "free")).strip().lower() == "free" and profile != "short":
            await show_managed_content(
                chat_id,
                "На Free доступен только короткий файл раз в 14 дней. Средняя и полная версии открываются на платных тарифах.",
                attachments=purchase_help_keyboard_for_row(user_profile(chat_id)),
                callback_id=callback_id,
                source_mid=source_mid,
                page=UI_PAGE_FILES_MENU,
                notification="Только короткий файл",
            )
            return True
        state.pending_file_kind[chat_id] = kind
        state.pending_file_profile[chat_id] = profile
        await show_managed_content(
            chat_id,
            build_file_prompt_text(kind, profile, chat_id),
            attachments=build_files_prompt_keyboard(),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_FILES_MENU,
            notification="Жду описание",
        )
        return True

    if payload in {"files:doc", "files:ppt", "files:sheet"}:
        kind = payload.split(":", 1)[1]
        state.pending_file_kind[chat_id] = kind
        state.pending_file_profile[chat_id] = default_file_profile_for_plan(str(user_profile(chat_id).get("plan", "free")))
        await show_managed_content(
            chat_id,
            build_file_prompt_text(kind, state.pending_file_profile[chat_id], chat_id),
            attachments=build_files_prompt_keyboard(),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_FILES_MENU,
            notification="Жду описание",
        )
        return True

    if payload == "files:cancel":
        clear_file_pending_input(chat_id)
        await show_ui_page(chat_id, UI_PAGE_FILES_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False)
        return True

    if payload == "image_mode:generate":
        set_image_mode(chat_id, "generate")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Режим генерации")
        return True

    if payload == "image_mode:edit":
        set_image_mode(chat_id, "edit")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Режим редактирования")
        return True

    if payload == "image_mode:back":
        set_image_mode(chat_id, "")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Выбор режима")
        return True

    if payload == "image_panel:scenario":
        set_image_panel(chat_id, "scenario")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Назад")
        return True

    if payload == "image_panel:style":
        set_image_panel(chat_id, "style")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Стиль и формат")
        return True

    if payload.startswith("image_style:"):
        style = payload.split(":", 1)[1].strip().lower()
        if style not in IMAGE_STYLE_OPTIONS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный стиль")
            return True
        prefs = get_image_prefs(chat_id)
        prefs["style"] = style
        prefs["preset"] = ""
        prefs["edit_preset"] = ""
        state.image_request_prefs[chat_id] = prefs
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False)
        return True

    if payload.startswith("image_aspect:"):
        aspect = payload.split(":", 1)[1].strip().lower()
        if aspect not in IMAGE_ASPECT_OPTIONS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный формат")
            return True
        prefs = get_image_prefs(chat_id)
        prefs["aspect"] = aspect
        prefs["preset"] = ""
        prefs["edit_preset"] = ""
        state.image_request_prefs[chat_id] = prefs
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False)
        return True

    if payload.startswith("image_preset:"):
        preset_key = payload.split(":", 1)[1].strip().lower()
        if preset_key == "none":
            clear_image_preset(chat_id)
            set_image_panel(chat_id, "style")
            await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Без сценария")
            return True
        if preset_key not in IMAGE_PRESET_OPTIONS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный сценарий")
            return True
        apply_image_preset(chat_id, preset_key)
        set_image_panel(chat_id, "style")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Сценарий выбран")
        return True

    if payload.startswith("image_edit_preset:"):
        preset_key = payload.split(":", 1)[1].strip().lower()
        if preset_key == "none":
            clear_image_preset(chat_id)
            set_image_panel(chat_id, "style")
            await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Без сценария")
            return True
        if preset_key not in IMAGE_EDIT_PRESET_OPTIONS:
            if callback_id:
                await answer_callback(callback_id, "Неизвестный сценарий")
            return True
        apply_image_edit_preset(chat_id, preset_key)
        set_image_panel(chat_id, "style")
        await show_ui_page(chat_id, UI_PAGE_IMAGE_MENU, callback_id=callback_id, source_mid=source_mid, push_history=False, notification="Сценарий выбран")
        return True

    if payload == "image_prompt:start":
        row = user_profile(chat_id)
        plan_name = str(row.get("plan", "free")).strip().lower()
        prompt_target_mid = state.ui_message_mid.get(chat_id) or source_mid
        if plan_name != "free" and not plan_allowed(plan_name, DEFAULT_IMAGE_MODEL.min_plan):
            if callback_id:
                await answer_callback(callback_id, "Недоступно на текущем тарифе")
            await show_managed_content(
                chat_id,
                f"Картинки доступны с тарифа {DEFAULT_IMAGE_MODEL.min_plan}. Открой «Тарифы».",
                attachments=build_tariffs_keyboard_pricing(),
                source_mid=prompt_target_mid,
                page=UI_PAGE_TARIFFS,
                push_history=False,
            )
            return True

        ok_limit, reason_limit = check_limit_only(chat_id, "images")
        if not ok_limit:
            if callback_id:
                await answer_callback(callback_id, "Лимит достигнут")
            await show_managed_content(
                chat_id,
                reason_limit,
                attachments=build_tariffs_keyboard_pricing(),
                callback_id=None,
                source_mid=prompt_target_mid,
                page=UI_PAGE_TARIFFS,
                push_history=False,
                notification="Лимит достигнут",
            )
            return True
        state.pending_image_prompt.add(chat_id)
        prefs = get_image_prefs(chat_id)
        preset_hint = ""
        if prefs.get("preset", "") in IMAGE_PRESET_OPTIONS:
            chosen = IMAGE_PRESET_OPTIONS[prefs["preset"]]
            preset_hint = f"Сценарий: {chosen['label']}\nПодсказка: {chosen['hint']}\n"
        generation_cost_line = (
            "Лимит Free: 1 генерация или редактирование фото в 7 дней с момента последнего использования."
            if str(row.get("plan", "free")).strip().lower() == "free"
            else f"Стоимость: {request_cost_text(image_credit_cost())} запросов."
        )
        await show_managed_content(
            chat_id,
            "Напиши, что нарисовать одним сообщением.\n\n"
            f"{preset_hint}"
            f"{image_params_summary(chat_id)}\n"
            f"{generation_cost_line}",
            attachments=build_image_prompt_keyboard(),
            callback_id=callback_id,
            source_mid=prompt_target_mid,
            page=UI_PAGE_IMAGE_MENU,
            notification="Жду описание",
        )
        return True

    if payload == "image_ref:start":
        row = user_profile(chat_id)
        plan_name = str(row.get("plan", "free")).strip().lower()
        prompt_target_mid = state.ui_message_mid.get(chat_id) or source_mid
        if plan_name != "free" and not plan_allowed(plan_name, DEFAULT_IMAGE_MODEL.min_plan):
            if callback_id:
                await answer_callback(callback_id, "Недоступно на текущем тарифе")
            await show_managed_content(
                chat_id,
                f"Режим «по фото» доступен с тарифа {DEFAULT_IMAGE_MODEL.min_plan}. Открой «Тарифы».",
                attachments=build_tariffs_keyboard_pricing(),
                source_mid=prompt_target_mid,
                page=UI_PAGE_TARIFFS,
                push_history=False,
            )
            return True
        state.pending_image_ref_prompt.add(chat_id)
        prefs = get_image_prefs(chat_id)
        preset_hint = ""
        if prefs.get("edit_preset", "") in IMAGE_EDIT_PRESET_OPTIONS:
            chosen = IMAGE_EDIT_PRESET_OPTIONS[prefs["edit_preset"]]
            preset_hint = f"Сценарий: {chosen['label']}\nПодсказка: {chosen['hint']}\n"
        edit_prompt_cost_line = (
            "Лимит Free: 1 генерация или редактирование фото в 7 дней с момента последнего использования.\n"
            if plan_name == "free"
            else f"Стоимость: {request_cost_text(image_edit_credit_cost())} запросов.\n"
        )
        await show_managed_content(
            chat_id,
            (
                "Пришли фото и коротко опиши, что сделать.\n"
                f"{preset_hint}"
                f"{edit_prompt_cost_line}"
                "Если фото уже отправлено — просто напиши описание (например: «нарисуй её в стиле аниме»)."
            ),
            attachments=build_image_prompt_keyboard(),
            callback_id=callback_id,
            source_mid=prompt_target_mid,
            page=UI_PAGE_IMAGE_MENU,
            notification="Жду фото",
        )
        return True

    if payload == "image_prompt:cancel":
        state.pending_image_prompt.discard(chat_id)
        state.pending_image_ref_prompt.discard(chat_id)
        if callback_id:
            await answer_callback(callback_id, "Отменено")
        await send_image_menu(chat_id)
        return True

    if payload == "sub_cancel:start":
        row = user_profile(chat_id)
        if row.get("plan") == "free" or not recurring_enabled_for_row(row):
            if callback_id:
                await answer_callback(callback_id, "Нечего отменять")
            await show_managed_content(
                chat_id,
                "Автопродление уже отключено или подписка не активна.",
                attachments=build_plan_keyboard(row),
                callback_id=callback_id,
                source_mid=source_mid,
                page=UI_PAGE_PLAN,
                notification="Нечего отменять",
            )
            return True

        expires_at = parse_iso_datetime(str(row.get("subscription_expires_at", "")))
        expires_text = format_msk_datetime(expires_at) if expires_at else "конца текущего периода"
        await show_managed_content(
            chat_id,
            "Подтвердить отмену автопродления?\n\n"
            f"Подписка будет работать до {expires_text}.\n"
            f"Отмена автосписаний вступит в силу с {expires_text}.",
            attachments=build_cancel_subscription_keyboard(),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_PLAN,
            notification="Подтверди отмену",
        )
        return True

    if payload == "sub_cancel:confirm":
        row = user_profile(chat_id)
        if row.get("plan") == "free" or not recurring_enabled_for_row(row):
            if callback_id:
                await answer_callback(callback_id, "Уже отключено")
            await send_plan(chat_id)
            return True

        expires_at = parse_iso_datetime(str(row.get("subscription_expires_at", "")))
        cancel_from = expires_at.replace(microsecond=0).isoformat() if expires_at else datetime.utcnow().replace(microsecond=0).isoformat()
        state.user_store.cancel_recurring(chat_id, cancel_from)
        row = user_profile(chat_id)
        cancel_dt = parse_iso_datetime(cancel_from)
        cancel_text = format_msk_datetime(cancel_dt) if cancel_dt else cancel_from
        await show_managed_content(
            chat_id,
            f"Автопродление отключено.\nПодписка действует до конца оплаченного периода.\nОтмена с {cancel_text}.",
            attachments=build_plan_keyboard(row),
            callback_id=callback_id,
            source_mid=source_mid,
            page=UI_PAGE_PLAN,
            notification="Отменено",
        )
        return True

    if payload == "sub_cancel:back":
        if callback_id:
            await answer_callback(callback_id, "Без отмены")
        await send_plan(chat_id)
        return True

    if payload.startswith("buy_consent:"):
        plan = payload.split(":", 1)[1].lower().strip()
        if plan not in BUYABLE_PLANS:
            if callback_id:
                await answer_callback(callback_id, "Неверный тариф")
            return True
        terms = recurring_terms_for_plan(plan)
        request_id, msg = await create_buy_request_v2(chat_id, plan, consent_text=terms)
        if request_id is None:
            await max_send_message(chat_id, msg, attachments=build_tariffs_keyboard_pricing(), notify=False)
            return True
        payment_url = str((state.user_store.get_payment(request_id) or {}).get("payment_url", "") or "")
        await show_managed_content(
            chat_id,
            msg,
            attachments=build_payment_request_keyboard(request_id, payment_url=payment_url, status="pending"),
            callback_id=callback_id,
            source_mid=source_mid,
            notification="Открываю оплату",
        )
        return True

    if payload.startswith("buy:"):
        plan = payload.split(":", 1)[1].lower()
        if plan == "free":
            user_profile(chat_id)
            state.user_store.set_plan(chat_id, "free")
            state.user_store.set_selected_model(chat_id, best_default_alias_for_plan("free"))
            if callback_id:
                await answer_callback(callback_id, "Переключено на Free")
            await max_send_message(chat_id, "Тариф переключен на free.", attachments=build_tariffs_keyboard_pricing(), notify=False)
            return True

        if plan not in BUYABLE_PLANS:
            if callback_id:
                await answer_callback(callback_id, "Неверный тариф")
            await max_send_message(chat_id, "Доступно: Lite, Start или Pro.", attachments=build_tariffs_keyboard_pricing(), notify=False)
            return True
        ok = await start_buy_flow(chat_id, plan, notify=False, callback_id=callback_id, source_mid=source_mid)
        if not ok:
            return True
        return True

    if payload.startswith("topup:"):
        code = payload.split(":", 1)[1].lower().strip()
        pack = topup_spec(code)
        if not pack:
            if callback_id:
                await answer_callback(callback_id, "Неверный пакет")
            await max_send_message(chat_id, "Неизвестный пакет запросов.", attachments=build_topups_keyboard(), notify=False)
            return True

        row = user_profile(chat_id)
        email, phone = effective_receipt_contact(row)
        if not (email or phone):
            state.pending_receipt_plan[chat_id] = f"topup_consent:{code}"
            await request_receipt_contact(
                chat_id,
                f"topup_consent:{code}",
                notify=False,
                callback_id=callback_id,
                source_mid=source_mid,
            )
            return True

        await send_topup_consent(chat_id, code, notify=False, callback_id=callback_id, source_mid=source_mid)
        return True

    if payload.startswith("topup_quick:"):
        code = payload.split(":", 1)[1].lower().strip() or TOPUP_QUICK_CODE
        row = user_profile(chat_id)
        email, phone = effective_receipt_contact(row)
        if not (email or phone):
            state.pending_receipt_plan[chat_id] = f"topup_consent:{code}"
            await request_receipt_contact(
                chat_id,
                f"topup_consent:{code}",
                notify=False,
                callback_id=callback_id,
                source_mid=source_mid,
            )
            return True

        request_id, msg = await create_topup_request_v2(chat_id, code)
        if request_id is None:
            await max_send_message(chat_id, msg, attachments=build_topups_keyboard(), notify=False)
            return True
        payment_url = str((state.user_store.get_payment(request_id) or {}).get("payment_url", "") or "")
        await show_managed_content(
            chat_id,
            msg,
            attachments=build_payment_request_keyboard(request_id, payment_url=payment_url, status="pending"),
            callback_id=callback_id,
            source_mid=source_mid,
            notification="Открываю быструю покупку",
        )
        return True

    if payload.startswith("topup_consent:"):
        code = payload.split(":", 1)[1].lower().strip()
        request_id, msg = await create_topup_request_v2(chat_id, code)
        if request_id is None:
            await max_send_message(chat_id, msg, attachments=build_topups_keyboard(), notify=False)
            return True
        payment_url = str((state.user_store.get_payment(request_id) or {}).get("payment_url", "") or "")
        await show_managed_content(
            chat_id,
            msg,
            attachments=build_payment_request_keyboard(request_id, payment_url=payment_url, status="pending"),
            callback_id=callback_id,
            source_mid=source_mid,
            notification="Открываю оплату",
        )
        return True

    if payload.startswith("payment_status:"):
        request_raw = payload.split(":", 1)[1].strip()
        if not request_raw.isdigit():
            if callback_id:
                await answer_callback(callback_id, "Ошибка номера заявки")
            return True
        request_id = int(request_raw)
        payment = state.user_store.get_payment(request_id)
        if not payment or int(payment.get("chat_id", 0) or 0) != chat_id:
            if callback_id:
                await answer_callback(callback_id, "Заявка не найдена")
            await max_send_message(chat_id, "Заявка не найдена.", attachments=build_tariffs_keyboard_pricing(), notify=False)
            return True

        bank_status = ""
        status = str(payment.get("status", "pending")).lower()
        provider_ref = str(payment.get("provider_ref", ""))
        if provider_ref.startswith("tbank:") and status in {"pending", "claimed"}:
            try:
                payment, bank_status = await refresh_payment_from_tbank(request_id, source="T-Bank GetState (status button)")
                payment = payment or state.user_store.get_payment(request_id) or payment
            except Exception:
                log.exception("T-Bank GetState failed from status button for request_id=%s", request_id)

        refreshed_status = str((payment or {}).get("status", "pending")).lower()
        payment_url = str((payment or {}).get("payment_url", "") or "")
        await show_managed_content(
            chat_id,
            payment_user_status_text(payment or {}, bank_status=bank_status),
            attachments=build_payment_request_keyboard(request_id, payment_url=payment_url, status=refreshed_status),
            callback_id=callback_id,
            source_mid=source_mid,
            notification=payment_status_label(refreshed_status),
        )
        return True

    if payload.startswith("paid:"):
        request_raw = payload.split(":", 1)[1].strip()
        if not request_raw.isdigit():
            if callback_id:
                await answer_callback(callback_id, "Ошибка номера заявки")
            return True
        request_id = int(request_raw)
        payment = state.user_store.get_payment(request_id)
        if not payment or int(payment["chat_id"]) != chat_id:
            await show_managed_content(
                chat_id,
                "Заявка не найдена.",
                attachments=build_tariffs_keyboard_pricing(),
                callback_id=callback_id,
                source_mid=source_mid,
                notification="Заявка не найдена",
            )
            return True

        async def render_payment_state(
            payment_row: dict[str, Any],
            *,
            bank_status: str = "",
            notification: str = "Обновил",
        ) -> None:
            current_status = str((payment_row or {}).get("status", "pending")).lower()
            payment_url = str((payment_row or {}).get("payment_url", "") or "")
            await show_managed_content(
                chat_id,
                payment_user_status_text(payment_row or {}, bank_status=bank_status),
                attachments=build_payment_request_keyboard(request_id, payment_url=payment_url, status=current_status),
                callback_id=callback_id,
                source_mid=source_mid,
                notification=notification,
            )

        status = str(payment["status"]).lower()
        provider_ref = str(payment.get("provider_ref", ""))
        if provider_ref.startswith("tbank:") and status == "pending":
            try:
                payment, bank_status = await refresh_payment_from_tbank(request_id, source="T-Bank GetState (user button)")
            except Exception:
                log.exception("T-Bank GetState failed from paid button for request_id=%s", request_id)
                payment = state.user_store.get_payment(request_id) or payment
                await show_managed_content(
                    chat_id,
                    (
                        "Не удалось мгновенно получить ответ от банка.\n"
                        "Платеж продолжает проверяться автоматически, обычно до 1-2 минут."
                    ),
                    attachments=build_payment_request_keyboard(
                        request_id,
                        payment_url=str((payment or {}).get("payment_url", "") or ""),
                        status=str((payment or {}).get("status", "pending")).lower(),
                    ),
                    callback_id=callback_id,
                    source_mid=source_mid,
                    notification="Проверка банка",
                )
                return True

            refreshed_status = str((payment or {}).get("status", "pending")).lower()
            if refreshed_status == "paid":
                await render_payment_state(payment or {}, bank_status=bank_status, notification="Оплата подтверждена")
                return True
            if refreshed_status == "refunded":
                await render_payment_state(payment or {}, bank_status=bank_status, notification="Возврат")
                return True
            if refreshed_status == "canceled":
                await render_payment_state(payment or {}, bank_status=bank_status, notification="Оплата отменена")
                return True
            await render_payment_state(payment or {}, bank_status=bank_status, notification="Ожидаем банк")
            return True
        if status == "paid":
            await render_payment_state(payment, notification="Уже подтверждено")
            return True
        if status == "claimed":
            await render_payment_state(payment, notification="Уже на проверке")
            return True
        if status == "canceled":
            await render_payment_state(payment, notification="Заявка отменена")
            return True
        state.user_store.set_payment_status(request_id, "claimed")
        payment = state.user_store.get_payment(request_id) or payment
        await notify_admin_about_payment_claim(request_id, payment)
        await render_payment_state(payment, notification="Передано админу")
        return True

    return False


async def handle_command(chat_id: int, text: str) -> bool:
    lowered = text.strip().lower()
    parts = text.strip().split(maxsplit=1)
    command = parts[0].lower()
    arg = parts[1].strip() if len(parts) > 1 else ""

    if lowered in {"gpt", "gpt4o", "gemini", "deepseek", "gpt54"}:
        command = "/model"
        arg = lowered
    elif lowered in {"старт", "start", "начать"}:
        command = "/start"
        arg = ""
    elif lowered in {"меню", "menu"}:
        command = "/menu"
        arg = ""
    elif command in {"/gpt", "/gpt4o", "/gemini", "/deepseek", "/gpt54"}:
        command = "/model"
        arg = command[1:]

    if command in {"/start", "/menu"}:
        row = user_profile(chat_id)
        if command == "/start" and arg:
            start_source, start_campaign = acquisition_meta_from_start_payload({"payload": arg})
            state.user_store.set_acquisition_meta(chat_id, source=start_source or ("referral" if is_referral_code(arg) else "direct"), campaign=start_campaign)
            start_ref = referral_code_from_start_payload({"payload": arg}) or (normalize_referral_code(arg) if is_referral_code(arg) else "")
            if start_ref:
                ok, info = state.user_store.apply_referral_code(chat_id, start_ref, REFERRAL_BONUS_CREDITS)
                if ok:
                    owner_chat_id = int(info)
                    log_referral_activation(chat_id, owner_chat_id, start_ref, "start_command")
                    await notify_referral_success(chat_id, owner_chat_id, "start")
                    row = user_profile(chat_id)
        if command == "/start" and CHANNEL_GATE_ENABLED:
            if not await ensure_channel_access(chat_id):
                return True
        if command == "/start" and int(row.get("onboarding_done", 0) or 0) == 0:
            await send_onboarding(chat_id, step=1)
            return True
        await send_menu(chat_id)
        return True

    if command == "/help":
        await send_help(chat_id)
        return True

    if command == "/models":
        await send_models(chat_id)
        return True

    if command == "/costs":
        if not is_admin(chat_id):
            await max_send_message(chat_id, "Команда доступна только администратору.")
            return True
        await send_costs(chat_id)
        return True

    if command == "/id":
        row = user_profile(chat_id)
        max_user_id = int(row.get("max_user_id", 0) or 0)
        text = f"Твой chat_id: {chat_id}"
        if max_user_id > 0:
            text += f"\nТвой max_user_id: {max_user_id}"
        else:
            text += "\nТвой max_user_id пока не определён."
        await max_send_message(chat_id, text)
        return True

    if command == "/tariffs":
        await send_managed_message(chat_id, build_tariffs_text(), attachments=build_tariffs_keyboard_pricing(), page=UI_PAGE_TARIFFS)
        return True

    if command == "/topup":
        await send_topups(chat_id)
        return True

    if command == "/plan":
        await send_plan(chat_id)
        return True

    if command in {"/credits", "/requests"}:
        await send_credits(chat_id)
        return True

    if command == "/support":
        await send_managed_message(chat_id, support_help_text(), attachments=build_keyboard(), page=UI_PAGE_SUPPORT)
        return True

    if command == "/files":
        await send_files_menu(chat_id)
        return True

    if command == "/channel":
        await send_channel(chat_id)
        return True

    if command == "/buy":
        if not arg:
            await max_send_message(chat_id, "Выбери тариф кнопками ниже.", attachments=build_tariffs_keyboard_pricing())
            return True
        plan = arg.lower().strip()
        if plan == "free":
            user_profile(chat_id)
            state.user_store.set_plan(chat_id, "free")
            state.user_store.set_selected_model(chat_id, best_default_alias_for_plan("free"))
            await max_send_message(chat_id, "Тариф переключен на free.", attachments=build_tariffs_keyboard_pricing())
            return True
        if plan not in BUYABLE_PLANS:
            await max_send_message(chat_id, "Доступно: Lite, Start или Pro.", attachments=build_tariffs_keyboard_pricing())
            return True
        ok = await start_buy_flow(chat_id, plan)
        if not ok:
            return True
        return True

    if command == "/payments":
        await send_payments(chat_id)
        return True

    if command == "/ref":
        row = user_profile(chat_id)
        if not arg:
            code = str(row.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
            invited = int(row.get("referrals_invited", 0) or 0)
            await max_send_message(
                chat_id,
                (
                    f"👥 Твой реф-код: {code}\n"
                    f"Приглашено друзей: {invited}\n"
                    f"Бонус за каждого друга: +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов вам обоим.\n\n"
                    f"Нажми «Поделиться» или попроси друга ввести: /ref {code}"
                ),
                attachments=build_growth_keyboard(chat_id),
            )
            return True
        if not arg:
            code = str(row.get("referral_code", "")).strip() or referral_code_for_chat(chat_id)
            invited = int(row.get("referrals_invited", 0) or 0)
            await max_send_message(
                chat_id,
                (
                    f"👥 Твой реф-код: {code}\n"
                    f"Приглашено друзей: {invited}\n"
                    f"Бонус за каждого друга: +{request_balance_text(REFERRAL_BONUS_CREDITS)} запросов вам обоим.\n\n"
                    "Нажми «Поделиться» или попроси друга открыть Меню → Бонусы → Ввести реф-код."
                ),
                attachments=build_growth_keyboard(chat_id),
            )
            return True

        ok, info = state.user_store.apply_referral_code(chat_id, arg, REFERRAL_BONUS_CREDITS)
        if not ok:
            await max_send_message(chat_id, info, attachments=build_growth_keyboard(chat_id))
            return True
        owner_chat_id = int(info)
        log_referral_activation(chat_id, owner_chat_id, arg, "command")
        await notify_referral_success(chat_id, owner_chat_id, "command")
        return True

    if command == "/promo":
        if not arg:
            state.pending_promo_code_input.add(chat_id)
            await max_send_message(chat_id, "Введи промокод одним сообщением.", attachments=build_growth_keyboard(chat_id))
            return True
        code = normalize_referral_code(arg)
        credits, bonus_ttl_days, reason = promo_offer_for_code(code)
        if credits <= 0:
            await max_send_message(chat_id, reason or "Такого промокода нет или он выключен.", attachments=build_growth_keyboard(chat_id))
            return True
        ok, info = state.user_store.redeem_promo_code(chat_id, code, credits, bonus_ttl_days=bonus_ttl_days)
        if not ok:
            await max_send_message(chat_id, info, attachments=build_growth_keyboard(chat_id))
            return True
        state.user_store.set_acquisition_meta(chat_id, source="promo", campaign=code)
        state.user_store.record_usage_event(
            chat_id=chat_id,
            event_type="promo_activation",
            plan=str(row.get("plan", "")),
            credits_spent=int(info),
            details=f"code={code};source=command",
        )
        ttl_tail = f" Срок действия бонуса: {bonus_ttl_days} дн." if bonus_ttl_days > 0 else ""
        await max_send_message(chat_id, f"Промокод активирован: +{request_balance_text(int(info))} запросов.{ttl_tail}", attachments=build_growth_keyboard(chat_id))
        return True

    if command == "/preset":
        if not arg:
            await max_send_message(
                chat_id,
                "Выбери режим: /preset fast|balanced|quality|expert",
                attachments=build_keyboard(),
            )
            return True
        preset = arg.lower().strip()
        preset_cfg = MODEL_PRESETS.get(preset)
        if not preset_cfg:
            await max_send_message(
                chat_id,
                "Неизвестный режим. Доступно: fast, balanced, quality, expert.",
                attachments=build_keyboard(),
            )
            return True
        alias = resolve_preset_alias_for_chat(chat_id, preset)
        label = await set_user_model(chat_id, alias)
        state.user_store.set_selected_preset(chat_id, preset)
        state.user_store.record_usage_event(
            chat_id=chat_id,
            event_type="preset_select",
            plan=str(user_profile(chat_id).get("plan", "free")),
            model_alias=alias,
            details=f"preset={preset};source=command",
        )
        await max_send_message(
            chat_id,
            f"Режим: {preset_cfg['label']}\nМодель: {label}",
            attachments=build_keyboard(),
        )
        return True

    if command == "/model":
        if not arg:
            await max_send_message(chat_id, "Укажи модель: /model deepseek|gpt|gpt4o|gemini|gpt54", attachments=build_keyboard())
            return True
        label = await set_user_model(chat_id, arg)
        state.user_store.set_selected_preset(chat_id, "")
        state.user_store.record_usage_event(
            chat_id=chat_id,
            event_type="model_select",
            plan=str(user_profile(chat_id).get("plan", "free")),
            model_alias=str(arg).strip().lower(),
            details="source=command",
        )
        await max_send_message(chat_id, f"Выбрана модель: {label}", attachments=build_keyboard())
        return True

    if command in {"/clear", "/reset"}:
        state.history(chat_id).clear()
        await max_send_message(chat_id, "Контекст диалога очищен.", attachments=build_keyboard())
        return True

    if command == "/image":
        if not arg:
            await send_image_menu(chat_id)
            return True
        prepared_prompt = build_image_prompt(arg, chat_id)
        return await process_image_generation(chat_id, arg, model_prompt=prepared_prompt)

    if command == "/image_ref":
        row = user_profile(chat_id)
        plan_name = str(row.get("plan", "free")).strip().lower()
        if plan_name != "free" and not plan_allowed(plan_name, DEFAULT_IMAGE_MODEL.min_plan):
            await max_send_message(
                chat_id,
                f"Режим «по фото» доступен с тарифа {DEFAULT_IMAGE_MODEL.min_plan}. Открой «Тарифы».",
                attachments=build_tariffs_keyboard_pricing(),
            )
            return True
        reference = get_recent_reference_image(chat_id)
        if not reference:
            state.pending_image_ref_prompt.add(chat_id)
            await max_send_message(
                chat_id,
                (
                    "Сначала отправь фото. "
                    + (
                        "На Free доступно 1 действие с картинкой каждые 7 дней."
                        if plan_name == "free"
                        else f"Стоимость редактирования фото: {request_cost_text(image_edit_credit_cost())} запросов."
                    )
                ),
                attachments=build_image_prompt_keyboard(),
            )
            return True
        if not arg:
            state.pending_image_ref_prompt.add(chat_id)
            await max_send_message(
                chat_id,
                "Опиши, что сделать с фото (например: «нарисуй её в стиле киберпанк»).",
                attachments=build_image_prompt_keyboard(),
            )
            return True
        return await process_image_edit_generation(chat_id, arg, reference)

    if command.startswith("/admin"):
        return await handle_admin(chat_id, text)

    return False


async def handle_pending_receipt_input(chat_id: int, text: str) -> bool:
    plan = state.pending_receipt_plan.get(chat_id)
    if not plan:
        return False

    lowered = text.strip().lower()
    if text.strip().startswith("/") and lowered != "/cancel":
        return False
    if lowered in {"отмена", "cancel", "/cancel"}:
        state.pending_receipt_plan.pop(chat_id, None)
        await show_ui_page(chat_id, receipt_return_page(plan), push_history=False)
        return True

    email, phone = parse_receipt_contact(text)
    if not (email or phone):
        target_label = "подписки" if not plan.startswith("topup") else "пакета запросов"
        await show_managed_content(
            chat_id,
            (
                f"Не удалось распознать контакт для {target_label}.\n"
                "Отправь email (user@example.com) или телефон (+79991234567).\n\n"
                "Можно нажать «Отмена» ниже."
            ),
            attachments=build_receipt_contact_keyboard(),
            page=UI_PAGE_SUPPORT,
        )
        return True

    user_profile(chat_id)
    state.user_store.set_receipt_contact(chat_id, email=email, phone=phone)
    state.pending_receipt_plan.pop(chat_id, None)
    if plan.startswith("topup_consent:"):
        code = plan.split(":", 1)[1].lower().strip()
        await send_topup_consent(chat_id, code)
        return True
    if plan.startswith("topup:"):
        code = plan.split(":", 1)[1].lower().strip()
        await send_topup_consent(chat_id, code)
        return True
    await send_buy_consent(chat_id, plan)
    return True


async def handle_pending_referral_input(chat_id: int, text: str) -> bool:
    if chat_id not in state.pending_referral_code_input:
        return False
    lowered = text.strip().lower()
    if lowered in {"отмена", "cancel", "/cancel"}:
        state.pending_referral_code_input.discard(chat_id)
        await show_ui_page(chat_id, UI_PAGE_GROWTH, push_history=False)
        return True
    if text.strip().startswith("/"):
        state.pending_referral_code_input.discard(chat_id)
        return False
    if not looks_like_bonus_code(text):
        state.pending_referral_code_input.discard(chat_id)
        return False

    state.pending_referral_code_input.discard(chat_id)
    ok, info = state.user_store.apply_referral_code(chat_id, text, REFERRAL_BONUS_CREDITS)
    if not ok:
        await show_managed_content(chat_id, info, attachments=build_growth_keyboard(chat_id), page=UI_PAGE_GROWTH)
        return True
    owner_chat_id = int(info)
    log_referral_activation(chat_id, owner_chat_id, text, "input")
    await notify_referral_success(chat_id, owner_chat_id, "input")
    return True


async def handle_pending_promo_input(chat_id: int, text: str) -> bool:
    if chat_id not in state.pending_promo_code_input:
        return False
    lowered = text.strip().lower()
    if lowered in {"отмена", "cancel", "/cancel"}:
        state.pending_promo_code_input.discard(chat_id)
        await show_ui_page(chat_id, UI_PAGE_GROWTH, push_history=False)
        return True
    if text.strip().startswith("/"):
        state.pending_promo_code_input.discard(chat_id)
        return False
    if not looks_like_bonus_code(text):
        state.pending_promo_code_input.discard(chat_id)
        return False
    state.pending_promo_code_input.discard(chat_id)

    code = normalize_referral_code(text)
    credits, bonus_ttl_days, reason = promo_offer_for_code(code)
    if credits <= 0:
        await show_managed_content(
            chat_id,
            reason or "Такого промокода нет или он выключен.",
            attachments=build_growth_keyboard(chat_id),
            page=UI_PAGE_GROWTH,
        )
        return True
    ok, info = state.user_store.redeem_promo_code(chat_id, code, credits, bonus_ttl_days=bonus_ttl_days)
    if not ok:
        await show_managed_content(chat_id, info, attachments=build_growth_keyboard(chat_id), page=UI_PAGE_GROWTH)
        return True
    state.user_store.set_acquisition_meta(chat_id, source="promo", campaign=code)
    row = user_profile(chat_id)
    state.user_store.record_usage_event(
        chat_id=chat_id,
        event_type="promo_activation",
        plan=str(row.get("plan", "")),
        credits_spent=int(info),
        details=f"code={code};source=input",
    )
    ttl_tail = f" Срок действия бонуса: {bonus_ttl_days} дн." if bonus_ttl_days > 0 else ""
    await show_managed_content(
        chat_id,
        f"Промокод активирован: +{request_balance_text(int(info))} запросов.{ttl_tail}",
        attachments=build_growth_keyboard(chat_id),
        page=UI_PAGE_GROWTH,
    )
    return True


async def process_update(update: dict[str, Any]) -> None:
    if not isinstance(update, dict) or not is_supported_update(update):
        return
    if is_channel_update(update):
        log.info("Channel update skipped")
        return
    if not remember_update(update):
        log.info("Duplicate update skipped")
        return

    update_type = update.get("update_type")
    if update_type == "message_callback":
        await handle_callback(update)
        return

    chat_id, text = parse_incoming_text(update)
    incoming_image_url = parse_incoming_image_url(update)
    if chat_id is None:
        return
    if chat_id <= 0:
        log.info("Non-dialog update skipped chat_id=%s update_type=%s", chat_id, update_type)
        return

    ensure_update_user_binding(chat_id, update)
    row = user_profile(chat_id)
    state.user_store.touch_last_active(chat_id)
    if row["is_blocked"]:
        return

    if update_type in {"bot_started", "user_added", "bot_added"}:
        source, campaign = acquisition_meta_from_start_payload(update)
        state.user_store.set_acquisition_meta(chat_id, source=source or "direct", campaign=campaign)
        referral_applied = await maybe_apply_start_referral(chat_id, update)
        if referral_applied:
            row = user_profile(chat_id)

    if CHANNEL_GATE_ENABLED and not channel_gate_allows_text(text or ""):
        if not await ensure_channel_access(chat_id):
            return

    if incoming_image_url:
        try:
            try:
                incoming_image = await fetch_image_bytes(incoming_image_url, use_max_auth=True)
            except Exception:
                incoming_image = await fetch_image_bytes(incoming_image_url, use_max_auth=False)
            remember_reference_image(chat_id, encode_data_url(incoming_image))
        except Exception as exc:
            log.exception("Failed to save incoming reference image for chat_id=%s", chat_id)
            with suppress(Exception):
                await max_send_message(chat_id, f"Не удалось обработать фото: {exc}", attachments=build_image_menu_keyboard(chat_id))
            return

        if text and (chat_id in state.pending_image_ref_prompt or looks_like_image_ref_request(text)):
            state.pending_image_ref_prompt.discard(chat_id)
            await process_image_edit_generation(chat_id, text, get_recent_reference_image(chat_id))
            return

        if not text:
            state.pending_image_ref_prompt.add(chat_id)
            row = user_profile(chat_id)
            plan_name = str(row.get("plan", "free")).strip().lower()
            await max_send_message(
                chat_id,
                (
                    "Фото получил ✅\n"
                    "Теперь напиши, что с ним сделать (например: «нарисуй её в стиле аниме»).\n"
                    + (
                        "На Free доступно 1 действие с картинкой каждые 7 дней."
                        if plan_name == "free"
                        else f"Стоимость редактирования фото: {request_cost_text(image_edit_credit_cost())} запросов."
                    )
                ),
                attachments=build_image_prompt_keyboard(),
            )
            return

    if update_type in {"bot_started", "user_added", "bot_added"} and not text:
        if int(row.get("onboarding_done", 0) or 0) == 0:
            await send_onboarding(chat_id, step=1)
        else:
            await send_menu(chat_id)
        return
    if not text:
        return

    log.info("Incoming update=%s chat_id=%s text=%r", update_type, chat_id, text[:120])
    first_word = text.strip().lower().split(maxsplit=1)[0] if text.strip() else ""
    if int(row.get("onboarding_done", 0) or 0) == 0 and first_word not in {"/start", "старт", "start", "начать"}:
        state.user_store.set_onboarding_done(chat_id, True)
    try:
        if await handle_pending_referral_input(chat_id, text):
            return
        if await handle_pending_promo_input(chat_id, text):
            return
        if await handle_pending_receipt_input(chat_id, text):
            return
        if await handle_pending_file_input(chat_id, text):
            return
        if await handle_pending_image_ref_prompt_input(chat_id, text):
            return
        if await handle_pending_image_prompt_input(chat_id, text):
            return
        if should_intercept_image_flow_text(chat_id, text):
            if get_image_prefs(chat_id).get("mode") == "edit":
                state.pending_image_ref_prompt.add(chat_id)
                if await handle_pending_image_ref_prompt_input(chat_id, text):
                    return
            else:
                state.pending_image_prompt.add(chat_id)
                if await handle_pending_image_prompt_input(chat_id, text):
                    return
        recent_reference = get_recent_reference_image(chat_id)
        if recent_reference and looks_like_image_ref_request(text) and not text.strip().startswith("/"):
            await process_image_edit_generation(chat_id, text, recent_reference)
            return
        if not text.strip().startswith("/"):
            file_kind = looks_like_file_request(text)
            if file_kind:
                file_profile = detect_file_profile_from_text(text, default_file_profile_for_plan(str(row.get("plan", "free"))))
                if str(row.get("plan", "free")).strip().lower() == "free" and file_profile != "short":
                    await max_send_message(
                        chat_id,
                        "На Free доступен только короткий файл раз в 14 дней. Средняя и полная версии открываются на платных тарифах.",
                        attachments=purchase_help_keyboard_for_row(row),
                    )
                    return
                if file_prompt_has_enough_detail(text):
                    await process_file_request(chat_id, file_kind, text, profile=file_profile)
                else:
                    state.pending_file_kind[chat_id] = file_kind
                    state.pending_file_profile[chat_id] = file_profile
                    await show_managed_content(
                        chat_id,
                        build_file_prompt_text(file_kind, file_profile, chat_id),
                        attachments=build_files_prompt_keyboard(),
                        page=UI_PAGE_FILES_MENU,
                        push_history=False,
                    )
                return
        if await handle_command(chat_id, text):
            return

        if len(text) > MAX_TEXT_INPUT_CHARS:
            await max_send_message(
                chat_id,
                f"Слишком длинное сообщение. Максимум {MAX_TEXT_INPUT_CHARS} символов.",
                attachments=build_keyboard(),
            )
            return

        ok_cd, reason_cd = check_cooldown(chat_id, "message")
        if not ok_cd:
            await max_send_message(chat_id, reason_cd, attachments=build_keyboard())
            return

        ok, reason = check_limit_only(chat_id, "messages")
        if not ok:
            await max_send_message(chat_id, reason, attachments=build_keyboard())
            return

        selected_alias = str(user_profile(chat_id).get("selected_model_alias") or DEFAULT_TEXT_MODEL.alias)
        model_label = TEXT_MODELS.get(selected_alias, DEFAULT_TEXT_MODEL).label
        fixed_text_cost = text_credit_cost(selected_alias)
        plan_name, _, _, estimated_messages = build_text_request(chat_id, text, selected_alias=selected_alias)
        estimated_prompt_tokens = estimate_tokens_from_messages(estimated_messages)
        estimated_total_tokens = estimated_prompt_tokens + completion_tokens_for_plan(plan_name)
        reserved_var_cost = variable_text_credits(selected_alias, estimated_total_tokens)
        reserved_total_cost = normalize_public_request_credit_cost(fixed_text_cost + reserved_var_cost)

        ok_credit, reason_credit = check_and_consume_credits(
            chat_id,
            reserved_total_cost,
            f"текст ({model_label})",
        )
        if not ok_credit:
            await max_send_message(chat_id, reason_credit, attachments=purchase_help_keyboard_for_row(user_profile(chat_id)))
            return

        ok, reason = check_and_consume_limit(chat_id, "messages")
        if not ok:
            state.user_store.refund_credits(chat_id, reserved_total_cost)
            await max_send_message(chat_id, reason, attachments=build_keyboard())
            return

        await max_send_message(chat_id, f"Думаю... Модель: {current_model_label(chat_id)}", notify=False)
        try:
            result = await ask_text_model(chat_id, text, selected_alias=selected_alias)
            actual_var_cost = variable_text_credits(selected_alias, result.total_tokens)
            actual_total_cost = normalize_public_request_credit_cost(fixed_text_cost + actual_var_cost)
            if reserved_total_cost > actual_total_cost:
                state.user_store.refund_credits(chat_id, reserved_total_cost - actual_total_cost)
            elif actual_total_cost > reserved_total_cost:
                extra_cost = actual_total_cost - reserved_total_cost
                ok_extra, _ = check_and_consume_credits(chat_id, extra_cost, f"сложность ({model_label})")
                if not ok_extra:
                    log.warning(
                        "Variable credits under-reserved chat_id=%s alias=%s reserved_var=%s actual_var=%s tokens=%s",
                        chat_id,
                        selected_alias,
                        reserved_var_cost,
                        actual_var_cost,
                        result.total_tokens,
                    )
            final_row = user_profile(chat_id)
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="text_request",
                plan=str(final_row.get("plan", "")),
                model_alias=selected_alias,
                credits_spent=actual_total_cost,
                tokens_total=int(result.total_tokens),
                details=f"prompt={int(result.prompt_tokens)};completion={int(result.completion_tokens)}",
            )
            await max_send_message(
                chat_id,
                result.text,
                attachments=build_reply_shortcuts_keyboard(chat_id),
                text_format="markdown",
            )
            with suppress(Exception):
                await maybe_send_low_credits_nudge(chat_id)
        except Exception:
            state.user_store.refund_credits(chat_id, reserved_total_cost)
            raise
    except Exception as exc:
        log.exception("Failed to process update")
        capture_exception_safe(exc)
        record_runtime_error()
        if is_expected_max_delivery_error(exc):
            log.info("Skipping process_update alert for suspended/denied chat_id=%s: %s", chat_id, exc)
        else:
            await notify_admin_alert("process_update", f"chat_id={chat_id}\nerror={exc}")
        with suppress(Exception):
            await max_send_message(chat_id, f"Ошибка: {exc}")


async def get_updates(marker: int | None = None) -> tuple[list[dict[str, Any]], int | None]:
    params: list[tuple[str, str]] = [("limit", "100"), ("timeout", "25")]
    if marker is not None:
        params.append(("marker", str(marker)))

    status, data, _ = await http_json_request_with_retries(
        "GET",
        f"{MAX_API}/updates",
        headers={"Authorization": MAX_TOKEN},
        params=params,
        semaphore=state.max_api_semaphore,
        request_name="max_get_updates",
    )
    if status >= 400:
        raise RuntimeError(f"MAX updates error {status}: {data}")

    updates = data.get("updates")
    next_marker = data.get("marker")
    if not isinstance(updates, list):
        updates = []
    if next_marker is not None and not isinstance(next_marker, int):
        try:
            next_marker = int(next_marker)
        except Exception:
            next_marker = marker
    return updates, next_marker


async def polling_loop() -> None:
    marker: int | None = None
    log.info("Polling mode started")
    while True:
        try:
            updates, marker = await get_updates(marker)
            for update in updates:
                await process_update(update)
            cpu_per_core = float(system.get("cpu_load_1m_per_core", 0.0) or 0.0)
            if cpu_per_core >= max(0.1, ALERT_CPU_LOAD_PER_CORE_THRESHOLD):
                await notify_admin_alert(
                    "cpu_pressure",
                    (
                        f"CPU load per core (1m): {cpu_per_core:.2f} "
                        f"(threshold {ALERT_CPU_LOAD_PER_CORE_THRESHOLD:.2f})."
                    ),
                )
            mem_used_pct = float(system.get("memory_used_pct", 0.0) or 0.0)
            if mem_used_pct >= max(1.0, ALERT_MEMORY_USED_PCT_THRESHOLD):
                await notify_admin_alert(
                    "memory_pressure",
                    (
                        f"Memory used: {mem_used_pct:.1f}% "
                        f"(threshold {ALERT_MEMORY_USED_PCT_THRESHOLD:.1f}%)."
                    ),
                )
            disk_used_pct = float(system.get("disk_used_pct", 0.0) or 0.0)
            if disk_used_pct >= max(1.0, ALERT_DISK_USED_PCT_THRESHOLD):
                await notify_admin_alert(
                    "disk_pressure",
                    (
                        f"Disk used: {disk_used_pct:.1f}% "
                        f"(threshold {ALERT_DISK_USED_PCT_THRESHOLD:.1f}%)."
                    ),
                )
            backup_age_hours_raw = report.get("latest_backup_age_hours")
            backup_age_hours = float(backup_age_hours_raw) if backup_age_hours_raw is not None else None
            if backup_age_hours is None:
                await notify_admin_alert(
                    "backup_missing",
                    "No DB backup file found in data/backups.",
                )
            elif backup_age_hours >= max(1, ALERT_BACKUP_STALE_HOURS):
                await notify_admin_alert(
                    "backup_stale",
                    (
                        f"Latest backup age is {backup_age_hours:.1f}h "
                        f"(threshold {ALERT_BACKUP_STALE_HOURS}h)."
                    ),
                )
        except asyncio.CancelledError:
            raise
        except Exception as exc:
            log.exception("Polling loop error")
            capture_exception_safe(exc)
            record_runtime_error()
            await notify_admin_alert("polling_loop", f"Polling loop error: {exc}")
            await asyncio.sleep(3)


async def backup_loop() -> None:
    interval_seconds = max(3600, AUTO_BACKUP_INTERVAL_HOURS * 3600)
    log.info("Backup loop started, interval=%ss", interval_seconds)
    while True:
        try:
            await asyncio.sleep(interval_seconds)
            backup_file = create_db_backup()
            log.info("Automatic DB backup created: %s", backup_file)
        except asyncio.CancelledError:
            raise
        except Exception as exc:
            log.exception("Automatic backup failed")
            capture_exception_safe(exc)
            record_runtime_error()
            await notify_admin_alert("db_backup", f"Automatic DB backup failed: {exc}")


async def monitor_loop() -> None:
    interval_seconds = max(300, SERVICE_MONITOR_INTERVAL_MINUTES * 60)
    log.info("Service monitor loop started, interval=%ss", interval_seconds)
    while True:
        try:
            await asyncio.sleep(interval_seconds)
            report = service_status_report()
            monitor = report["monitor"]
            system = report.get("system", {}) if isinstance(report, dict) else {}
            recent_errors = int(report["recent_runtime_errors"] or 0)
            if recent_errors >= max(1, ALERT_HIGH_ERRORS_THRESHOLD):
                await notify_admin_alert(
                    "high_errors",
                    f"За последние {ALERT_HIGH_ERRORS_WINDOW_MINUTES} минут ошибок: {recent_errors}",
                )
            active_users = int(monitor.get("active_users", 0) or 0)
            payments_count = int(monitor.get("payments_count", 0) or 0)
            if active_users >= max(1, ALERT_LOW_PAYMENTS_MIN_ACTIVE_USERS) and payments_count <= max(0, ALERT_LOW_PAYMENTS_MAX_PAYMENTS):
                await notify_admin_alert(
                    "low_payments",
                    (
                        f"За последние {ALERT_LOW_PAYMENTS_LOOKBACK_HOURS} ч активных пользователей: {active_users}, "
                        f"оплат: {payments_count}. Проверь платежный UX и конверсию."
                    ),
                )
            recent_cost = float(monitor.get("recent_text_cost_rub", 0.0) or 0.0)
            baseline_cost = float(monitor.get("baseline_text_cost_rub", 0.0) or 0.0)
            if (
                recent_cost >= max(0.0, ALERT_SPEND_SPIKE_MIN_RUB)
                and baseline_cost > 0
                and recent_cost >= baseline_cost * max(1.1, ALERT_SPEND_SPIKE_MULTIPLIER)
            ):
                await notify_admin_alert(
                    "spend_spike",
                    (
                        f"Текстовая себестоимость за последние {ALERT_SPEND_SPIKE_LOOKBACK_HOURS} ч выросла до "
                        f"{recent_cost:.0f} ₽ против ожидаемых {baseline_cost:.0f} ₽."
                    ),
                )
        except asyncio.CancelledError:
            raise
        except Exception as exc:
            log.exception("Service monitor loop failed")
            capture_exception_safe(exc)
            record_runtime_error()
            await notify_admin_alert("service_monitor", f"Service monitor loop failed: {exc}")


def render_status_html() -> str:
    report = service_status_report()
    esc = html.escape
    smoke_rows = []
    for row in report["smoke_checks"]:
        ok = row["ok"] == "ok"
        smoke_rows.append(
            "<tr>"
            f"<td>{'✅' if ok else '❌'} {esc(str(row['name']))}</td>"
            f"<td>{esc(str(row['details']))}</td>"
            "</tr>"
        )
    monitor = report["monitor"]
    backup_at = report["latest_backup_at"] or "-"
    backup_path = report["latest_backup_path"] or "-"
    return f"""<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>Состояние сервиса</title>
  <link rel="stylesheet" href="/assets/style.css"/>
  <style>
    table {{ width:100%; border-collapse:collapse; }}
    th, td {{ border-bottom:1px solid var(--line); padding:10px 8px; text-align:left; font-size:14px; vertical-align:top; }}
    .mini-grid {{ display:grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap:10px; margin:14px 0; }}
    .mini-metric {{ border:1px solid var(--line); border-radius:12px; padding:12px; background:#fcfdff; }}
    .mini-metric strong {{ display:block; font-size:22px; margin-top:4px; }}
    .muted {{ color:var(--muted); }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="card">
      <div class="top"><span class="badge">Состояние сервиса</span><span class="badge">Срез: {esc(str(report['generated_at']))}</span></div>
      <h1>Статус бота</h1>
      <div class="actions">
        <a class="btn" href="/analytics">Аналитика</a>
        <a class="btn" href="/admin/panel">Админка</a>
        <a class="btn btn-primary" href="/status">Обновить</a>
        <a class="btn" href="/health/deep">JSON</a>
      </div>
      <div class="mini-grid">
        <div class="mini-metric">Run mode<strong>{esc(str(report['run_mode']))}</strong></div>
        <div class="mini-metric">Uptime<strong>{esc(str(report['uptime']))}</strong></div>
        <div class="mini-metric">DB size<strong>{int(report['db_size_bytes']) // 1024} KB</strong></div>
        <div class="mini-metric">Ошибки за окно<strong>{int(report['recent_runtime_errors'])}</strong></div>
        <div class="mini-metric">Активные за {ALERT_LOW_PAYMENTS_LOOKBACK_HOURS}ч<strong>{int(monitor.get('active_users', 0) or 0)}</strong></div>
        <div class="mini-metric">Оплаты за {ALERT_LOW_PAYMENTS_LOOKBACK_HOURS}ч<strong>{int(monitor.get('payments_count', 0) or 0)}</strong></div>
        <div class="mini-metric">Text cost {ALERT_SPEND_SPIKE_LOOKBACK_HOURS}ч<strong>{float(monitor.get('recent_text_cost_rub', 0.0) or 0.0):.0f} ₽</strong></div>
        <div class="mini-metric">Baseline cost<strong>{float(monitor.get('baseline_text_cost_rub', 0.0) or 0.0):.0f} ₽</strong></div>
      </div>
      <p class="muted">Последний бэкап: {esc(str(backup_at))}<br/>{esc(str(backup_path))}</p>
      <p class="muted">Polling: {'on' if report['polling_task'] else 'off'} • Backup loop: {'on' if report['backup_task'] else 'off'} • Monitor loop: {'on' if report['monitor_task'] else 'off'} • HTTP session: {'ok' if report['session_ready'] else 'closed'}</p>
    </div>
    <div class="card">
      <h2>Smoke-check после деплоя</h2>
      <table>
        <tr><th>Проверка</th><th>Детали</th></tr>
        {''.join(smoke_rows)}
      </table>
    </div>
  </div>
</body>
</html>"""


@asynccontextmanager
async def lifespan(_: FastAPI):
    require_env()
    validate_pricing_sanity()
    init_sentry_if_enabled()
    await get_session()
    worker_count = max(1, OPENROUTER_IMAGE_CONCURRENCY)
    state.image_worker_tasks = [
        asyncio.create_task(image_worker_loop(index + 1))
        for index in range(worker_count)
    ]
    if RUN_MODE == "polling":
        state.polling_task = asyncio.create_task(polling_loop())
    if AUTO_BACKUP_ENABLED:
        state.backup_task = asyncio.create_task(backup_loop())
    if SERVICE_MONITOR_ENABLED:
        state.monitor_task = asyncio.create_task(monitor_loop())
    try:
        yield
    finally:
        for task in state.image_worker_tasks:
            task.cancel()
        for task in state.image_worker_tasks:
            with suppress(asyncio.CancelledError):
                await task
        state.image_worker_tasks.clear()
        if state.polling_task:
            state.polling_task.cancel()
            with suppress(asyncio.CancelledError):
                await state.polling_task
        if state.backup_task:
            state.backup_task.cancel()
            with suppress(asyncio.CancelledError):
                await state.backup_task
        if state.monitor_task:
            state.monitor_task.cancel()
            with suppress(asyncio.CancelledError):
                await state.monitor_task
        if state.session and not state.session.closed:
            await state.session.close()


app = FastAPI(title="MAX Multi AI Bot", lifespan=lifespan)
app.mount("/assets", StaticFiles(directory=str(SITE_DIR / "assets")), name="assets")


@app.get("/health")
async def health() -> dict[str, Any]:
    report = service_status_report()
    return {
        "status": "ok" if report["db_exists"] and report["session_ready"] else "degraded",
        "run_mode": report["run_mode"],
        "uptime": report["uptime"],
        "db_exists": report["db_exists"],
        "backup_task": report["backup_task"],
        "monitor_task": report["monitor_task"],
    }


@app.get("/health/deep")
async def health_deep(request: Request, token: str = "") -> dict[str, Any]:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        raise HTTPException(status_code=401, detail="auth required")
    return service_status_report()


@app.get("/status", response_class=HTMLResponse)
async def status_page(request: Request, token: str = "") -> HTMLResponse:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        response = HTMLResponse(render_admin_login_html("Нужен пароль администратора."), status_code=401)
        clear_admin_cookie(response)
        return response
    response = HTMLResponse(render_status_html())
    set_admin_cookie(response, session_id)
    return response


@app.get("/", response_class=FileResponse)
async def landing_index() -> FileResponse:
    return FileResponse(site_file("index.html"))


@app.get("/offer", response_class=FileResponse)
async def landing_offer() -> FileResponse:
    return FileResponse(site_file("offer.html"))


@app.get("/privacy", response_class=FileResponse)
async def landing_privacy() -> FileResponse:
    return FileResponse(site_file("privacy.html"))


@app.get("/refund", response_class=FileResponse)
async def landing_refund() -> FileResponse:
    return FileResponse(site_file("refund.html"))


@app.get("/contacts", response_class=FileResponse)
async def landing_contacts() -> FileResponse:
    return FileResponse(site_file("contacts.html"))


@app.get("/contacts/meta")
async def contacts_meta() -> dict[str, str]:
    return {
        "email": CONTACT_EMAIL,
        "phone": CONTACT_PHONE,
        "support_url": support_url_value(),
        "support_text": SUPPORT_TEXT,
    }


@app.get("/support", response_class=FileResponse)
async def landing_support() -> FileResponse:
    return FileResponse(site_file("support.html"))


@app.get("/support/meta")
async def support_meta() -> dict[str, str]:
    return {"url": support_url_value(), "text": SUPPORT_TEXT, "email": CONTACT_EMAIL}


@app.get("/mailru-domainMB5PESlCeJQEXuoC.html", response_class=FileResponse)
async def mailru_domain_verify() -> FileResponse:
    return FileResponse(site_file("mailru-domainMB5PESlCeJQEXuoC.html"))


def render_admin_login_html(error: str = "") -> str:
    esc = html.escape
    error_block = ""
    if error:
        error_block = f"<p class='error-box'>{esc(error)}</p>"
    return f"""<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>Вход в аналитику</title>
  <link rel="stylesheet" href="/assets/style.css"/>
  <style>
    .login-card {{ max-width: 520px; margin: 48px auto 0; }}
    .form-grid {{ display: grid; gap: 12px; margin-top: 18px; }}
    .field-label {{ font-size: 14px; font-weight: 700; color: var(--text); }}
    .input {{ width: 100%; border: 1px solid var(--line); border-radius: 12px; padding: 12px 14px; font: inherit; }}
    .error-box {{ border: 1px solid #ffd1d5; background: #fff0f1; color: #b82b3d; border-radius: 12px; padding: 10px 12px; }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="card login-card">
      <div class="top">
        <span class="badge">Закрытый раздел</span>
      </div>
      <h1>Аналитика и управление</h1>
      <p>Введи пароль администратора, чтобы открыть денежную аналитику и панель управления ботом.</p>
      {error_block}
      <form class="form-grid" method="post" action="/analytics/login">
        <label class="field-label" for="password">Пароль</label>
        <input class="input" id="password" name="password" type="password" autocomplete="current-password" placeholder="Пароль администратора"/>
        <button class="btn btn-primary" type="submit">Войти</button>
      </form>
    </div>
  </div>
</body>
</html>"""


def render_admin_analytics_html_v2(token: str, days: int = 30) -> str:
    report = state.user_store.kpi_report(days=days)
    esc = html.escape

    def money(value: int | float) -> str:
        return f"{int(round(value)):,}".replace(",", " ")

    def requests_per_rub(credits: int, price_rub: int) -> str:
        if price_rub <= 0:
            return "0.00"
        return f"{public_requests_from_credits(credits) / float(price_rub):.2f}"

    def fmt_msk(value: Any) -> str:
        text = str(value or "").strip()
        dt = parse_iso_datetime(text)
        return format_msk_datetime(dt) if dt else "-"

    active_users = int(report.get("active_users", 0) or 0)
    payers = int(report.get("payers", 0) or 0)
    revenue_rub = int(report.get("revenue_rub", 0) or 0)
    refunds_raw = int(report.get("refunds_rub", 0) or 0)
    refunds_abs = abs(refunds_raw)
    payments_count = int(report.get("payments_count", 0) or 0)
    refunds_count = int(report.get("refunds_count", 0) or 0)
    net_rub = revenue_rub + refunds_raw if refunds_raw < 0 else revenue_rub - refunds_abs
    pay_share = (payers * 100.0 / active_users) if active_users else 0.0
    arpu = (net_rub / active_users) if active_users else 0.0
    arppu = (net_rub / payers) if payers else 0.0
    avg_check = (revenue_rub / payments_count) if payments_count else 0.0
    refund_rate_pct = (refunds_count * 100.0 / payments_count) if payments_count else 0.0
    text_requests = int(report.get("text_requests", 0) or 0)
    image_requests = int(report.get("image_requests", 0) or 0)
    text_credits = int(report.get("text_credits", 0) or 0)
    image_credits = int(report.get("image_credits", 0) or 0)
    total_credits = int(report.get("total_credits_spent", 0) or 0)
    text_tokens = int(report.get("text_tokens", 0) or 0)
    avg_tokens = (text_tokens / text_requests) if text_requests else 0.0
    estimated_text_cost_rub = float(report.get("estimated_text_cost_rub", 0.0) or 0.0)
    estimated_text_contribution_rub = net_rub - estimated_text_cost_rub
    estimated_text_margin_pct = (estimated_text_contribution_rub * 100.0 / net_rub) if net_rub > 0 else 0.0
    text_cost_share_pct = (estimated_text_cost_rub * 100.0 / net_rub) if net_rub > 0 else 0.0
    referral_activations = int(report.get("referral_activations", 0) or 0)
    referred_payers = int(report.get("referred_payers", 0) or 0)
    referral_bonus_credits = int(report.get("referral_bonus_credits", 0) or 0)
    referral_conversion_pct = (referred_payers * 100.0 / referral_activations) if referral_activations else 0.0
    referred_share_of_payers_pct = (referred_payers * 100.0 / payers) if payers else 0.0
    new_users = int(report.get("new_users", 0) or 0)
    new_paid_users = int(report.get("new_paid_users", 0) or 0)
    free_to_paid_pct = (new_paid_users * 100.0 / new_users) if new_users else 0.0
    d1_eligible = int(report.get("d1_eligible", 0) or 0)
    d1_retained = int(report.get("d1_retained", 0) or 0)
    d1_retention_pct = (d1_retained * 100.0 / d1_eligible) if d1_eligible else 0.0
    d7_eligible = int(report.get("d7_eligible", 0) or 0)
    d7_retained = int(report.get("d7_retained", 0) or 0)
    d7_retention_pct = (d7_retained * 100.0 / d7_eligible) if d7_eligible else 0.0

    period_links = " ".join(
        f"<a class='btn {'btn-primary' if days == option else ''}' href='{esc(admin_url('/analytics', token, days=option))}'>{option} дней</a>"
        for option in (7, 30, 90, 180)
    )

    plan_rows: list[str] = []
    for row in report.get("plans", []):
        plan_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('plan', '-') or '-').title())}</td>"
            f"<td>{int(row.get('payments', 0) or 0)}</td>"
            f"<td>{money(int(row.get('revenue', 0) or 0))} ₽</td>"
            "</tr>"
        )
    if not plan_rows:
        plan_rows.append("<tr><td colspan='3'>Платежей за выбранный период пока нет.</td></tr>")

    top_referrer_rows: list[str] = []
    for row in report.get("top_referrers", []):
        cid = int(row.get("chat_id", 0) or 0)
        top_referrer_rows.append(
            "<tr>"
            f"<td>{cid}</td>"
            f"<td>{int(row.get('max_user_id', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('referral_code', '') or '-'))}</td>"
            f"<td>{int(row.get('referrals_invited', 0) or 0)}</td>"
            f"<td>{int(row.get('paid_referrals', 0) or 0)}</td>"
            f"<td>{esc(fmt_msk(row.get('last_referral_at')))}</td>"
            f"<td><a href='{esc(admin_url('/admin/panel', chat_id=cid))}'>Открыть</a></td>"
            "</tr>"
        )
    if not top_referrer_rows:
        top_referrer_rows.append("<tr><td colspan='7'>Пока нет успешных реферальных активаций.</td></tr>")

    suspicious_referral_rows: list[str] = []
    for row in report.get("suspicious_referrals", []):
        cid = int(row.get("chat_id", 0) or 0)
        suspicious_referral_rows.append(
            "<tr>"
            f"<td>{cid}</td>"
            f"<td>{int(row.get('max_user_id', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('referral_code', '') or '-'))}</td>"
            f"<td>{int(row.get('referrals_invited', 0) or 0)}</td>"
            f"<td>{int(row.get('active_recent_referrals', 0) or 0)}</td>"
            f"<td>{int(row.get('paid_referrals', 0) or 0)}</td>"
            f"<td><a href='{esc(admin_url('/admin/panel', chat_id=cid))}'>Открыть</a></td>"
            "</tr>"
        )
    if not suspicious_referral_rows:
        suspicious_referral_rows.append("<tr><td colspan='7'>Подозрительных реферальных кластеров пока нет.</td></tr>")

    source_rows: list[str] = []
    for row in report.get("sources", []):
        source_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('source', 'direct') or 'direct'))}</td>"
            f"<td>{esc(str(row.get('campaign', '-') or '-'))}</td>"
            f"<td>{int(row.get('users_count', 0) or 0)}</td>"
            f"<td>{int(row.get('paid_users', 0) or 0)}</td>"
            f"<td>{money(int(row.get('revenue_rub', 0) or 0))} ₽</td>"
            "</tr>"
        )
    if not source_rows:
        source_rows.append("<tr><td colspan='5'>По источникам пока нет данных.</td></tr>")

    promo_code_rows: list[str] = []
    for row in report.get("promo_codes", []):
        activations = int(row.get("activations", 0) or 0)
        paid_users = int(row.get("paid_users", 0) or 0)
        conversion = (paid_users * 100.0 / activations) if activations else 0.0
        promo_code_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('promo_code', '-') or '-'))}</td>"
            f"<td>{activations}</td>"
            f"<td>{paid_users}</td>"
            f"<td>{conversion:.1f}%</td>"
            f"<td>{request_balance_text(int(row.get('credits', 0) or 0))}</td>"
            f"<td>{money(int(row.get('revenue_rub', 0) or 0))} ₽</td>"
            "</tr>"
        )
    if not promo_code_rows:
        promo_code_rows.append("<tr><td colspan='6'>По промокодам пока нет активаций.</td></tr>")

    economics_plan_rows: list[str] = []
    unit_margin_rows: list[str] = []
    for plan in ("lite", "start", "pro"):
        amount, period_days = plan_price_and_days(plan)
        credits = credits_for_plan(plan)
        ratio = requests_per_rub(credits, amount)
        daily_requests = public_requests_from_credits(credits) / max(1, period_days)
        models_text = {
            "lite": "DeepSeek, GPT-4.1 Nano, GPT-4o Mini, Gemini 2.5 Flash",
            "start": f"DeepSeek, GPT-4.1 Nano, GPT-4o Mini, Gemini 2.5 Flash, GPT-5.4 в режиме «Эксперт» до {PLAN_CONFIGS['start'].daily_gpt54_limit}/день",
            "pro": "DeepSeek, GPT-4.1 Nano, GPT-4o Mini, Gemini 2.5 Flash, GPT-5.4",
        }[plan]
        economics_plan_rows.append(
            "<tr>"
            f"<td>{esc(plan.title())}</td>"
            f"<td>{money(amount)} ₽</td>"
            f"<td>{period_days}</td>"
            f"<td>{request_cost_text(credits)}</td>"
            f"<td>{daily_requests:.0f}</td>"
            f"<td>{ratio}</td>"
            f"<td>{esc(models_text)}</td>"
            "</tr>"
        )
        econ = expected_unit_economics(amount, credits)
        unit_margin_rows.append(
            "<tr>"
            f"<td>{esc(plan.title())}</td>"
            f"<td>{money(amount)} ₽</td>"
            f"<td>{request_cost_text(credits)}</td>"
            f"<td>{money(econ['expected_cost_rub'])} ₽</td>"
            f"<td>{money(econ['payment_fee_rub'])} ₽</td>"
            f"<td>{money(econ['receipt_fee_rub'])} ₽</td>"
            f"<td>{money(econ['tax_rub'])} ₽</td>"
            f"<td>{money(econ['margin_rub'])} ₽</td>"
            f"<td>{econ['margin_pct']:.1f}%</td>"
            "</tr>"
        )

    economics_pack_rows: list[str] = []
    for code in ("small", "medium", "large"):
        pack = TOPUP_PACKS[code]
        credits = int(pack["credits"])
        price_rub = int(pack["price_rub"])
        economics_pack_rows.append(
            "<tr>"
            f"<td>{esc(str(pack['label']))}</td>"
            f"<td>{money(price_rub)} ₽</td>"
            f"<td>{request_balance_text(credits)}</td>"
            f"<td>{requests_per_rub(credits, price_rub)}</td>"
            "</tr>"
        )
        econ = expected_unit_economics(price_rub, credits)
        unit_margin_rows.append(
            "<tr>"
            f"<td>{esc(str(pack['label']))}</td>"
            f"<td>{money(price_rub)} ₽</td>"
            f"<td>{request_balance_text(credits)}</td>"
            f"<td>{money(econ['expected_cost_rub'])} ₽</td>"
            f"<td>{money(econ['payment_fee_rub'])} ₽</td>"
            f"<td>{money(econ['receipt_fee_rub'])} ₽</td>"
            f"<td>{money(econ['tax_rub'])} ₽</td>"
            f"<td>{money(econ['margin_rub'])} ₽</td>"
            f"<td>{econ['margin_pct']:.1f}%</td>"
            "</tr>"
        )

    margin_rows: list[str] = []
    for row in report.get("margins", []):
        revenue = int(row.get("revenue_rub", 0) or 0)
        text_cost = float(row.get("estimated_text_cost_rub", 0.0) or 0.0)
        contribution = float(row.get("contribution_rub", 0.0) or 0.0)
        if revenue <= 0 and text_cost <= 0:
            continue
        margin_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('plan', '-') or '-').title())}</td>"
            f"<td>{money(revenue)} ₽</td>"
            f"<td>{money(text_cost)} ₽</td>"
            f"<td>{money(contribution)} ₽</td>"
            f"<td>{float(row.get('margin_pct', 0.0) or 0.0):.1f}%</td>"
            f"<td>{int(row.get('text_requests', 0) or 0)}</td>"
            f"<td>{int(row.get('image_requests', 0) or 0)}</td>"
            "</tr>"
        )
    if not margin_rows:
        margin_rows.append("<tr><td colspan='7'>Пока недостаточно данных для оценки маржи по тарифам.</td></tr>")

    models_rows: list[str] = []
    for row in report.get("models", []):
        label = str(row.get("label", "") or "-")
        kind = str(row.get("kind", "") or "—")
        requests_count = int(row.get("requests", 0) or 0)
        tokens = int(row.get("tokens", 0) or 0)
        credits = int(row.get("credits", 0) or 0)
        avg_model_tokens = (tokens / requests_count) if requests_count else 0.0
        estimated_cost_rub = float(row.get("estimated_cost_usd", 0.0) or 0.0) * ANALYTICS_USD_TO_RUB
        estimated_cost_cell = f"{money(estimated_cost_rub)} ₽" if kind == "Текст" else "—"
        models_rows.append(
            "<tr>"
            f"<td>{esc(label)}</td>"
            f"<td>{esc(kind)}</td>"
            f"<td>{requests_count}</td>"
            f"<td>{money(tokens)}</td>"
            f"<td>{money(avg_model_tokens)}</td>"
            f"<td>{request_cost_text(credits)}</td>"
            f"<td>{estimated_cost_cell}</td>"
            "</tr>"
        )
    if not models_rows:
        models_rows.append("<tr><td colspan='7'>Пока нет данных по моделям.</td></tr>")

    daily_rows: list[str] = []
    for row in report.get("daily", []):
        daily_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('day', '-') or '-'))}</td>"
            f"<td>{money(int(row.get('revenue', 0) or 0))} ₽</td>"
            f"<td>{int(row.get('active_users', 0) or 0)}</td>"
            "</tr>"
        )
    if not daily_rows:
        daily_rows.append("<tr><td colspan='3'>За период пока нет дневных данных.</td></tr>")

    screen_rows: list[str] = []
    for row in report.get("top_screens", []):
        screen_rows.append(
            "<tr>"
            f"<td>{esc(str(row.get('screen', '-') or '-'))}</td>"
            f"<td>{int(row.get('views', 0) or 0)}</td>"
            "</tr>"
        )
    if not screen_rows:
        screen_rows.append("<tr><td colspan='2'>Пока нет данных по экранам.</td></tr>")

    preset_rows: list[str] = []
    preset_names = {
        "fast": "⚡ Быстро",
        "balanced": "⚖ Баланс",
        "quality": "🧠 Качество",
        "expert": "🚀 Эксперт",
        "unknown": "Неизвестно",
    }
    for row in report.get("top_presets", []):
        preset_key = str(row.get("preset", "unknown") or "unknown")
        preset_rows.append(
            "<tr>"
            f"<td>{esc(preset_names.get(preset_key, preset_key))}</td>"
            f"<td>{int(row.get('uses', 0) or 0)}</td>"
            "</tr>"
        )
    if not preset_rows:
        preset_rows.append("<tr><td colspan='2'>Пока нет данных по режимам.</td></tr>")

    return f"""<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>Аналитика бота</title>
  <link rel="stylesheet" href="/assets/style.css"/>
  <style>
    .dash-grid {{ display:grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap:12px; margin-top:16px; }}
    .metric {{ border:1px solid var(--line); border-radius:16px; padding:16px; background:#fcfdff; }}
    .metric-label {{ color:var(--muted); font-size:13px; margin-bottom:8px; }}
    .metric-value {{ font-size:28px; font-weight:800; color:var(--text); }}
    .metric-note {{ color:var(--muted); font-size:13px; margin-top:6px; }}
    .panel-nav {{ display:flex; gap:10px; flex-wrap:wrap; margin-top:14px; }}
    .table-card {{ margin-top:14px; }}
    table {{ width:100%; border-collapse:collapse; }}
    th, td {{ border-bottom:1px solid var(--line); padding:10px 8px; text-align:left; font-size:14px; vertical-align:top; }}
    .muted {{ color:var(--muted); }}
    code {{ background:#f3f6fb; border-radius:6px; padding:2px 6px; }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="card">
      <div class="top">
        <span class="badge">Закрытая аналитика</span>
        <span class="badge">Период: {int(report.get('days', days) or days)} дней</span>
      </div>
      <h1>Деньги и метрики</h1>
      <p>Здесь удобно смотреть выручку, расход по моделям, рефералку и бумажную экономику проекта.</p>
      <div class="panel-nav">
        <a class="btn btn-primary" href="{esc(admin_url('/analytics', token, days=days))}">Аналитика</a>
        <a class="btn" href="{esc(admin_url('/admin/panel', token))}">Админка</a>
        <a class="btn" href="{esc(admin_url('/analytics/export/media-plan.xlsx', token, days=days))}">Excel медиаплан</a>
        <a class="btn" href="{esc(admin_url('/analytics/export/campaigns.csv', token, days=days))}">CSV кампаний</a>
        <a class="btn" href="{esc(admin_url('/analytics/export/campaigns.json', token, days=days))}">JSON кампаний</a>
        <a class="btn" href="{esc(admin_url('/analytics/logout', token))}">Выйти</a>
      </div>
      <div class="actions">{period_links}</div>
      <div class="dash-grid">
        <div class="metric"><div class="metric-label">Чистая выручка</div><div class="metric-value">{money(net_rub)} ₽</div><div class="metric-note">Выручка {money(revenue_rub)} ₽, возвраты {money(refunds_abs)} ₽</div></div>
        <div class="metric"><div class="metric-label">Плательщики</div><div class="metric-value">{payers}</div><div class="metric-note">{pay_share:.1f}% от активных</div></div>
        <div class="metric"><div class="metric-label">Активные пользователи</div><div class="metric-value">{active_users}</div><div class="metric-note">За выбранный период</div></div>
        <div class="metric"><div class="metric-label">ARPU / ARPPU</div><div class="metric-value">{money(arpu)} / {money(arppu)} ₽</div><div class="metric-note">На активного и на платящего</div></div>
        <div class="metric"><div class="metric-label">Средний чек</div><div class="metric-value">{money(avg_check)} ₽</div><div class="metric-note">{payments_count} оплат, refund rate {refund_rate_pct:.1f}%</div></div>
        <div class="metric"><div class="metric-label">Текстовые запросы</div><div class="metric-value">{text_requests}</div><div class="metric-note">Среднее {avg_tokens:.0f} токенов на текст</div></div>
        <div class="metric"><div class="metric-label">Картинки</div><div class="metric-value">{image_requests}</div><div class="metric-note">Списано {request_cost_text(image_credits)} запросов</div></div>
        <div class="metric"><div class="metric-label">Все запросы</div><div class="metric-value">{request_cost_text(total_credits)}</div><div class="metric-note">Текст {request_cost_text(text_credits)} / картинки {request_cost_text(image_credits)}</div></div>
        <div class="metric"><div class="metric-label">Себестоимость текста</div><div class="metric-value">{money(estimated_text_cost_rub)} ₽</div><div class="metric-note">{text_cost_share_pct:.1f}% от чистой выручки</div></div>
        <div class="metric"><div class="metric-label">Оценочная маржа</div><div class="metric-value">{money(estimated_text_contribution_rub)} ₽</div><div class="metric-note">{estimated_text_margin_pct:.1f}% от чистой выручки, пока без image-cost</div></div>
        <div class="metric"><div class="metric-label">Реф. активации</div><div class="metric-value">{referral_activations}</div><div class="metric-note">Реф. плательщики: {referred_payers}, конверсия {referral_conversion_pct:.1f}%</div></div>
        <div class="metric"><div class="metric-label">Бонусы рефералки</div><div class="metric-value">{request_balance_text(referral_bonus_credits)}</div><div class="metric-note">Выдано запросов по реферальной механике</div></div>
      </div>
    </div>
    <div class="grid">
      <div class="card table-card">
        <h2>Воронка</h2>
        <table>
          <tr><th>Этап</th><th>Пользователей</th><th>Доля</th></tr>
          <tr><td>Активные пользователи</td><td>{active_users}</td><td>100%</td></tr>
          <tr><td>Плательщики</td><td>{payers}</td><td>{pay_share:.1f}% от активных</td></tr>
          <tr><td>Новые пользователи</td><td>{new_users}</td><td>—</td></tr>
          <tr><td>Новые платные</td><td>{new_paid_users}</td><td>{free_to_paid_pct:.1f}% от новых</td></tr>
        </table>
      </div>
      <div class="card table-card">
        <h2>Retention и конверсия</h2>
        <table>
          <tr><th>Метрика</th><th>Значение</th><th>Комментарий</th></tr>
          <tr><td>D1 retention</td><td>{d1_retention_pct:.1f}%</td><td>{d1_retained} из {d1_eligible} пользователей</td></tr>
          <tr><td>D7 retention</td><td>{d7_retention_pct:.1f}%</td><td>{d7_retained} из {d7_eligible} пользователей</td></tr>
          <tr><td>Новые → платные</td><td>{free_to_paid_pct:.1f}%</td><td>{new_paid_users} из {new_users} новых пользователей</td></tr>
          <tr><td>Реферальная конверсия</td><td>{referral_conversion_pct:.1f}%</td><td>{referred_payers} платящих из {referral_activations} реф-активаций</td></tr>
        </table>
      </div>
    </div>
    <div class="grid">
      <div class="card table-card">
        <h2>Навигация и режимы</h2>
        <div class="grid">
          <div class="card table-card">
            <h3>Популярные экраны</h3>
            <table>
              <tr><th>Экран</th><th>Просмотров</th></tr>
              {''.join(screen_rows)}
            </table>
          </div>
          <div class="card table-card">
            <h3>Популярные режимы</h3>
            <table>
              <tr><th>Режим</th><th>Выборов</th></tr>
              {''.join(preset_rows)}
            </table>
          </div>
        </div>
      </div>
      <div class="card table-card">
        <h2>Монетизация</h2>
        <table>
          <tr><th>Метрика</th><th>Значение</th><th>Комментарий</th></tr>
          <tr><td>ARPU</td><td>{money(arpu)} ₽</td><td>На одного активного пользователя</td></tr>
          <tr><td>ARPPU</td><td>{money(arppu)} ₽</td><td>На одного платящего пользователя</td></tr>
          <tr><td>Средний чек</td><td>{money(avg_check)} ₽</td><td>{payments_count} оплат за период</td></tr>
          <tr><td>Refund rate</td><td>{refund_rate_pct:.1f}%</td><td>{refunds_count} возвратов</td></tr>
          <tr><td>Доля текстовой себестоимости</td><td>{text_cost_share_pct:.1f}%</td><td>{money(estimated_text_cost_rub)} ₽ от чистой выручки</td></tr>
          <tr><td>Маржа после текста</td><td>{money(estimated_text_contribution_rub)} ₽</td><td>{estimated_text_margin_pct:.1f}% от чистой выручки</td></tr>
        </table>
      </div>
    </div>
    <div class="card table-card">
      <h2>Выручка по тарифам</h2>
      <table>
        <tr><th>Тариф</th><th>Оплат</th><th>Выручка</th></tr>
        {''.join(plan_rows)}
      </table>
    </div>
    <div class="grid">
      <div class="card table-card">
        <h2>Топ рефереров</h2>
        <table>
          <tr><th>chat_id</th><th>user_id</th><th>Код</th><th>Пригласил</th><th>Платных</th><th>Последний реф</th><th></th></tr>
          {''.join(top_referrer_rows)}
        </table>
      </div>
      <div class="card table-card">
        <h2>Подозрительные реф-кластеры</h2>
        <p class="muted">Это не бан-лист, а очередь на проверку: много приглашений, но ноль платных, либо слишком плотная активность по рефералке.</p>
        <table>
          <tr><th>chat_id</th><th>user_id</th><th>Код</th><th>Пригласил</th><th>Активно за 14 дней</th><th>Платных</th><th></th></tr>
          {''.join(suspicious_referral_rows)}
        </table>
      </div>
      <div class="card table-card">
        <h2>Источники и кампании</h2>
        <p class="muted">Автоматически сюда попадает только то, что бот реально видит при первом запуске. Для рекламы в канал главный измеритель ниже — промокоды кампаний.</p>
        <table>
          <tr><th>Источник</th><th>Кампания</th><th>Пользователей</th><th>Платных</th><th>Выручка</th></tr>
          {''.join(source_rows)}
        </table>
      </div>
      <div class="card table-card">
        <h2>Промокоды кампаний</h2>
        <p class="muted">Для рекламы в канал проще давать в посте один код кампании. Лид в боте считается после активации кода, дальше видно оплаты и выручку.</p>
        <table>
          <tr><th>Код</th><th>Активаций</th><th>Платных</th><th>Конверсия</th><th>Запросов выдано</th><th>Выручка</th></tr>
          {''.join(promo_code_rows)}
        </table>
      </div>
    </div>
    <div class="card table-card">
      <h2>Экономика проекта</h2>
      <p>Шпаргалка по текущей сетке: цены, запросы, модели, бумажная маржа и допущения.</p>
      <div class="grid">
        <div class="card table-card">
          <h3>Подписки</h3>
          <table>
            <tr><th>Тариф</th><th>Цена</th><th>Дней</th><th>Запросов</th><th>Средне в день</th><th>Запр/₽</th><th>Модели</th></tr>
            {''.join(economics_plan_rows)}
          </table>
        </div>
        <div class="card table-card">
          <h3>Пакеты запросов</h3>
          <table>
            <tr><th>Пакет</th><th>Цена</th><th>Запросов</th><th>Запр/₽</th></tr>
            {''.join(economics_pack_rows)}
          </table>
        </div>
      </div>
      <div class="grid">
        <div class="card table-card">
          <h3>Free и лимиты</h3>
          <p>Free: {request_balance_text(FREE_DAILY_CREDITS)} запросов в день и 1 картинка каждые 7 дней.</p>
          <p>Картинка: {request_cost_text(image_credit_cost())} запросов. Редактирование фото: {request_cost_text(image_edit_credit_cost())} запросов.</p>
          <p>Макс. переменная доплата за текст: {request_cost_text(MAX_VARIABLE_CREDITS_PER_TEXT)} запрос.</p>
        </div>
        <div class="card table-card">
          <h3>Расчётные допущения</h3>
          <p>Курс для аналитики: {money(ANALYTICS_USD_TO_RUB)} ₽ за $1.</p>
          <p>Текстовая себестоимость выше считается по фактическим prompt/completion токенам и ценам из реестра моделей.</p>
          <p>Для бумажной экономики используем: эквайринг {ANALYTICS_PAYMENT_FEE_PCT:.2f}%, Т-Чеки {ANALYTICS_RECEIPT_FEE_PCT:.2f}%, налог {ANALYTICS_TAX_PCT:.2f}%, ожидаемая себестоимость 1 запроса {(ANALYTICS_EXPECTED_COST_PER_CREDIT_RUB * PUBLIC_REQUEST_UNIT_CREDITS):.3f} ₽.</p>
          <p class="muted">Настраивается через <code>ANALYTICS_USD_TO_RUB</code>, <code>ANALYTICS_PAYMENT_FEE_PCT</code>, <code>ANALYTICS_RECEIPT_FEE_PCT</code>, <code>ANALYTICS_TAX_PCT</code> и <code>ANALYTICS_EXPECTED_COST_PER_CREDIT_RUB</code>.</p>
        </div>
      </div>
      <div class="card table-card">
        <h3>Ожидаемая юнит-экономика</h3>
        <p>Стрессовый сценарий: считаем, что купленные запросы будут полностью выбраны. Это помогает видеть, остаётся ли прибыль на бумаге ещё до масштабирования.</p>
        <table>
          <tr><th>Продукт</th><th>Цена</th><th>Запросов</th><th>Себестоимость</th><th>Эквайринг</th><th>Т-Чеки</th><th>Налог</th><th>Маржа</th><th>Маржа %</th></tr>
          {''.join(unit_margin_rows)}
        </table>
      </div>
      <div class="card table-card">
        <h3>Маржа по тарифам</h3>
        <p>Это управленческая оценка: выручка минус текстовая себестоимость по фактическим токенам. Картинки и платёжные комиссии в этот блок пока не включены.</p>
        <table>
          <tr><th>Тариф</th><th>Выручка</th><th>Себестоимость текста</th><th>Маржа</th><th>Маржа %</th><th>Текстов</th><th>Картинок</th></tr>
          {''.join(margin_rows)}
        </table>
      </div>
    </div>
    <div class="grid">
      <div class="card table-card">
        <h2>Расход по моделям</h2>
        <p>Для текстовых моделей считаем токены, запросы и оценочную себестоимость. Для картинок в этой таблице стоимость не считаем, поэтому стоит «—».</p>
        <table>
          <tr><th>Модель</th><th>Тип</th><th>Событий</th><th>Токенов</th><th>Ср./событие</th><th>Запросов списано</th><th>Себестоимость</th></tr>
          {''.join(models_rows)}
        </table>
      </div>
      <div class="card table-card">
        <h2>Последние дни</h2>
        <table>
          <tr><th>День</th><th>Выручка</th><th>Активные</th></tr>
          {''.join(daily_rows)}
        </table>
      </div>
    </div>
  </div>
</body>
</html>"""


def campaign_export_rows(days: int = 30) -> list[dict[str, Any]]:
    report = state.user_store.kpi_report(days=days)
    source_rows = report.get("sources", []) or []
    promo_rows = report.get("promo_codes", []) or []

    by_campaign: dict[str, dict[str, Any]] = {}

    for row in source_rows:
        source = str(row.get("source", "") or "direct").strip().lower() or "direct"
        campaign = str(row.get("campaign", "") or "-").strip().lower() or "-"
        key = campaign if campaign != "-" else f"source:{source}"
        item = by_campaign.setdefault(
            key,
            {
                "campaign": campaign,
                "source": source,
                "users_count": 0,
                "promo_activations": 0,
                "paid_users": 0,
                "revenue_rub": 0,
            },
        )
        item["users_count"] += int(row.get("users_count", 0) or 0)
        item["paid_users"] += int(row.get("paid_users", 0) or 0)
        item["revenue_rub"] += int(row.get("revenue_rub", 0) or 0)

    for row in promo_rows:
        promo_code = str(row.get("promo_code", "") or "").strip().lower()
        if not promo_code:
            continue
        item = by_campaign.setdefault(
            promo_code,
            {
                "campaign": promo_code,
                "source": "promo",
                "users_count": 0,
                "promo_activations": 0,
                "paid_users": 0,
                "revenue_rub": 0,
            },
        )
        item["promo_activations"] += int(row.get("activations", 0) or 0)
        item["paid_users"] = max(item["paid_users"], int(row.get("paid_users", 0) or 0))
        item["revenue_rub"] = max(item["revenue_rub"], int(row.get("revenue_rub", 0) or 0))

    rows = sorted(
        by_campaign.values(),
        key=lambda x: (
            -int(x.get("revenue_rub", 0) or 0),
            -int(x.get("paid_users", 0) or 0),
            -int(x.get("promo_activations", 0) or 0),
            -int(x.get("users_count", 0) or 0),
            str(x.get("campaign", "")),
        ),
    )

    for item in rows:
        promo_activations = int(item.get("promo_activations", 0) or 0)
        paid_users = int(item.get("paid_users", 0) or 0)
        item["conversion_paid_from_promo_pct"] = (
            round(paid_users * 100.0 / promo_activations, 2) if promo_activations > 0 else 0.0
        )

    return rows


def campaign_export_csv(days: int = 30) -> str:
    rows = campaign_export_rows(days=days)
    buf = StringIO()
    writer = csv.writer(buf)
    writer.writerow(
        [
            "campaign",
            "source",
            "users_count",
            "promo_activations",
            "paid_users",
            "revenue_rub",
            "conversion_paid_from_promo_pct",
        ]
    )
    for row in rows:
        writer.writerow(
            [
                str(row.get("campaign", "") or ""),
                str(row.get("source", "") or ""),
                int(row.get("users_count", 0) or 0),
                int(row.get("promo_activations", 0) or 0),
                int(row.get("paid_users", 0) or 0),
                int(row.get("revenue_rub", 0) or 0),
                float(row.get("conversion_paid_from_promo_pct", 0.0) or 0.0),
            ]
        )
    return buf.getvalue()


def campaign_media_plan_xlsx(days: int = 30) -> bytes:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.datavalidation import DataValidation
    except Exception as exc:
        raise RuntimeError(f"xlsx_builder_unavailable: {exc}")

    rows = campaign_export_rows(days=days)
    by_campaign = {
        str(item.get("campaign", "") or "").strip().lower(): dict(item)
        for item in rows
        if str(item.get("campaign", "") or "").strip()
    }

    wb = Workbook()
    ws = wb.active
    ws.title = "Media Plan"
    headers = [
        "Priority",
        "Ad Link",
        "Campaign",
        "Channel",
        "Placement URL",
        "Planned Budget (RUB)",
        "Promo Status",
        "Channel Link",
        "Post Date",
        "Spent (RUB)",
        "Starts (users_count)",
        "Promo Activations",
        "Paid Users",
        "Revenue (RUB)",
        "CPL",
        "CPA Paid",
        "ROMI",
    ]
    ws.append(headers)
    header_fill = PatternFill("solid", fgColor="1F4E78")
    header_font = Font(color="FFFFFF", bold=True)
    thin = Side(style="thin", color="D9D9D9")
    for col, title in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=title)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)

    channel_link = CHANNEL_URL if CHANNEL_URL else "https://max.ru"
    for idx, (campaign, channel_name, budget, placement_url) in enumerate(ADS_MEDIA_CHANNELS, start=1):
        row_idx = idx + 1
        live = by_campaign.get(campaign.lower(), {})
        users_count = int(live.get("users_count", 0) or 0)
        promo_activations = int(live.get("promo_activations", 0) or 0)
        paid_users = int(live.get("paid_users", 0) or 0)
        revenue_rub = int(live.get("revenue_rub", 0) or 0)
        spent_rub = int(budget)
        ad_link = campaign_deep_link(campaign)

        ws.cell(row=row_idx, column=1, value=idx)
        ad_link_cell = ws.cell(row=row_idx, column=2, value=ad_link or "Set BOT_PUBLIC_URL in .env")
        if ad_link:
            ad_link_cell.hyperlink = ad_link
            ad_link_cell.style = "Hyperlink"
        ws.cell(row=row_idx, column=3, value=campaign)
        ws.cell(row=row_idx, column=4, value=channel_name)
        placement_cell = ws.cell(row=row_idx, column=5, value="Open placement")
        placement_cell.hyperlink = placement_url
        placement_cell.style = "Hyperlink"
        ws.cell(row=row_idx, column=6, value=budget)
        ws.cell(row=row_idx, column=7, value="planned")
        channel_cell = ws.cell(row=row_idx, column=8, value="Open MAX channel")
        channel_cell.hyperlink = channel_link
        channel_cell.style = "Hyperlink"
        ws.cell(row=row_idx, column=9, value="")
        ws.cell(row=row_idx, column=10, value=spent_rub)
        ws.cell(row=row_idx, column=11, value=users_count)
        ws.cell(row=row_idx, column=12, value=promo_activations)
        ws.cell(row=row_idx, column=13, value=paid_users)
        ws.cell(row=row_idx, column=14, value=revenue_rub)
        ws.cell(row=row_idx, column=15, value=f"=IF(L{row_idx}>0,J{row_idx}/L{row_idx},0)")
        ws.cell(row=row_idx, column=16, value=f"=IF(M{row_idx}>0,J{row_idx}/M{row_idx},0)")
        ws.cell(row=row_idx, column=17, value=f"=IF(J{row_idx}>0,(N{row_idx}-J{row_idx})/J{row_idx},0)")

    total_row = len(ADS_MEDIA_CHANNELS) + 3
    ws.cell(row=total_row, column=1, value="TOTAL").font = Font(bold=True)
    for col in (6, 10, 11, 12, 13, 14):
        letter = get_column_letter(col)
        ws.cell(row=total_row, column=col, value=f"=SUM({letter}2:{letter}{total_row-1})")
    ws.cell(row=total_row, column=15, value=f"=IF(L{total_row}>0,J{total_row}/L{total_row},0)")
    ws.cell(row=total_row, column=16, value=f"=IF(M{total_row}>0,J{total_row}/M{total_row},0)")
    ws.cell(row=total_row, column=17, value=f"=IF(J{total_row}>0,(N{total_row}-J{total_row})/J{total_row},0)")

    for r in range(2, total_row + 1):
        for c in range(1, 18):
            cell = ws.cell(row=r, column=c)
            cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
            if c in (1, 3, 6, 7, 10, 11, 12, 13, 14, 15, 16, 17):
                cell.alignment = Alignment(horizontal="center", vertical="center")
            else:
                cell.alignment = Alignment(vertical="center")

    for col in (6, 10, 14, 15, 16):
        for r in range(2, total_row + 1):
            ws.cell(row=r, column=col).number_format = "#,##0.00"
    for r in range(2, total_row + 1):
        ws.cell(row=r, column=17).number_format = "0.00%"

    widths = {
        1: 10, 2: 48, 3: 12, 4: 42, 5: 20, 6: 20, 7: 14, 8: 18, 9: 14,
        10: 12, 11: 16, 12: 16, 13: 12, 14: 14, 15: 10, 16: 12, 17: 10,
    }
    for col, width in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:Q{total_row-1}"
    status_validation = DataValidation(type="list", formula1='"planned,booked,posted,done,cancelled"', allow_blank=True)
    ws.add_data_validation(status_validation)
    status_validation.add(f"G2:G{total_row-1}")

    guide = wb.create_sheet("Guide")
    guide["A1"] = "Обновление данных"
    guide["A1"].font = Font(size=14, bold=True)
    guide["A3"] = "Файл генерируется автоматически из /analytics/export/campaigns.*."
    guide["A4"] = "Чтобы обновить цифры, просто скачай Excel снова из аналитики."
    guide["A5"] = "Лиды считаем по Promo Activations, деньги по Revenue."
    guide["A6"] = f"Период среза: {max(1, min(int(days), 365))} дней."
    guide.column_dimensions["A"].width = 120

    output = BytesIO()
    wb.save(output)
    return output.getvalue()


@app.get("/analytics", response_class=HTMLResponse)
async def analytics_page(request: Request, token: str = "", days: int = 30) -> HTMLResponse:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        return HTMLResponse(render_admin_login_html())

    response = HTMLResponse(render_admin_analytics_html_v2(token="", days=days))
    set_admin_cookie(response, session_id)
    return response


@app.get("/analytics/export/campaigns.json")
async def analytics_export_campaigns_json(request: Request, token: str = "", days: int = 30) -> dict[str, Any]:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        raise HTTPException(status_code=403, detail="forbidden")
    return {
        "days": max(1, min(int(days), 365)),
        "generated_at": datetime.utcnow().isoformat(),
        "rows": campaign_export_rows(days=days),
    }


@app.get("/analytics/export/campaigns.csv")
async def analytics_export_campaigns_csv(request: Request, token: str = "", days: int = 30) -> PlainTextResponse:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        raise HTTPException(status_code=403, detail="forbidden")
    csv_body = campaign_export_csv(days=days)
    response = PlainTextResponse(csv_body, media_type="text/csv; charset=utf-8")
    response.headers["Content-Disposition"] = "attachment; filename=campaigns_export.csv"
    set_admin_cookie(response, session_id)
    return response


@app.get("/analytics/export/media-plan.xlsx")
async def analytics_export_media_plan_xlsx(request: Request, token: str = "", days: int = 30) -> Response:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        raise HTTPException(status_code=403, detail="forbidden")
    try:
        content = campaign_media_plan_xlsx(days=days)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"xlsx_export_failed: {exc}")
    filename = f"media_plan_{max(1, min(int(days), 365))}d.xlsx"
    response = Response(
        content=content,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    response.headers["Content-Disposition"] = f'attachment; filename="{filename}"'
    set_admin_cookie(response, session_id)
    return response


@app.post("/analytics/login", response_class=HTMLResponse)
async def analytics_login(request: Request) -> HTMLResponse:
    allowed, reason = admin_login_allowed(request)
    if not allowed:
        return HTMLResponse(render_admin_login_html(reason), status_code=429)

    raw_body = (await request.body()).decode("utf-8", errors="ignore")
    form_data = dict(parse_qsl(raw_body, keep_blank_values=True))
    password = str(form_data.get("password", "")).strip()
    if not admin_panel_authorized(password):
        admin_login_register_failure(request)
        return HTMLResponse(render_admin_login_html("Неверный пароль."), status_code=401)

    admin_login_register_success(request)
    session_id = issue_admin_session()
    response = RedirectResponse(url="/analytics", status_code=303)
    set_admin_cookie(response, session_id)
    return response


@app.get("/analytics/logout")
async def analytics_logout(request: Request) -> RedirectResponse:
    session_id = request.cookies.get(ADMIN_SESSION_COOKIE, "").strip()
    if session_id:
        state.admin_sessions.pop(session_id, None)
    response = RedirectResponse(url="/analytics", status_code=303)
    clear_admin_cookie(response)
    return response


@app.get("/admin/panel", response_class=HTMLResponse)
async def admin_panel(
    request: Request,
    token: str = "",
    chat_id: int | None = None,
    request_id: int | None = None,
    q: str = "",
    payment_status: str = "",
) -> HTMLResponse:
    session_id = resolve_admin_session(request, token)
    if not session_id:
        raise HTTPException(status_code=403, detail="forbidden")
    response = HTMLResponse(
        render_admin_panel_html_v2(
            csrf_token=admin_csrf_token(session_id),
            chat_id=chat_id,
            request_id=request_id,
            q=q,
            payment_status=payment_status,
        )
    )
    set_admin_cookie(response, session_id)
    return response


def payment_status_signing_secret() -> str:
    return WEBHOOK_SECRET or TBANK_PASSWORD or ADMIN_PANEL_TOKEN


def payment_status_signature(request_id: int, status_ts: int) -> str:
    secret = payment_status_signing_secret()
    if not secret:
        return ""
    payload = f"{int(request_id)}:{int(status_ts)}"
    return hmac.new(secret.encode("utf-8"), payload.encode("utf-8"), hashlib.sha256).hexdigest()


def payment_status_signature_valid(request_id: int, status_ts_raw: str, status_sig: str) -> bool:
    if not status_sig:
        return False
    if not status_ts_raw.strip().isdigit():
        return False
    status_ts = int(status_ts_raw.strip())
    now_ts = int(datetime.utcnow().timestamp())
    # Allow small clock skew and enforce finite link lifetime.
    if status_ts < now_ts - 120:
        return False
    if status_ts > now_ts + max(60, PAYMENT_STATUS_TOKEN_TTL_SECONDS):
        return False
    expected = payment_status_signature(request_id, status_ts)
    if not expected:
        return False
    return hmac.compare_digest(expected, status_sig.strip())


@app.get("/admin/panel/action", response_class=HTMLResponse)
async def admin_panel_action_get() -> HTMLResponse:
    return HTMLResponse("Method Not Allowed", status_code=405)


@app.post("/admin/panel/action", response_class=HTMLResponse)
async def admin_panel_action(
    request: Request,
) -> HTMLResponse:
    session_id = resolve_admin_session(request, "")
    if not session_id:
        raise HTTPException(status_code=403, detail="forbidden")
    raw_body = (await request.body()).decode("utf-8", errors="ignore")
    form_data = dict(parse_qsl(raw_body, keep_blank_values=True))
    csrf = str(form_data.get("csrf", "")).strip()
    if not admin_csrf_valid(session_id, csrf):
        raise HTTPException(status_code=403, detail="csrf failed")

    action = str(form_data.get("type", "")).strip().lower()
    chat_raw = str(form_data.get("chat_id", "")).strip()
    request_raw = str(form_data.get("request_id", "")).strip()
    plan = str(form_data.get("plan", "")).strip()
    value = str(form_data.get("value", "")).strip()
    amount_raw = str(form_data.get("amount", "")).strip()

    chat_id = int(chat_raw) if chat_raw.isdigit() else None
    request_id = int(request_raw) if request_raw.isdigit() else None
    amount = int(amount_raw) if re.fullmatch(r"-?\d+", amount_raw) else 0
    message = "Готово"
    try:
        if action == "backup":
            backup_file = create_db_backup()
            message = f"Бэкап создан: {backup_file}"
        elif action == "nudge":
            sent, total = await send_reengage_nudges(days=REENGAGE_DORMANT_DAYS, limit=REENGAGE_BATCH_LIMIT)
            message = f"Реактивация: отправлено {sent}/{total}"
        elif action == "set_plan" and chat_id is not None and plan in PLAN_CONFIGS:
            user_profile(chat_id)
            selected = best_default_alias_for_plan(plan)
            if plan in PAID_PLANS:
                expires_at = state.user_store.set_subscription(
                    chat_id,
                    plan,
                    30,
                    selected,
                    recurring_enabled=False,
                )
                state.user_store.record_usage_event(
                    chat_id=chat_id,
                    event_type="admin_plan_change",
                    plan=plan,
                    details=f"source=admin_panel;mode=set_plan;expires_at={expires_at}",
                )
                message = f"Подписка пользователя {chat_id} -> {plan} до {format_msk_datetime(parse_iso_datetime(expires_at))}"
            else:
                state.user_store.set_plan(chat_id, plan)
                state.user_store.set_selected_model(chat_id, selected)
                state.user_store.set_credits(chat_id, credits_for_plan(plan))
                state.user_store.record_usage_event(
                    chat_id=chat_id,
                    event_type="admin_plan_change",
                    plan=plan,
                    details=f"source=admin_panel;mode=set_plan",
                )
                message = f"План пользователя {chat_id} -> {plan}"
        elif action == "set_sub" and chat_id is not None and plan in PAID_PLANS:
            user_profile(chat_id)
            expires_at = state.user_store.set_subscription(
                chat_id,
                plan,
                30,
                best_default_alias_for_plan(plan),
                recurring_enabled=False,
            )
            state.user_store.record_usage_event(
                chat_id=chat_id,
                event_type="admin_plan_change",
                plan=plan,
                details=f"source=admin_panel;mode=set_sub;expires_at={expires_at}",
            )
            message = f"Подписка пользователя {chat_id} -> {plan} до {format_msk_datetime(parse_iso_datetime(expires_at))}"
        elif action == "add_credits" and chat_id is not None and amount != 0:
            user_profile(chat_id)
            balance = state.user_store.adjust_credits(chat_id, amount)
            message = f"Баланс пользователя {chat_id}: {request_balance_text(balance)} запросов"
        elif action == "cancel_recurring" and chat_id is not None:
            row = user_profile(chat_id)
            expires_at = str(row.get("subscription_expires_at", "") or "")
            cancel_from = expires_at if expires_at else datetime.utcnow().replace(microsecond=0).isoformat()
            state.user_store.cancel_recurring(chat_id, cancel_from)
            message = f"Автопродление пользователя {chat_id} отключено"
        elif action == "block" and chat_id is not None and value in {"on", "off"}:
            state.user_store.set_blocked(chat_id, value == "on")
            message = f"Пользователь {chat_id}: block={value}"
        elif action == "payment" and request_id is not None and value in {"paid", "cancel"}:
            if value == "paid":
                ok, info = await activate_payment_request(request_id, source="admin panel")
                message = f"Платеж #{request_id}: {info}" if ok else f"Платеж #{request_id}: {info}"
            else:
                changed, info = cancel_payment_if_open(request_id)
                message = f"Платеж #{request_id}: {'отменен' if changed else info}"
        else:
            message = "Некорректные параметры действия."
    except Exception as exc:
        capture_exception_safe(exc)
        message = f"Ошибка: {exc}"

    response = HTMLResponse(
        render_admin_panel_html_v2(
            csrf_token=admin_csrf_token(session_id),
            chat_id=chat_id,
            request_id=request_id,
            message=message,
        )
    )
    set_admin_cookie(response, session_id)
    return response


def payment_status_view(request_id: int | None) -> dict[str, Any]:
    if request_id is None:
        return {
            "known": False,
            "status": "unknown",
            "status_label": payment_status_label("unknown"),
            "title": "Заявка не найдена",
            "message": "Не удалось определить номер заявки. Вернись в бот и нажми «Тарифы».",
        }

    payment = state.user_store.get_payment(request_id)
    if not payment:
        return {
            "known": False,
            "request_id": request_id,
            "status": "unknown",
            "status_label": payment_status_label("unknown"),
            "title": "Заявка не найдена",
            "message": "Такой заявки нет. Создай новую оплату в боте.",
        }

    status = str(payment.get("status", "pending")).lower()
    title, message = payment_status_title_message(status)
    return {
        "known": True,
        "request_id": request_id,
        "status": status,
        "status_label": payment_status_label(status),
        "title": title,
        "message": message,
        "plan": str(payment.get("plan", "")),
        "amount_rub": int(payment.get("amount_rub", 0)),
    }


def render_admin_panel_html_v2(
    csrf_token: str,
    chat_id: int | None = None,
    request_id: int | None = None,
    message: str = "",
    q: str = "",
    payment_status: str = "",
) -> str:
    query = str(q or "").strip()
    status_filter = str(payment_status or "").strip().lower()
    users = state.user_store.search_users(query, limit=50) if query else state.user_store.list_recent_users(limit=25)
    payments = state.user_store.search_payments(query=query, status=status_filter, limit=50) if (query or status_filter) else state.user_store.list_recent_payments(limit=25)
    suspicious_users = state.user_store.suspicious_users_report(limit=20)
    selected_user = state.user_store.get_user(chat_id) if chat_id else None
    selected_payment = state.user_store.get_payment(request_id) if request_id else None
    selected_user_payments = state.user_store.list_user_payments(chat_id, limit=8) if chat_id else []
    selected_user_events = state.user_store.list_user_usage_events(chat_id, limit=12) if chat_id else []
    esc = html.escape

    def action_form(label: str, **params: Any) -> str:
        hidden = [f"<input type='hidden' name='csrf' value='{esc(csrf_token)}'/>"]
        for key, value in params.items():
            if value is None or value == "":
                continue
            hidden.append(f"<input type='hidden' name='{esc(str(key))}' value='{esc(str(value))}'/>")
        return (
            "<form method='post' action='/admin/panel/action' style='display:inline-block; margin:0 6px 6px 0;'>"
            + "".join(hidden)
            + f"<button class='btn' type='submit'>{esc(label)}</button>"
            + "</form>"
        )

    def yes_no(flag: Any) -> str:
        return "да" if int(flag or 0) == 1 else "нет"

    info_block = ""
    if message:
        info_block += f"<p class='notice'>{esc(repair_mojibake(message))}</p>"
    if selected_user:
        selected_plan = str(selected_user.get("plan", "free") or "free").title()
        selected_alias = str(selected_user.get("selected_model_alias", "") or "").strip()
        selected_model_label = TEXT_MODELS.get(selected_alias, DEFAULT_TEXT_MODEL).label if selected_alias else "-"
        acquisition_source = str(selected_user.get("acquisition_source", "") or "").strip() or "direct"
        acquisition_campaign = str(selected_user.get("acquisition_campaign", "") or "").strip() or "-"
        expires_at = format_msk_datetime(parse_iso_datetime(str(selected_user.get("subscription_expires_at", "") or ""))) if str(selected_user.get("plan", "free")) != "free" else "—"
        referred_by = int(selected_user.get("referred_by_chat_id", 0) or 0)
        selected_user_summary = (
            "<div class='mini-grid'>"
            f"<div class='mini-metric'>Тариф<strong>{esc(selected_plan)}</strong></div>"
            f"<div class='mini-metric'>Запросы<strong>{request_balance_text(int(selected_user.get('credits_balance', 0) or 0))}</strong></div>"
            f"<div class='mini-metric'>Модель<strong>{esc(selected_model_label)}</strong></div>"
            f"<div class='mini-metric'>Автопродление<strong>{yes_no(selected_user.get('recurring_enabled', 0))}</strong></div>"
            f"<div class='mini-metric'>Блок<strong>{yes_no(selected_user.get('is_blocked', 0))}</strong></div>"
            f"<div class='mini-metric'>Реф-код<strong>{esc(str(selected_user.get('referral_code', '') or '-'))}</strong></div>"
            f"<div class='mini-metric'>Пригласил<strong>{int(selected_user.get('referrals_invited', 0) or 0)}</strong></div>"
            f"<div class='mini-metric'>Доступ до<strong>{esc(expires_at)}</strong></div>"
            f"<div class='mini-metric'>Источник<strong>{esc(acquisition_source)}</strong></div>"
            f"<div class='mini-metric'>Кампания<strong>{esc(acquisition_campaign)}</strong></div>"
            f"<div class='mini-metric'>Пригласил его<strong>{referred_by if referred_by > 0 else '—'}</strong></div>"
            "</div>"
            f"<p class='muted'>chat_id: {int(selected_user.get('chat_id', 0) or 0)} • user_id: {int(selected_user.get('max_user_id', 0) or 0)} • активность: {esc(str(selected_user.get('last_active_at', '') or '-'))}</p>"
        )
        payment_history_rows = []
        for item in selected_user_payments:
            payment_history_rows.append(
                "<tr>"
                f"<td>#{int(item.get('id', 0) or 0)}</td>"
                f"<td>{esc(payment_item_human_name(str(item.get('plan', '') or '')))}</td>"
                f"<td>{int(item.get('amount_rub', 0) or 0)} ₽</td>"
                f"<td>{esc(payment_status_label(str(item.get('status', '') or '')))}</td>"
                f"<td>{esc(format_msk_datetime(parse_iso_datetime(str(item.get('created_at', '') or ''))))}</td>"
                "</tr>"
            )
        if not payment_history_rows:
            payment_history_rows.append("<tr><td colspan='5'>Платежей по пользователю пока нет.</td></tr>")
        event_history_rows = []
        event_names = {
            "text_request": "Текстовый запрос",
            "image_request": "Картинка",
            "payment": "Оплата",
            "refund": "Возврат",
            "screen_view": "Экран",
            "preset_select": "Режим",
            "model_select": "Модель",
            "referral_activation": "Реф-активация",
            "referral_reward": "Реф-бонус",
            "promo_activation": "Промокод",
        }
        for item in selected_user_events:
            event_type = str(item.get("event_type", "") or "")
            details = str(item.get("details", "") or "")
            label = event_names.get(event_type, event_type or "Событие")
            event_history_rows.append(
                "<tr>"
                f"<td>{esc(label)}</td>"
                f"<td>{esc(str(item.get('model_alias', '') or '—'))}</td>"
                f"<td>{request_cost_text(int(item.get('credits_spent', 0) or 0))}</td>"
                f"<td>{int(item.get('rub_amount', 0) or 0)} ₽</td>"
                f"<td>{esc(details[:90] + ('…' if len(details) > 90 else ''))}</td>"
                f"<td>{esc(format_msk_datetime(parse_iso_datetime(str(item.get('created_at', '') or ''))))}</td>"
                "</tr>"
            )
        if not event_history_rows:
            event_history_rows.append("<tr><td colspan='6'>Истории действий пока нет.</td></tr>")
        info_block += (
            "<h3>Пользователь</h3>"
            + selected_user_summary
            + "<h3>Последние оплаты пользователя</h3>"
            + "<table><tr><th>ID</th><th>Продукт</th><th>Сумма</th><th>Статус</th><th>Когда</th></tr>"
            + "".join(payment_history_rows)
            + "</table>"
            + "<h3>Последние действия пользователя</h3>"
            + "<table><tr><th>Событие</th><th>Модель</th><th>Запросы</th><th>Сумма</th><th>Детали</th><th>Когда</th></tr>"
            + "".join(event_history_rows)
            + "</table>"
            + f"<details><summary>Сырой JSON пользователя</summary><pre>{esc(json.dumps(selected_user, ensure_ascii=False, indent=2))}</pre></details>"
        )
    if selected_payment:
        payment_status_text = payment_status_label(str(selected_payment.get("status", "") or ""))
        payment_summary = (
            "<div class='mini-grid'>"
            f"<div class='mini-metric'>Статус<strong>{esc(payment_status_text)}</strong></div>"
            f"<div class='mini-metric'>Сумма<strong>{int(selected_payment.get('amount_rub', 0) or 0)} ₽</strong></div>"
            f"<div class='mini-metric'>Продукт<strong>{esc(str(selected_payment.get('plan', '') or '-'))}</strong></div>"
            f"<div class='mini-metric'>Провайдер<strong>{esc(str(selected_payment.get('provider_payment_id', '') or '-'))}</strong></div>"
            "</div>"
            f"<p class='muted'>Заявка #{int(selected_payment.get('id', 0) or 0)} • chat_id: {int(selected_payment.get('chat_id', 0) or 0)} • создан: {esc(str(selected_payment.get('created_at', '') or '-'))}</p>"
        )
        info_block += "<h3>Платёж</h3>" + payment_summary + f"<details><summary>Сырой JSON платежа</summary><pre>{esc(json.dumps(selected_payment, ensure_ascii=False, indent=2))}</pre></details>"

    user_rows = []
    for row in users:
        cid = int(row.get("chat_id", 0) or 0)
        recurring = "on" if int(row.get("recurring_enabled", 0) or 0) == 1 else "off"
        user_rows.append(
            "<tr>"
            f"<td>{cid}</td>"
            f"<td>{int(row.get('max_user_id', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('plan', '')))}</td>"
            f"<td>{request_balance_text(int(row.get('credits_balance', 0) or 0))}</td>"
            f"<td>{recurring}</td>"
            f"<td>{int(row.get('is_blocked', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('last_active_at', '') or '-'))}</td>"
            f"<td><a href='{esc(admin_url('/admin/panel', chat_id=cid))}'>Открыть</a></td>"
            "</tr>"
        )

    payment_rows = []
    for row in payments:
        rid = int(row.get("id", 0) or 0)
        payment_rows.append(
            "<tr>"
            f"<td>{rid}</td>"
            f"<td>{int(row.get('chat_id', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('plan', '')))}</td>"
            f"<td>{int(row.get('amount_rub', 0) or 0)}</td>"
            f"<td>{esc(payment_status_label(str(row.get('status', ''))))}</td>"
            f"<td><a href='{esc(admin_url('/admin/panel', request_id=rid))}'>Открыть</a></td>"
            "</tr>"
        )

    suspicious_rows = []
    for row in suspicious_users:
        suspicious_rows.append(
            "<tr>"
            f"<td>{int(row.get('chat_id', 0) or 0)}</td>"
            f"<td>{int(row.get('max_user_id', 0) or 0)}</td>"
            f"<td>{esc(str(row.get('plan', '') or ''))}</td>"
            f"<td>{request_balance_text(int(row.get('credits_balance', 0) or 0))}</td>"
            f"<td>{esc(str(row.get('risk_reason', '') or ''))}</td>"
            f"<td>{esc(str(row.get('last_active_at', '') or '-'))}</td>"
            f"<td><a href='{esc(admin_url('/admin/panel', chat_id=int(row.get('chat_id', 0) or 0)))}'>Открыть</a></td>"
            "</tr>"
        )
    if not suspicious_rows:
        suspicious_rows.append("<tr><td colspan='7'>Явных подозрительных кейсов пока нет.</td></tr>")

    users_count = len(users)
    payments_count = len(payments)
    suspicious_count = len(suspicious_users)
    paid_users_count = sum(1 for row in users if str(row.get("plan", "free") or "free") != "free")
    recurring_count = sum(1 for row in users if int(row.get("recurring_enabled", 0) or 0) == 1)
    blocked_count = sum(1 for row in users if int(row.get("is_blocked", 0) or 0) == 1)
    pending_count = sum(1 for row in payments if str(row.get("status", "") or "").lower() == "pending")
    claimed_count = sum(1 for row in payments if str(row.get("status", "") or "").lower() == "claimed")
    paid_count = sum(1 for row in payments if str(row.get("status", "") or "").lower() == "paid")
    refunded_count = sum(1 for row in payments if str(row.get("status", "") or "").lower() == "refunded")
    search_form = (
        "<form class='actions' method='get' action='/admin/panel'>"
        f"<input class='input' type='text' name='q' value='{esc(query)}' placeholder='chat_id / user_id / тариф / заявка'/>"
        "<select class='input' name='payment_status'>"
        f"<option value='' {'selected' if not status_filter else ''}>Все статусы платежей</option>"
        f"<option value='pending' {'selected' if status_filter == 'pending' else ''}>Ожидает оплату</option>"
        f"<option value='claimed' {'selected' if status_filter == 'claimed' else ''}>Проверка оплаты</option>"
        f"<option value='paid' {'selected' if status_filter == 'paid' else ''}>Оплачено</option>"
        f"<option value='canceled' {'selected' if status_filter == 'canceled' else ''}>Отменено</option>"
        f"<option value='refunded' {'selected' if status_filter == 'refunded' else ''}>Возврат</option>"
        "</select>"
        "<button class='btn btn-primary' type='submit'>Искать</button>"
        "<a class='btn' href='/admin/panel'>Сбросить</a>"
        "</form>"
    )

    action_block = ""
    if selected_user:
        cid = int(selected_user["chat_id"])
        recurring_on = int(selected_user.get("recurring_enabled", 0) or 0) == 1
        action_block = (
            "<h3>Ручная корректировка</h3>"
            "<p>" + action_form("Вернуть на free", type="set_plan", chat_id=cid, plan="free") + "</p>"
            "<p>"
            + action_form("Lite на 30 дней", type="set_sub", chat_id=cid, plan="lite")
            + action_form("Start на 30 дней", type="set_sub", chat_id=cid, plan="start")
            + action_form("Pro на 30 дней", type="set_sub", chat_id=cid, plan="pro")
            + "</p>"
            "<p>"
            + action_form("+20 запросов", type="add_credits", chat_id=cid, amount=100)
            + action_form("+100 запросов", type="add_credits", chat_id=cid, amount=500)
            + action_form("-20 запросов", type="add_credits", chat_id=cid, amount=-100)
            + action_form("-100 запросов", type="add_credits", chat_id=cid, amount=-500)
            + "</p>"
            "<p>"
            + action_form("Отключить автопродление", type="cancel_recurring", chat_id=cid)
            + f"{' ✅' if recurring_on else ' (уже выключено)'}"
            + "</p>"
            "<p>"
            + action_form("Block ON", type="block", chat_id=cid, value="on")
            + action_form("Block OFF", type="block", chat_id=cid, value="off")
            + "</p>"
        )
    if selected_payment:
        rid = int(selected_payment["id"])
        action_block += (
            "<h3>Платёж</h3>"
            "<p>"
            + action_form("Подтвердить оплату", type="payment", request_id=rid, value="paid")
            + action_form("Отменить", type="payment", request_id=rid, value="cancel")
            + "</p>"
        )

    return f"""<!doctype html>
<html lang="ru">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>Админка бота</title>
  <link rel="stylesheet" href="/assets/style.css"/>
  <style>
    table {{ width:100%; border-collapse:collapse; }}
    th, td {{ border-bottom:1px solid var(--line); padding:10px 8px; text-align:left; font-size:14px; }}
    pre {{ white-space:pre-wrap; word-break:break-word; background:#fbfdff; border:1px solid #e1e8f5; border-radius:12px; padding:12px; }}
    .notice {{ padding:10px 12px; border:1px solid var(--line); border-radius:12px; background:#f8fbff; }}
    .mini-grid {{ display:grid; grid-template-columns: repeat(auto-fit, minmax(160px, 1fr)); gap:10px; margin:14px 0; }}
    .mini-metric {{ border:1px solid var(--line); border-radius:12px; padding:12px; background:#fcfdff; }}
    .mini-metric strong {{ display:block; font-size:22px; margin-top:4px; }}
    .input {{ min-height:42px; border:1px solid var(--line); border-radius:12px; padding:0 12px; background:#fff; color:var(--text); }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="card">
      <div class="top">
        <span class="badge">Закрытая админка</span>
      </div>
      <h1>Управление ботом</h1>
      <div class="actions">
        <a class="btn" href="/analytics">Аналитика</a>
        <a class="btn btn-primary" href="/admin/panel">Админка</a>
        <a class="btn" href="/status">Статус</a>
        {action_form("Создать бэкап БД", type="backup")}
        {action_form("Реактивировать free", type="nudge")}
        <a class="btn" href="/analytics/logout">Выйти</a>
      </div>
      {search_form}
      <div class="mini-grid">
        <div class="mini-metric">Пользователи в выдаче<strong>{users_count}</strong></div>
        <div class="mini-metric">Платежи в выдаче<strong>{payments_count}</strong></div>
        <div class="mini-metric">Подозрительные кейсы<strong>{suspicious_count}</strong></div>
        <div class="mini-metric">Платные в выдаче<strong>{paid_users_count}</strong></div>
        <div class="mini-metric">Автопродление<strong>{recurring_count}</strong></div>
        <div class="mini-metric">Заблокированы<strong>{blocked_count}</strong></div>
        <div class="mini-metric">Pending / Claimed<strong>{pending_count} / {claimed_count}</strong></div>
        <div class="mini-metric">Paid / Refunded<strong>{paid_count} / {refunded_count}</strong></div>
      </div>
      {info_block}
      {action_block}
    </div>
    <div class="card">
      <h2>Последние пользователи</h2>
      <table><tr><th>chat_id</th><th>user_id</th><th>plan</th><th>requests</th><th>recur</th><th>blocked</th><th>last_active_at</th><th></th></tr>
      {''.join(user_rows)}
      </table>
    </div>
    <div class="card">
      <h2>Подозрительные случаи</h2>
      <p>Здесь показываются платные аккаунты без успешной оплаченной заявки и случаи, где оплата подтверждалась вручную через админку. Это не приговор, а очередь на проверку.</p>
      <table><tr><th>chat_id</th><th>user_id</th><th>plan</th><th>requests</th><th>Причина</th><th>last_active_at</th><th></th></tr>
      {''.join(suspicious_rows)}
      </table>
    </div>
    <div class="card">
      <h2>Последние платежи</h2>
      <table><tr><th>id</th><th>chat_id</th><th>plan</th><th>amount</th><th>status</th><th></th></tr>
      {''.join(payment_rows)}
      </table>
    </div>
  </div>
</body>
</html>"""


@app.get("/payment/status")
async def payment_status(
    request: Request,
    request_id: int | None = None,
    status_ts: str = "",
    status_sig: str = "",
) -> dict[str, Any]:
    session_id = resolve_admin_session(request, "")
    is_admin_request = bool(session_id)
    if is_admin_request:
        # Keep admin session alive while inspecting statuses manually.
        pass
    elif request_id is None or not payment_status_signature_valid(request_id, status_ts, status_sig):
        raise HTTPException(status_code=401, detail="auth required")

    view = payment_status_view(request_id)
    if request_id is None or not view.get("known"):
        return view

    if view["status"] not in {"pending", "claimed", "paid"}:
        return view

    try:
        _, bank_status = await refresh_payment_from_tbank(request_id, source="T-Bank GetState")
        refreshed = payment_status_view(request_id)
        refreshed["bank_status"] = bank_status
        return refreshed
    except Exception:
        log.exception("T-Bank GetState failed for request_id=%s", request_id)
        return view


@app.get("/payment/success", response_class=FileResponse)
async def payment_success_page() -> FileResponse:
    return FileResponse(site_file("payment_success.html"))


@app.get("/payment/fail", response_class=FileResponse)
async def payment_fail_page() -> FileResponse:
    return FileResponse(site_file("payment_fail.html"))


@app.post("/webhook/max")
async def max_webhook(request: Request) -> dict[str, bool]:
    if WEBHOOK_SECRET:
        header_secret = request.headers.get("X-Max-Bot-Api-Secret", "")
        if header_secret != WEBHOOK_SECRET:
            log.warning("Webhook secret mismatch")
            raise HTTPException(status_code=403, detail="forbidden")

    payload = await request.json()
    updates = payload if isinstance(payload, list) else [payload]
    for update in updates:
        try:
            await process_update(update)
        except Exception as exc:
            if is_expected_max_delivery_error(exc):
                log.info("Skipping max_webhook alert for suspended/denied dialog: %s", exc)
                continue
            log.exception("Unhandled webhook processing error")
            capture_exception_safe(exc)
            record_runtime_error()
            await notify_admin_alert("max_webhook", f"Unhandled webhook error: {exc}")
    return {"ok": True}


@app.post("/webhook/tbank")
async def tbank_webhook(request: Request) -> PlainTextResponse:
    payload = await request.json()
    if not isinstance(payload, dict):
        raise HTTPException(status_code=400, detail="invalid payload")

    if not tbank_notification_is_valid(payload):
        log.warning("Invalid T-Bank webhook signature or terminal")
        with suppress(Exception):
            await notify_admin_alert("tbank_webhook", "Invalid webhook signature or terminal key.")
        raise HTTPException(status_code=403, detail="forbidden")

    order_id = scalar_string(payload.get("OrderId"))
    status = scalar_string(payload.get("Status")).upper()
    success_raw = payload.get("Success")
    success = success_raw is True or scalar_string(success_raw).lower() == "true"

    request_id = parse_request_id_from_order_id(order_id)
    if request_id is not None and success and status == "CONFIRMED":
        activated, info = await activate_payment_request(request_id, source="T-Bank webhook")
        if activated:
            log.info("T-Bank payment activated request_id=%s until=%s", request_id, info)
        else:
            log.info("T-Bank payment skipped request_id=%s reason=%s", request_id, info)
    elif request_id is not None and tbank_is_refund_status(status):
        changed, info = await process_refund_payment_request(request_id, source="T-Bank webhook", bank_status=status)
        if changed:
            log.info("T-Bank payment refunded request_id=%s result=%s", request_id, info)
        else:
            log.info("T-Bank refund ignored request_id=%s reason=%s", request_id, info)
    elif request_id is not None and tbank_is_cancel_status(status):
        changed, info = cancel_payment_if_open(request_id)
        if changed:
            log.info("T-Bank payment canceled request_id=%s status=%s", request_id, status)
        else:
            log.info("T-Bank cancel ignored request_id=%s reason=%s", request_id, info)

    return PlainTextResponse("OK")


def run() -> None:
    if RUN_MODE == "webhook":
        require_env()
        uvicorn.run("main:app", host=HOST, port=PORT, reload=False)
        return
    require_env()
    asyncio.run(polling_loop())


if __name__ == "__main__":
    run()
